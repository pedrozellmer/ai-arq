# -*- coding: utf-8 -*-
"""Compara múltiplas quotes de fornecedores (já parseadas) e gera:
  - dict com análise comparativa (totais, rankings, discrepâncias)
  - planilha XLSX (comparativo detalhado)
  - PPT executivo

Input: lista de quotes no formato produzido por supplier_quote_parser.

Também pode comparar contra o quantitativo ORIGINAL do projeto (items do
project_items) pra identificar itens que fornecedor esqueceu ou acrescentou.
"""
import os
from collections import defaultdict
from difflib import SequenceMatcher
from statistics import mean, stdev
from typing import List, Dict, Optional, Tuple
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter


# Threshold pra considerar dois textos "o mesmo item" (0-1, maior = mais estrito).
# Era 0.78 e NÃO casava reescrita real de fornecedor (2/8 nos testes de 19/07):
# "Piso vinílico em régua 3mm" vs "Fornecimento e instalação de piso vinílico
# 3mm" ficava fora. Baixado pra 0.62 + limpeza de stopwords/dimensões + trava
# de números incompatíveis (abaixo) — 0.62 sem a trava fundiria "porta 80cm"
# com "porta 90cm".
FUZZY_MATCH_THRESHOLD = 0.62

# Unidades equivalentes: fornecedor escreve "m2", nós "m²"; "und"/"unid"/"pç"
# são tudo unidade. Sem normalizar, a exigência de "mesma unidade" zerava o
# pareamento (0 pares nos testes de 19/07) → ranking "🏆 R$ 0,00".
_UN_MAP = {
    "m2": "m2", "m²": "m2", "m^2": "m2", "mq": "m2",
    "m3": "m3", "m³": "m3", "m^3": "m3",
    "un": "un", "und": "un", "unid": "un", "unidade": "un", "u": "un",
    "pc": "un", "pç": "un", "pca": "un", "pça": "un", "peca": "un", "peça": "un",
    "vb": "vb", "verba": "vb", "gb": "vb", "glb": "vb", "global": "vb",
    "ml": "m", "m": "m", "mt": "m", "metro": "m",
    "cj": "cj", "conj": "cj", "conjunto": "cj", "jg": "cj", "jogo": "cj",
    "pt": "pt", "ponto": "pt",
    "h": "h", "hr": "h", "hora": "h", "hs": "h",
    "kg": "kg", "ton": "t", "t": "t", "l": "l", "lt": "l",
    "mes": "mes", "mês": "mes", "dia": "dia", "diaria": "dia", "diária": "dia",
}


def _norm_un(un: str) -> str:
    """Unidade canônica pra comparação ('M².' → 'm2'). Vazio fica vazio."""
    u = (un or "").strip().lower().replace(".", "").replace(" ", "")
    return _UN_MAP.get(u, u)


_STOPWORDS = {"de", "da", "do", "das", "dos", "em", "com", "c/", "c",
              "para", "p/", "p", "e", "a", "o", "no", "na", "tipo",
              "fornecimento", "instalacao", "execucao", "incluso", "inclusa"}

def _clean_tokens(s: str):
    """Tokens limpos pra matching: colapsa dimensões ('60x60cm'→'60x60',
    vale pra 2 ou 3 medidas: '14x19x39cm'→'14x19x39'), remove stopwords/verbos
    de serviço e plural simples."""
    import re as _re
    s = _re.sub(r"(\d+(?:\s*[x×]\s*\d+)+)\s*(cm|mm|m)\b", r"\1", s or "")
    s = _re.sub(r"\s*[x×]\s*", "x", s)
    out = []
    for t in s.split():
        if t in _STOPWORDS:
            continue
        if len(t) > 3 and t.endswith("s"):
            t = t[:-1]
        out.append(t)
    return out


def _similarity(a: str, b: str) -> float:
    """Similaridade combinada: SequenceMatcher + Jaccard de tokens limpos.
    Trava de números: se ambos têm números e NENHUM coincide ('porta 80cm' vs
    'porta 90cm'), penaliza — dimensão diferente costuma ser item diferente."""
    if not a or not b:
        return 0.0
    if a == b:
        return 1.0

    ta, tb = _clean_tokens(a), _clean_tokens(b)
    ca, cb = " ".join(ta), " ".join(tb)
    if not ca or not cb:
        return 0.0
    if ca == cb:
        return 1.0

    seq = SequenceMatcher(None, ca, cb).ratio()
    sa, sb = set(ta), set(tb)
    jaccard = len(sa & sb) / len(sa | sb) if (sa and sb) else 0.0
    score = 0.5 * seq + 0.5 * jaccard

    import re as _re
    nums_a = {t for t in ta if _re.search(r"\d", t)}
    nums_b = {t for t in tb if _re.search(r"\d", t)}
    if nums_a and nums_b and not (nums_a & nums_b):
        score *= 0.7

    return score


def _fuzzy_merge(
    quotes: List[Dict],
    threshold: float = FUZZY_MATCH_THRESHOLD,
) -> List[Dict]:
    """Merge fuzzy: agrupa itens similares entre fornecedores.

    Estratégia:
    1) Acumula itens do fornecedor 1 como "âncoras"
    2) Pra cada item dos fornecedores 2..N, acha a melhor âncora (>=threshold)
       com MESMA UNIDADE; se achou, vincula. Se não, vira nova âncora.
    3) Resultado: lista de "buckets" com itens dos 3 fornecedores alinhados.
    """
    suppliers = [q["supplier_name"] for q in quotes]

    # Cada bucket: {desc, un, <sup_name>: item_or_None}
    buckets = []

    for q in quotes:
        sup = q["supplier_name"]
        for it in q.get("items", []):
            if it.get("is_section"):
                continue
            if not it.get("un") or not it.get("qtd") or not it.get("total"):
                continue  # só itens precificados
            if it["total"] <= 0:
                continue

            desc_norm = it.get("desc_norm", "")
            un = (it.get("un") or "").lower()
            un_norm = _norm_un(un)
            if not desc_norm:
                continue

            # Procura bucket com maior similaridade (mesma unidade NORMALIZADA:
            # "m²" do fornecedor A casa com "m2" do B)
            best_bucket_idx = -1
            best_score = threshold
            for idx, b in enumerate(buckets):
                if b.get(sup) is not None:
                    continue  # bucket já tem item desse fornecedor
                if (b.get("un_norm") or "") != un_norm:
                    continue  # unidade tem que bater
                score = _similarity(b["desc_norm"], desc_norm)
                if score > best_score:
                    best_score = score
                    best_bucket_idx = idx

            if best_bucket_idx >= 0:
                buckets[best_bucket_idx][sup] = it
                # Atualiza desc do bucket pro mais curto/limpo se aplicável
                if len(it["desc"]) < len(buckets[best_bucket_idx]["desc"]):
                    buckets[best_bucket_idx]["desc"] = it["desc"]
            else:
                # Cria novo bucket
                new_bucket = {
                    "desc": it["desc"],
                    "desc_norm": desc_norm,
                    "un": un,
                    "un_norm": un_norm,
                    **{s: None for s in suppliers},
                }
                new_bucket[sup] = it
                buckets.append(new_bucket)

    return buckets


# ═══════════════════════════════════════════════════════════════
#  COMPARAÇÃO / ANÁLISE
# ═══════════════════════════════════════════════════════════════

def compare_quotes(quotes: List[Dict],
                    reference_items: Optional[List[Dict]] = None) -> Dict:
    """Compara N quotes entre si (e opcionalmente contra reference).

    Args:
        quotes: lista de dicts no formato do parser (com supplier_name + items)
        reference_items: opcional — itens do quantitativo original do AI.arq
            (pra identificar o que foi cotado vs o que ficou de fora)

    Returns: dict {
        "suppliers": [...nomes...],
        "totals_by_supplier": {name: {total, mat, mo, n_items}},
        "coverage": {name: {n_priced, n_total, pct}},
        "merged_items": [...],  # itens unificados com total de cada fornecedor
        "paired_totals": {...},  # total só dos itens que TODOS orçaram
        "ranking": [...],
        "biggest_discrepancies": [...],
        "reference_check": {...}  # se reference_items fornecido
    }
    """
    suppliers = [q["supplier_name"] for q in quotes]

    # Totais brutos
    totals = {}
    for q in quotes:
        totals[q["supplier_name"]] = {
            "total": q["total_bruto"],
            "mat": q["total_material"],
            "mo": q["total_mao_obra"],
            "n_items": q["n_items_quoted"],
        }

    # Merge FUZZY — usa similaridade de string + Jaccard pra agrupar itens
    # descritos de formas ligeiramente diferentes entre fornecedores.
    # Exemplo: "Demolição de drywall em placas" casa com "Demolicao de drywall
    # com placas" (score ~0.85).
    merged = _fuzzy_merge(quotes)

    # Cobertura: quantos itens únicos cada fornecedor orçou
    n_unique = len(merged)
    coverage = {}
    for sup in suppliers:
        n = sum(1 for m in merged if m[sup])
        coverage[sup] = {
            "n_priced": n,
            "n_total_unique": n_unique,
            "pct": (n / n_unique * 100) if n_unique else 0,
        }

    # TOTAL PAREADO (só itens que TODOS orçaram)
    paired_totals = {sup: 0 for sup in suppliers}
    paired_count = 0
    for m in merged:
        if all(m[sup] and m[sup].get("total") for sup in suppliers):
            paired_count += 1
            for sup in suppliers:
                paired_totals[sup] += m[sup]["total"] or 0

    # Ranking por total pareado
    ranking = sorted(
        paired_totals.items(),
        key=lambda x: x[1] or float("inf")
    )

    # Maiores discrepâncias item a item
    discrepancies = []
    for m in merged:
        vals = [(sup, m[sup]["total"]) for sup in suppliers
                if m[sup] and m[sup].get("total")]
        if len(vals) < 2:
            continue
        tots = [v[1] for v in vals]
        vmin, vmax = min(tots), max(tots)
        if vmin <= 0:
            continue
        pct = (vmax - vmin) / vmin * 100
        discrepancies.append({
            "desc": m["desc"][:120],
            "un": m["un"],
            "totals": dict(vals),
            "min": vmin,
            "max": vmax,
            "pct_diff": pct,
            "cheapest": min(vals, key=lambda x: x[1])[0],
            "most_expensive": max(vals, key=lambda x: x[1])[0],
        })
    discrepancies.sort(key=lambda x: -x["pct_diff"])

    # Comparação contra quantitativo original (se fornecido)
    reference_check = None
    if reference_items:
        ref_check = _compare_against_reference(merged, reference_items, suppliers)
        reference_check = ref_check

    return {
        "suppliers": suppliers,
        "totals_by_supplier": totals,
        "coverage": coverage,
        "merged_items": merged,
        "paired_totals": paired_totals,
        "paired_count": paired_count,
        # Ranking só é honesto com itens de fato pareados: com <3 pares o
        # "menor preço" compara quase nada (chegava a coroar 🏆 R$ 0,00 com
        # zero pares). Quem exibe (XLSX/PPT/tela) DEVE checar esta flag.
        "ranking_confiavel": paired_count >= 3,
        "ranking": ranking,
        "biggest_discrepancies": discrepancies,
        "n_unique_items": n_unique,
        "reference_check": reference_check,
    }


def _compare_against_reference(merged, reference_items, suppliers):
    """Identifica:
    - quais items do quantitativo original NENHUM fornecedor cotou (esquecido)
    - quais items orçados NÃO existem no quantitativo original (adicionado)
    - divergências grandes de quantidade vs o CAD
    """
    # Normaliza referência pra comparar
    import unicodedata, re as re_mod
    def _n(s):
        if not s: return ""
        s = str(s).strip().lower()
        s = unicodedata.normalize("NFKD", s)
        s = "".join(c for c in s if not unicodedata.combining(c))
        s = re_mod.sub(r"[^\w\s]", " ", s)
        s = re_mod.sub(r"\s+", " ", s).strip()
        return s[:80]

    # Pareia REFERÊNCIA × ITENS ORÇADOS pelo mesmo critério fuzzy que o
    # _fuzzy_merge usa entre fornecedores. Antes era igualdade exata de texto —
    # mas fornecedor NUNCA copia a descrição igual ("Luminária LED de embutir
    # 60x60" vira "Luminaria LED embutir 60x60"). Com igualdade, o MESMO item
    # entrava nas duas listas: "nenhum fornecedor orçou" E "fornecedor
    # acrescentou". O painel mentia nos dois sentidos, sempre.
    refs = [it for it in reference_items if it.get("description")]
    ref_norm = [(_n(it.get("description", "")), _norm_un(it.get("unit") or ""), it)
                for it in refs]
    merged_norm = [(_n(m["desc"]), m.get("un_norm") or _norm_un(m["un"] or ""), m)
                   for m in merged]

    def _acha(alvo_desc, alvo_un, candidatos):
        """Melhor candidato acima do threshold, exigindo MESMA unidade (m² × m
        muda o item). Devolve None se ninguém passar."""
        melhor, melhor_score = None, 0.0
        for c_desc, c_un, c_obj in candidatos:
            if c_un != alvo_un:
                continue
            sc = _similarity(alvo_desc, c_desc)
            if sc >= FUZZY_MATCH_THRESHOLD and sc > melhor_score:
                melhor, melhor_score = c_obj, sc
        return melhor

    missing_from_suppliers = []
    for r_desc, r_un, ref in ref_norm:
        if _acha(r_desc, r_un, merged_norm) is None:
            missing_from_suppliers.append({
                "description": ref.get("description"),
                "unit": ref.get("unit"),
                "quantity": ref.get("quantity"),
                "note": "Item do quantitativo AI.arq que NENHUM fornecedor orçou",
            })

    extra_from_suppliers = []
    for m_desc, m_un, m in merged_norm:
        if _acha(m_desc, m_un, ref_norm) is None:
            cotou = [sup for sup in suppliers if m[sup]]
            extra_from_suppliers.append({
                "description": m["desc"],
                "unit": m["un"],
                "cotado_por": cotou,
                "note": "Item cotado por fornecedor que não existe no quantitativo AI.arq",
            })

    # Divergência de quantidade — mesmo pareamento fuzzy. Com igualdade exata
    # esta checagem quase nunca achava par, então nunca avisava que o fornecedor
    # estava orçando uma quantidade diferente da medida no CAD (que é o alerta
    # mais caro dos três).
    qty_divergence = []
    for m_desc, m_un, m in merged_norm:
        ref = _acha(m_desc, m_un, ref_norm)
        if not ref or not ref.get("quantity"):
            continue
        ref_qty = float(ref["quantity"])
        for sup in suppliers:
            q = m[sup]
            if not q or not q.get("qtd"):
                continue
            if ref_qty <= 0:
                continue
            diff_pct = abs(q["qtd"] - ref_qty) / ref_qty * 100
            if diff_pct > 20:  # > 20% de diferença
                qty_divergence.append({
                    "description": m["desc"],
                    "unit": m["un"],
                    "qty_aiarq": ref_qty,
                    "qty_supplier": q["qtd"],
                    "supplier": sup,
                    "diff_pct": diff_pct,
                })

    return {
        "missing_from_suppliers": missing_from_suppliers[:50],
        "extra_from_suppliers": extra_from_suppliers[:50],
        "qty_divergences": sorted(qty_divergence, key=lambda x: -x["diff_pct"])[:30],
        "summary": {
            "n_missing": len(missing_from_suppliers),
            "n_extra": len(extra_from_suppliers),
            "n_qty_divergence": len(qty_divergence),
        },
    }


# ═══════════════════════════════════════════════════════════════
#  XLSX COMPARATIVO
# ═══════════════════════════════════════════════════════════════

F_TITLE = Font(name="Arial", bold=True, size=14, color="FFFFFF")
F_HDR = Font(name="Arial", bold=True, size=10, color="FFFFFF")
F_N = Font(name="Arial", size=9)
F_BOLD = Font(name="Arial", bold=True, size=9)
F_SMALL = Font(name="Arial", size=8)
F_WIN = Font(name="Arial", size=9, bold=True, color="15803D")
F_LOSS = Font(name="Arial", size=9, color="DC2626")
F_MISS = Font(name="Arial", size=8, italic=True, color="9CA3AF")

P_TITLE = PatternFill("solid", fgColor="1E3A8A")
P_SEC = PatternFill("solid", fgColor="2563EB")
P_HDR = PatternFill("solid", fgColor="3B82F6")
P_WIN = PatternFill("solid", fgColor="DCFCE7")
P_MISS = PatternFill("solid", fgColor="F3F4F6")

AC = Alignment(horizontal="center", vertical="center", wrap_text=True)
AL = Alignment(horizontal="left", vertical="center", wrap_text=True)
BD = Border(left=Side("thin"), right=Side("thin"),
             top=Side("thin"), bottom=Side("thin"))


def generate_comparative_xlsx(analysis: Dict, output_path: str,
                                 project_name: str = "",
                                 architect_name: str = "",
                                 client_name: str = "") -> str:
    """Gera planilha comparativa dos fornecedores."""
    suppliers = analysis["suppliers"]
    n_sup = len(suppliers)

    wb = Workbook()
    ws = wb.active
    ws.title = "Comparativo"

    # Colunas: #, Descrição, UN, 1-per-supplier (R$), Δ, %dif, Menor, Obs
    base_cols = 3  # #, desc, un
    sup_cols = n_sup
    extra_cols = 4  # Δ, %dif, menor, obs
    total_cols = base_cols + sup_cols + extra_cols

    widths = [5, 45, 6] + [14] * n_sup + [14, 9, 14, 22]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # Título
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=total_cols)
    title = "COMPARATIVO DE FORNECEDORES"
    if project_name:
        title += f" — {project_name}"
    ws.cell(row=1, column=1, value=title).font = F_TITLE
    ws.cell(row=1, column=1).fill = P_TITLE
    ws.cell(row=1, column=1).alignment = AC
    ws.row_dimensions[1].height = 28

    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=total_cols)
    info = []
    if architect_name: info.append(f"Escritório: {architect_name}")
    if client_name: info.append(f"Cliente: {client_name}")
    ws.cell(row=2, column=1, value=" | ".join(info) or "Análise comparativa").font = Font(
        name="Arial", italic=True, size=9, color="6B7280")
    ws.cell(row=2, column=1).alignment = AC

    # RESUMO EXECUTIVO
    ro = 4
    ws.merge_cells(start_row=ro, start_column=1, end_row=ro, end_column=total_cols)
    ws.cell(row=ro, column=1, value="RESUMO EXECUTIVO").font = F_HDR
    ws.cell(row=ro, column=1).fill = P_SEC
    ws.cell(row=ro, column=1).alignment = AC
    ro += 1

    # Header totais
    hdr = ["", "Fornecedor", "", "Total MAT", "Total M.O.",
           "TOTAL BRUTO", "Cobertura", "TOTAL PAREADO", "Ranking*"]
    # Encaixa no layout (adapta pra n colunas)
    for c, h in enumerate(hdr[:total_cols], 1):
        ws.cell(row=ro, column=c, value=h).font = F_HDR
        ws.cell(row=ro, column=c).fill = P_HDR
        ws.cell(row=ro, column=c).alignment = AC
        ws.cell(row=ro, column=c).border = BD
    ro += 1

    ranking_list = sorted(
        analysis["paired_totals"].items(),
        key=lambda x: x[1] or float("inf")
    )
    # Com <3 itens pareados o ranking não compara nada — sem vencedor/destaque
    _rk_ok = bool(analysis.get("ranking_confiavel", analysis.get("paired_count", 0) >= 3))
    for rank, (sup, paired_tot) in enumerate(ranking_list, 1):
        data = analysis["totals_by_supplier"][sup]
        cov = analysis["coverage"][sup]
        is_winner = rank == 1 and _rk_ok

        ws.cell(row=ro, column=2, value=sup).font = F_WIN if is_winner else F_BOLD
        ws.cell(row=ro, column=4, value=data["mat"] or 0).number_format = "#,##0.00"
        ws.cell(row=ro, column=5, value=data["mo"] or 0).number_format = "#,##0.00"
        ws.cell(row=ro, column=6, value=data["total"] or 0).number_format = "#,##0.00"
        ws.cell(row=ro, column=7,
                value=f"{cov['n_priced']}/{cov['n_total_unique']} ({cov['pct']:.0f}%)"
                ).font = F_LOSS if cov["pct"] < 50 else F_N
        ws.cell(row=ro, column=8, value=paired_tot).number_format = "#,##0.00"
        if is_winner:
            ws.cell(row=ro, column=8).font = F_WIN
        ws.cell(row=ro, column=9,
                value=(f"{rank}º {'✔' if is_winner else ''}" if _rk_ok else "—"))
        if is_winner:
            ws.cell(row=ro, column=9).font = F_WIN
        for c in range(2, min(10, total_cols + 1)):
            ws.cell(row=ro, column=c).border = BD
            if c != 2:
                ws.cell(row=ro, column=c).alignment = AC
            if is_winner:
                ws.cell(row=ro, column=c).fill = P_WIN
        ro += 1

    # Nota metodológica
    ro += 1
    ws.merge_cells(start_row=ro, start_column=2, end_row=ro, end_column=total_cols)
    ws.cell(row=ro, column=2,
            value=(f"* Ranking por TOTAL PAREADO ({analysis['paired_count']} "
                   f"itens que todos os {n_sup} fornecedores orçaram). "
                   f"Comparação justa."
                   if _rk_ok else
                   f"* Poucos itens pareados entre os fornecedores "
                   f"({analysis['paired_count']}). Ranking suprimido — compare "
                   f"pelo TOTAL BRUTO e confira as descrições item a item.")
            ).font = Font(name="Arial", italic=True, size=9, color="6B7280")
    ro += 2

    # DETALHE ITEM A ITEM
    ws.merge_cells(start_row=ro, start_column=1, end_row=ro, end_column=total_cols)
    ws.cell(row=ro, column=1, value="DETALHE ITEM A ITEM").font = F_HDR
    ws.cell(row=ro, column=1).fill = P_SEC
    ws.cell(row=ro, column=1).alignment = AC
    ro += 1

    # Header detalhe
    hdrs = ["#", "DESCRIÇÃO", "UN"] + [f"{s} (R$)" for s in suppliers] + \
           ["Δ", "% dif", "MENOR", "OBS"]
    for c, h in enumerate(hdrs, 1):
        ws.cell(row=ro, column=c, value=h).font = F_HDR
        ws.cell(row=ro, column=c).fill = P_HDR
        ws.cell(row=ro, column=c).alignment = AC
        ws.cell(row=ro, column=c).border = BD
    ro += 1

    for idx, m in enumerate(analysis["merged_items"], 1):
        ws.cell(row=ro, column=1, value=idx).font = F_SMALL
        ws.cell(row=ro, column=2, value=m["desc"]).font = F_N
        ws.cell(row=ro, column=3, value=m["un"]).font = F_SMALL

        totals_row = {}
        for offset, sup in enumerate(suppliers):
            col = 4 + offset
            it = m[sup]
            if it and it.get("total"):
                ws.cell(row=ro, column=col, value=it["total"]
                        ).number_format = "#,##0.00"
                totals_row[sup] = it["total"]
            else:
                ws.cell(row=ro, column=col, value="—").font = F_MISS
                ws.cell(row=ro, column=col).fill = P_MISS

        vals = [v for v in totals_row.values() if v and v > 0]
        col_delta = 4 + n_sup
        col_pct = 5 + n_sup
        col_menor = 6 + n_sup
        col_obs = 7 + n_sup
        if len(vals) >= 2:
            vmin, vmax = min(vals), max(vals)
            delta = vmax - vmin
            pct = (delta / vmin * 100) if vmin > 0 else 0
            ws.cell(row=ro, column=col_delta, value=delta
                    ).number_format = "#,##0.00"
            ws.cell(row=ro, column=col_pct, value=f"{pct:.0f}%"
                    ).font = F_LOSS if pct > 50 else F_N
            winner = min(totals_row.items(), key=lambda x: x[1])[0]
            ws.cell(row=ro, column=col_menor, value=winner).font = F_WIN
            win_col = 4 + suppliers.index(winner)
            ws.cell(row=ro, column=win_col).fill = P_WIN
            if pct > 100:
                ws.cell(row=ro, column=col_obs,
                        value="⚠ dif >100% — escopo?").font = Font(
                    name="Arial", size=8, italic=True, color="DC2626")
        elif len(vals) == 1:
            ws.cell(row=ro, column=col_obs,
                    value="Só 1 orçou").font = F_MISS

        for c in range(1, total_cols + 1):
            ws.cell(row=ro, column=c).border = BD
            if c != 2:
                ws.cell(row=ro, column=c).alignment = AC
            else:
                ws.cell(row=ro, column=c).alignment = AL
        ro += 1

    ws.freeze_panes = "A12"

    # Insere logo do escritório (se URL fornecida)
    # Logo é tratado fora desta função — ver generate_comparative_xlsx_with_logo

    # Aba 2: reference check se tiver
    ref_check = analysis.get("reference_check")
    if ref_check:
        ws2 = wb.create_sheet("Vs. Quantitativo AI.arq")
        r = 1
        ws2.cell(row=r, column=1, value="ANÁLISE vs. QUANTITATIVO ORIGINAL"
                 ).font = F_TITLE
        ws2.cell(row=r, column=1).fill = P_TITLE
        ws2.merge_cells("A1:F1")
        r += 2

        ws2.cell(row=r, column=1,
                 value="Itens do seu quantitativo que NENHUM fornecedor cotou:"
                 ).font = F_BOLD
        r += 1
        for item in ref_check["missing_from_suppliers"][:20]:
            ws2.cell(row=r, column=1, value=item.get("description", "")).font = F_N
            ws2.cell(row=r, column=2, value=item.get("quantity") or "").font = F_N
            ws2.cell(row=r, column=3, value=item.get("unit", "")).font = F_SMALL
            r += 1

        r += 2
        ws2.cell(row=r, column=1,
                 value="Divergências de quantidade (fornecedor ≠ seu quantitativo):"
                 ).font = F_BOLD
        r += 1
        for qd in ref_check["qty_divergences"][:20]:
            ws2.cell(row=r, column=1, value=qd["description"]).font = F_N
            ws2.cell(row=r, column=2, value=qd["supplier"]).font = F_SMALL
            ws2.cell(row=r, column=3, value=f"AI.arq: {qd['qty_aiarq']}").font = F_N
            ws2.cell(row=r, column=4, value=f"Forn: {qd['qty_supplier']}").font = F_N
            ws2.cell(row=r, column=5, value=f"{qd['diff_pct']:.0f}% dif"
                     ).font = F_LOSS
            r += 1

        ws2.column_dimensions["A"].width = 60
        ws2.column_dimensions["B"].width = 15
        ws2.column_dimensions["C"].width = 20
        ws2.column_dimensions["D"].width = 20
        ws2.column_dimensions["E"].width = 15

    wb.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════════
#  PPT EXECUTIVO
# ═══════════════════════════════════════════════════════════════

def generate_comparative_pptx(analysis: Dict, output_path: str,
                                 project_name: str = "",
                                 architect_name: str = "",
                                 client_name: str = "",
                                 logo_path: Optional[str] = None,
                                 brand_color_hex: Optional[str] = None) -> str:
    """Gera PPT executivo com logo do escritório (co-branding).

    brand_color_hex: cor primária do escritório (#RRGGBB). Se não fornecido,
    usa azul AI.arq como default.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def _hex_to_rgb(hex_str):
        """Converte '#RRGGBB' ou 'RRGGBB' em RGBColor."""
        if not hex_str:
            return None
        s = hex_str.lstrip("#").strip()
        if len(s) != 6:
            return None
        try:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        except ValueError:
            return None

    # Aplica cor do escritório (se fornecida) em AZUL_ESCURO/AZUL.
    # Se não, usa azul AI.arq.
    brand_primary = _hex_to_rgb(brand_color_hex)
    if brand_primary:
        AZUL_ESCURO = brand_primary
        # Deriva versão "clara" (sobe em 25% cada canal pra contraste)
        r, g, b = brand_primary[0], brand_primary[1], brand_primary[2]
        AZUL = RGBColor(
            min(255, int(r * 1.2)),
            min(255, int(g * 1.2)),
            min(255, int(b * 1.2)),
        )
    else:
        AZUL_ESCURO = RGBColor(0x1E, 0x3A, 0x8A)
        AZUL = RGBColor(0x25, 0x63, 0xEB)

    VERDE = RGBColor(0x15, 0x80, 0x3D)
    VERMELHO = RGBColor(0xDC, 0x26, 0x26)
    CINZA = RGBColor(0x6B, 0x72, 0x80)
    CINZA_CLARO = RGBColor(0x9C, 0xA3, 0xAF)
    BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
    PRETO = RGBColor(0x1F, 0x29, 0x37)

    def _blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _rect(slide, x, y, w, h, fill):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        return shp

    def _txt(slide, x, y, w, h, text, size=16, bold=False, color=PRETO,
             align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return tb

    def _add_logo(slide, logo_path, x=Inches(0.3), y=Inches(0.2),
                    max_height=Inches(0.8)):
        if not logo_path or not os.path.exists(logo_path):
            return
        try:
            slide.shapes.add_picture(logo_path, x, y, height=max_height)
        except Exception:
            pass

    # ── SLIDE 1: CAPA ──
    slide = _blank()
    _rect(slide, 0, 0, prs.slide_width, prs.slide_height, AZUL_ESCURO)
    if logo_path:
        _add_logo(slide, logo_path, x=Inches(0.5), y=Inches(0.5),
                   max_height=Inches(1.2))
    _txt(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1),
         "COMPARATIVO DE FORNECEDORES", size=36, bold=True,
         color=BRANCO, align=PP_ALIGN.CENTER)
    _txt(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
         project_name or "Análise Comparativa", size=24,
         color=BRANCO, align=PP_ALIGN.CENTER)
    info_line = []
    if architect_name: info_line.append(architect_name)
    if client_name: info_line.append(f"Cliente: {client_name}")
    if info_line:
        _txt(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
             " | ".join(info_line), size=14,
             color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)
    from datetime import datetime as _dt
    _txt(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
         f"Análise gerada em {_dt.now().strftime('%d/%m/%Y')}", size=11,
         color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

    # ── SLIDE 2: RESUMO EXECUTIVO — CARDS ──
    slide = _blank()
    _add_logo(slide, logo_path)
    _txt(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Resumo Executivo — Totais por Fornecedor", size=28, bold=True,
         color=AZUL_ESCURO)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.03), AZUL)

    ranking_list = sorted(
        analysis["paired_totals"].items(),
        key=lambda x: x[1] or float("inf")
    )
    n_sup = len(ranking_list)
    card_w = Inches(min(3.9, 12.3 / n_sup - 0.2))
    card_gap = Inches(0.2)
    total_w = n_sup * card_w + (n_sup - 1) * card_gap
    start_x = (prs.slide_width - total_w) / 2
    card_y = Inches(1.6)
    card_h = Inches(3.5)

    # Com <3 itens pareados o "menor preço" não compara nada: sem troféu/verde
    _rk_ok_ppt = bool(analysis.get("ranking_confiavel",
                                   analysis.get("paired_count", 0) >= 3))
    for rank, (sup, paired_tot) in enumerate(ranking_list, 1):
        data = analysis["totals_by_supplier"][sup]
        cov = analysis["coverage"][sup]
        bg = (VERDE if rank == 1 else (AZUL if rank == 2 else CINZA)) \
            if _rk_ok_ppt else CINZA
        x = start_x + (rank - 1) * (card_w + card_gap)

        _rect(slide, x, card_y, card_w, card_h, bg)

        rank_label = ("🏆 MENOR PREÇO" if rank == 1 else f"{rank}º") \
            if _rk_ok_ppt else "—"
        _txt(slide, x + Inches(0.1), card_y + Inches(0.1),
              card_w - Inches(0.2), Inches(0.5),
              rank_label, size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        _txt(slide, x + Inches(0.1), card_y + Inches(0.7),
              card_w - Inches(0.2), Inches(0.6),
              sup, size=26, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        total_str = (f"R$ {paired_tot:,.2f}" if paired_tot
                     else "—")
        _txt(slide, x + Inches(0.1), card_y + Inches(1.5),
              card_w - Inches(0.2), Inches(0.8),
              total_str, size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        _txt(slide, x + Inches(0.1), card_y + Inches(2.3),
              card_w - Inches(0.2), Inches(0.4),
              f"(total pareado)", size=10,
              color=BRANCO, align=PP_ALIGN.CENTER)
        _txt(slide, x + Inches(0.1), card_y + Inches(2.7),
              card_w - Inches(0.2), Inches(0.4),
              f"Bruto: R$ {data['total']:,.0f}", size=11,
              color=BRANCO, align=PP_ALIGN.CENTER)
        _txt(slide, x + Inches(0.1), card_y + Inches(3.1),
              card_w - Inches(0.2), Inches(0.3),
              f"Cobertura: {cov['pct']:.0f}% ({cov['n_priced']} itens)", size=10,
              color=BRANCO, align=PP_ALIGN.CENTER)

    cheapest_paired = ranking_list[0][1] or 0
    most_exp_paired = ranking_list[-1][1] or 0
    if _rk_ok_ppt and cheapest_paired > 0 and most_exp_paired > cheapest_paired:
        delta = most_exp_paired - cheapest_paired
        pct = (delta / cheapest_paired) * 100
        _txt(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
              f"▲ DIFERENÇA: R$ {delta:,.2f}   ({pct:.1f}%)", size=18,
              bold=True, color=VERMELHO, align=PP_ALIGN.CENTER)
    elif not _rk_ok_ppt:
        _txt(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
              f"⚠ Poucos itens pareados ({analysis.get('paired_count', 0)}) — "
              f"compare pelo total bruto e revise as descrições.", size=14,
              bold=True, color=CINZA, align=PP_ALIGN.CENTER)

    _txt(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6),
         "Total pareado = só os itens que TODOS os fornecedores orçaram (comparação justa).",
         size=10, color=CINZA, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: TOP DISCREPÂNCIAS ──
    slide = _blank()
    _add_logo(slide, logo_path)
    _txt(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Top Itens — Maiores Divergências", size=26, bold=True, color=AZUL_ESCURO)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.03), AZUL)

    suppliers = analysis["suppliers"]
    top = analysis["biggest_discrepancies"][:12]

    table_y = Inches(1.5)
    row_h = Inches(0.42)

    # Header
    _rect(slide, Inches(0.5), table_y, Inches(12.3), row_h, AZUL_ESCURO)
    _txt(slide, Inches(0.6), table_y + Inches(0.05), Inches(5.5), row_h,
          "Item", size=10, bold=True, color=BRANCO)
    col_width_sup = Inches(1.5)
    for i, sup in enumerate(suppliers):
        _txt(slide, Inches(6.2 + i * 1.5), table_y + Inches(0.05),
              col_width_sup, row_h, sup, size=10, bold=True,
              color=BRANCO, align=PP_ALIGN.RIGHT)
    _txt(slide, Inches(6.2 + len(suppliers) * 1.5), table_y + Inches(0.05),
          Inches(1.3), row_h, "% dif", size=10, bold=True,
          color=BRANCO, align=PP_ALIGN.RIGHT)

    table_y += row_h
    for idx, d in enumerate(top):
        row_bg = RGBColor(0xF9, 0xFA, 0xFB) if idx % 2 == 0 else BRANCO
        _rect(slide, Inches(0.5), table_y, Inches(12.3), row_h, row_bg)
        _txt(slide, Inches(0.6), table_y + Inches(0.05), Inches(5.5),
              row_h, d["desc"][:75], size=9, color=PRETO)
        for i, sup in enumerate(suppliers):
            v = d["totals"].get(sup)
            if v is not None:
                val_str = f"R$ {v:,.0f}"
                is_min = abs(v - d["min"]) < 0.01
                is_max = abs(v - d["max"]) < 0.01
                col = VERDE if is_min else (VERMELHO if is_max else PRETO)
                bold = is_min or is_max
            else:
                val_str = "—"
                col = CINZA_CLARO
                bold = False
            _txt(slide, Inches(6.2 + i * 1.5), table_y + Inches(0.05),
                  col_width_sup, row_h, val_str, size=9, bold=bold,
                  color=col, align=PP_ALIGN.RIGHT)
        _txt(slide, Inches(6.2 + len(suppliers) * 1.5),
              table_y + Inches(0.05), Inches(1.3), row_h,
              f"{d['pct_diff']:.0f}%", size=9, bold=True,
              color=VERMELHO, align=PP_ALIGN.RIGHT)
        table_y += row_h

    # ── SLIDE 4: vs. QUANTITATIVO (se houver) ──
    ref_check = analysis.get("reference_check")
    if ref_check and ref_check.get("summary"):
        slide = _blank()
        _add_logo(slide, logo_path)
        _txt(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "vs. Seu Quantitativo AI.arq", size=26, bold=True,
              color=AZUL_ESCURO)
        _rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.03), AZUL)

        summary = ref_check["summary"]
        stats = [
            (f"{summary['n_missing']}",
              "itens do seu quantitativo que NENHUM fornecedor cotou",
              VERMELHO),
            (f"{summary['n_extra']}",
              "itens acrescentados pelos fornecedores (não estavam no seu quantitativo)",
              RGBColor(0xEA, 0x58, 0x0C)),
            (f"{summary['n_qty_divergence']}",
              "divergências de quantidade (>20%) entre fornecedor e seu CAD",
              RGBColor(0xEA, 0x58, 0x0C)),
        ]
        y = Inches(1.7)
        for num, label, color in stats:
            _rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.3),
                   RGBColor(0xF3, 0xF4, 0xF6))
            _rect(slide, Inches(0.5), y, Inches(0.15), Inches(1.3), color)
            _txt(slide, Inches(0.8), y + Inches(0.2), Inches(2.5), Inches(1),
                  num, size=44, bold=True, color=color)
            _txt(slide, Inches(3.5), y + Inches(0.35), Inches(9), Inches(0.9),
                  label, size=14, color=PRETO)
            y += Inches(1.4)

    # ── SLIDE 5: METODOLOGIA ──
    slide = _blank()
    _add_logo(slide, logo_path)
    _txt(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Metodologia", size=26, bold=True, color=AZUL_ESCURO)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.03), AZUL)

    metodo = (
        "• Comparação item-a-item entre todas as cotações recebidas, normalizando descrições "
        "(sem acento/case) pra casar o mesmo serviço escrito de formas diferentes.\n\n"
        "• TOTAL PAREADO: soma apenas dos itens que TODOS os fornecedores orçaram — "
        "garante comparação justa, mesmo que cada um tenha um escopo ligeiramente diferente.\n\n"
        "• COBERTURA: % de itens únicos que cada fornecedor cotou. Cobertura baixa "
        "sugere escopo incompleto (não 'mais barato').\n\n"
        "• Diferenças acima de 100% entre fornecedores foram sinalizadas — normalmente "
        "indicam que estão orçando coisas diferentes (material premium vs básico, "
        "escopo com/sem BDI, etc), não simplesmente 'mais caro'.\n\n"
        "• Cruzamento contra o quantitativo AI.arq destaca itens esquecidos ou "
        "divergências de medição.\n\n"
        "• Critério de decisão: MENOR PREÇO ≠ melhor escolha. Considere prazo, "
        "garantia, referências, capacidade técnica e reputação do fornecedor."
    )
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, linha in enumerate(metodo.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = linha
        r.font.name = "Arial"
        r.font.size = Pt(13)
        r.font.color.rgb = PRETO

    # Rodapé em todas as slides
    footer_left = architect_name or ""
    footer_right = "Análise preparada por AI.arq (ai.arq.br)"
    for s in prs.slides:
        if footer_left:
            _txt(s, Inches(0.3), Inches(7.2), Inches(6), Inches(0.25),
                  footer_left, size=8, color=CINZA_CLARO)
        _txt(s, Inches(7), Inches(7.2), Inches(6), Inches(0.25),
              footer_right, size=8, color=CINZA_CLARO, align=PP_ALIGN.RIGHT)

    prs.save(output_path)
    return output_path
