# -*- coding: utf-8 -*-
"""API Backend AI.arq — Processamento de pranchas de arquitetura."""
import os
import time
import uuid
import shutil
import asyncio
import tempfile
import threading
from pathlib import Path
from typing import Optional
from datetime import datetime, timedelta

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from pydantic import BaseModel
from dotenv import load_dotenv

# Carregar .env do mesmo diretório do script
_script_dir = os.path.dirname(os.path.abspath(__file__))
_env_path = os.path.join(_script_dir, '.env')
if os.path.exists(_env_path):
    with open(_env_path, 'r') as _f:
        for _line in _f:
            _line = _line.strip()
            if '=' in _line and not _line.startswith('#'):
                _k, _v = _line.split('=', 1)
                os.environ[_k.strip()] = _v.strip()
else:
    load_dotenv()

from models import ProcessingStatus
from processor import process_pdfs
from analyzer import analyze_all_sheets
from spreadsheet import generate_spreadsheet
from instagram_webhook import router as instagram_router
from whatsapp_notify import router as whatsapp_router, send_whatsapp_template
from engine_rules import (
    salvage_truncated_json as _salvage_truncated_json,
    extract_balanced_obj as _extract_balanced_obj,
    normalize_items_payload as _normalize_items_payload,
    should_force_steel_kg as _should_force_steel_kg,
    is_likely_wrong_type as _is_likely_wrong_type,
    extraction_has_quality_caveat as _extraction_has_quality_caveat,
    extract_block_name as _extract_block_name,
    is_nonsense_item as _is_nonsense_item,
    extract_type_code as _extract_type_code,
    response_truncated as _response_truncated,
    is_floor_surface as _is_floor_surface,
    is_unit_mismatch_countable as _is_unit_mismatch_countable,
    AREA_UNITS_HONESTY as _AREA_UNITS_HONESTY,
    FLOOR_M2_UNITS as _FLOOR_M2_UNITS,
)
# calibrator.py foi desativado: o modelo de "fator absoluto" (real/ai) não
# respeita o isolamento entre projetos. A calibração agora é 100% por
# densidade (density_calibration.py) e só gera alertas.

# Supabase client para salvar projetos
# Env vars recomendados no Render:
#   - SUPABASE_URL
#   - SUPABASE_KEY (ou SUPABASE_ANON_KEY) — usada APENAS como `apikey` no header
#     do PostgREST (exigência do gateway). Não deve mais ir no Authorization.
#   - SUPABASE_SERVICE_ROLE_KEY — usada em queries não autenticadas (webhook,
#     jobs background, contato, leads, admin). Bypassa RLS. NUNCA no frontend.
# Decidido 2026-06-02 após auditoria RLS revelar policies abertas pra anon.
# Fluxo atualizado:
#   - JWT do user → _supa_rest_as_user(request, ...) — pra queries onde o user
#     é dono dos dados (RLS valida `auth.uid() = user_id`).
#   - service_role → _supa_rest_service(...) — pra queries sem request (jobs
#     background, webhooks, formulários públicos, admin RPCs etc.).
SUPABASE_URL = os.getenv("SUPABASE_URL", "https://kqjabzwgbfuivzlcfvvu.supabase.co")
SUPABASE_KEY = os.getenv("SUPABASE_KEY") or os.getenv("SUPABASE_ANON_KEY")
if not SUPABASE_KEY:
    print("[WARN] SUPABASE_KEY não configurado no ambiente — usando fallback hardcoded (remover em breve)")
    SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6ImtxamFiendnYmZ1aXZ6bGNmdnZ1Iiwicm9sZSI6ImFub24iLCJpYXQiOjE3NzYwMDg5NzcsImV4cCI6MjA5MTU4NDk3N30.48xSenZlDV0LfD94ZxwGvX41Kf9Je2n-ouZpJrrCSKI"

# Service role key pra queries não autenticadas. Fallback pra SUPABASE_KEY
# (anon) com WARN — mantém compat enquanto o Render não tem a env setada.
# Quando a env for setada e a Onda B do RLS aplicada, anon não passa mais por
# nenhuma policy (revoga GRANT pra anon nas tabelas internas).
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY") or SUPABASE_KEY
if SUPABASE_SERVICE_ROLE_KEY == SUPABASE_KEY:
    print("[WARN] SUPABASE_SERVICE_ROLE_KEY não setada — usando anon como fallback. "
          "Setar em Render env vars antes de aplicar a migration que revoga anon.")

# Email do admin (tem acesso a /api/admin/*). Pode ser overridado por env.
ADMIN_EMAIL = os.getenv("ADMIN_EMAIL", "zarelalopes@gmail.com").lower()

# Log persistente de operações Supabase (só erros + último sucesso por operação)
# pra poder investigar via /api/debug/supa-log quando o log do Render tá fora de alcance.
_SUPA_LOG_PATH = os.path.join(tempfile.gettempdir(), "aiarq_supa_log.txt")


def _supa_log(line: str):
    try:
        with open(_SUPA_LOG_PATH, "a", encoding="utf-8") as f:
            f.write(f"{datetime.utcnow().isoformat()}Z {line}\n")
    except Exception:
        pass


def _supabase_insert(table, data):
    """Insere registro no Supabase via REST API."""
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/{table}"
        body = json.dumps(data).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        req.add_header('Prefer', 'return=minimal')
        urllib.request.urlopen(req, timeout=20)
        _supa_log(f"INSERT {table} OK  data={json.dumps(data)[:200]}")
        return True
    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode('utf-8', errors='replace')[:500]
        except Exception:
            resp_body = '(unreadable)'
        msg = f"INSERT {table} HTTP {e.code}: {resp_body}  data={json.dumps(data)[:200]}"
        print(f"Supabase insert HTTP {e.code} ({table}): {resp_body}")
        _supa_log(msg)
        return False
    except Exception as e:
        msg = f"INSERT {table} ERR {type(e).__name__}: {e}  data={json.dumps(data)[:200]}"
        print(f"Supabase insert error ({table}): {type(e).__name__}: {e}")
        _supa_log(msg)
        return False


def _log_error(stage, message, job_id=None, severity="error"):
    """Grava um erro técnico do motor na tabela error_log do Supabase, pra ser
    lido via MCP/admin SEM precisar abrir o log do Render. Best-effort, NUNCA
    levanta (não pode atrapalhar quem já está num except)."""
    try:
        _supabase_insert("error_log", {
            "stage": str(stage)[:80],
            "message": str(message)[:2000],
            "job_id": (str(job_id)[:40] if job_id else None),
            "severity": severity,
        })
    except Exception:
        pass


def _error_log_causa_real(job_id: str, limit: int = 3) -> str:
    """Causa TÉCNICA real de um job (o que o cliente NÃO viu — ele viu o rótulo).
    Junta as entradas mais recentes do error_log por job_id pra o alerta do Pedro
    mostrar a raiz ao lado do rótulo (QW3, 20/07). Best-effort — nunca levanta."""
    if not job_id:
        return ""
    import urllib.request as _u, json as _j
    try:
        q = (f"{SUPABASE_URL}/rest/v1/error_log?job_id=eq.{job_id}"
             f"&severity=in.(error,critical)"
             f"&select=stage,message,created_at&order=created_at.desc&limit={limit}")
        req = _u.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _j.loads(_u.urlopen(req, timeout=8).read().decode("utf-8"))
        partes = []
        for r in rows or []:
            _stage = (r.get("stage") or "").strip()
            _msg = (r.get("message") or "").strip()
            if _msg:
                partes.append(f"[{_stage}] {_msg[:280]}" if _stage else _msg[:280])
        return " · ".join(partes)
    except Exception:
        return ""


def _supabase_update(table, match_field, match_value, data):
    """Atualiza registro no Supabase via REST API.

    IMPORTANTE: pra table='projects' com match_field='job_id', usa RPC
    `update_project_status` (SECURITY DEFINER) que bypassa RLS. Sem isso, o
    PATCH volta 200 OK mas 0 rows afetadas — causa do bug silencioso onde
    projetos de usuários logados ficavam 'queued' pra sempre no admin mesmo
    com processamento concluído."""
    import urllib.request, urllib.error, json

    # Caminho especial: atualizar projects por job_id via RPC.
    # Roteia conforme os campos sendo atualizados:
    # - meta editável (project_name/typology/address/phase) → update_project_meta
    # - status/area/warnings/etc → update_project_status (legado)
    if table == "projects" and match_field == "job_id":
        if any(k in _META_FIELDS for k in data.keys()):
            return _rpc_update_project_meta(match_value, data)
        return _rpc_update_project_status(match_value, data)

    try:
        url = f"{SUPABASE_URL}/rest/v1/{table}?{match_field}=eq.{match_value}"
        body = json.dumps(data).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='PATCH')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        req.add_header('Prefer', 'return=minimal')
        urllib.request.urlopen(req, timeout=20)
        _supa_log(f"UPDATE {table} {match_field}={match_value} OK  data={json.dumps(data)[:200]}")
        return True
    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode('utf-8', errors='replace')[:500]
        except Exception:
            resp_body = '(unreadable)'
        msg = f"UPDATE {table} {match_field}={match_value} HTTP {e.code}: {resp_body}  data={json.dumps(data)[:200]}"
        print(f"Supabase update HTTP {e.code} ({table} where {match_field}={match_value}): {resp_body}")
        _supa_log(msg)
        return False
    except Exception as e:
        msg = f"UPDATE {table} {match_field}={match_value} ERR {type(e).__name__}: {e}  data={json.dumps(data)[:200]}"
        print(f"Supabase update error ({table} where {match_field}={match_value}): {type(e).__name__}: {e}")
        _supa_log(msg)
        return False


# ═══════════════════════════════════════════════════════════════
#  AUTH HELPERS — JWT do Supabase para admin e ownership
# ═══════════════════════════════════════════════════════════════

def _get_user_from_request(request):
    """Valida header Authorization: Bearer <jwt> contra /auth/v1/user do Supabase.

    Retorna dict com {"id", "email"} se válido, ou None se inválido/ausente."""
    import urllib.request, urllib.error, json as _j
    try:
        auth_header = request.headers.get("Authorization", "") or request.headers.get("authorization", "")
        if not auth_header.lower().startswith("bearer "):
            return None
        token = auth_header[7:].strip()
        if not token:
            return None
        url = f"{SUPABASE_URL}/auth/v1/user"
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {token}")
        resp = urllib.request.urlopen(req, timeout=10)
        data = _j.loads(resp.read().decode("utf-8"))
        uid = data.get("id")
        email = (data.get("email") or "").lower()
        if not uid:
            return None
        return {"id": uid, "email": email}
    except Exception as _e:
        _supa_log(f"AUTH validate fail: {type(_e).__name__}: {_e}")
        return None


def _require_admin(request):
    """Valida JWT + email == ADMIN_EMAIL. Raise HTTPException se falha.

    Retorna o dict do usuário em caso de sucesso."""
    user = _get_user_from_request(request)
    if not user:
        raise HTTPException(401, "Autenticação requerida (Bearer token ausente ou inválido)")
    if user.get("email", "").lower() != ADMIN_EMAIL:
        raise HTTPException(403, "Acesso restrito a administradores")
    return user


def _supa_rest_service(method: str, path: str, body=None, params=None,
                        prefer: str = None, timeout: int = 15):
    """Faz uma chamada Supabase REST usando a SERVICE_ROLE_KEY (bypassa RLS).

    Use SOMENTE em rotas/funções não autenticadas, onde não há JWT do user:
      - Webhooks (Meta/Instagram, Stripe etc.)
      - Jobs background (recovery de jobs travados, cleanup, processamento)
      - Formulários públicos (contato, leads do chat widget)
      - Admin/RPCs que já têm checagem própria (_require_admin)
      - Helpers internos chamados de threads sem `request` em escopo

    Para queries onde o user é dono dos dados e o request está em escopo, use
    `_supa_rest_as_user(request, ...)` — preserva RLS por usuário.

    Comportamento:
      - `apikey` sempre é SUPABASE_KEY (PostgREST exige).
      - `Authorization: Bearer <service_role>` bypassa toda RLS.
      - `path` aceita formato absoluto (`/rest/v1/...`) ou só a tabela
        (ex.: `chat_leads?email=eq.X`); prefixo é adicionado se faltar.
      - `body` aceita dict/list; serializa pra JSON. None = sem body.
      - `params` é dict opcional → query string (concatenado).
      - `prefer` é o header Prefer (`return=minimal`, `resolution=...`, etc.).

    Retorna tupla (status_code, parsed_json_or_None). Erros HTTP devolvem
    (code, None) com log. Falha total devolve (0, None).
    """
    import urllib.request, urllib.error, json as _j

    # Monta URL
    if path.startswith("http"):
        url = path
    else:
        if not path.startswith("/"):
            path = "/" + path
        if not path.startswith("/rest/v1/") and not path.startswith("/storage/"):
            path = "/rest/v1" + path
        url = f"{SUPABASE_URL}{path}"

    if params:
        from urllib.parse import urlencode
        sep = "&" if "?" in url else "?"
        url = f"{url}{sep}{urlencode(params)}"

    # Serializa body se houver
    data_bytes = None
    if body is not None:
        data_bytes = _j.dumps(body, default=str).encode("utf-8")

    try:
        req = urllib.request.Request(url, data=data_bytes, method=method.upper())
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        if data_bytes is not None:
            req.add_header("Content-Type", "application/json")
        if prefer:
            req.add_header("Prefer", prefer)
        resp = urllib.request.urlopen(req, timeout=timeout)
        status = resp.getcode()
        raw = resp.read()
        if not raw:
            return status, None
        try:
            return status, _j.loads(raw.decode("utf-8"))
        except Exception:
            return status, None
    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode("utf-8", errors="replace")[:300]
        except Exception:
            resp_body = "(unreadable)"
        _supa_log(f"REST_SERVICE {method} {path} HTTP {e.code}: {resp_body}")
        print(f"[supa_rest_service] HTTP {e.code} {method} {path}: {resp_body}")
        return e.code, None
    except Exception as e:
        _supa_log(f"REST_SERVICE {method} {path} ERR {type(e).__name__}: {e}")
        print(f"[supa_rest_service] err {method} {path}: {type(e).__name__}: {e}")
        return 0, None


def _supa_rest_as_user(request, method: str, path: str, body=None, params=None,
                       prefer: str = None, timeout: int = 15):
    """Faz uma chamada Supabase REST repassando o JWT do usuário.

    Substitui o padrão antigo `Authorization: Bearer SUPABASE_KEY` (anon),
    que causa dois problemas:
      1) RLS bloqueia: anon não casa com `auth.uid()` → rows=[] mesmo com
         usuário autorizado (bug Daniela 2026-05-18 download, Pedro 2026-05-25
         reprocessar);
      2) RLS permissiva: anon vaza dado de outros usuários se a policy de
         SELECT for ampla demais.

    Comportamento:
      - `apikey` SEMPRE é SUPABASE_KEY (PostgREST exige).
      - Authorization usa o Bearer do header do request. Se ausente, cai pra
         SUPABASE_KEY (anon) — preserva compat com projetos `user_id='anonymous'`
         onde o JWT pode não vir.
      - `path` deve começar com `/rest/v1/...` ou só com a tabela (ex:
         `project_clients?job_id=eq.X`); prefixo é adicionado se faltar.
      - `body` pode ser dict ou lista; serializa pra JSON. None = sem body.
      - `params` é dict opcional → query string adicional (concatenado).
      - `prefer` é o header Prefer do PostgREST (ex: `return=representation`,
         `return=minimal`, `resolution=merge-duplicates`).

    Retorna tupla (status_code, parsed_json_or_None). Erros HTTP devolvem
    (code, None) com log. Falha total devolve (0, None).
    """
    import urllib.request, urllib.error, json as _j

    # Monta URL
    if path.startswith("http"):
        url = path
    else:
        if not path.startswith("/"):
            path = "/" + path
        if not path.startswith("/rest/v1/"):
            path = "/rest/v1" + path
        url = f"{SUPABASE_URL}{path}"

    if params:
        from urllib.parse import urlencode
        sep = "&" if "?" in url else "?"
        url = f"{url}{sep}{urlencode(params)}"

    # Extrai Bearer do request (fallback anon)
    user_token = None
    try:
        auth_header = request.headers.get("Authorization", "") or \
                       request.headers.get("authorization", "")
        if auth_header.lower().startswith("bearer "):
            tok = auth_header[7:].strip()
            if tok:
                user_token = tok
    except Exception:
        user_token = None

    bearer = user_token or SUPABASE_KEY

    # Serializa body se houver
    data_bytes = None
    if body is not None:
        data_bytes = _j.dumps(body, default=str).encode("utf-8")

    try:
        req = urllib.request.Request(url, data=data_bytes, method=method.upper())
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {bearer}")
        if data_bytes is not None:
            req.add_header("Content-Type", "application/json")
        if prefer:
            req.add_header("Prefer", prefer)
        resp = urllib.request.urlopen(req, timeout=timeout)
        status = resp.getcode()
        raw = resp.read()
        if not raw:
            return status, None
        try:
            return status, _j.loads(raw.decode("utf-8"))
        except Exception:
            # PATCH/DELETE com Prefer:return=minimal pode devolver texto vazio
            return status, None
    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode("utf-8", errors="replace")[:300]
        except Exception:
            resp_body = "(unreadable)"
        _supa_log(f"REST_AS_USER {method} {path} HTTP {e.code}: {resp_body}")
        print(f"[supa_rest_as_user] HTTP {e.code} {method} {path}: {resp_body}")
        return e.code, None
    except Exception as e:
        _supa_log(f"REST_AS_USER {method} {path} ERR {type(e).__name__}: {e}")
        print(f"[supa_rest_as_user] err {method} {path}: {type(e).__name__}: {e}")
        return 0, None


def _get_project_owner(job_id: str):
    """Retorna o user_id registrado no projeto (ou None se não existe).

    Usa RPC `get_project_owner` (SECURITY DEFINER) pra bypassar RLS.
    Antes (até 2026-05-17) chamava REST /projects?job_id=eq.X com a anon key,
    mas a RLS de SELECT só autoriza 'auth.uid() = user_id' OU admin. Anon ficava
    sem autorização e retornava 0 rows, levando ao falso 404 'Projeto não
    encontrado' em todos os endpoints que chamavam _require_project_owner —
    download, items, cronograma, quotes, agent/ask, projects/{id}/client.
    Bug detectado pela Daniela em 2026-05-18 ao tentar baixar a planilha.
    """
    import urllib.request, urllib.error, json as _j
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/get_project_owner"
        body = _j.dumps({"p_job_id": job_id}).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        # _get_project_owner usa RPC SECURITY DEFINER — bypassa RLS sem precisar
        # de service_role. Mantém anon como Bearer por contrato (a RPC checa
        # internamente; mudar pra service_role aqui é desnecessário e arrisca
        # mascarar bugs de policy).
        req.add_header("Authorization", f"Bearer {SUPABASE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=10)
        result = _j.loads(resp.read().decode("utf-8"))
        # RPC SQL retorna text — pode vir como string direto, null, ou list[str]
        if result is None:
            return None
        if isinstance(result, str):
            return result or "anonymous"
        if isinstance(result, list) and result:
            return (result[0] or "anonymous") if isinstance(result[0], str) else None
        return None
    except urllib.error.HTTPError as e:
        try:
            print(f"[get_project_owner] HTTP {e.code}: {e.read().decode('utf-8', errors='replace')[:200]}")
        except Exception:
            pass
        return None
    except Exception as e:
        print(f"[get_project_owner] erro: {type(e).__name__}: {e}")
        return None


def _require_project_owner(request, job_id: str):
    """Valida que quem chamou é dono do projeto.

    Regra:
    - Projeto com user_id='anonymous' (ou vazio/null): SÓ admin. Ver nota abaixo.
    - Projeto com user_id real: exige JWT válido e user.id == project.user_id.
      JWT inválido → 401; válido mas não dono → 403.

    Admin (ADMIN_EMAIL) tem acesso a qualquer projeto.
    Retorna o user_id do projeto."""
    owner = _get_project_owner(job_id)
    if owner is None:
        raise HTTPException(404, "Projeto não encontrado")

    # 🚨 FECHADO EM 28/07/2026. Antes: `if not owner or owner == "anonymous":
    # return owner` — ou seja, projeto sem dono era LIBERADO sem login nenhum.
    # São 53 projetos do começo do beta (17/04 a 22/05/2026), 9 deles com
    # planilha pronta, contendo prancha e quantitativo de cliente real. O id é
    # UUID (não dá pra adivinhar), mas quem tivesse o link — recebido por
    # WhatsApp, e-mail ou histórico — abria sem senha.
    # Como esses projetos não têm dono, não existe a quem liberar: só admin.
    if not owner or owner == "anonymous":
        user = _get_user_from_request(request)
        if not user or user.get("email", "").lower() != ADMIN_EMAIL:
            raise HTTPException(
                403,
                "Este projeto é de uma versão antiga do AI.arq e não está mais "
                "acessível por link. Fale com a gente em contato@ai.arq.br."
            )
        return owner

    # Projetos de usuário logado: exigem JWT
    user = _get_user_from_request(request)
    if not user:
        raise HTTPException(401, "Autenticação requerida para acessar este projeto")
    # Admin liberado
    if user.get("email", "").lower() == ADMIN_EMAIL:
        return owner
    if user.get("id") != owner:
        raise HTTPException(403, "Acesso restrito ao dono do projeto")
    return owner


_DISCIPLINE_TO_SECTION = {
    "Estrutura":                   "0. Estrutura",
    "Serviços Preliminares":       "1. Serviços Preliminares",
    "Demolição e Remoção":         "2. Demolição e Remoção",
    "Fechamentos Verticais":       "3. Fechamentos Verticais",
    "Revestimentos":               "4. Revestimentos",
    "Pisos e Rodapés":             "5. Pisos e Rodapés",
    "Forros":                      "6. Forros",
    "Iluminação":                  "7. Iluminação",
    "Instalações Elétricas e Dados":"8. Instalações Elétricas e Dados",
    "Instalações Hidráulicas":     "9. Instalações Hidráulicas",
    "Instalações de Gás":          "10. Instalações de Gás",
    "Ar-Condicionado":             "11. Ar-Condicionado",
    "Incêndio e Segurança":        "12. Incêndio e Segurança",
    "Portas e Ferragens":          "13. Portas e Ferragens",
    "Divisórias e Vidros":         "14. Divisórias e Vidros",
    "Persianas e Cortinas":        "15. Persianas e Cortinas",
    "Marcenaria":                  "16. Marcenaria",
    "Mobiliário":                  "17. Mobiliário",
    "Complementares":              "18. Complementares",
}


def _persist_items_to_supabase(job_id: str, items: list) -> int:
    """Insere cada BudgetItem como row em project_items.
    Permite revisão inline no navegador via endpoint /api/items/{job_id}.
    Retorna quantos foram persistidos com sucesso."""
    import urllib.request, urllib.error, json

    import math as _math
    rows = []
    for idx, it in enumerate(items):
        disc = getattr(it, "discipline", "") or "Complementares"
        section = _DISCIPLINE_TO_SECTION.get(disc, f"99. {disc}")
        # Sanitiza quantity: NaN/inf viram 0. json.dumps emitiria `NaN` (JSON
        # inválido) → o PostgREST rejeitava o BATCH INTEIRO (400) → 0 itens salvos
        # mas items_count=len(all_items). Bug do job a2b9316b (14/07, diagrama sem medidas).
        try:
            _q = float(getattr(it, "quantity", 0) or 0)
        except (TypeError, ValueError):
            _q = 0.0
        if not _math.isfinite(_q):
            _q = 0.0
        rows.append({
            "job_id": job_id,
            "item_num": str(getattr(it, "item_num", "") or ""),
            "description": (getattr(it, "description", "") or "")[:500],
            "unit": (getattr(it, "unit", "") or "vb")[:20],
            "quantity": _q,
            "observations": (getattr(it, "observations", "") or "")[:1000],
            "ref_sheet": (getattr(it, "ref_sheet", "") or "")[:200],
            "confidence": str(getattr(getattr(it, "confidence", None), "value", "estimado"))
                          if hasattr(getattr(it, "confidence", None), "value")
                          else str(getattr(it, "confidence", "estimado") or "estimado"),
            "discipline": disc,
            "section": section,
            "sort_order": idx,
        })

    if not rows:
        _supa_log(f"PERSIST items job={job_id} SKIP (empty)")
        return 0

    # SWAP no caminho de SUCESSO: apaga os itens antigos deste job SÓ agora (com os
    # novos itens já prontos) e insere. Assim um reprocesso que FALHA antes daqui
    # (0 itens/erro) NÃO destrói a planilha anterior — crítico pro /add-file (troca
    # PDF→CAD sem risco de perder o estimado). Idempotente pro run normal (1º upload
    # não tem itens; job filho de reprocesso usa job_id novo → delete é no-op).
    try:
        _del = urllib.request.Request(
            f"{SUPABASE_URL}/rest/v1/project_items?job_id=eq.{job_id}", method='DELETE')
        _del.add_header('apikey', SUPABASE_KEY)
        _del.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        urllib.request.urlopen(_del, timeout=20)
    except Exception as _de:
        _supa_log(f"PERSIST pre-delete job={job_id} WARN (segue pro insert): {_de}")

    # Batch insert — REST do Supabase aceita array
    try:
        url = f"{SUPABASE_URL}/rest/v1/project_items"
        body = json.dumps(rows).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        req.add_header('Prefer', 'return=minimal')
        urllib.request.urlopen(req, timeout=30)
        _supa_log(f"PERSIST items job={job_id} OK n={len(rows)}")
        return len(rows)
    except urllib.error.HTTPError as e:
        try:
            resp = e.read().decode('utf-8', errors='replace')[:500]
        except Exception:
            resp = '(unreadable)'
        _supa_log(f"PERSIST items job={job_id} HTTP {e.code}: {resp}")
        print(f"[persist_items] HTTP {e.code}: {resp}")
        return 0
    except Exception as e:
        _supa_log(f"PERSIST items job={job_id} ERR {type(e).__name__}: {e}")
        print(f"[persist_items] err: {e}")
        return 0


def _rpc_update_project_status(job_id: str, data: dict) -> bool:
    """Atualiza uma row de projects via RPC `update_project_status`.
    Retorna True só se a RPC retornou rows_updated >= 1 (detecta silent fail)."""
    import urllib.request, urllib.error, json

    # Converte o dict de colunas → argumentos da RPC
    payload = {
        "p_job_id": job_id,
        "p_status":        data.get("status"),
        "p_items_count":   data.get("items_count"),
        "p_total_area":    data.get("total_area"),
        "p_layout_area":   data.get("layout_area"),
        "p_error_message": data.get("error_message"),
        "p_completed_at":  data.get("completed_at"),
        "p_warnings":      data.get("warnings"),
    }

    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/update_project_status"
        body = json.dumps(payload).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=20)
        resp_body = resp.read().decode('utf-8', errors='replace')
        rows = 0
        try:
            rows = int(resp_body.strip())
        except (ValueError, TypeError):
            # Pode vir como string JSON: "1" ou como array
            try:
                parsed = json.loads(resp_body)
                if isinstance(parsed, int):
                    rows = parsed
                elif isinstance(parsed, list) and parsed:
                    rows = int(parsed[0]) if isinstance(parsed[0], int) else 0
            except Exception:
                pass

        if rows >= 1:
            _supa_log(f"RPC update_project_status job_id={job_id} OK rows={rows}  data={json.dumps(data)[:200]}")
            return True

        # 0 rows: a row não existe (provável erro de job_id ou race condition)
        _supa_log(f"RPC update_project_status job_id={job_id} ZERO_ROWS  data={json.dumps(data)[:200]}")
        print(f"[supabase] RPC update_project_status job={job_id}: row não encontrada")
        return False

    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode('utf-8', errors='replace')[:500]
        except Exception:
            resp_body = '(unreadable)'
        msg = f"RPC update_project_status job_id={job_id} HTTP {e.code}: {resp_body}  data={json.dumps(data)[:200]}"
        print(f"[supabase] RPC HTTP {e.code}: {resp_body}")
        _supa_log(msg)
        return False
    except Exception as e:
        msg = f"RPC update_project_status job_id={job_id} ERR {type(e).__name__}: {e}  data={json.dumps(data)[:200]}"
        print(f"[supabase] RPC err: {e}")
        _supa_log(msg)
        return False


# Campos editáveis via update_project_meta (RPC criada na migration
# add_project_address_phase_and_meta_rpc_v2). Mantém alinhado com a RPC.
_META_FIELDS = {"project_name", "typology", "address", "phase"}


def _rpc_update_project_meta(job_id: str, data: dict) -> bool:
    """Atualiza metadados editáveis (project_name/typology/address/phase) via
    RPC `update_project_meta`. Mesmo padrão da update_project_status.
    Retorna True só se rows_updated >= 1."""
    import urllib.request, urllib.error, json
    payload = {
        "p_job_id":       job_id,
        "p_project_name": data.get("project_name"),
        "p_typology":     data.get("typology"),
        "p_address":      data.get("address"),
        "p_phase":        data.get("phase"),
    }
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/update_project_meta"
        body = json.dumps(payload).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=20)
        resp_body = resp.read().decode('utf-8', errors='replace')
        rows = 0
        try:
            rows = int(resp_body.strip())
        except (ValueError, TypeError):
            try:
                parsed = json.loads(resp_body)
                if isinstance(parsed, int):
                    rows = parsed
                elif isinstance(parsed, list) and parsed:
                    rows = int(parsed[0]) if isinstance(parsed[0], int) else 0
            except Exception:
                pass
        if rows >= 1:
            _supa_log(f"RPC update_project_meta job_id={job_id} OK rows={rows}  data={json.dumps(data)[:200]}")
            return True
        _supa_log(f"RPC update_project_meta job_id={job_id} ZERO_ROWS  data={json.dumps(data)[:200]}")
        return False
    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode('utf-8', errors='replace')[:500]
        except Exception:
            resp_body = '(unreadable)'
        _supa_log(f"RPC update_project_meta job_id={job_id} HTTP {e.code}: {resp_body}  data={json.dumps(data)[:200]}")
        return False
    except Exception as e:
        _supa_log(f"RPC update_project_meta job_id={job_id} ERR {type(e).__name__}: {e}  data={json.dumps(data)[:200]}")
        return False

# ═══════════════════════════════════════════════════════════════
#  Supabase Storage helpers (bucket aiarq-planilhas)
#  Persiste planilhas geradas pra sobreviverem a redeploys do Render.
# ═══════════════════════════════════════════════════════════════

PLANILHAS_BUCKET = "aiarq-planilhas"
PRANCHAS_BUCKET = "aiarq-pranchas"  # PDFs originais pra "abrir prancha" na revisão inline


def _supabase_storage_upload(local_path: str, remote_key: str) -> bool:
    """Faz upload de um arquivo pro Supabase Storage. Sobrescreve se existe."""
    import urllib.request, urllib.error
    try:
        with open(local_path, "rb") as f:
            body = f.read()
        url = f"{SUPABASE_URL}/storage/v1/object/{PLANILHAS_BUCKET}/{remote_key}"
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type",
                       "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        req.add_header("x-upsert", "true")
        urllib.request.urlopen(req, timeout=30)
        _supa_log(f"STORAGE upload {remote_key} OK ({len(body)} bytes)")
        return True
    except urllib.error.HTTPError as e:
        try:
            resp_body = e.read().decode("utf-8", errors="replace")[:300]
        except Exception:
            resp_body = "(unreadable)"
        _supa_log(f"STORAGE upload {remote_key} HTTP {e.code}: {resp_body}")
        print(f"Storage upload {remote_key} HTTP {e.code}: {resp_body}")
        return False
    except Exception as e:
        _supa_log(f"STORAGE upload {remote_key} ERR {type(e).__name__}: {e}")
        print(f"Storage upload error: {e}")
        return False


_PRANCHA_MIME = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".dxf": "application/acad",
    ".dwg": "application/acad",
}


def _sanitize_filename_for_storage(filename: str) -> str:
    """Remove acentos e caracteres especiais que quebram upload pro Supabase Storage.

    Bug detectado: PDF com nome 'Planta 1 - Galpão.pdf' dava HTTP 400 no upload
    mesmo com URL-encoding correto (Galp%C3%A3o). O Storage rejeita certos
    bytes UTF-8 multi-byte. Solução: ASCII-only filename pro storage.
    """
    import unicodedata
    # Decompõe acentos (ã = a + ̃) e remove os combinadores
    nfkd = unicodedata.normalize("NFKD", filename)
    ascii_only = "".join(c for c in nfkd if not unicodedata.combining(c))
    # Mantém só ASCII printáveis seguros pra URL
    safe = "".join(c if (c.isalnum() or c in " ._-()") else "_" for c in ascii_only)
    # Compacta múltiplos underscores
    while "__" in safe:
        safe = safe.replace("__", "_")
    return safe.strip("_ ") or "file"


def _safe_local_filename(filename: str) -> str:
    """Sanitiza filename pra uso em os.path.join local — anti-traversal.

    Aplica _sanitize_filename_for_storage E garante que o resultado:
    - É só basename (sem barras / nem backslashes)
    - Não é '.' nem '..' nem começa com '.'
    - Não é vazio (fallback 'upload')

    Use isso em TODO os.path.join(dir, filename) onde filename vem
    de UploadFile.filename. Sem isso, cliente manda '../../etc/passwd'
    e escreve fora do work_dir.
    """
    import os as _os
    if not filename:
        return "upload"
    # Pega só o último componente (cobre '/foo/bar' e '\\foo\\bar')
    base = _os.path.basename(filename.replace("\\", "/"))
    safe = _sanitize_filename_for_storage(base)
    # Rejeita resultados que ainda seriam traversal ou vazios
    if not safe or safe in (".", "..") or safe.startswith("."):
        return "upload"
    return safe


# Motivo da última falha de upload de prancha, pra sobreviver até o _log_error
# do chamador (ver docstring de _supabase_storage_upload_prancha).
_ULTIMA_FALHA_UPLOAD_PRANCHA = ""


def _supabase_storage_upload_prancha(local_path: str, job_id: str, filename: str) -> bool:
    """Upload de prancha (PDF, PNG, JPG) pro bucket aiarq-pranchas.
    Key: {job_id}/{filename_sanitizado}. Content-Type derivado da extensão.

    Filename é sanitizado pra remover acentos/especiais — Supabase Storage
    rejeita certos UTF-8 multi-byte mesmo URL-encoded.

    🪤 Em falha, grava o MOTIVO e o TAMANHO em _ULTIMA_FALHA_UPLOAD_PRANCHA.
    Antes só devolvia False: o painel dizia "CAD não guardado" sem dizer por quê,
    e o motivo real (413? timeout? tamanho?) morria no log do Render. Caso Ana
    31/07: o DXF que MEDIU não foi guardado e não deu pra saber a causa."""
    import urllib.request, urllib.error
    global _ULTIMA_FALHA_UPLOAD_PRANCHA
    _tam_mb = 0.0
    try:
        _tam_mb = os.path.getsize(local_path) / 1048576.0
    except Exception:
        pass
    try:
        with open(local_path, "rb") as f:
            body = f.read()
        import urllib.parse as _up
        safe_name = _sanitize_filename_for_storage(filename)
        remote_key = f"{job_id}/{_up.quote(safe_name)}"
        url = f"{SUPABASE_URL}/storage/v1/object/{PRANCHAS_BUCKET}/{remote_key}"
        ext = os.path.splitext(safe_name.lower())[1]
        mime = _PRANCHA_MIME.get(ext, "application/octet-stream")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", mime)
        req.add_header("x-upsert", "true")
        # Timeout proporcional ao tamanho: 60s dava conta de PDF, mas DXF de
        # planta grande passa de 100 MB e estourava calado. ~2s por MB, teto 5min.
        _tmo = max(60, min(300, int(_tam_mb * 2)))
        urllib.request.urlopen(req, timeout=_tmo)
        if safe_name != filename:
            _supa_log(f"STORAGE upload prancha OK (saneado: '{filename}' → '{safe_name}')")
        _ULTIMA_FALHA_UPLOAD_PRANCHA = ""
        return True
    except Exception as e:
        _detalhe = ""
        if isinstance(e, urllib.error.HTTPError):
            try:
                _corpo = e.read().decode("utf-8", "replace")[:200]
            except Exception:
                _corpo = ""
            _detalhe = f"HTTP {e.code} {_corpo}".strip()
        else:
            _detalhe = f"{type(e).__name__}: {e}"
        _ULTIMA_FALHA_UPLOAD_PRANCHA = f"{_detalhe} — arquivo de {_tam_mb:.1f} MB"
        _supa_log(f"STORAGE upload prancha {filename} ERR {_ULTIMA_FALHA_UPLOAD_PRANCHA}")
        print(f"[storage pranchas] upload {filename} error: {_ULTIMA_FALHA_UPLOAD_PRANCHA}")
        return False


# Alias de compatibilidade com código antigo
_supabase_storage_upload_pdf = _supabase_storage_upload_prancha


def _supabase_storage_download_prancha(job_id: str, filename: str) -> Optional[bytes]:
    """Baixa prancha (PDF/PNG/etc) do Storage. Retorna bytes ou None."""
    import urllib.request, urllib.parse as _up
    try:
        remote_key = f"{job_id}/{_up.quote(filename)}"
        url = f"{SUPABASE_URL}/storage/v1/object/{PRANCHAS_BUCKET}/{remote_key}"
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = urllib.request.urlopen(req, timeout=30)
        return resp.read()
    except Exception as e:
        print(f"[storage pranchas] download {filename}: {e}")
        return None


# Alias de compatibilidade
_supabase_storage_download_pdf = _supabase_storage_download_prancha


# ─── CHECKPOINT por prancha (19/07) ───────────────────────────────────
# A retomada pós-restart refazia o job INTEIRO do zero — num projeto de 22
# pranchas isso paga a IA de novo e reabre a janela pra outra queda. Agora o
# resultado da análise de cada prancha vira um JSON em
# {job_id}/checkpoint/{stem}.json no bucket de pranchas; na retomada, prancha
# com checkpoint NÃO chama a IA de novo. O filtro do _retomar (.pdf/.dwg/.dxf)
# ignora esses .json. Tudo best-effort: sem checkpoint = comportamento antigo.

def _ckpt_save(job_id: str, stem: str, result: dict) -> None:
    import urllib.request, urllib.parse as _up, json as _j
    try:
        stem = _sanitize_filename_for_storage(stem)  # acentos quebram o Storage
        body = _j.dumps(result, ensure_ascii=False, default=str).encode("utf-8")
        key = f"{job_id}/checkpoint/{_up.quote(stem)}.json"
        url = f"{SUPABASE_URL}/storage/v1/object/{PRANCHAS_BUCKET}/{key}"
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        req.add_header("x-upsert", "true")
        urllib.request.urlopen(req, timeout=30)
    except Exception as e:
        print(f"[ckpt] save {stem} falhou (segue sem): {e}")


def _ckpt_load_all(job_id: str) -> dict:
    """{stem: result} de todos os checkpoints do job. Vazio se não há."""
    import urllib.request, urllib.parse as _up, json as _j
    out = {}
    try:
        list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        body = _j.dumps({"prefix": f"{job_id}/checkpoint/", "limit": 300}).encode("utf-8")
        req = urllib.request.Request(list_url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        objs = _j.loads(urllib.request.urlopen(req, timeout=15).read().decode("utf-8"))
        for o in objs or []:
            name = o.get("name", "")
            if not name.endswith(".json"):
                continue
            try:
                url = (f"{SUPABASE_URL}/storage/v1/object/{PRANCHAS_BUCKET}/"
                       f"{job_id}/checkpoint/{_up.quote(name)}")
                r = urllib.request.Request(url, method="GET")
                r.add_header("apikey", SUPABASE_KEY)
                r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                data = _j.loads(urllib.request.urlopen(r, timeout=20).read().decode("utf-8"))
                out[_up.unquote(name[:-5])] = data
            except Exception:
                continue
    except Exception as e:
        print(f"[ckpt] load_all {job_id}: {e}")
    return out


def _ckpt_limpar(job_id: str) -> None:
    """Apaga os checkpoints do job (chamado no done — não precisam mais)."""
    import urllib.request, urllib.parse as _up, json as _j
    try:
        list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        body = _j.dumps({"prefix": f"{job_id}/checkpoint/", "limit": 300}).encode("utf-8")
        req = urllib.request.Request(list_url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        objs = _j.loads(urllib.request.urlopen(req, timeout=15).read().decode("utf-8"))
        for o in objs or []:
            name = o.get("name", "")
            if not name:
                continue
            try:
                url = (f"{SUPABASE_URL}/storage/v1/object/{PRANCHAS_BUCKET}/"
                       f"{job_id}/checkpoint/{_up.quote(name)}")
                r = urllib.request.Request(url, method="DELETE")
                r.add_header("apikey", SUPABASE_KEY)
                r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                urllib.request.urlopen(r, timeout=15)
            except Exception:
                continue
    except Exception:
        pass


def _supabase_storage_download(remote_key: str, local_path: str) -> bool:
    """Baixa arquivo do Supabase Storage pra path local. Cria diretório se preciso."""
    import urllib.request
    try:
        url = f"{SUPABASE_URL}/storage/v1/object/{PLANILHAS_BUCKET}/{remote_key}"
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = urllib.request.urlopen(req, timeout=30)
        os.makedirs(os.path.dirname(local_path), exist_ok=True)
        with open(local_path, "wb") as f:
            f.write(resp.read())
        return True
    except Exception as e:
        print(f"Storage download error ({remote_key}): {e}")
        return False


def get_planilha_path(job_id: str) -> Optional[str]:
    """Retorna o path local da planilha de um job. Se sumiu (Render
    /tmp volátil), tenta baixar do Supabase Storage. None se falhou."""
    local = os.path.join(WORK_DIR, job_id, f"orcamento_{job_id}.xlsx")
    if os.path.exists(local):
        return local
    if _supabase_storage_download(f"{job_id}.xlsx", local):
        return local
    return None


# Swagger/OpenAPI só quando LIGADO explicitamente (env EXPOSE_API_DOCS=1). Em
# produção fica OFF: /docs, /redoc e /openapi.json eram públicos e expunham toda
# a superfície da API pra reconhecimento (achado auditoria 27/07). Os endpoints
# seguem protegidos por auth — isto só tira o mapa de bandeja do atacante.
_EXPOSE_DOCS = os.environ.get("EXPOSE_API_DOCS", "0") == "1"
app = FastAPI(
    title="AI.arq API",
    description="Motor de processamento de pranchas de arquitetura com IA",
    version="1.0.0",
    docs_url="/docs" if _EXPOSE_DOCS else None,
    redoc_url="/redoc" if _EXPOSE_DOCS else None,
    openapi_url="/openapi.json" if _EXPOSE_DOCS else None,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    # Segurança (16/07): a API se autentica por Bearer token (Authorization header),
    # NÃO por cookie. Com credentials=False, o '*' não reflete origem pra requisição
    # credenciada — fecha o combo inseguro '*'+credentials sem quebrar nada (Bearer
    # segue via allow_headers). Nenhum fluxo usa cookie cross-site.
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["X-Filename", "Content-Disposition"],
)

# ── Instagram Agent (desativado por padrão, ativar manualmente via /api/instagram/toggle) ──
app.include_router(instagram_router)
app.include_router(whatsapp_router)  # WhatsApp Cloud API (webhook); envio dormente até setar env

# Armazenamento de jobs em arquivo JSON (sobrevive a restarts)
import json as _json
WORK_DIR = os.path.join(tempfile.gettempdir(), "aiarq_jobs")
os.makedirs(WORK_DIR, exist_ok=True)
JOBS_FILE = os.path.join(WORK_DIR, "_jobs.json")

# Trava de escrita do store de jobs. Recovery, request e processamento escrevem
# o MESMO JSON. Sem ela, dois read-modify-write concorrentes se atropelam
# (lost-update): um 'error' pode ser sobrescrito de volta pra 'queued'/'processing'
# por uma thread com snapshot velho — a raiz do job órfão. RLock (reentrante)
# porque _load_jobs/_save_jobs são chamados de dentro de seções já travadas.
_JOBS_LOCK = threading.RLock()

def _load_jobs() -> dict:
    with _JOBS_LOCK:
        try:
            if os.path.exists(JOBS_FILE):
                with open(JOBS_FILE, 'r') as f:
                    return _json.load(f)
        except Exception as _e:
            # NÃO fica calado: _jobs.json corrompido → retornar {} apagaria TODOS
            # os jobs de uma vez (o cenário catastrófico que a escrita atômica
            # minimiza). Deixa rastro no log do Render pra não sumir em silêncio.
            print(f"[jobs] FALHA ao ler {JOBS_FILE}: {type(_e).__name__}: {_e}")
        return {}

def _save_jobs(jobs_dict):
    # Escrita atômica: grava num tmp e faz rename. Se o processo morrer no meio,
    # o _jobs.json antigo continua íntegro em vez de virar arquivo truncado —
    # que o _load_jobs engoliria como {} e apagaria TODOS os jobs de uma vez.
    with _JOBS_LOCK:
        try:
            tmp = f"{JOBS_FILE}.{os.getpid()}.tmp"
            with open(tmp, 'w') as f:
                _json.dump(jobs_dict, f)
            os.replace(tmp, JOBS_FILE)
        except Exception as _e:
            print(f"[jobs] FALHA ao gravar {JOBS_FILE}: {type(_e).__name__}: {_e}")

# ─── Email transacional (SMTP — Google Workspace) ───────────────────
# Configurado por env vars SMTP_* no Render. Se não estiver configurado,
# vira no-op silencioso (loga e segue) — NUNCA derruba o fluxo que chamou.
def _send_email_smtp(to_email: str, subject: str, html_body: str, text_body: str = "", log_kind: str = "email") -> bool:
    host = os.getenv("SMTP_HOST", "")
    user = os.getenv("SMTP_USER", "")
    password = os.getenv("SMTP_PASSWORD", "")
    if not (host and user and password and to_email):
        print(f"[email] SMTP não configurado ou sem destino — pulando '{subject}'")
        return False
    try:
        port = int(os.getenv("SMTP_PORT", "587") or "587")
    except ValueError:
        port = 587
    from_name = os.getenv("SMTP_FROM_NAME", "AI.arq")
    from_email = os.getenv("SMTP_FROM", user)
    try:
        import smtplib
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText
        from email.utils import formataddr
        msg = MIMEMultipart("alternative")
        msg["Subject"] = subject
        msg["From"] = formataddr((from_name, from_email))
        msg["To"] = to_email
        msg["Reply-To"] = from_email
        if text_body:
            msg.attach(MIMEText(text_body, "plain", "utf-8"))
        msg.attach(MIMEText(html_body, "html", "utf-8"))
        with smtplib.SMTP(host, port, timeout=20) as server:
            server.ehlo()
            server.starttls()
            server.login(user, password)
            server.sendmail(from_email, [to_email], msg.as_string())
        print(f"[email] OK -> {to_email}: {subject}")
        # Registra TODO email enviado pro usuario (medir volume/pessoa, nao poluir).
        # Pula contas internas (Pedro/aliases) e o alerta interno (NOTIFY_EMAIL).
        try:
            _to = (to_email or "").lower()
            if not _email_eh_interno(_to) and _to != (NOTIFY_EMAIL or "").lower():
                _supabase_insert("email_sent_log", {
                    "email": to_email,
                    "kind": (log_kind or "email"),
                    "subject": (subject or "")[:200],
                })
        except Exception as _e:
            print(f"[email] log de envio falhou (nao critico): {_e}")
        return True
    except Exception as e:
        print(f"[email] FALHA -> {to_email}: {type(e).__name__}: {e}")
        return False


# Alertas internos pro Pedro (novo cliente / novo projeto) — vão pro gmail
# pessoal dele. Configurável por env; default é o pessoal informado.
NOTIFY_EMAIL = os.getenv("NOTIFY_EMAIL", "pedro.zellmer@gmail.com")


def _notify_admin(subject: str, body_html: str) -> bool:
    """Alerta interno simples pro Pedro. Best-effort (não derruba nada)."""
    return _send_email_smtp(
        NOTIFY_EMAIL, f"[AI.arq] {subject}",
        '<div style="font-family:Arial,sans-serif;font-size:15px;color:#0F172A;line-height:1.6;">'
        + body_html +
        '<div style="margin-top:16px;color:#94a3b8;font-size:12px;">Alerta automático do AI.arq '
        '&middot; ai.arq.br</div></div>')


def _saudacao() -> str:
    """bom dia / boa tarde / boa noite no horário de Brasília (UTC-3)."""
    h = (datetime.utcnow().hour - 3) % 24
    if 5 <= h < 12:
        return "bom dia"
    if 12 <= h < 18:
        return "boa tarde"
    return "boa noite"


def _first_name(full: str) -> str:
    full = (full or "").strip()
    return full.split(" ")[0] if full else ""


def _greeting_line(full_name: str) -> str:
    """'Fulano, boa noite,' — com nome; ou 'Boa noite,' sem nome."""
    fn = _first_name(full_name)
    s = _saudacao()
    return f"{fn}, {s}," if fn else f"{s.capitalize()},"


# Número dedicado do AI.arq no WhatsApp Business (app, atendido pelo Pedro).
# 🪤 NÃO é o pessoal dele. Se mudar, mudar também em `whatsapp-button.js`
# (const NUMERO), que é o botão flutuante do site.
# O cliente é sempre quem INICIA a conversa — foi o contrário disso (cold-DM do
# número pessoal) que restringiu a conta dele em 23/07.
_WHATSAPP_NUM = "551151968034"
_WHATSAPP_LINK = (f"https://wa.me/{_WHATSAPP_NUM}"
                  "?text=Ol%C3%A1%21%20Vim%20pelo%20e-mail%20do%20AI.arq.")


def _email_img(arquivo: str, alt: str, margem: str = "14px 0 4px") -> str:
    """Imagem dos e-mails (assets/email/ no site). Largura fixa + alt sempre
    (Gmail bloqueia imagem por padrão — o alt segura a mensagem)."""
    return (f'<img src="https://ai.arq.br/assets/email/{arquivo}" width="460" alt="{alt}" '
            f'style="width:100%;max-width:460px;height:auto;border-radius:12px;'
            f'display:block;margin:{margem};border:0;">')


def _email_wrap(title: str, body_html: str, cta_text: str = "", cta_url: str = "", badge: str = "",
                badge_color: str = "green",
                reason: str = "Você está recebendo este e-mail porque tem uma conta no AI.arq.",
                signoff: bool = True, preheader: str = "") -> str:
    """Layout moderno e acessível dos emails (table-based + estilo inline, do
    jeito que Gmail/Outlook exigem). Logo = ícone hospedado em ai.arq.br.

    Assinatura e rodapé ficam DENTRO do card (contíguos ao conteúdo) de
    propósito: quando o rodapé é um bloco solto no fim, o Gmail o reconhece
    como boilerplate repetido e colapsa atrás do "•••". Integrado, não colapsa.
    """
    cta = ""
    if cta_text and cta_url:
        cta = ('<tr><td style="padding:22px 30px 8px;">'
               f'<a href="{cta_url}" style="background:#4F46E5;color:#ffffff;text-decoration:none;'
               'padding:14px 26px;border-radius:10px;font-size:15px;font-weight:600;'
               'font-family:Arial,sans-serif;display:inline-block;">'
               f'{cta_text} &rarr;</a></td></tr>')
    badge_html = ""
    if badge:
        _bbg, _bfg = ("#fef3c7", "#b45309") if badge_color == "amber" else ("#dcfce7", "#15803d")
        badge_html = ('<tr><td style="padding:18px 30px 0;"><span style="display:inline-block;'
                      f'background:{_bbg};color:{_bfg};font-size:12px;font-weight:700;'
                      'font-family:Arial,sans-serif;padding:4px 10px;border-radius:20px;">'
                      f'{badge}</span></td></tr>')
    # Assinatura pessoal (dentro do card)
    sig_html = ""
    if signoff:
        sig_html = ('<tr><td style="padding:22px 30px 0;">'
                    '<div style="border-top:1px solid #eef2f7;padding-top:18px;font-size:15px;'
                    'line-height:1.55;color:#475569;font-family:Arial,sans-serif;">'
                    'Um abraço,<br>'
                    '<span style="font-weight:700;color:#0F172A;">Pedro</span><br>'
                    '<span style="color:#94a3b8;font-size:13px;">AI.arq &middot; ai.arq.br</span>'
                    f'<div style="margin-top:12px;"><a href="{_WHATSAPP_LINK}" '
                    'style="display:inline-block;border:1px solid #25D366;color:#128C4A;'
                    'text-decoration:none;padding:8px 14px;border-radius:8px;font-size:13px;'
                    'font-weight:600;font-family:Arial,sans-serif;">'
                    '&#128172; Falar no WhatsApp</a></div>'
                    '</div></td></tr>')
    # Preheader: linha de preview que o Gmail/Apple Mail mostram ao lado do
    # assunto. Invisível no corpo; o &zwnj; empurra pra fora o texto real do
    # e-mail que apareceria no preview. (Upgrade 02/08 — vale pra todos.)
    pre_html = ""
    if preheader:
        pre_html = ('<div style="display:none;max-height:0;overflow:hidden;mso-hide:all;">'
                    f'{preheader}' + ('&zwnj;&nbsp;' * 30) + '</div>')
    return (
        f'{pre_html}'
        '<div style="background:#eaeef3;padding:28px 14px;font-family:Arial,sans-serif;">'
        '<table role="presentation" width="100%" cellpadding="0" cellspacing="0" '
        'style="max-width:520px;margin:0 auto;"><tr><td>'
        '<table role="presentation" width="100%" cellpadding="0" cellspacing="0" '
        'style="background:#ffffff;border-radius:16px;overflow:hidden;border:1px solid #e2e8f0;">'
        '<tr><td style="height:5px;background:#4F46E5;background:linear-gradient(90deg,#4F46E5,#22D3EE);'
        'font-size:0;line-height:0;">&nbsp;</td></tr>'
        '<tr><td style="padding:26px 30px 4px;"><table role="presentation" cellpadding="0" cellspacing="0"><tr>'
        '<td><img src="https://ai.arq.br/email-logo.png" width="48" height="48" alt="AI.arq" '
        'style="width:48px;height:48px;border-radius:12px;display:block;"></td>'
        '<td style="padding-left:12px;font-size:21px;font-weight:700;color:#0F172A;'
        'letter-spacing:-.3px;font-family:Arial,sans-serif;">AI.arq</td></tr></table></td></tr>'
        f'{badge_html}'
        f'<tr><td style="padding:16px 30px 0;font-size:22px;font-weight:700;color:#0F172A;'
        f'line-height:1.3;font-family:Arial,sans-serif;">{title}</td></tr>'
        f'<tr><td style="padding:12px 30px 4px;font-size:15px;line-height:1.65;color:#475569;'
        f'font-family:Arial,sans-serif;">{body_html}</td></tr>'
        f'{cta}'
        f'{sig_html}'
        # Rodapé LGPD DENTRO do card (contíguo -> não vira "•••")
        '<tr><td style="padding:18px 30px 24px;">'
        '<div style="border-top:1px solid #eef2f7;padding-top:14px;font-size:11px;color:#aab4c0;'
        'line-height:1.6;font-family:Arial,sans-serif;">'
        f'{reason}<br>'
        '<a href="https://ai.arq.br/privacidade.html" style="color:#8b93f6;text-decoration:none;">Política de Privacidade</a>'
        ' &middot; <a href="https://ai.arq.br" style="color:#8b93f6;text-decoration:none;">ai.arq.br</a>'
        ' &middot; Para remover seus dados, é só responder este e-mail.'
        '</div></td></tr>'
        '</table>'
        '</td></tr></table></div>'
    )


# Dedup em memória (vida do processo) pra não mandar 2x o email de falha do
# MESMO job — ex.: except do process_job + recovery no boot. Em produção o
# disco é efêmero; isto cobre o processo atual, e o parent_job_id cobre os
# reprocessamentos manuais (que criam job novo com pai).
_falha_emailed = set()


def _build_falha_email(name: str, project_name: str, reprocessavel: bool, error_hint: str = ""):
    """Monta (subject, html) do email de falha. Separado pra reuso no preview.
    error_hint = mensagem do erro; usada pra dar orientação ESPECÍFICA quando
    reprocessavel=False (DWG não abre vs arquivo grande vs sem cotas)."""
    import html as _hf
    pn = _hf.escape(project_name or "seu projeto")
    _pn_raw = (project_name or "").strip()
    greet = _greeting_line(_hf.escape(name or ""))
    if reprocessavel:
        body = (f"{greet}<br><br>"
                f"Tivemos um problema ao processar o projeto <b>{pn}</b> e ele não "
                f"foi concluído. Quase sempre é coisa passageira e <b>reprocessar "
                f"resolve</b> — e o reprocessamento é grátis."
                + _email_img("falha-retry.png", "Reprocessar quase sempre resolve — é um clique, e é grátis")
                + f"Se continuar dando problema, é só responder este e-mail que a gente "
                f"te ajuda pessoalmente. 🙂")
        subject = (f"{_pn_raw} — tivemos um problema no AI.arq"
                   if _pn_raw else "Tivemos um problema com seu projeto no AI.arq")
        html = _email_wrap("Não conseguimos concluir seu projeto", body,
                           "Reprocessar no painel", "https://ai.arq.br/dashboard.html",
                           badge="⚠ Precisa reprocessar", badge_color="amber",
                           reason="Você está recebendo este e-mail porque enviou um projeto ao AI.arq.")
    else:
        # Orientação ESPECÍFICA por tipo de problema de arquivo (mesmo diagnóstico do
        # erro técnico). Reprocessar o MESMO arquivo não resolve em nenhum destes.
        _eh = (error_hint or "").lower()
        if "dwg" in _eh and ("abrir" in _eh or "convert" in _eh):
            motivo = ("não conseguimos <b>abrir o seu DWG</b> automaticamente — costuma "
                      "acontecer com arquivo salvo numa versão muito recente do AutoCAD, ou "
                      "com objetos especiais (comum em incêndio, hidráulica e elétrica feitos "
                      "em software MEP).")
            fix = ("O ideal é <b>reenviar em DXF ou PDF vetorial</b>, ou salvar o DWG numa "
                   "<b>versão mais antiga</b> do AutoCAD (ex.: 2013) e mandar de novo.")
        elif "grande demais" in _eh or ("limite" in _eh and "mb" in _eh):
            motivo = "o arquivo ficou <b>grande demais</b> pra processar com segurança."
            fix = ("Exporte <b>só a prancha necessária</b> (não o projeto inteiro), ou "
                   "<b>divida o arquivo em partes</b> e reenvie.")
        else:
            motivo = ("não conseguimos ler as quantidades nesse arquivo. Quase sempre é "
                      "porque o PDF é uma imagem escaneada/fotografada, ou a prancha tem só "
                      "o desenho, sem cotas e quadros de áreas.")
            fix = ("O ideal é <b>reenviar a planta completa exportada direto do CAD</b> "
                   "(PDF vetorial, DWG ou DXF).")
        body = (f"{greet}<br><br>"
                f"Recebemos o projeto <b>{pn}</b>, mas {motivo} Ou seja, <b>reprocessar o "
                f"mesmo arquivo não vai resolver</b>."
                + _email_img("falha-arquivo.png", "Um ajuste no arquivo resolve — exporte em DXF")
                + f"{fix}<br><br>"
                f"Se quiser, responda este e-mail com o arquivo que a gente te ajuda a "
                f"preparar. 🙂")
        subject = (f"{_pn_raw} — precisamos de outro arquivo (AI.arq)"
                   if _pn_raw else "Sobre o seu projeto no AI.arq — precisamos de outro arquivo")
        html = _email_wrap("Precisamos de outro arquivo pra continuar", body,
                           "Enviar outra prancha", "https://ai.arq.br/dashboard.html",
                           badge="⚠ Revisar o arquivo", badge_color="amber",
                           reason="Você está recebendo este e-mail porque enviou um projeto ao AI.arq.")
    return subject, html


def _email_falha_cliente(job_id: str, reprocessavel: bool = True) -> bool:
    """Avisa o cliente que o projeto falhou. Best-effort, NUNCA levanta.

    - reprocessavel=True  -> falha passageira (IA sobrecarregada / reinício do
      servidor / erro técnico): "reprocessar resolve, é grátis".
    - reprocessavel=False -> arquivo não-quantificável (PDF escaneado, prancha
      só de layout, vetorial sem texto): reprocessar o MESMO arquivo NÃO
      resolve -> orienta a reenviar a planta completa exportada do CAD.

    Dedup: pula se já avisou este job_id neste processo, ou se o job é fruto de
    reprocessamento (parent_job_id setado -> o job-pai já avisou)."""
    try:
        if job_id in _falha_emailed:
            return False
        import html as _hf, urllib.request as _urf
        _qf = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
               f"&select=user_email,user_name,project_name,parent_job_id,reprocess_count,created_at,error_message")
        _rf = _urf.Request(_qf, method="GET")
        _rf.add_header("apikey", SUPABASE_KEY)
        _rf.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _rows = _json.loads(_urf.urlopen(_rf, timeout=10).read().decode("utf-8"))
        if not _rows:
            return False
        _email = _rows[0].get("user_email") or ""
        if not _email:
            return False
        if _rows[0].get("parent_job_id"):
            _falha_emailed.add(job_id)  # filho de reprocessamento: pai já avisou
            return False
        if (_rows[0].get("reprocess_count") or 0) > 0:
            # reprocesso (usuário clicou reprocessar e acompanha, OU revisão interna):
            # não re-emailar falha — só o 1º processamento notifica o cliente.
            _falha_emailed.add(job_id)
            return False
        # Freio anti-spam PERSISTENTE (sobrevive a restart — o dedup em memoria
        # furava quando um deploy reiniciava o processo e o cliente levava varios
        # "falhou" no mesmo incidente, caso Luciano). Como o semaforo processa 1
        # por vez na ORDEM, basta avisar a falha MAIS ANTIGA do usuario na janela:
        # se ja existe um projeto DESTE usuario que falhou ANTES deste (criado
        # antes) nos ultimos 15 min, aquele ja avisou -> nao manda de novo.
        try:
            from urllib.parse import quote as _quote
            from datetime import datetime as _dt, timedelta as _td, timezone as _tz
            _self_created = _rows[0].get("created_at") or ""
            _since = (_dt.now(_tz.utc) - _td(minutes=15)).isoformat()
            if _self_created:
                _tq = (f"{SUPABASE_URL}/rest/v1/projects"
                       f"?user_email=eq.{_quote(_email, safe='')}"
                       f"&status=eq.error"
                       f"&created_at=lt.{_quote(_self_created, safe='')}"
                       f"&created_at=gte.{_quote(_since, safe='')}"
                       f"&job_id=neq.{job_id}&select=job_id&limit=1")
                _tr = _urf.Request(_tq, method="GET")
                _tr.add_header("apikey", SUPABASE_KEY)
                _tr.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                if _json.loads(_urf.urlopen(_tr, timeout=8).read().decode("utf-8")):
                    _falha_emailed.add(job_id)
                    print(f"[email] falha-cliente FREADO ({_email}: ja avisado nesta janela de 15min)")
                    return False
        except Exception as _te:
            print(f"[email] checagem de freio falhou (segue e manda): {_te}")
        _nm = _resolve_client_name(_email, hint=_rows[0].get("user_name") or "")
        _subject, _html = _build_falha_email(
            _nm,
            _rows[0].get("project_name") or "seu projeto",
            reprocessavel,
            error_hint=(_rows[0].get("error_message") or ""))
        ok = _send_email_smtp(_email, _subject, _html,
                              log_kind="erro_reprocessar" if reprocessavel else "erro_trocar")
        _falha_emailed.add(job_id)
        return ok
    except Exception as _e:
        print(f"[email] falha-cliente nao enviado (nao-fatal): {_e}")
        return False


def _build_reading_diagnostic(all_items, n_pdf, n_cad, project_type, project_data):
    """Diagnóstico de leitura PERSONALIZADO pro email da planilha pronta: explica
    COMO a IA leu o projeto e POR QUE a planilha ficou daquele jeito, a partir dos
    sinais que o motor já calcula (tipo de arquivo, medido vs estimado, warnings).
    Honesto — nunca promete o que não entrega. Cor sempre com ícone+texto (Pedro é
    daltônico): ✓ medido / ⚠ estimado. Best-effort, nunca levanta."""
    import html as _hd
    try:
        from models import Confidence as _Conf
    except Exception:
        _Conf = None
    try:
        total = len(all_items or [])
        if not total:
            return ""

        def _is_medido(it):
            c = getattr(it, "confidence", None)
            ok_conf = (c == _Conf.CONFIRMADO) if _Conf else str(c).endswith("confirmado")
            return ok_conf and getattr(it, "origem", "") != "vision_pdf"

        medidos = sum(1 for it in all_items if _is_medido(it))
        estimados = total - medidos
        is_estrut = (project_type or "").strip().lower() == "estrutura"

        placar = (f"<b>&#10003; {medidos} medido(s)</b> direto do CAD (em branco na planilha) "
                  f"e <b>&#9888; {estimados} pra você confirmar</b> (em laranja).")

        if n_cad == 0 and n_pdf > 0:
            porque = ("Seu projeto veio em <b>PDF</b>, então a IA leu pela imagem — por isso "
                      "saiu tudo como <b>estimativa pra conferir</b> (a gente nunca marca como "
                      "medido um número que não veio da geometria). Pra medição exata, mande o "
                      "mesmo projeto em <b>DWG ou DXF</b> que a gente mede de verdade.")
        elif n_cad > 0 and medidos > 0:
            porque = ("Seu arquivo veio em <b>CAD (DWG/DXF)</b>, então medimos boa parte direto "
                      "da geometria do desenho (os itens em branco). Os em laranja dependem da "
                      "sua conferência.")
        elif n_cad > 0:
            porque = ("Seu arquivo é CAD, mas a IA não conseguiu medir a geometria diretamente "
                      "(comum quando os elementos foram desenhados como linhas soltas, não como "
                      "blocos) — identificamos os itens, mas a quantidade ficou pra você confirmar.")
        else:
            porque = ("Os itens em branco foram medidos do desenho; os em laranja são pra você "
                      "confirmar a quantidade.")

        if is_estrut:
            porque += (" Como é projeto <b>estrutural</b>, lemos concreto em m&sup3;, fôrma em "
                       "m&sup2; e aço em kg — o peso de aço sai medido quando a prancha tem um "
                       "<b>quadro/resumo de aço</b>.")

        extra = ""
        for w in (getattr(project_data, "warnings", None) or [])[:2]:
            ws = str(w).strip()
            if ws:
                extra += f"<br>&bull; {_hd.escape(ws)}"

        return (f"<div style='margin-top:14px;padding:12px 14px;background:#F8FAFC;"
                f"border-left:3px solid #4F46E5;border-radius:6px;font-size:14px;line-height:1.55;'>"
                f"<b>Como lemos o seu projeto</b><br>{placar}<br>{porque}{extra}</div>")
    except Exception as _de:
        print(f"[email] diagnostico de leitura nao montado (nao-fatal): {_de}")
        return ""


def _next_steps_html(job_id: str, n_medido: int = 0, n_total: int = 0,
                     veio_pdf: bool = False, tem_cronograma: bool = False) -> str:
    """Bloco 'o que fazer agora' PERSONALIZADO pro email de planilha pronta — os
    caminhos dependem do resultado real do projeto, pra não empurrar passo que não
    faz sentido. Tudo mora na página do projeto (o botão principal do email leva lá).

    - veio de PDF e mediu 0 → 1º passo é 'complemente com o CAD' (caso Diana 08/07).
    - mediu bem → pula o CAD e lidera com revisão (citando quantos ficaram em laranja).
    - chat sempre; cronograma vira 'ver' se já existe, 'montar' se não."""
    n_est = max(0, n_total - n_medido)
    passos = []
    # 1) Complementar com CAD — só quando veio de PDF e NADA foi medido
    if veio_pdf and n_medido == 0 and n_total > 0:
        passos.append(("&#128208;", "Me&ccedil;a de verdade: complemente com o CAD",
            f"Sua prancha veio em <b>PDF</b>, ent&atilde;o os {n_total} itens sa&iacute;ram como "
            f"<b>estimativa</b>. Suba o <b>DWG ou DXF</b> da mesma prancha no pr&oacute;prio projeto "
            f"que a gente <b>refaz medindo de verdade</b> &mdash; de gra&ccedil;a."))
    # 2) Revisar — sempre; se há estimativa, cita o número
    if n_est > 0:
        _it = "item" if n_est == 1 else "itens"
        passos.append(("&#128221;", f"Revise os {n_est} {_it} em laranja",
            "Esses sa&iacute;ram como <b>estimativa</b> pra voc&ecirc; conferir. Ajuste o que precisar &mdash; "
            "e ao subir a planilha revisada voc&ecirc; <b>afina o motor</b>: os pr&oacute;ximos projetos saem medindo melhor."))
    else:
        passos.append(("&#128221;", "Revise e ajuste",
            "Confira os itens e ajuste o que precisar. Ao subir a planilha revisada voc&ecirc; <b>afina o motor</b> pros pr&oacute;ximos projetos."))
    # 3) Memorial descritivo (novo entregável 01/08) — logo depois da revisão
    #    de propósito: "revise antes, o texto sai melhor" empurra a revisão.
    passos.append(("&#128196;", "Baixe o memorial descritivo (rascunho em Word)",
        "Escrito a partir dos itens do seu CAD, na ordem das etapas da obra &mdash; o documento "
        "que prefeitura, banco e incorpora&ccedil;&atilde;o pedem. Estimativas rotuladas e campos "
        "<b>[A PREENCHER]</b> pro que o desenho n&atilde;o informa. &Eacute; base pra voc&ecirc; editar "
        "e assinar &mdash; e sai melhor se voc&ecirc; revisar as quantidades antes."))
    # 4) Chat — sempre útil
    passos.append(("&#128172;", "Tire d&uacute;vidas no chat",
        "Pergunte sobre o seu quantitativo direto na p&aacute;gina do projeto: o que &eacute; medido, "
        "o que revisar, onde est&aacute; cada item."))
    # 5) Cronograma — 'ver' se já existe, senão 'montar'
    if tem_cronograma:
        passos.append(("&#128197;", "Veja o cronograma",
            "Seu cronograma f&iacute;sico j&aacute; est&aacute; gerado &mdash; abra pra acompanhar a obra."))
    else:
        passos.append(("&#128197;", "Monte o cronograma",
            "Gere o cronograma f&iacute;sico da obra a partir do quantitativo &mdash; as etapas "
            "distribu&iacute;das no tempo, com o peso de cada uma. Tamb&eacute;m &eacute; de gra&ccedil;a."))
    cards = ""
    for emoji, titulo, desc in passos:
        cards += (f'<div style="margin:10px 0;padding:12px 14px;border:1px solid #e5e7eb;border-radius:10px;background:#ffffff;">'
                  f'<div style="font-weight:600;color:#111827;font-size:14px;">{emoji} {titulo}</div>'
                  f'<div style="font-size:13px;color:#4b5563;margin-top:3px;line-height:1.5;">{desc}</div>'
                  f'</div>')
    return (f'<div style="margin-top:20px;">'
            f'<div style="font-weight:700;color:#111827;margin-bottom:4px;font-size:15px;">O que voc&ecirc; pode fazer agora</div>'
            f'<div style="font-size:12px;color:#6b7280;margin-bottom:8px;">Tudo isso est&aacute; na p&aacute;gina do seu projeto &mdash; &eacute; s&oacute; abrir no bot&atilde;o acima.</div>'
            f'{cards}'
            f'</div>')


def _build_planilha_pronta_email(name: str, project_name: str, job_id: str,
                                 n_total: int, extra_body_html: str = ""):
    """Monta (subject, html) do email 'planilha pronta'. Separado do envio pra
    reuso no preview do painel. `extra_body_html` traz o diagnóstico de leitura +
    próximos passos + avisos (montados com dados reais no envio; exemplo no
    preview) — anexados ao fim do corpo, exatamente como no envio real."""
    import html as _hp
    _pn = _hp.escape(project_name or "seu projeto")
    _greet = _greeting_line(_hp.escape(name or ""))
    _body = (f"{_greet}<br><br>"
             f"O quantitativo do projeto <b>{_pn}</b> terminou de processar "
             f"({n_total} itens). Abra seu projeto pra revisar e baixar a planilha."
             + _email_img("pronta-hero.png", "Planilha de quantitativos pronta, itens rotulados medido ou estimativa")
             + f"{extra_body_html}")
    _pn_subj = (project_name or "").strip()
    subject = (f"{_pn_subj} — sua planilha do AI.arq está pronta"
               if _pn_subj else "Sua planilha do AI.arq está pronta")
    html = _email_wrap("Sua planilha está pronta", _body,
                       "Abrir meu projeto", f"https://ai.arq.br/projeto.html?job_id={job_id}",
                       badge="&#10003; Concluído",
                       reason="Você está recebendo este e-mail porque processou um projeto no AI.arq.")
    return subject, html


def _build_welcome_email(name: str = ""):
    """Monta (subject, html) do email de boas-vindas. Separado do envio pra
    reuso no preview do painel (Central de Emails).

    Redesenhado 02/08/2026 (pedido do Pedro, benchmark = welcome de
    concorrente bem produzido): passos numerados + o que sai do CAD + a linha
    de honestidade medido×estimado como assinatura de marca. Tudo table-based
    inline (Gmail/Outlook). Copy segue as 4 regras (nada interno, PT revisado,
    nada inventado) e a regra do beta: "grátis durante o beta, quantos quiser"."""
    import html as _hw
    greet = _greeting_line(_hw.escape(name or ""))

    def _passo(n, titulo, texto):
        return ('<tr><td style="padding:7px 0;vertical-align:top;width:40px;">'
                '<div style="width:28px;height:28px;border-radius:50%;background:#4F46E5;'
                'color:#ffffff;font-size:14px;font-weight:700;text-align:center;'
                f'line-height:28px;font-family:Arial,sans-serif;">{n}</div></td>'
                '<td style="padding:7px 0 7px 10px;font-size:14px;line-height:1.55;'
                'color:#475569;font-family:Arial,sans-serif;">'
                f'<b style="color:#0F172A;">{titulo}</b><br>{texto}</td></tr>')

    def _check(texto):
        return ('<tr><td style="padding:3px 0;vertical-align:top;width:22px;color:#15803d;'
                'font-size:14px;font-family:Arial,sans-serif;font-weight:700;">&#10003;</td>'
                '<td style="padding:3px 0;font-size:14px;line-height:1.5;color:#475569;'
                f'font-family:Arial,sans-serif;">{texto}</td></tr>')

    _img = _email_img
    body = (
        _img("welcome-foto.jpg", "Arquiteto trabalhando sobre a prancha do projeto", margem="2px 0 14px")
        + f"{greet}<br><br>"
        "Que bom ter você aqui! O AI.arq <b>lê o seu projeto e mede</b> — a planilha de "
        "quantitativos que levaria horas no Excel sai em minutos, direto do seu arquivo, "
        "com cada item dizendo se foi <b>medido</b> ou <b>estimado</b>:"
        + _img("welcome-hero.png", "Planilha de quantitativos com itens medidos e estimativas rotuladas")
        + "<br>"
        '<b style="color:#0F172A;font-size:15px;">Como funciona</b>'
        '<table role="presentation" cellpadding="0" cellspacing="0" width="100%" style="margin-top:6px;">'
        + _passo(1, "Envie o projeto",
                 "DXF é o formato que dá o melhor resultado (medimos a geometria exata). "
                 "DWG e PDF também funcionam.")
        + _passo(2, "A IA mede e monta a planilha",
                 "Itens por disciplina, com referências SINAPI onde há correspondência.")
        + _passo(3, "Revise e baixe",
                 "Confirme as quantidades na tela e exporte em Excel — e os documentos "
                 "abaixo saem do mesmo projeto, sem custo extra.")
        + "</table><br>"
        '<b style="color:#0F172A;font-size:15px;">E não para na planilha</b><br>'
        '<span style="font-size:14px;">O mesmo projeto também gera:</span>'
        + _img("welcome-cronograma.png", "Cronograma físico com barras por etapa e curva de avanço")
        + '<span style="font-size:14px;"><b style="color:#0F172A;">Cronograma físico</b> — '
        'etapas na ordem da obra, com durações calculadas das suas quantidades e curva de avanço. '
        'Baixa em PDF ou apresentação com a sua marca.</span>'
        + _img("welcome-memorial.png", "Memorial descritivo em rascunho com campos a preencher")
        + '<span style="font-size:14px;"><b style="color:#0F172A;">Memorial descritivo (rascunho)</b> — '
        'escrito a partir dos itens do seu CAD, editável na tela, sai em Word ou PDF. '
        'Você completa, o responsável técnico assina.</span>'
        + _img("welcome-comparativo.png", "Comparativo de cotações de fornecedores lado a lado")
        + '<span style="font-size:14px;"><b style="color:#0F172A;">Comparativo de cotações</b> — '
        'suba as planilhas dos fornecedores e veja lado a lado, com comparação justa '
        'e quem esqueceu o quê.</span><br><br>'
        # A linha de honestidade É a marca — vai no primeiro e-mail de propósito.
        '<div style="background:#FFF7ED;border:1px solid #fed7aa;border-radius:10px;'
        'padding:12px 14px;font-size:13px;line-height:1.55;color:#7c4a12;'
        'font-family:Arial,sans-serif;">'
        '<b>Nosso compromisso:</b> cada número diz de onde veio. Em <b>branco</b>, o que '
        'foi <b>medido</b> do seu arquivo; em <b>laranja</b>, o que é <b>estimativa</b> '
        'pra você confirmar. A gente nunca vende chute como certeza.</div><br>'
        "Estamos em <b>beta — grátis, quantos projetos você quiser</b>, sem cartão. "
        "Aproveite pra testar com um projeto real."
    )
    subject = "Bem-vindo ao AI.arq — seu projeto vira planilha medida"
    html = _email_wrap(
        "Bem-vindo ao AI.arq", body,
        "Enviar meu primeiro projeto", "https://ai.arq.br/dashboard.html",
        reason="Você está recebendo este e-mail porque criou sua conta no AI.arq.",
        preheader="Envie o CAD, receba a planilha medida — grátis durante o beta, sem cartão.")
    return subject, html


def _send_welcome_email(email: str, name: str = "") -> bool:
    """Monta + envia o email de boas-vindas. Usado no 1º acesso ao dashboard
    (gated por created_at) e no reenvio manual pelo admin. Best-effort."""
    name = _resolve_client_name(email, hint=name)
    subject, html = _build_welcome_email(name)
    return _send_email_smtp(email, subject, html, log_kind="boas_vindas")


def _generate_magic_link(email: str, redirect_to: str = "https://ai.arq.br/login.html") -> str:
    """Gera um link de login de 1 clique (magic link) via admin API do GoTrue.
    Precisa SUPABASE_SERVICE_ROLE_KEY. Retorna o link ou "" se falhar."""
    try:
        import urllib.request as _urm
        _payload = _json.dumps({
            "type": "magiclink",
            "email": email,
            "redirect_to": redirect_to,
        }).encode("utf-8")
        _r = _urm.Request(f"{SUPABASE_URL}/auth/v1/admin/generate_link", data=_payload, method="POST")
        _r.add_header("apikey", SUPABASE_KEY)
        _r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _r.add_header("Content-Type", "application/json")
        _resp = _json.loads(_urm.urlopen(_r, timeout=10).read().decode("utf-8"))
        return (_resp.get("action_link")
                or (_resp.get("properties") or {}).get("action_link")
                or "")
    except Exception as _e:
        print(f"[magiclink] falha pra {email}: {_e}")
        return ""


def _name_from_auth(user_id: str) -> str:
    """Pega o nome do metadata do auth (full_name/name) por user_id. Útil pra
    incompletos: deram o nome no cadastro mas não criaram o profile. Precisa
    service_role. Retorna "" se não achar."""
    if not user_id:
        return ""
    try:
        import urllib.request as _urn
        _r = _urn.Request(f"{SUPABASE_URL}/auth/v1/admin/users/{user_id}", method="GET")
        _r.add_header("apikey", SUPABASE_KEY)
        _r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _u = _json.loads(_urn.urlopen(_r, timeout=8).read().decode("utf-8"))
        _m = _u.get("user_metadata") or _u.get("raw_user_meta_data") or {}
        return (_m.get("full_name") or _m.get("name") or "").strip()
    except Exception as _e:
        print(f"[name] auth lookup falhou pra {user_id}: {_e}")
        return ""


def _resolve_client_name(email: str = "", user_id: str = "", hint: str = "") -> str:
    """Nome do cliente pra saudação personalizada dos e-mails — regra do Pedro
    (19/07): TODO e-mail pro cliente tem que vir com o nome dele. Busca na ordem
    do que MAIS acha: hint > profiles.full_name (por user_id, depois por email) >
    projects.user_name (por email) > metadata do auth. O nome real quase sempre
    vive em profiles.full_name (o metadata do auth costuma vir VAZIO — foi o que
    fez o e-mail do Tiago sair 'Boa tarde,' sem nome). Best-effort, nunca lança;
    só devolve "" se realmente não houver nome em lugar nenhum."""
    h = (hint or "").strip()
    if h:
        return h
    import urllib.request as _ur, urllib.parse as _up
    def _get(url):
        try:
            r = _ur.Request(url, method="GET")
            r.add_header("apikey", SUPABASE_KEY)
            r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            return _json.loads(_ur.urlopen(r, timeout=8).read().decode("utf-8")) or []
        except Exception:
            return []
    # 1) profiles por user_id (exato)
    if user_id:
        rows = _get(f"{SUPABASE_URL}/rest/v1/profiles?user_id=eq.{_up.quote(str(user_id), safe='')}"
                    f"&select=full_name&limit=1")
        if rows and (rows[0].get("full_name") or "").strip():
            return rows[0]["full_name"].strip()
    # 2) profiles por email (cobre 100% dos usuários hoje — profiles.email sempre preenchido)
    em = (email or "").strip()
    _cands = [em, em.lower()] if (em and em != em.lower()) else ([em] if em else [])
    for cand in _cands:
        rows = _get(f"{SUPABASE_URL}/rest/v1/profiles?email=eq.{_up.quote(cand, safe='')}"
                    f"&select=full_name&limit=1")
        if rows and (rows[0].get("full_name") or "").strip():
            return rows[0]["full_name"].strip()
    # 3) projects por email (nome mais recente que não veio vazio)
    if em:
        rows = _get(f"{SUPABASE_URL}/rest/v1/projects?user_email=eq.{_up.quote(em.lower(), safe='')}"
                    f"&select=user_name&order=created_at.desc&limit=10")
        for row in rows:
            if (row.get("user_name") or "").strip():
                return row["user_name"].strip()
    # 4) metadata do auth (último recurso — costuma vir vazio)
    if user_id:
        return _name_from_auth(user_id)
    return ""


def _build_nudge_email(name: str, kind: str, magic_link: str):
    """Monta (subject, html) do email de lembrete por kind. Separado do envio
    pra reuso no preview. kinds: 'cadastro', 'onboarding', 'feedback'."""
    import html as _hn
    greet = _greeting_line(_hn.escape(name or ""))
    if kind == "cadastro":
        title = "Falta pouco pra começar"
        # 02/08: saiu o "primeiro projeto é por nossa conta" (regra velha) —
        # no beta é "grátis, quantos projetos quiser" (feedback do Pedro 17/07).
        body = (_email_img("nudge-foto.jpg", "Edifício em construção", margem="2px 0 14px")
                + f"{greet}<br><br>Você começou seu cadastro no AI.arq mas não chegou a terminar — "
                f"e falta <b>só um passo</b>. Termine pra subir sua primeira prancha: o "
                f"levantamento de quantitativos sai em minutos, e no beta está <b>grátis — "
                f"quantos projetos você quiser</b>.<br><br>É só clicar abaixo pra entrar direto, sem precisar de senha:")
        cta = "Terminar meu cadastro"
        subject = "Falta pouco pra terminar seu cadastro no AI.arq"
    elif kind == "onboarding":
        title = "Vem subir sua primeira prancha"
        body = (_email_img("nudge-foto.jpg", "Edifício em construção", margem="2px 0 14px")
                + f"{greet}<br><br>Você já tem conta no AI.arq, mas ainda não testou com uma prancha. "
                f"Que tal agora? Manda um PDF, DWG ou DXF e em minutos você recebe a planilha de "
                f"quantitativos — e no beta está <b>grátis, quantos projetos você quiser</b>.<br><br>"
                f"Clica abaixo pra entrar direto e subir:")
        cta = "Subir minha primeira prancha"
        subject = "Sua primeira prancha no AI.arq — grátis no beta"
    else:  # feedback
        title = "Como foi seu projeto no AI.arq?"
        body = (_email_img("feedback-foto.jpg", "Mesa de trabalho de arquitetura", margem="2px 0 14px")
                + f"{greet}<br><br>Vi que você usou o AI.arq pra levantar quantitativos — e eu "
                f"queria muito saber a sua opinião. Leva <b>1 minutinho</b>: você dá uma nota "
                f"pra cada etapa (subir a prancha, processamento, precisão, a planilha) e "
                f"deixa um comentário, se quiser.<br><br>Seu feedback vale ouro pra deixar o "
                f"AI.arq cada vez melhor. 🙂")
        cta = "Avaliar meu projeto"
        subject = "Como foi seu projeto no AI.arq? (1 min)"
    html = _email_wrap(title, body, cta, magic_link,
                       reason="Você está recebendo este e-mail porque criou uma conta no AI.arq.")
    return subject, html


# Mapeia kind do nudge -> coluna 'kind' do email_sent_log (volume por tipo).
_NUDGE_LOG_KIND = {"cadastro": "nudge_cadastro", "onboarding": "nudge_onboarding", "feedback": "feedback"}


def _send_nudge_email(email: str, name: str, kind: str, magic_link: str) -> bool:
    """Email de lembrete com login de 1 clique. kind:
    - 'cadastro'   -> incompleto: 'falta pouco, termine o cadastro'.
    - 'onboarding' -> tem conta mas 0 projetos: 'vem subir sua 1ª prancha'.
    - 'feedback'   -> tem projeto: 'como foi?'."""
    name = _resolve_client_name(email, hint=name)
    subject, html = _build_nudge_email(name, kind, magic_link)
    return _send_email_smtp(email, subject, html, log_kind=_NUDGE_LOG_KIND.get(kind, "nudge"))


def _build_calibracao_email(name: str, project_name: str, job_id: str = ""):
    """Monta (subject, html) do email de calibração: pede as correções reais do
    cliente pra afinar o motor. SEM mencionar dinheiro/cashback (beta grátis).

    Reescrito 02/08/2026 (revisão da esteira): o caminho principal agora é a
    REVISÃO NA TELA (revisao.html), não o upload de planilha revisada. Motivo:
    subir XLSX corrigido é um gesto que quase ninguém faz (o painel de
    revision_feedback vivia vazio enquanto as revisões inline se acumulavam),
    e a tela de revisão ficou muito mais rápida em 01/08 (clique na linha,
    confirmar/editar por palavra, salvamento automático). O upload continua
    citado como alternativa pra quem já trabalha no Excel."""
    import html as _hc
    pn = _hc.escape(project_name or "seu projeto")
    greet = _greeting_line(_hc.escape(name or ""))
    _link = (f"https://ai.arq.br/revisao.html?job_id={job_id}" if job_id
             else "https://ai.arq.br/dashboard.html")
    body = (f"{greet}<br><br>"
            f"Você chegou a conferir o quantitativo do projeto <b>{pn}</b>? Os itens em "
            f"<b>laranja</b> são estimativas esperando a sua palavra final — e cada "
            f"correção sua <b>vale ouro pra gente</b>."
            + _email_img("calibracao-art.png", "Correção do cliente: estimativa vira número confirmado")
            + f"Dá pra fazer <b>direto na tela</b>, item por item: confirmar o que está "
            f"certo, corrigir a quantidade, tirar o que não existe. Salva sozinho, "
            f"e leva poucos minutos.<br><br>"
            f"Cada ajuste seu <b>afina o nosso motor</b>, e o resultado volta pra você: "
            f"os <b>seus próximos projetos saem medindo melhor</b> exatamente o que hoje "
            f"ainda sai como estimativa.<br><br>"
            f"<span style=\"color:#94a3b8;font-size:13px;\">Prefere trabalhar no Excel? "
            f"Você também pode subir a planilha corrigida na página do projeto — "
            f"funciona igual pra nós.</span>")
    subject = "Que tal ajudar a afinar seu quantitativo?"
    html = _email_wrap("Suas correções deixam o motor mais preciso", body,
                       "Revisar meu quantitativo", _link,
                       reason="Você está recebendo este e-mail porque processou um projeto no AI.arq.",
                       preheader="Confirme as estimativas na tela — leva poucos minutos e melhora seus próximos projetos.")
    return subject, html


def _job_tem_planilha_revisada(job_id: str) -> bool:
    """True se o cliente JÁ deu o feedback deste job — por qualquer caminho:
    upload da planilha revisada (evento 'planilha_upload') OU revisão inline
    na tela (item_reviews). Se já revisou, o e-mail de calibração cala a boca.

    02/08/2026: a checagem de item_reviews entrou junto com a virada do e-mail
    pra revisão na tela — sem ela, quem revisou clicando continuaria recebendo
    "que tal revisar?" (o gesto que a gente pede virou outro).
    Best-effort: na dúvida diz True (não pede)."""
    if not job_id:
        return True
    import urllib.request
    def _tem(url):
        try:
            req = urllib.request.Request(url, method="GET")
            req.add_header("apikey", SUPABASE_KEY)
            req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            return bool(_json.loads(urllib.request.urlopen(req, timeout=8).read().decode("utf-8")))
        except Exception:
            return None  # erro: não decide aqui
    up = _tem(f"{SUPABASE_URL}/rest/v1/project_cashback_events?job_id=eq.{job_id}"
              f"&event_type=eq.planilha_upload&select=id&limit=1")
    if up is None:
        return True
    if up:
        return True
    inline = _tem(f"{SUPABASE_URL}/rest/v1/item_reviews?job_id=eq.{job_id}&select=id&limit=1")
    if inline is None:
        return True
    return bool(inline)


def _send_email_calibracao(email: str, name: str, project_name: str, job_id: str = "") -> bool:
    """Pede as correções do cliente (loop de aprendizado do motor). Ligado no
    tick automático em 19/07 (janela 12-30 dias pós-done, 1x por projeto).
    job_id manda o CTA direto pra revisão daquele projeto (02/08)."""
    name = _resolve_client_name(email, hint=name)
    subject, html = _build_calibracao_email(name, project_name, job_id)
    return _send_email_smtp(email, subject, html, log_kind="calibracao")


class JobsStore:
    """Armazena jobs em arquivo JSON."""
    def __getitem__(self, key):
        jobs = _load_jobs()
        if key not in jobs:
            raise KeyError(key)
        return ProcessingStatus(**jobs[key])

    def __setitem__(self, key, value):
        with _JOBS_LOCK:
            jobs = _load_jobs()
            if isinstance(value, ProcessingStatus):
                jobs[key] = value.model_dump()
            else:
                jobs[key] = value
            _save_jobs(jobs)

    def __contains__(self, key):
        return key in _load_jobs()

    def update_field(self, key, **kwargs):
        with _JOBS_LOCK:
            jobs = _load_jobs()
            if key in jobs:
                jobs[key].update(kwargs)
                _save_jobs(jobs)
            else:
                # Atualização de status pra um job que sumiu do store vira no-op
                # MUDO — deixa rastro pra não esconder "status não gravou".
                print(f"[jobs] update_field no-op: job '{key}' fora do store {list(kwargs)}")

jobs = JobsStore()


# ═══════════════════════════════════════════════════════════════
#  Startup recovery: jobs travados após redeploy
#
#  O /tmp do Render é volátil: num redeploy, a thread de processamento
#  morre no meio e o `jobs[job_id]` em memória/JSON some, mas o registro
#  no Supabase fica pendurado em status='queued'/'processing' pra sempre.
#  Admin via acumular jobs-fantasma e o usuário final via "processando…"
#  sem nenhuma thread de verdade rodando.
#
#  Fix: no startup, varrer o Supabase atrás de jobs em aberto e marcá-los
#  como error com mensagem clara. Trades-offs:
#  - Não tenta retomar o processamento (seria complexo e caro); só encerra
#    o status pra o usuário poder reenviar.
#  - Só toca em linhas cujo `created_at` é mais velho que RECOVERY_GRACE_MIN
#    (3min) — evita marcar como erro um job que acabou de nascer no mesmo
#    tick do startup.
# ═══════════════════════════════════════════════════════════════

RECOVERY_GRACE_MIN = 3  # minutos: jobs mais novos que isso são ignorados
RECOVERY_SWEEP_MIN = 5  # minutos: intervalo da varredura periódica de recovery


def _retomar_job_do_storage(job_id: str, typology: str = "office",
                            project_type: str = "arquitetura",
                            conta_retomada: bool = True) -> bool:
    """RESILIÊNCIA (16/06): RETOMA um job interrompido por restart.

    Baixa os arquivos originais do Storage (que agora sobem no upload, antes
    de processar), limpa itens parciais (idempotência) e re-dispara o
    processamento. Retorna True se conseguiu re-disparar; False se não há
    arquivos no Storage (job antigo, pré-resiliência) — aí o caller marca erro.

    Antes (bug que derrubou o Adriano 2x): qualquer restart no meio do
    processamento matava o job pra sempre — 92% dos erros históricos."""
    import urllib.request, json as _j
    from urllib.parse import unquote
    try:
        list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        body = _j.dumps({"prefix": f"{job_id}/", "limit": 200}).encode("utf-8")
        req = urllib.request.Request(list_url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=15)
        objs = _j.loads(resp.read().decode("utf-8"))
        names = [unquote(o.get("name", "")) for o in objs if o.get("name")]
        names = [n for n in names if n.lower().endswith(('.pdf', '.dwg', '.dxf'))]
        if not names:
            return False
        work_dir = os.path.join(WORK_DIR, job_id)
        os.makedirs(work_dir, exist_ok=True)
        file_paths = []
        for n in names:
            data = _supabase_storage_download_prancha(job_id, os.path.basename(n))
            if not data:
                continue
            lp = os.path.join(work_dir, os.path.basename(n))
            with open(lp, "wb") as f:
                f.write(data)
            file_paths.append(lp)
        if not file_paths:
            return False
        # Idempotência: limpa itens parciais antes de reprocessar (DELETE REST)
        try:
            del_url = f"{SUPABASE_URL}/rest/v1/project_items?job_id=eq.{job_id}"
            del_req = urllib.request.Request(del_url, method="DELETE")
            del_req.add_header("apikey", SUPABASE_KEY)
            del_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            urllib.request.urlopen(del_req, timeout=15)
        except Exception as _de:
            print(f"[recovery-retomar] limpeza itens {job_id}: {_de}")
        _supabase_update("projects", "job_id", job_id,
                         {"status": "queued", "error_message": None})
        # Anti-loop: incrementa AUTO_RESUME_COUNT ANTES de disparar (contador
        # PRÓPRIO da auto-retomada, separado do reprocess_count do reprocesso
        # manual). Por que separado: assim a auto-retomada (a) NÃO gasta a cota de
        # reprocesso grátis do usuário e (b) NÃO cai no gate "pular email de
        # reprocesso" — job retomado que conclui/falha volta a NOTIFICAR o cliente.
        # CRÍTICO: NÃO usar _supabase_update (roteia pro update_project_status,
        # que não tem a coluna e descartaria o incremento → loop de retomada).
        if conta_retomada:
            try:
                inc_url = f"{SUPABASE_URL}/rest/v1/rpc/increment_auto_resume_count"
                inc_body = _j.dumps({"p_job_id": job_id}).encode("utf-8")
                inc_req = urllib.request.Request(inc_url, data=inc_body, method="POST")
                inc_req.add_header("apikey", SUPABASE_KEY)
                inc_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                inc_req.add_header("Content-Type", "application/json")
                urllib.request.urlopen(inc_req, timeout=10)
            except Exception as _ie:
                print(f"[recovery-retomar] {job_id}: falha ao incrementar contador: {_ie}")
        else:
            print(f"[recovery-retomar] {job_id}: restart por DEPLOY — retomada "
                  f"não desconta do orçamento (auto_resume_count intacto)")
        # Re-semeia a entrada local (o JSON volátil foi zerado no restart) pra
        # que /api/status volte a responder a barra de progresso durante o
        # resume — sem isso o usuário via 404/"não encontrado" no meio.
        try:
            jobs[job_id] = ProcessingStatus(
                job_id=job_id, status="queued", progress=0,
                current_step="Retomando após reinício do servidor...",
                total_steps=3,
            )
        except Exception as _se:
            print(f"[recovery-retomar] {job_id}: falha ao semear status local: {_se}")
        import threading as _t
        _t.Thread(target=_process_job_throttled,
                  args=(job_id, file_paths, work_dir),
                  kwargs={"typology": typology, "project_type": project_type}, daemon=True).start()
        print(f"[recovery-retomar] {job_id}: RETOMADO com {len(file_paths)} arquivo(s)")
        return True
    except Exception as e:
        print(f"[recovery-retomar] {job_id} falhou: {e}")
        return False


# Erros PASSAGEIROS (infra/IA) que valem re-tentativa automática. Erro de
# ARQUIVO (DWG inválido, PDF escaneado, 0 itens) NÃO casa — precisa do cliente.
import re as _re_auto
_TRANSIENT_ERR_RX = _re_auto.compile(
    r"reinici|interrompid|sobrecarregad|timeout|tempo\s+limite|excedeu|conex|momentane",
    _re_auto.IGNORECASE)


def _auto_retry_erros_transitorios():
    """REVISÃO AUTOMÁTICA (decisão Pedro 07/07): projeto que caiu por causa
    passageira re-tenta SOZINHO na varredura de 5min — o que antes era resgate
    manual (casos sumi/Lia/Thamiry 06/07). Mesma trava do recovery:
    auto_resume_count < 2. Quando o erro vira TERMINAL (esgotou tentativas ou
    é problema do arquivo), alerta interno pro Pedro com o diagnóstico — 1x
    por job (dedup em email_auto_log) — pra ele nunca mais descobrir fuçando."""
    import urllib.request as _u, json as _j
    from datetime import timedelta as _td
    try:
        _cut = (datetime.utcnow() - _td(hours=24)).isoformat()
        q = (f"{SUPABASE_URL}/rest/v1/projects?status=eq.error&archived=not.is.true"
             f"&is_eval=not.is.true"  # avaliações (teste) ficam fora do auto-retry/alerta
             f"&created_at=gte.{_cut}"
             f"&select=job_id,user_email,project_name,error_message,typology,project_type,auto_resume_count"
             f"&limit=20")
        req = _u.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _j.loads(_u.urlopen(req, timeout=15).read().decode("utf-8"))
    except Exception as e:
        print(f"[auto-retry] busca de erros falhou: {e}")
        return

    for row in rows or []:
        job_id = row.get("job_id")
        if not job_id:
            continue
        msg = row.get("error_message") or ""
        count = int(row.get("auto_resume_count") or 0)
        transitorio = bool(_TRANSIENT_ERR_RX.search(msg))

        if transitorio and count < 2:
            typ = row.get("typology") or "office"
            ptype = row.get("project_type") or "arquitetura"
            if _retomar_job_do_storage(job_id, typ, ptype):
                print(f"[auto-retry] {job_id}: erro passageiro → re-tentando sozinho "
                      f"(tentativa {count + 1}/2)")
                _log_error("auto-retry", f"re-tentativa automática {count + 1}/2 "
                           f"(erro passageiro: {msg[:120]})", job_id, severity="info")
                continue
            # sem arquivo no Storage → cai pro alerta terminal abaixo

        # TERMINAL: esgotou tentativas, não é transitório, ou não tem arquivo.
        # Alerta interno 1x por job — o diagnóstico chega no email do Pedro.
        if not _email_auto_ja_enviado(NOTIFY_EMAIL, "alerta_erro_terminal", ref=job_id):
            _causa = ("esgotou as 2 re-tentativas automáticas" if transitorio
                      else "problema no arquivo do cliente (não re-tentável)")
            # QW3 (20/07): a causa TÉCNICA real (error_log) ao lado do rótulo que
            # o cliente viu — pro Pedro parar de investigar às cegas.
            _causa_real = _error_log_causa_real(job_id)
            _bloco_real = (f"<b>Causa técnica real:</b> {_causa_real[:600]}<br>"
                           if _causa_real else "")
            _ok = _notify_admin(
                f"Projeto com erro terminal: {job_id}",
                f"<b>Projeto:</b> {row.get('project_name') or job_id}<br>"
                f"<b>Cliente:</b> {row.get('user_email') or '—'}<br>"
                f"<b>Classificação:</b> {_causa}<br>"
                f"<b>Rótulo que o cliente viu:</b> {msg[:400]}<br>"
                f"{_bloco_real}<br>"
                f"O cliente já recebeu o email de falha com orientação. "
                f"Se for caso de resgate manual, o arquivo está no Storage "
                f"(job <code>{job_id}</code>).")
            if _ok:
                _email_auto_registrar(NOTIFY_EMAIL, "alerta_erro_terminal", ref=job_id)
            else:
                # QW6 (20/07): o alerta é canal único best-effort. Se o SMTP está
                # fora, o aviso sumiria em silêncio — registra crítico no error_log
                # pra ficar rastreável via MCP mesmo sem email. NÃO registra o
                # dedup (assim re-tenta o alerta na próxima varredura, quando o
                # SMTP voltar). A causa real vai junto pra não perder o diagnóstico.
                _log_error(
                    "alert:admin",
                    f"alerta de erro terminal NÃO entregue (SMTP fora) — job {job_id}; "
                    f"cliente {row.get('user_email') or '—'}; rótulo: {msg[:200]}"
                    + (f"; causa real: {_causa_real[:400]}" if _causa_real else ""),
                    job_id, severity="critical")


def _versao_build() -> str:
    """Versão do código em execução (commit do deploy, injetado pelo Render).
    'dev' fora do Render. Usada pra distinguir restart por DEPLOY (versão
    mudou) de crash por OOM/plataforma (mesma versão)."""
    return (os.environ.get("RENDER_GIT_COMMIT") or "").strip()[:8] or "dev"


def _restart_foi_deploy() -> bool:
    """True se o boot ANTERIOR rodava outra versão → este restart foi deploy.
    Deploy não deve gastar o orçamento de retomada do job (2×) nem contar
    pro disjuntor de crash-loop. Best-effort: na dúvida, False (trata como
    crash — mais conservador)."""
    import urllib.request, json, re as _re
    try:
        q = (f"{SUPABASE_URL}/rest/v1/error_log?stage=eq.boot"
             f"&order=created_at.desc&limit=2&select=message")
        req = urllib.request.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = json.loads(urllib.request.urlopen(req, timeout=8).read().decode("utf-8"))
        # rows[0] = o boot ATUAL (registrado antes desta chamada); rows[1] = anterior
        if len(rows) < 2:
            return False
        m = _re.search(r"\(v ([0-9a-z]+)\)", rows[1].get("message") or "")
        if not m:
            return False
        return m.group(1) != _versao_build()
    except Exception:
        return False


def _boots_recentes(minutos: int = 15) -> int:
    """Quantos boots o servidor teve nos últimos N minutos (inclui o atual,
    registrado no startup). 2+ boots em janela curta = crash-loop: algum
    arquivo está derrubando o servidor a cada retomada (incidente 06/07:
    5 boots em 35min, 1 arquivo venenoso matou 5 jobs). Best-effort: em
    dúvida devolve 1 (sem quarentena indevida)."""
    import urllib.request, urllib.parse, json
    try:
        desde = (datetime.utcnow() - timedelta(minutes=minutos)).isoformat() + "+00:00"
        q = (f"{SUPABASE_URL}/rest/v1/error_log?stage=eq.boot"
             f"&created_at=gte.{urllib.parse.quote(desde, safe='')}&select=id")
        req = urllib.request.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = json.loads(urllib.request.urlopen(req, timeout=8).read().decode("utf-8"))
        return len(rows or []) or 1
    except Exception:
        return 1


def _recover_stuck_jobs_on_startup(skip_local_active: bool = False,
                                   crash_loop: bool = False,
                                   deploy_restart: bool = False):
    """RESILIÊNCIA: varre jobs em queued/processing que sobreviveram a um
    restart e tenta RETOMAR (baixa do Storage + reprocessa). Só marca 'error'
    se não houver arquivo no Storage (job antigo) ou se já auto-retomou antes
    (proteção anti-loop via reprocess_count).

    Roda em 2 modos:
    - startup (skip_local_active=False): no boot, o JobsStore local está vazio.
    - periódico (skip_local_active=True): a cada RECOVERY_SWEEP_MIN; pula jobs
      que ESTE processo já está tocando (entrada viva no JobsStore) pra não
      reprocessar em dobro. Pega o caso raro de job que encalhou entre dois
      restarts (ex.: restart dentro da janela de graça).

    Usa RPC `list_stuck_jobs` (SECURITY DEFINER) pra enxergar rows de
    qualquer usuário — o SELECT direto via anon esbarra em RLS."""
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_stuck_jobs"
        body = json.dumps({"p_older_than_minutes": RECOVERY_GRACE_MIN}).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=15)
        rows = json.loads(resp.read().decode("utf-8"))
    except Exception as e:
        print(f"[recovery] falha ao buscar jobs pendentes via RPC: {e}")
        _supa_log(f"RECOVERY fetch ERR {type(e).__name__}: {e}")
        return

    if not rows:
        print("[recovery] nenhum job pendente no Supabase — startup limpo")
        return

    now = datetime.utcnow()
    recovered = 0
    skipped_grace = 0

    for row in rows:
        job_id = row.get("job_id")
        created_at = row.get("created_at")
        if not job_id:
            continue

        # Grace period: não marcar jobs que acabaram de nascer
        try:
            if created_at:
                # formato ISO: 2026-04-20T12:34:56.789Z ou +00:00
                clean = created_at.replace("Z", "+00:00")
                dt = datetime.fromisoformat(clean)
                age_min = (now.replace(tzinfo=dt.tzinfo) - dt).total_seconds() / 60
                if age_min < RECOVERY_GRACE_MIN:
                    skipped_grace += 1
                    continue
        except Exception:
            pass  # se data malformada, marca como erro mesmo assim

        # Modo periódico: se ESTE processo já cuida do job (entrada viva no
        # JobsStore local), não mexe — evita reprocessar em dobro um job que
        # está rodando aqui agora. No startup o store local está vazio (passa).
        if skip_local_active:
            try:
                if job_id in jobs and getattr(jobs[job_id], "status", None) in ("queued", "processing"):
                    continue
            except Exception:
                pass

        # RESILIÊNCIA (16/06): tenta RETOMAR antes de marcar erro.
        # Busca reprocess_count + typology do job (proteção anti-loop:
        # auto-retoma até 2× — job longo multi-arquivo (ex.: 22 DWGs) pode
        # pegar 2 reinícios de deploy no meio; ampliado de 1 pra 2 na
        # auditoria 06/07. Após 2, vira erro pra não entrar em loop).
        prev_count = None   # sentinela: None = leitura FALHOU → fail-closed (não retoma)
        typ = "office"
        ptype = "arquitetura"
        try:
            q = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
                 f"&select=auto_resume_count,typology,project_type")
            qreq = urllib.request.Request(q, method="GET")
            qreq.add_header("apikey", SUPABASE_KEY)
            qreq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            prow = json.loads(urllib.request.urlopen(qreq, timeout=10).read().decode("utf-8"))
            prev_count = int(prow[0].get("auto_resume_count") or 0) if prow else 0
            if prow:
                typ = prow[0].get("typology") or "office"
                ptype = prow[0].get("project_type") or "arquitetura"
        except Exception:
            # Fail-CLOSED (fix 2026-07-22): antes o default 0 anulava OS DOIS freios
            # (quarentena exige >=1, retomada exige <2) — uma única falha de leitura
            # fazia retomar sem limite. Agora leitura falha = prev_count None = não
            # retoma (trata como "já retomou").
            _log_error("recovery:read-fail",
                       "Falha ao ler auto_resume_count — fail-closed, não retoma",
                       job_id, severity="error")
            prev_count = None

        # DISJUNTOR anti-crash-loop (19/07): em crash-loop (2+ boots em 15min),
        # o job que estava PROCESSANDO no momento da morte é o provável veneno
        # (arquivo pesado que re-OOMa a cada retomada). Se ele já gastou uma
        # retomada (prev_count>=1), NÃO retomar de novo — quarentena: erro
        # terminal com orientação de arquivo + alerta pro Pedro. Jobs 'queued'
        # (inocentes na fila) seguem o fluxo normal. Teria transformado o
        # incidente de 06/07 (5 jobs mortos) em 1.
        # Fail-closed: se não deu pra ler o contador de retomadas, não dá pra saber
        # quantas já gastou — não retoma (evita loop infinito); marca erro pro
        # cliente reenviar.
        if prev_count is None:
            _supabase_update("projects", "job_id", job_id, {
                "status": "error",
                "error_message": "Processamento interrompido por reinício do servidor. Reenvie o projeto.",
                "completed_at": now.isoformat(),
            })
            recovered += 1
            continue

        quarentenado = (crash_loop and not deploy_restart
                        and row.get("status") == "processing"
                        and prev_count >= 1)
        if quarentenado:
            _supabase_update("projects", "job_id", job_id, {
                "status": "error",
                "error_message": ("Esse arquivo é pesado demais e derrubou o "
                                  "processamento mais de uma vez. Envie só a "
                                  "prancha de arquitetura (sem 3D/imagens), ou "
                                  "divida o arquivo em partes menores."),
                "completed_at": now.isoformat(),
            })
            _log_error("recovery:quarentena",
                       f"Crash-loop detectado — job em quarentena após "
                       f"{prev_count} retomada(s)", job_id, severity="error")
            try:
                _email_falha_cliente(job_id, reprocessavel=False)
            except Exception:
                pass
            try:
                import threading as _thq
                _thq.Thread(target=_notify_admin, args=(
                    "Disjuntor: job em quarentena (crash-loop)",
                    f"O servidor reiniciou 2+ vezes em 15min e o job "
                    f"<b>{job_id}</b> estava processando nas quedas — provável "
                    f"arquivo venenoso (OOM). Ele foi marcado como erro e o "
                    f"cliente orientado a enviar um arquivo menor."),
                    daemon=True).start()
            except Exception:
                pass
            recovered += 1
            continue

        # Restart por DEPLOY não gasta o orçamento de retomada (o job não tem
        # culpa da atualização); crash real (OOM/plataforma) conta como antes.
        if prev_count < 2 and _retomar_job_do_storage(
                job_id, typ, ptype, conta_retomada=not deploy_restart):
            # _retomar incrementou auto_resume_count (quando conta_retomada)
            # e re-disparou o processamento.
            recovered += 1
            continue

        # Não deu pra retomar (sem arquivo no Storage ou já retomou): marca erro
        ok = _supabase_update("projects", "job_id", job_id, {
            "status": "error",
            "error_message": "Processamento interrompido por reinício do servidor. Reenvie o projeto.",
            "completed_at": now.isoformat(),
        })
        if ok:
            recovered += 1
            # Registra no error_log (visível via MCP/admin) — antes o modo de
            # falha nº1 (job morto por reinício) nunca escrevia aqui, dando a
            # falsa sensação de "tudo verde". De quebra valida a escrita em prod.
            try:
                _log_error("recovery:restart",
                           "Job não-retomável marcado como erro após reinício do servidor",
                           job_id, severity="warning")
            except Exception:
                pass
            # Também atualiza o JobsStore local se ainda estiver lá
            try:
                if job_id in jobs:
                    jobs.update_field(job_id,
                                       status="error",
                                       error_message="Processamento interrompido por reinício do servidor.",
                                       current_step="Erro: reinício do servidor")
            except Exception:
                pass
            # Avisa o cliente que o projeto falhou por reinício (reprocessar
            # resolve). Best-effort; o helper tem dedup interno por job.
            try:
                _email_falha_cliente(job_id, reprocessavel=True)
            except Exception:
                pass

    print(f"[recovery] startup: {recovered} jobs marcados como erro, "
          f"{skipped_grace} pulados (dentro da janela de graça)")
    _supa_log(f"RECOVERY startup recovered={recovered} grace_skip={skipped_grace}")


@app.on_event("startup")
async def _on_startup_recover_jobs():
    """Hook de startup do FastAPI: (1) registra o boot (sinal do disjuntor
    anti-crash-loop), (2) recupera jobs travados do redeploy anterior e
    (3) sobe uma varredura periódica que pega jobs encalhados entre dois
    restarts (ex.: restart dentro da janela de graça)."""
    crash_loop = False
    deploy_restart = False
    try:
        _log_error("boot", f"Servidor iniciou (v {_versao_build()})", severity="info")
        deploy_restart = _restart_foi_deploy()
        n_boots = _boots_recentes(15)
        crash_loop = n_boots >= 2
        if crash_loop and not deploy_restart:
            print(f"[recovery] ⚠ {n_boots} boots em 15min — modo crash-loop "
                  f"(job ativo com retomada gasta vai pra quarentena)")
        elif deploy_restart:
            print(f"[recovery] restart por deploy (v {_versao_build()}) — "
                  f"retomadas não descontam do orçamento")
    except Exception:
        pass
    try:
        _recover_stuck_jobs_on_startup(crash_loop=crash_loop,
                                       deploy_restart=deploy_restart)
    except Exception as e:
        # Nunca deixar o startup falhar por causa da recuperação
        print(f"[recovery] exceção não-fatal no startup: {e}")

    # Varredura periódica (modo skip_local_active: não toca em job que este
    # processo já está rodando). Thread daemon — morre junto com o processo.
    def _periodic_recovery_loop():
        while True:
            try:
                time.sleep(RECOVERY_SWEEP_MIN * 60)
                _recover_stuck_jobs_on_startup(skip_local_active=True)
                _auto_retry_erros_transitorios()
            except Exception as _e:
                print(f"[recovery] varredura periódica falhou (segue ativa): {_e}")
    try:
        import threading as _tr
        _tr.Thread(target=_periodic_recovery_loop, daemon=True).start()
        print(f"[recovery] varredura periódica ativa (a cada {RECOVERY_SWEEP_MIN}min)")
    except Exception as _e:
        print(f"[recovery] não subiu a varredura periódica: {_e}")


# _extract_balanced_obj e _salvage_truncated_json movidos pra engine_rules.py
# (fonte única, sem duplicação; testados em tests/test_engine_rules.py).
# Importados no topo do arquivo.


# Regras determinísticas de unidade por tipo de serviço (pós-IA).
# Se a IA retornar unidade errada pra descrição específica, o código força a
# unidade correta e marca o item como "estimado" (laranja) pra o usuário revisar.
import re as _re
_UNIT_SURFACE_KEYWORDS = _re.compile(
    r"\b(pisos?|forros?|pinturas?|revestimentos?|azulejos?|cer[âa]micas?|"
    r"porcelanatos?|carpetes?|viníl|vinílicos?|tapetes?)\b",
    _re.IGNORECASE,
)
_UNIT_LINEAR_KEYWORDS = _re.compile(
    r"\b(rodap[eé]s?|tabicas?|soleiras?|perfi(?:l|s)|perfilados?|molduras?|"
    r"eletrocalhas?|eletrodutos?|trilhos?|cord[aã]o|cord[õo]es|"
    r"cornijas?)\b",
    _re.IGNORECASE,
)
_UNIT_COUNT_KEYWORDS = _re.compile(
    r"\b(lumin[áa]rias?|spots?|projetores?|pendentes?|arandelas?|"
    r"portas?|janelas?|esquadrias?|tomadas?|interruptores?|sensores?|"
    r"difusores?|grelhas?|sprinklers?|detectores?|c[âa]meras?|cftv|"
    r"tvs?|televisor(?:es)?|monitor(?:es)?|"  # eletrônicos: TV/monitor são UN, nunca m²
    r"quadro|qdf|caixa(?:\s+de\s+som)?|al[çc]ap[ãa]o|al[çc]ap[õo]es|"
    r"chuveiros?|torneiras?)\b",
    _re.IGNORECASE,
)


_GENERIC_WORDS = {
    "de", "do", "da", "dos", "das", "para", "com", "em",
    "nova", "novo", "existente", "existentes",
    "conforme", "especificacao", "especificacoes", "projeto",
    "instalacao", "execucao", "fornecimento", "fornecida",
    "tipo", "tipos", "cor", "modelo", "padrao",
    "altura", "comprimento", "largura", "espessura",
    "ceramico", "ceramica", "ceramicos", "ceramicas",
    "metalico", "metalica", "plastico", "plastica",
    "area", "areas",
    "m2", "m", "un", "ml",
}


def _normalize_description_key(desc: str) -> str:
    """Reduz a descrição a uma chave de comparação — remove acentos, sufixos
    por departamento, números extras, e faz lowercase. Usado pra detectar
    itens similares em consolidação."""
    if not desc:
        return ""
    s = desc.lower()
    # Remover acentos
    import unicodedata
    s = unicodedata.normalize('NFD', s)
    s = ''.join(c for c in s if unicodedata.category(c) != 'Mn')
    # Remover sufixos de departamento/variante
    # Ex.: "painel divisorio - contabilidade" → "painel divisorio"
    # Ex.: "demarcacao de area departamento contabilidade" → "demarcacao de area"
    for sep in (' - ', ' — ', ' / ', ' departamento ', ' deptos ', ' do depto ',
                ' da sala ', ' sala ', ' para sala '):
        s = s.split(sep)[0]
    # Normalizar espaços e pontuação
    s = _re.sub(r"[^a-z0-9]+", " ", s)
    s = _re.sub(r"\s+", " ", s).strip()
    # Remover palavras genéricas que não mudam o significado
    tokens = [t for t in s.split() if t not in _GENERIC_WORDS and len(t) > 1]
    return " ".join(sorted(tokens[:6]))  # primeiras 6 palavras ordenadas


def _primary_noun(desc: str) -> str:
    """Retorna o primeiro token significativo (>3 chars, não-genérico) da
    descrição. Usado pra detectar mesma 'família' (alvenaria, luminária, etc.)
    mesmo quando descrições divergem bastante."""
    if not desc:
        return ""
    import unicodedata
    s = unicodedata.normalize('NFD', desc.lower())
    s = ''.join(c for c in s if unicodedata.category(c) != 'Mn')
    s = _re.sub(r"[^a-z0-9]+", " ", s)
    for t in s.split():
        if len(t) > 3 and t not in _GENERIC_WORDS:
            return t
    return ""


def _pick_best_unit(units: list[str], description: str) -> str:
    """Entre vários units candidatos, escolhe o mais coerente com a descrição.
    Se o normalizer semântico diz algo, usa. Senão, prefere contagem (un) >
    linear (ml) > superfície (m²) > vb — ordem de especificidade crescente."""
    norm, corrected = _normalize_unit_for_item(description, units[0] if units else "vb")
    if corrected:
        return norm
    # Sem opinião semântica: primeiro unidade não-m² presente (m² é o problema usual)
    priority = ["un", "ml", "m²", "m", "vb"]
    for u in priority:
        if u in units:
            return u
    return units[0] if units else "vb"


def _pick_area_consensus(readings: list) -> float:
    """Consenso de área entre leituras de várias pranchas.

    Antes: cada prancha sobrescrevia o valor anterior — bug onde uma prancha
    de detalhe com valor errado (135m² em vez de 270m² ou vice-versa)
    "ganhava" simplesmente por ser a última processada. Resultado: área
    dependia da ordem de processamento.

    Agora: coleta todas as leituras, agrupa por similaridade (±5%), e
    retorna a MODA (valor mais frequente). Empate resolvido pelo maior valor.

    Exemplo:
      leituras = [135.4, 135.4, 270.0, 135.0]  → winner bucket = [135.4, 135.4, 135.0] (3 pranchas)
      → retorna ~135.3
    """
    vals = [float(v) for v in readings if v and float(v) > 0]
    if not vals:
        return 0.0
    if len(vals) == 1:
        return vals[0]

    buckets: list[list[float]] = []
    for v in vals:
        placed = False
        for bucket in buckets:
            median = sum(bucket) / len(bucket)
            if median > 0 and abs(v - median) / median <= 0.05:
                bucket.append(v)
                placed = True
                break
        if not placed:
            buckets.append([v])

    # Mais leituras = mais confiável. Empate: maior valor (laje bruta
    # costuma ser o maior dos candidatos na dúvida).
    buckets.sort(key=lambda b: (-len(b), -max(b)))
    winner = buckets[0]
    return round(sum(winner) / len(winner), 2)


def _dedupe_rooms(rooms: list) -> list:
    """Dedup de ambientes em kept_elements/new_rooms.

    A IA extrai "BANHEIRO" de uma prancha e "Banheiro suíte" de outra; sem
    dedup, ambos aparecem no resumo. Regra:
    - Case-insensitive + remove acentos pra comparar
    - Se um ambiente é substring do outro (ex.: "banheiro" ⊂ "banheiro
      suíte"), mantém só o MAIS específico
    - Preserva a forma original de escrita do item mantido
    - Aceita itens como strings OU dicts (formato usado por new_rooms)
    """
    if not rooms:
        return rooms

    import unicodedata

    def _key(r) -> str:
        if isinstance(r, dict):
            name = r.get("name", "")
        else:
            name = str(r or "")
        s = name.lower().strip()
        s = unicodedata.normalize("NFD", s)
        s = "".join(c for c in s if unicodedata.category(c) != "Mn")
        s = _re.sub(r"[^a-z0-9]+", " ", s).strip()
        return s

    # Ordena por length decrescente — "banheiro suíte" vem antes de "banheiro"
    # pra garantir que a versão mais específica seja mantida
    sorted_rooms = sorted(rooms, key=lambda r: -len(_key(r)))
    seen_keys: list[str] = []
    out = []
    for r in sorted_rooms:
        k = _key(r)
        if not k:
            continue
        # Se já vimos key igual ou key atual é substring de alguma já vista
        # (ex.: "banheiro" enquanto já temos "banheiro suite"), pula
        is_subset = any(k == s or k in s.split() and len(k.split()) == 1
                        for s in seen_keys)
        if is_subset:
            continue
        # Também: se key atual tem substring em comum de >=2 palavras com
        # alguma já vista, considera duplicata (ex.: "sala de estar e jantar"
        # vs "sala de estar" — ambos devem virar 1)
        cur_tokens = set(k.split())
        is_dup = False
        for s in seen_keys:
            overlap = cur_tokens & set(s.split())
            if len(overlap) >= 2 and len(overlap) >= len(cur_tokens) * 0.6:
                is_dup = True
                break
        if is_dup:
            continue
        seen_keys.append(k)
        out.append(r)
    return out


def _dedupe_by_block(items: list) -> list:
    """Funde itens que citam o MESMO bloco CAD (nome entre aspas na descrição).
    A IA às vezes cria o mesmo bloco em 2-3 disciplinas (ex: 'cad-escr-02' como
    'Cadeira para escritório' E 'Mobiliário de escritório') → a contagem inflava
    (14 viravam 28). Mesmo bloco = MESMA contagem física: mantém 1 item (o de
    maior confiança/descrição), usa a MAIOR qty (NÃO soma), descarta o resto."""
    try:
        from models import Confidence
        _conf_ok = Confidence("confirmado")
    except Exception:
        _conf_ok = None
    by_block: dict = {}
    passthrough: list = []
    for it in items:
        bk = _extract_block_name(getattr(it, "description", "") or "")
        if not bk:
            passthrough.append(it)
        else:
            by_block.setdefault(bk, []).append(it)
    merged = list(passthrough)
    fundidos = 0
    for bk, group in by_block.items():
        if len(group) == 1:
            merged.append(group[0])
            continue
        best = max(group, key=lambda it: (
            1 if (_conf_ok is not None and it.confidence == _conf_ok) else 0,
            len(getattr(it, "description", "") or ""),
        ))
        try:
            best.quantity = max((float(it.quantity or 0) for it in group),
                                default=float(best.quantity or 0))
        except Exception:
            pass
        best.observations = ((getattr(best, "observations", "") or "") +
            f" | Bloco '{bk}' aparecia em {len(group)} itens — fundido pra não "
            f"duplicar a contagem (mesma peça do CAD)").strip(" |")
        merged.append(best)
        fundidos += len(group) - 1
    if fundidos:
        print(f"[dedup-bloco] {fundidos} itens duplicados de bloco fundidos")
    return merged


def _drop_nonsense_items(items: list) -> list:
    """Remove itens-artefato que não são quantitativo de verdade (ex: 'área de
    seção transversal' de parede — a hachura da ESPESSURA virou item de m²).
    Caso Thamiry (projeto drywall multi-prancha)."""
    kept = [it for it in items if not _is_nonsense_item(getattr(it, "description", "") or "")]
    n = len(items) - len(kept)
    if n:
        print(f"[limpeza] {n} itens-artefato (seção transversal) removidos")
    return kept


def _consolidate_by_type_code(items: list) -> list:
    """Funde o MESMO tipo de divisória/parede (DRY 07, DW-12...) que aparece em
    VÁRIAS pranchas: SOMA as qtys (paredes diferentes em ambientes diferentes),
    1 linha por (tipo, unidade, disciplina). Marca estimado (soma cross-prancha =
    revisar). Caso Thamiry: 191 itens fragmentados de drywall, 156 zerados."""
    from models import Confidence
    by_key: dict = {}
    passthrough: list = []
    for it in items:
        tc = _extract_type_code(getattr(it, "description", "") or "")
        if not tc:
            passthrough.append(it)
        else:
            by_key.setdefault((tc, it.unit, it.discipline), []).append(it)
    merged = list(passthrough)
    fundidos = 0
    for (tc, _u, _d), group in by_key.items():
        if len(group) == 1:
            merged.append(group[0])
            continue
        best = max(group, key=lambda it: len(getattr(it, "description", "") or ""))
        try:
            best.quantity = round(sum(float(it.quantity or 0) for it in group), 2)
        except Exception:
            pass
        try:
            best.confidence = Confidence("estimado")
        except Exception:
            pass
        best.observations = ((getattr(best, "observations", "") or "") +
            f" | Tipo '{tc}' somado de {len(group)} entradas em várias pranchas — confira o total").strip(" |")
        merged.append(best)
        fundidos += len(group) - 1
    if fundidos:
        print(f"[consolida-tipo] {fundidos} entradas de tipo de divisória fundidas")
    return merged


def _consolidate_items(items: list) -> list:
    """Consolida itens redundantes em múltiplas passadas:

    PASSADA 1 — por (chave_normalizada, unidade):
    - Mesma chave + mesma qty + mesma unidade → mantém 1 (desc mais completa).
    - Mesma chave + mesma unidade + qtys diferentes → mantém todos.
    - Réplica por departamento (4+ itens, qty<2) → funde em 1 estimado.

    PASSADA 2 — fusões qty+discipline+noun (qty >= 2):
    - Mesma qty_arredondada + mesma discipline + units diferentes →
      funde. Escolhe melhor unit.
    - Mesma qty_arredondada + mesma discipline + mesma unit + primary_noun
      igual → funde.

    PASSADA 3 — fusões por FAMÍLIA em qty pequena (unit=un, vb, ml):
    - Itens com mesmo primary_noun + mesma discipline + qty <= 2 +
      descrições "quase iguais" (variantes de legenda "conforme
      especificação XX") → consolida em 1 linha "NOUN — N un (várias
      variantes)" somando qtys.
    - Casos típicos: 6 portas "conforme especificação X" qty=1, 5 ralos,
      4 luminárias LM1/LM2 duplicadas cross-prancha.

    PASSADA 4 — agrupar luminárias por código LM*/LN*:
    - Descrições tipo "Luminária LM1" extraídas de pranchas diferentes
      viram 1 linha por código, somando qtys.

    PASSADA 5 — filtrar itens "a definir" sem fonte real:
    - Descrições/obs com "a definir", "por definir", "conforme projeto"
      + qty=1 default + 3+ casos similares → consolida em 1 linha única
      "Itens a especificar em projeto executivo".
    """
    from models import BudgetItem, Confidence

    # ── Passada 1 ──
    groups: dict = {}
    for item in items:
        key = (_normalize_description_key(item.description), item.unit)
        groups.setdefault(key, []).append(item)

    pass1 = []
    for key, group in groups.items():
        if len(group) == 1:
            pass1.append(group[0])
            continue

        quantities = [round(float(it.quantity), 2) for it in group]
        unique_qtys = set(quantities)

        # Réplica por departamento tem prioridade (4+ itens com qty < 2) —
        # independente de qtys serem idênticas, porque itens "Contabilidade"
        # e "RH" costumam bater mesmo quando são na verdade áreas distintas.
        if max(quantities) < 2.0 and len(group) >= 4:
            best = max(group, key=lambda x: len(x.description or ""))
            clean_desc = best.description
            for sep in (' - ', ' — ', ' departamento ', ' deptos ', ' da sala '):
                clean_desc = clean_desc.split(sep)[0]
            total_qty = round(sum(quantities), 2)
            consolidated = BudgetItem(
                item_num=best.item_num,
                description=f"{clean_desc.strip()} (várias variantes)",
                unit=best.unit,
                quantity=total_qty,
                observations=(
                    f"Consolidado de {len(group)} entradas replicadas por "
                    f"departamento/variante — soma de qtys: {total_qty} {best.unit}. "
                    f"Revisar se faz sentido tratar como item único."
                ),
                ref_sheet=best.ref_sheet,
                confidence=Confidence("estimado"),
                discipline=best.discipline,
            )
            pass1.append(consolidated)
        elif len(unique_qtys) == 1:
            best = max(group, key=lambda x: len(x.description or ""))
            pass1.append(best)
        else:
            pass1.extend(group)

    # ── Passada 2 ──
    # Cada item é reapresentado com fingerprint (discipline, qty_arred). Se
    # mais de um item compartilha fingerprint, avaliamos noun+overlap pra
    # decidir se funde.
    # Qty pequena (< 2) é frequentemente um "un" que se repete por acaso — não
    # fundimos por coincidência de qty+discipline nesse range pra evitar falso
    # positivo (ex.: 1 porta de emergência + 1 porta comum ambas qty=1).
    MIN_QTY_PASS2 = 2.0

    buckets: dict = {}
    for it in pass1:
        try:
            qty_r = round(float(it.quantity or 0), 2)
        except Exception:
            qty_r = 0.0
        if qty_r < MIN_QTY_PASS2:
            buckets.setdefault(("__solo__", id(it)), []).append(it)
            continue
        buckets.setdefault((it.discipline or "", qty_r), []).append(it)

    pass2 = []
    for (disc, qty_r), group in buckets.items():
        if disc == "__solo__" or len(group) == 1:
            pass2.extend(group)
            continue

        # Tenta fundir itens do mesmo bucket em "famílias". Critério de fusão:
        # mesmo primary_noun OU interseção de >= 2 tokens significativos.
        # (1 token só gera FP — ex.: dois itens com "porta" mas sentidos distintos.)
        families: list[list] = []
        for it in group:
            noun = _primary_noun(it.description)
            key_tokens = set(_normalize_description_key(it.description).split())
            placed = False
            for fam in families:
                fam_noun = _primary_noun(fam[0].description)
                fam_tokens = set(_normalize_description_key(fam[0].description).split())
                overlap = key_tokens & fam_tokens
                if noun and noun == fam_noun:
                    fam.append(it); placed = True; break
                if len(overlap) >= 2:
                    fam.append(it); placed = True; break
            if not placed:
                families.append([it])

        for fam in families:
            if len(fam) == 1:
                pass2.append(fam[0])
                continue
            # Fundir família: melhor descrição, melhor unidade, obs combinada
            best = max(fam, key=lambda x: len(x.description or ""))
            units = [it.unit for it in fam]
            chosen_unit = _pick_best_unit(units, best.description)
            unit_changed = len(set(units)) > 1
            variant_count = len(fam)
            obs_parts = [best.observations or ""]
            if unit_changed:
                obs_parts.append(
                    f"Fundido de {variant_count} entradas com units divergentes "
                    f"({'/'.join(sorted(set(units)))}) — mesma qty {qty_r}"
                )
            else:
                obs_parts.append(
                    f"Fundido de {variant_count} entradas com mesma qty "
                    f"{qty_r} {chosen_unit} — descrições similares"
                )
            merged_item = BudgetItem(
                item_num=best.item_num,
                description=best.description,
                unit=chosen_unit,
                quantity=float(qty_r),
                observations=" | ".join(p for p in obs_parts if p).strip(),
                ref_sheet=best.ref_sheet,
                confidence=Confidence("estimado"),
                discipline=best.discipline,
            )
            pass2.append(merged_item)

    # ═══════════════════════════════════════════════════════════════
    #  PASSADA 3 — Famílias em qty pequena (un/ml/vb com qty<=2)
    # ═══════════════════════════════════════════════════════════════
    # Alvo: "Porta conforme especificação 08" qty=1 × 6 duplicadas;
    # "Ralo sifonado" qty=1 × 5 vezes; "Mobilização" qty=1 × 3 vezes.
    # A IA analisa cada prancha em isolamento e cria variantes da mesma
    # família que passada 2 ignora (corte em qty>=2).
    #
    # Regra pra evitar falso-positivo: só funde se as descrições são
    # "quase iguais" (começam com o mesmo noun + mesmas primeiras 3
    # palavras OU todas contêm "especificação"/"conforme" como marca
    # de item-de-legenda).
    def _is_legend_variant(descs: list[str]) -> bool:
        """True se as descrições parecem variações numéricas da mesma
        linha de legenda (ex.: 'Porta conforme especificação 08', 09, 10)."""
        markers = ("conforme especifica", "conforme projeto", "especifica",
                   "a definir", "por definir", "conforme detalhe",
                   "conforme indicado")
        hits = sum(1 for d in descs if any(m in (d or "").lower() for m in markers))
        # Se >=60% têm marca de legenda, trata como variante
        return hits >= max(2, int(len(descs) * 0.6))

    def _same_opening(descs: list[str], n_tokens: int = 3) -> bool:
        """True se todas descrições começam com os mesmos n primeiros tokens
        significativos (>3 chars, não-genéricos)."""
        def _open(d):
            norm = _normalize_description_key(d or "")
            return " ".join(norm.split()[:n_tokens])
        openings = {_open(d) for d in descs}
        return len(openings) == 1 and next(iter(openings)) != ""

    # Código de luminária (LM1, LM2, etc.) — se presente, separa o bucket
    # pra passada 3 não misturar luminárias de códigos diferentes.
    _lum_code_re_p3 = _re.compile(r"\b([LMN][MN]\d{1,2})\b", _re.IGNORECASE)

    buckets_p3: dict = {}
    keep_as_is: list = []
    for it in pass2:
        try:
            qty_r = round(float(it.quantity or 0), 2)
        except Exception:
            qty_r = 0.0
        # Só considera qty pequena (1, 1.5, 2); qty grandes já foram
        # tratadas em passada 2
        if qty_r > 2.0 or qty_r <= 0:
            keep_as_is.append(it)
            continue
        if it.unit not in ("un", "ml", "vb", "cj"):
            keep_as_is.append(it)
            continue
        noun = _primary_noun(it.description)
        if not noun:
            keep_as_is.append(it)
            continue
        # Código de luminária entra no key pra não misturar LM1 com LM2
        code = _lum_code_re_p3.search(it.description or "")
        code_tag = code.group(1).upper() if code else ""
        # Bucket: discipline + primary_noun + unit + código LM (se houver)
        key = (it.discipline or "", noun, it.unit, code_tag)
        buckets_p3.setdefault(key, []).append(it)

    pass3 = list(keep_as_is)
    for (disc, noun, unit, code_tag), group in buckets_p3.items():
        descs = [it.description for it in group]

        # Critério de consolidação por nível de segurança:
        #
        # SEGURO — sempre consolida, mesmo com só 2 itens:
        #   - unit=vb + mesmo noun (ex.: 3 "Mobilização" vb=1 é sempre erro;
        #     vb não faz sentido duplicar)
        #
        # FORTE — consolida com 3+ itens se mesmo noun/disc/unit:
        #   - un + mesmo noun ("Ralo" × 5, "Porta" × 6) são quase sempre
        #     variantes da mesma família em pranchas diferentes
        #
        # BORDERLINE — precisa reforço (legenda-variant ou mesma abertura):
        #   - ml + mesmo noun (rodapé de cozinha ≠ rodapé do banheiro)
        #   - un com só 2 itens (pode ser 2 portas distintas legítimas)
        should_merge = False
        if unit == "vb" and len(group) >= 2:
            should_merge = True  # vb duplicado é sempre erro
        elif unit == "un" and len(group) >= 3 and (
                _is_legend_variant(descs) or _same_opening(descs, 3)):
            # 3+ "un" SÓ funde COM reforço (mesma legenda OU mesma abertura de 3
            # tokens) — igual ao ramo genérico abaixo. Sem isso, 3 produtos
            # DISTINTOS de mesmo substantivo (porta de madeira PM-01 + de vidro
            # PV-01 + de emergência P-EM) virariam 1 linha "qtd 3", violando a
            # regra dura nº4 (tipo específico nunca some) e entregando
            # quantitativo errado. Bug achado na auditoria 27/07
            # (test_no_false_positive_portas). Ralo/Porta IGUAIS repetidos têm a
            # mesma abertura → o _same_opening deixa fundir; tipos diferentes não.
            should_merge = True
        elif unit == "cj" and len(group) >= 2:
            should_merge = True  # conjuntos duplicados
        elif len(group) >= 3 and (_is_legend_variant(descs) or _same_opening(descs, 3)):
            should_merge = True  # ml/etc só com reforço
        elif len(group) >= 2 and _is_legend_variant(descs):
            # 2 itens ambos com "conforme especificação" também funde
            should_merge = True

        if not should_merge:
            pass3.extend(group)
            continue

        # Consolida
        best = max(group, key=lambda x: len(x.description or ""))
        total_qty = round(sum(float(it.quantity or 0) for it in group), 2)
        # Remove sufixo numérico da legenda pra descrição limpa
        clean = _re.sub(r"\s*(conforme\s+)?(especifica[çc][aã]o\s+\d+|especifica[çc][aã]o\b).*$",
                        "", best.description, flags=_re.IGNORECASE).strip()
        if not clean or len(clean) < 5:
            clean = noun.capitalize()
        consolidated = BudgetItem(
            item_num=best.item_num,
            description=f"{clean} — {len(group)} variantes consolidadas",
            unit=unit,
            quantity=total_qty,
            observations=(
                f"Consolidado de {len(group)} entradas com mesma família "
                f"({noun}) — soma: {total_qty} {unit}. "
                f"Ver legenda do projeto pra especificações individuais."
            ),
            ref_sheet=best.ref_sheet,
            confidence=Confidence("estimado"),
            discipline=disc,
        )
        pass3.append(consolidated)

    # ═══════════════════════════════════════════════════════════════
    #  PASSADA 4 — Luminárias por código LM*/LN*/LU*
    # ═══════════════════════════════════════════════════════════════
    # "Luminária LM1" extraída de 3 pranchas como 3 itens separados
    # (qty 3 + 1 + 1) vira 1 linha (qty 5). Só aplicável em iluminação.
    lum_code_re = _re.compile(r"\b([LMN][MN]\d{1,2})\b", _re.IGNORECASE)

    def _lum_code(desc: str) -> str | None:
        m = lum_code_re.search(desc or "")
        return m.group(1).upper() if m else None

    by_code: dict[str, list] = {}
    pass4 = []
    for it in pass3:
        if (it.discipline or "").lower() != "iluminação".lower():
            pass4.append(it)
            continue
        code = _lum_code(it.description)
        if not code or it.unit != "un":
            pass4.append(it)
            continue
        by_code.setdefault(code, []).append(it)

    for code, group in by_code.items():
        if len(group) == 1:
            pass4.append(group[0])
            continue
        best = max(group, key=lambda x: len(x.description or ""))
        total_qty = round(sum(float(it.quantity or 0) for it in group), 2)
        total_qty_int = int(total_qty) if total_qty == int(total_qty) else total_qty
        consolidated = BudgetItem(
            item_num=best.item_num,
            description=best.description,
            unit="un",
            quantity=float(total_qty_int),
            observations=(
                f"Consolidado — código {code} aparece em {len(group)} pranchas, "
                f"somando {total_qty_int} unidades. Ver quadro de cargas "
                f"do projeto pra confirmar total."
            ),
            ref_sheet=best.ref_sheet,
            confidence=best.confidence,  # preserva confirmado se todos eram
            discipline=best.discipline,
        )
        pass4.append(consolidated)

    # ═══════════════════════════════════════════════════════════════
    #  PASSADA 5 — Consolidar itens "a definir" sem fonte clara
    # ═══════════════════════════════════════════════════════════════
    # "Acabamento elétrico a definir" qty=1 × 3 vezes é noise —
    # orçamentista já sabe que precisa de projeto executivo. Consolida
    # em 1 linha única por disciplina.
    def _is_vague(it) -> bool:
        text = f"{it.description or ''} {it.observations or ''}".lower()
        vague_markers = (
            "a definir", "por definir", "a ser definida",
            "acabamento elétrico a definir", "acabamentos elétricos a definir",
            "especificação a definir",
        )
        has_vague = any(m in text for m in vague_markers)
        short = (it.description or "").strip()
        # Só marca como vago se qty=1/0 e unit vb/un e o texto é genérico
        try:
            qty_r = float(it.quantity or 0)
        except Exception:
            qty_r = 0.0
        return has_vague and qty_r <= 1.0 and it.unit in ("un", "vb", "cj")

    vague_by_disc: dict[str, list] = {}
    pass5 = []
    for it in pass4:
        if _is_vague(it):
            vague_by_disc.setdefault(it.discipline or "Complementares", []).append(it)
        else:
            pass5.append(it)

    for disc, group in vague_by_disc.items():
        if len(group) == 1:
            pass5.append(group[0])
            continue
        # 2+ itens vagos mesma disciplina → consolida em 1 linha
        best = max(group, key=lambda x: len(x.description or ""))
        consolidated = BudgetItem(
            item_num=best.item_num,
            description=f"Itens de {disc.lower()} a especificar em projeto executivo",
            unit="vb",
            quantity=1.0,
            observations=(
                f"Consolidado de {len(group)} entradas genéricas "
                f"(\"a definir\", \"conforme projeto\") sem especificação na "
                f"legenda. Quantificar e especificar no projeto executivo."
            ),
            ref_sheet=best.ref_sheet,
            confidence=Confidence("estimado"),
            discipline=disc,
        )
        pass5.append(consolidated)

    # ═══════════════════════════════════════════════════════════════
    #  PASSADA 6 — DEDUP CROSS-PRANCHA (mesmo elemento em pranchas diferentes)
    # ═══════════════════════════════════════════════════════════════
    # Caso real Weslei (2026-05-08): "Escada metálica" apareceu 2 vezes:
    #   - "PROJETO AMPLIACAO" → qty=8.6, unit=m² (medida REAL da legenda)
    #   - "eletrico.pdf" → qty=1, unit=un (vista da mesma escada na elétrica)
    # Resultado: usuário recebia 2 escadas, sendo a mesma física.
    #
    # Regra: agrupar por (primary_noun, discipline). Pra grupos com 2+ itens
    # vindos de pranchas DIFERENTES (ref_sheet) — provável duplicação.
    # Mantém o item de MAIOR ESPECIFICIDADE (qty>0 vence qty=0, m²>m>un>vb,
    # confirmado vence estimado).

    # Prioridade de unidade (maior = mais específica). m²/m³ vencem un/vb.
    _UNIT_PRIORITY = {
        "m²": 100, "m2": 100, "m³": 100, "m3": 100,
        "m": 80, "ml": 80,
        "kg": 60,
        "un": 40, "cj": 35,
        "mês": 30, "dia": 25,
        "vb": 10, "%": 5,
    }

    def _unit_priority(u: str) -> int:
        return _UNIT_PRIORITY.get((u or "").lower().strip(), 50)

    def _item_score(it) -> tuple:
        """Pra ordenar candidatos a manter — maior = melhor."""
        try:
            qty = float(it.quantity or 0)
        except Exception:
            qty = 0.0
        return (
            1 if qty > 0 else 0,                          # qty>0 vence qty=0
            1 if it.confidence == Confidence("confirmado") else 0,
            _unit_priority(it.unit),                       # m² vence un vence vb
            qty,                                           # qty maior desempata
            len(it.description or ""),                     # descrição mais completa
        )

    buckets_p6: dict[tuple, list] = {}
    for it in pass5:
        noun = _primary_noun(it.description or "")
        if not noun:
            buckets_p6.setdefault(("__solo__", id(it)), []).append(it)
            continue
        key = (noun, it.discipline or "")
        buckets_p6.setdefault(key, []).append(it)

    pass6 = []
    for key, group in buckets_p6.items():
        if key[0] == "__solo__" or len(group) == 1:
            pass6.extend(group)
            continue

        # Pra grupo com 2+ itens, verificar se vêm de pranchas diferentes
        ref_sheets = set((it.ref_sheet or "").strip().lower() for it in group)
        if len(ref_sheets) < 2:
            # Mesma prancha — não é cross-prancha (provavelmente passadas anteriores
            # já trataram). Manter como está.
            pass6.extend(group)
            continue

        # Verificar similaridade de descrições (overlap >= 2 tokens)
        descs_keys = [set(_normalize_description_key(it.description or "").split()) for it in group]
        # Pra todos do grupo, ver se há overlap forte entre TODOS pares
        # (alternativa simplificada: ver se o primary_noun é o mesmo + alguma palavra
        # significativa em comum entre todos)
        common_tokens = set.intersection(*descs_keys) if descs_keys else set()
        sig_common = {t for t in common_tokens if len(t) > 3}  # ignora "de", "e", etc

        if not sig_common:
            # Mesmo noun mas sem outras palavras significativas em comum — provavelmente
            # itens distintos (ex: "porta" em 2 pranchas diferentes = 2 portas diferentes)
            pass6.extend(group)
            continue

        # Confirmada duplicação cross-prancha. Manter o item de MAIOR especificidade.
        winner = max(group, key=_item_score)
        losers = [it for it in group if id(it) != id(winner)]
        total_pranchas = len(ref_sheets)

        # Anota no obs do vencedor que dedupliquei
        loser_summary = "; ".join(
            f"{it.quantity} {it.unit} ({it.ref_sheet[:30] if it.ref_sheet else 'sem ref'})"
            for it in losers
        )
        merged_obs = (winner.observations or "").rstrip(". ") + (". " if winner.observations else "")
        merged_obs += (
            f"Deduplicado: este item aparecia em {total_pranchas} pranchas diferentes "
            f"do mesmo projeto. Versões descartadas: {loser_summary}. "
            f"Versão mantida tem maior especificidade (unidade física, qty>0)."
        )
        merged_item = BudgetItem(
            item_num=winner.item_num,
            description=winner.description,
            unit=winner.unit,
            quantity=winner.quantity,
            observations=merged_obs,
            ref_sheet=winner.ref_sheet,
            confidence=winner.confidence,
            discipline=winner.discipline,
        )
        pass6.append(merged_item)

    return pass6


# ═══════════════════════════════════════════════════════════════════════════
#  REGRAS PÓS-CONSOLIDAÇÃO (fixes aplicados em 2026-05-17 após auditoria
#  do projeto do Yuri / job b82a72ed):
#
#  🅐 Detector multifamiliar — flag quando contar 5+ bacias OU porta-elevador
#     OU porta corta-fogo, sinaliza no warning do projeto pra cliente decidir.
#  🅑 Dedup por layer m² — se 2+ itens m² compartilham o mesmo layer fonte,
#     mantém só o de maior qty (que representa o total do layer) e marca os
#     outros como estimado + warning de possível sobreposição.
#  🅒 "Conforme projeto" / "a definir" → estimado obrigatório (regra dura #1).
#  🅓 "Fundido"/"Consolidado de N entradas" / "(várias variantes)" → estimado
#     obrigatório quando ainda está confirmado.
# ═══════════════════════════════════════════════════════════════════════════

_LAYER_RE = _re.compile(r"layer\s+'?([A-Z][A-Z0-9_-]+)'?", _re.IGNORECASE)


def _extract_layer_from_obs(obs: str) -> str:
    """Extrai o nome do layer CAD da observação (ex.: 'A-FLOR-PATT'). Vazio se não achar."""
    if not obs:
        return ""
    m = _LAYER_RE.search(obs)
    return m.group(1).upper() if m else ""


def _detect_multifamiliar_signal(items: list, current_typology: str) -> tuple[bool, list]:
    """Retorna (is_multifamiliar_provavel, lista_de_evidencias).

    Triggers — qualquer um basta:
      - 5+ bacias sanitárias (residência unifamiliar tem 2-4 max)
      - Porta de elevador OU bloco 'elevador' detectado
      - 3+ portas corta-fogo (sinal de escada de incêndio = prédio)
      - 6+ portas-janela de correr OU 15+ janelas (multi-pavimento)
    """
    from collections import Counter
    if (current_typology or "").lower() != "residential":
        return False, []  # Só dispara o aviso pra residential — outras já não confundem

    bacia_count = 0
    porta_corta_fogo_count = 0
    elevador_detected = False
    janela_count = 0
    porta_janela_count = 0

    for it in items:
        desc = (it.description or "").lower()
        obs = (it.observations or "").lower()
        try:
            qty = float(it.quantity or 0)
        except Exception:
            qty = 0.0

        if "bacia sanit" in desc or "bacia sanit" in obs:
            bacia_count += qty
        if "elevador" in desc or "elevador" in obs:
            elevador_detected = True
        if "corta-fogo" in desc or "corta fogo" in desc or ("rf-60" in desc or "rf 60" in desc):
            porta_corta_fogo_count += qty
        if "janela" in desc and "porta" not in desc:
            janela_count += qty
        if "porta-janela" in desc or "porta janela" in desc:
            porta_janela_count += qty

    evidencias = []
    if bacia_count >= 5:
        evidencias.append(f"{int(bacia_count)} bacias sanitárias detectadas")
    if elevador_detected:
        evidencias.append("elevador detectado nas plantas")
    if porta_corta_fogo_count >= 3:
        evidencias.append(f"{int(porta_corta_fogo_count)} portas corta-fogo (escada de incêndio)")
    if porta_janela_count >= 6:
        evidencias.append(f"{int(porta_janela_count)} portas-janela de correr")
    if janela_count >= 15:
        evidencias.append(f"{int(janela_count)} janelas")

    return (len(evidencias) >= 1, evidencias)


def _apply_post_consolidation_rules(items: list) -> tuple[list, int]:
    """Aplica regras 🅑+🅒+🅓 nos items pós-consolidate.

    Retorna (items_modificados, n_alterados).
    Mutate os items in-place (mais simples — preserva referências).
    """
    from models import Confidence
    n_changed = 0

    # ── 🅒 + 🅓: forçar estimado quando vago ou somatório ──
    VAGUE_MARKERS = (
        "conforme projeto", "conforme especifica", "especificação a definir",
        "a definir", "por definir", "conforme detalhe",
    )
    FUSED_MARKERS = (
        "fundido de", "consolidado de", "(várias variantes)", "(varias variantes)",
        "várias variações", "varias variacoes", "(várias variações)",
        "contagem de blocos",  # quando obs diz "contagem de blocos X (várias variações)"
    )

    for it in items:
        if not it:
            continue
        desc = (it.description or "").lower()
        obs = (it.observations or "").lower()
        text = f"{desc} {obs}"

        is_vague = any(m in text for m in VAGUE_MARKERS)
        is_fused = any(m in text for m in FUSED_MARKERS)

        try:
            current_conf = str(it.confidence) if it.confidence else "estimado"
        except Exception:
            current_conf = "estimado"

        if current_conf == "confirmado" and (is_vague or is_fused):
            try:
                it.confidence = Confidence("estimado")
            except Exception:
                pass
            tag = []
            if is_vague:
                tag.append("descrição genérica")
            if is_fused:
                tag.append("agregação de variantes")
            it.observations = (
                (it.observations or "") +
                f" | ⚠ Estimado: {', '.join(tag)} — confirmar contra projeto"
            ).strip(" |")
            n_changed += 1

    # ── 🅑 Dedup por layer m² ──
    by_layer: dict[str, list] = {}
    for it in items:
        if (it.unit or "").lower() not in ("m²", "m2", "m³", "m3"):
            continue
        layer = _extract_layer_from_obs(it.observations or "")
        if not layer:
            continue
        by_layer.setdefault(layer, []).append(it)

    for layer, group in by_layer.items():
        if len(group) < 2:
            continue
        # Mantém o de maior qty como "principal"; os outros viram estimado com warning
        sorted_g = sorted(group, key=lambda x: float(x.quantity or 0), reverse=True)
        winner = sorted_g[0]
        others = sorted_g[1:]
        other_summary = "; ".join(
            f"{round(float(it.quantity or 0), 2)} {it.unit}" for it in others
        )
        for it in others:
            try:
                it.confidence = Confidence("estimado")
            except Exception:
                pass
            it.observations = (
                (it.observations or "") +
                f" | ⚠ Possível sobreposição: outro item ({round(float(winner.quantity or 0),2)} {winner.unit}) "
                f"compartilha o mesmo layer '{layer}'. Verifique se não há duplicação."
            ).strip(" |")
            n_changed += 1
        # Marca o winner também com aviso (mais brando — informativo)
        winner.observations = (
            (winner.observations or "") +
            f" | ⚠ {len(others)} outros itens m² do mesmo layer '{layer}' marcados como estimado (possível sobreposição: {other_summary})"
        ).strip(" |")

    return items, n_changed


# Ranges plausíveis por unidade — valores fora disso indicam erro provável.
# Valores em contexto de reforma de escritório corporativo (nosso nicho atual).
_PLAUSIBILITY_RANGES = {
    "un": (0, 5000),   # 5000+ un em uma obra é raro (ex.: difusores AC em torre)
    "ml": (0, 50000),  # 50km de perfil é improvável; mas casos grandes permitidos
    "m²": (0, 50000),  # 50000m² = escritório de 5 andares grandes
    "m": (0, 50000),
    "mês": (0, 60),    # 5 anos de obra é improvável
    "dia": (0, 1000),  # 3 anos em dias
    "%": (0, 100),     # percentual
    "vb": (0, 1),      # verba por definição é 0 ou 1
}
# Combinações disciplina × unidade que não fazem sentido.
# (disciplina, unidade) -> mensagem descritiva do problema
_DISCIPLINE_UNIT_MISMATCHES = {
    ("Iluminação", "m²"): "luminárias se contam em 'un', não em área",
    ("Iluminação", "m"):  "luminárias se contam em 'un', não em metros",
    ("Pisos e Rodapés", "un"): "piso é superfície, deve ser m² — exceção: elementos pontuais como grelhas",
    ("Forros", "un"): "forro é superfície, deve ser m²",
    ("Ar-Condicionado", "m²"): "equipamentos de AC são unidades, não área",
    ("Incêndio e Segurança", "m²"): "sprinklers/detectores são unidades",
    ("Portas e Ferragens", "m²"): "portas são unidades, não área",
}


def _check_plausibility(item, project_total_area_m2: float = 0) -> tuple[bool, str]:
    """Retorna (é_plausível, motivo_se_não). Só avalia — decisão é do caller."""
    if item is None:
        return True, ""
    try:
        qty = float(item.quantity or 0)
    except Exception:
        return False, "quantidade não-numérica"

    # 1. Range plausível pela unidade
    unit = item.unit or "vb"
    if unit in _PLAUSIBILITY_RANGES:
        lo, hi = _PLAUSIBILITY_RANGES[unit]
        if qty > hi:
            return False, f"{qty:.0f} {unit} é alto demais pro tipo (max típico {hi})"
        if qty < lo:
            return False, f"{qty} {unit} é negativo"

    # 2. Disciplina × unidade
    disc = item.discipline or ""
    mismatch_key = (disc, unit)
    if mismatch_key in _DISCIPLINE_UNIT_MISMATCHES:
        return False, _DISCIPLINE_UNIT_MISMATCHES[mismatch_key]

    # 3. Área vs laje (só pra superfícies)
    if unit == "m²" and project_total_area_m2 > 0 and qty > project_total_area_m2 * 1.5:
        return False, (
            f"área {qty:.1f} m² é maior que 1.5× área da laje "
            f"({project_total_area_m2:.0f} m²) — possível dupla contagem"
        )

    return True, ""


def _validate_quantity_for_unit(item) -> tuple[float, bool]:
    """Garante consistência entre unidade e quantidade.
    - 'un' só aceita inteiros positivos (arredonda se frac, ou zera + marca estimado)
    - 'ml' / 'm²' aceita qualquer número >= 0
    Retorna (qty_ajustada, foi_ajustada)"""
    qty = float(item.quantity) if item.quantity is not None else 0
    if item.unit == "un":
        if qty != int(qty):
            # un com decimal é suspeito (ex.: IA confunde um comprimento
            # em ml com contagem e devolve un=NNN.NN).
            # Se for "quase inteiro" (ex.: 9.0001), arredonda. Senão zera.
            if abs(qty - round(qty)) < 0.01:
                return float(round(qty)), True
            # Valor claramente não é contagem — descartar e marcar estimado
            return 0.0, True
    return qty, False


def _normalize_unit_for_item(description: str, current_unit: str) -> tuple[str, bool]:
    """Ajusta a unidade baseada na descrição do item.
    Retorna (unidade_nova, foi_corrigida).
    - Item com palavra-chave de SUPERFÍCIE (piso/forro/pintura) → m²
    - Item com palavra-chave LINEAR (rodapé/perfil/eletrocalha) → ml
    - Item com palavra-chave CONTÁVEL (luminária/porta) → un
    - Senão, mantém a unidade atual"""
    if not description:
        return current_unit, False

    # Vb / % / mês / dia são especiais — não corrigir
    if current_unit in ("vb", "%", "mês", "mes", "dia", "h"):
        return current_unit, False

    # Decide pela IDENTIDADE do item (o NOME), não pelo contexto entre parênteses.
    # Ex.: "MONITOR/TV 42\" (…visíveis na planta de forro)" NÃO é m² só porque a
    # observação cita "forro" — ali "forro" é LOCALIZAÇÃO, não o tipo do item.
    # Casa as palavras-chave só no trecho antes do 1º "(" (caso Roberta 23/07).
    head_lower = description.split("(", 1)[0].lower()

    # Ordem de precedência: contável > linear > superfície (senão piso vira superfície erroneamente)
    if _UNIT_COUNT_KEYWORDS.search(head_lower):
        if current_unit != "un":
            return "un", True
        return "un", False
    if _UNIT_LINEAR_KEYWORDS.search(head_lower):
        if current_unit != "ml":
            return "ml", True
        return "ml", False
    if _UNIT_SURFACE_KEYWORDS.search(head_lower):
        if current_unit != "m²":
            return "m²", True
        return "m²", False
    return current_unit, False


def _item_geo_category(description: str, discipline: str = "") -> str | None:
    """Mapeia um ITEM (descrição+disciplina) para UMA categoria física geométrica
    {piso, forro, pintura, paredes, demolicao} ou None.

    Usado SÓ pelo cross-check determinístico: casar o número da IA contra a medida
    da PRÓPRIA categoria do item (nunca confirmar forro com área de piso — furo C da
    revisão 15/07). CONSERVADOR de propósito: só retorna categoria quando a descrição
    tem a PALAVRA da categoria (piso/forro/parede/pintura/demolição). Material ambíguo
    ('cerâmica', 'porcelanato' — pode ser piso OU parede) retorna None → não promove.
    Recall menor, zero classificação errada (regra nº1)."""
    t = f"{description} {discipline}".lower()
    # Ordem: mais específico primeiro (demolição e forro antes de piso/parede).
    if "demoli" in t or "demolir" in t:
        return "demolicao"
    if "forro" in t or "sanca" in t:
        return "forro"
    if "pintura" in t or "pintar" in t:
        return "pintura"
    if "contrapiso" in t or "piso" in t or "assoalho" in t:
        return "piso"
    if ("parede" in t or "alvenaria" in t or "divisória" in t
            or "divisoria" in t or "drywall" in t):
        return "paredes"
    return None


# ══════════════════════════════════════════════════════════════════
#  🛡️ FREIO DE MEMÓRIA (fix 2026-07-22) — garante que o servidor NÃO CAI
#  por causa de UM projeto pesado. Lê a RAM real do container (cgroup) e
#  aborta o job de forma limpa ANTES do OOM killer, mantendo o serviço de
#  pé pra todos os outros clientes. Confiabilidade = cliente volta.
# ══════════════════════════════════════════════════════════════════
def _container_mem_frac():
    """Fração de RAM usada do CONTAINER (cgroup), 0..1, ou None se não medir.
    Lê o cgroup (v2 e v1) — reflete o limite REAL do container (4GB no Render),
    ao contrário de psutil.virtual_memory(), que lê o HOST inteiro e mente."""
    for used_p, lim_p in (
        ("/sys/fs/cgroup/memory.current", "/sys/fs/cgroup/memory.max"),
        ("/sys/fs/cgroup/memory/memory.usage_in_bytes",
         "/sys/fs/cgroup/memory/memory.limit_in_bytes"),
    ):
        try:
            with open(used_p) as _uf:
                used = int(_uf.read().strip())
            with open(lim_p) as _lf:
                raw = _lf.read().strip()
            if raw == "max":
                continue
            limit = int(raw)
            if 0 < limit < 200 * 1024 ** 3:   # ignora "sem limite" (número gigante)
                return used / limit
        except Exception:
            continue
    return None


def _mem_pressure(threshold: float = 0.85) -> bool:
    """True se o container passou de `threshold` do limite de RAM. Se não der pra
    medir (dev/local sem cgroup), retorna False — não bloqueia."""
    frac = _container_mem_frac()
    return frac is not None and frac >= threshold


def _abort_job_mem(job_id: str, done: int, total: int):
    """Aborta um job por pressão de memória, de forma LIMPA — mantém o servidor de
    pé pra todos os outros. Marca erro com orientação de dividir em lotes + alerta."""
    _msg = ("Seu projeto é grande demais pra processar de uma vez e chegou perto do "
            "limite de memória do servidor. Divida em 2-3 envios menores (ex.: "
            "metade das pranchas por vez) que processa tranquilo.")
    try:
        jobs.update_field(job_id, status="error")
        jobs.update_field(job_id, error_message=_msg)
    except Exception:
        pass
    try:
        _supabase_update("projects", "job_id", job_id, {
            "status": "error", "error_message": _msg,
            "completed_at": datetime.utcnow().isoformat(),
        })
    except Exception:
        pass
    try:
        _log_error("mem:freio",
                   f"Freio de memória: abortado em {done}/{total} pranchas "
                   f"(uso={_container_mem_frac()})", job_id, severity="error")
    except Exception:
        pass
    try:
        import threading as _thm
        _thm.Thread(target=_notify_admin, args=(
            "🛡️ Freio de memória disparou",
            f"O job <b>{job_id}</b> chegou perto do limite de RAM em {done}/{total} "
            f"pranchas e foi abortado ANTES de derrubar o servidor. O cliente foi "
            f"orientado a dividir em lotes menores. (Servidor seguiu de pé.)"),
            daemon=True).start()
    except Exception:
        pass
    try:
        _email_falha_cliente(job_id, reprocessavel=False)
    except Exception:
        pass


def _extract_dxf_inprocess(dxf_path, unit_consensus):
    """Extração IN-PROCESS (fallback): emagrece + extrai + gera o prompt. Mesmo
    resultado do worker isolado, mas roda no processo do servidor (sem teto de RAM)."""
    p = dxf_path
    try:
        from dxf_slim import emagrecer_dxf_se_preciso
        _s = emagrecer_dxf_se_preciso(p)
        if _s:
            p = _s
    except Exception as _sl_e:
        print(f"[dxf-slim] pulando ({_sl_e})")
    from dwg_extractor import extract_from_file
    _ext = extract_from_file(p, unit_factor_override=unit_consensus)
    return _ext, _ext.to_structured_prompt(), p


def _extract_dxf_isolated(dxf_path, unit_consensus, timeout_s=900):
    """Extrai UMA prancha DXF num SUBPROCESSO matável com teto de memória (fix
    confiabilidade 2026-07-22). Uma prancha densa que estouraria a RAM mata só o
    filho — o servidor fica de pé. Retorna (extraction, structured_text, path).

    - timeout → concurrent.futures.TimeoutError (o loop já trata: pula a prancha)
    - filho morto por RAM / erro de parse → RuntimeError (o loop trata: pula)
    - NÃO deu pra lançar o subprocesso (local/ambiente restrito) → fallback in-process
    """
    import concurrent.futures as _cf
    import subprocess
    import sys
    import tempfile
    import pickle as _pk
    _worker = os.path.join(os.path.dirname(os.path.abspath(__file__)), "dxf_extract_worker.py")
    if not os.path.exists(_worker):
        return _extract_dxf_inprocess(dxf_path, unit_consensus)   # sem worker → in-process
    _fd, _out = tempfile.mkstemp(suffix=".pkl", prefix="dxfext_")
    os.close(_fd)
    # Teto de memória do filho (POSIX). Prancha legítima cabe; densa patológica morre
    # aqui em vez de derrubar o container. RLIMIT_AS é VIRTUAL (> RSS real), então o
    # padrão é folgado (2,5GB) pra não matar prancha grande boa por engano — pai ~1GB
    # + filho 2,5GB = 3,5GB < 4GB do container. Env DXF_EXTRACT_MEM_MB ajusta: se
    # prancha legítima começar a falhar, sobe; se ainda estourar o container, desce.
    _preexec = None
    try:
        import resource as _res
        _cap = int(os.environ.get("DXF_EXTRACT_MEM_MB", "2560")) * 1024 * 1024
        def _cap_mem():
            _res.setrlimit(_res.RLIMIT_AS, (_cap, _cap))
        _preexec = _cap_mem
    except Exception:
        _preexec = None
    try:
        _kw = dict(cwd=os.path.dirname(os.path.abspath(__file__)),
                   timeout=timeout_s, capture_output=True)
        if _preexec is not None:
            _kw["preexec_fn"] = _preexec
        _proc = subprocess.run(
            [sys.executable, _worker, dxf_path, _out,
             ("" if unit_consensus is None else str(unit_consensus))], **_kw)
    except subprocess.TimeoutExpired:
        try: os.remove(_out)
        except OSError: pass
        raise _cf.TimeoutError()          # subprocess.run já matou o filho
    except Exception as _le:
        # Não deu pra LANÇAR o subprocesso (ex.: local sem 'resource'/preexec) →
        # cai pro in-process (mesmo resultado, sem isolamento). NÃO é o caso OOM.
        try: os.remove(_out)
        except OSError: pass
        print(f"[dxf-isolado] subprocesso não lançou ({type(_le).__name__}: {_le}) — fallback in-process")
        return _extract_dxf_inprocess(dxf_path, unit_consensus)
    if _proc.returncode != 0:
        _err = (_proc.stderr or b"").decode("utf-8", "replace")[-600:]
        try: os.remove(_out)
        except OSError: pass
        # rc != 0 = filho morto pelo teto de RAM OU erro de parse → NÃO faz fallback
        # (re-rodar in-process uma prancha que estourou a RAM derrubaria o servidor).
        raise RuntimeError(f"extração isolada falhou (rc={_proc.returncode}): {_err}")
    try:
        with open(_out, "rb") as _f:
            _result = _pk.load(_f)
    finally:
        try: os.remove(_out)
        except OSError: pass
    return _result


def _measure_unambiguous(value: float, cat: str, by_cat: dict) -> bool:
    """True se `value` (medida arredondada) aparece SÓ na categoria `cat` dentro de
    `by_cat` — não colide com nenhuma outra categoria do MESMO tipo de medida.

    Fecha o furo C: forro≈piso (mesma área). Se a área bate tanto em piso quanto em
    forro, o número é AMBÍGUO → não promove nenhum dos dois. Só confirma quando a
    geometria aponta inequivocamente pra categoria do item."""
    for _c, _s in by_cat.items():
        if _c != cat and value in _s:
            return False
    return True


# ─── Throttle de processamento concorrente ──────────────────────────
# Cada upload dispara process_job() numa thread daemon. SEM limite, 2-3
# projetos processando ao mesmo tempo somam picos de RAM (render PDF +
# Vision + listas de itens + XLSX) e estouram os 2GB do Render → OOM +
# restart → derruba quem está processando. Foi o que pegou o Adriano em
# 16/06/2026. O semáforo força 1 job por vez: quem chega depois espera na
# fila em vez de competir por memória. Reduz o pico de ~2GB pra 400-600MB.
import threading as _threading_sem
_JOB_SEMAPHORE = _threading_sem.Semaphore(1)


def _process_job_throttled(*args, **kwargs):
    """Wrapper que serializa process_job via semáforo (1 por vez)."""
    with _JOB_SEMAPHORE:
        process_job(*args, **kwargs)


def _job_medidos_count(job_id: str) -> int:
    """Quantos itens MEDIDOS (confidence=confirmado) este job já tem salvos.

    Retorna -1 se não deu pra saber (soluço do banco). Usado pela trava
    anti-perda do complemento: refazer um projeto que já mediu, usando só PDF,
    trocaria medição por estimativa — proibido pela regra dura nº1.
    """
    try:
        import urllib.request as _urm
        _q = (f"{SUPABASE_URL}/rest/v1/project_items?job_id=eq.{job_id}"
              f"&confidence=eq.confirmado&select=id")
        _r = _urm.Request(_q, method="GET")
        _r.add_header("apikey", SUPABASE_KEY)
        _r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _resp = _urm.urlopen(_r, timeout=10)
        return len(_json.loads(_resp.read().decode("utf-8")))
    except Exception as _e:
        print(f"[add-file] contagem de medidos falhou: {_e}")
        return -1


def _complement_base_has_items(job_id: str) -> bool:
    """Este job já tem itens salvos (= planilha anterior viva)?

    DEFAULT SEGURO (board 15/07): um complemento SEMPRE existe sobre um projeto-base.
    Se a contagem falhar por soluço do banco (timeout/5xx do Supabase, Render sob
    carga), NÃO pode concluir "base vazia" e derrubar a planilha. Assume base=existe
    por padrão; só devolve False se a query SUCEDER e vier 0. 2 tentativas.
    """
    for _try in range(2):
        try:
            import urllib.request as _urc
            _cq = f"{SUPABASE_URL}/rest/v1/project_items?job_id=eq.{job_id}&select=id&limit=1"
            _cr = _urc.Request(_cq, method="GET")
            _cr.add_header("apikey", SUPABASE_KEY)
            _cr.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            return len(_json.loads(_urc.urlopen(_cr, timeout=10).read().decode("utf-8"))) > 0
        except Exception as _cerr:
            print(f"[add-file] contagem de itens base falhou (tentativa {_try+1}, assumindo base existe): {_cerr}")
    return True


def _salvage_layout_esquadrias(file_paths):
    """Última cartada antes de declarar "0 itens" numa planta de LAYOUT vetorial:
    lê as ESQUADRIAS das cotas escritas no texto do PDF (largura×altura/peitoril,
    ex. "160x150/86"). 100% DETERMINÍSTICO — não chama IA, não inventa dimensão,
    só conta o que está escrito na prancha. Caso Catarina (20/07): estudo de layout
    de interiores sem quadro de áreas dava erro "nenhum item"; agora entrega as
    esquadrias medidas na prancha, marcadas 'estimado' (é leitura de layout, não
    medição geométrica nossa — regra dura nº1). Só PDF; retorna [] se nada casar.

    Nota: BudgetItem/Confidence são importados LOCALMENTE (o main.py não os traz
    no topo — só ProcessingStatus); por isso a assinatura não os anota."""
    from models import BudgetItem, Confidence
    import collections, re as _re_esq
    # Cota de esquadria no texto vetorial: "160x150/86" = largura×altura/peitoril (cm).
    # A forma COM peitoril (3 números) é notação inconfundível de esquadria — evita
    # falso-positivo com tamanho de piso ("60x60", que não tem "/"). Conservador.
    _rx = _re_esq.compile(r'^(\d{2,3})x(\d{2,3})/(\d{2,3})$')
    counts: "collections.Counter" = collections.Counter()
    for p in file_paths or []:
        if not str(p).lower().endswith(".pdf"):
            continue
        try:
            import pdfplumber
            with pdfplumber.open(p) as pdf:
                for page in pdf.pages:
                    for w in page.extract_words():
                        m = _rx.match((w.get("text") or "").replace(" ", ""))
                        if m:
                            counts[(m.group(1), m.group(2), m.group(3))] += 1
        except Exception as e:
            print(f"[salvage-esq] {os.path.basename(str(p))}: {e}")
            continue
    items = []
    _ref = os.path.basename(str(file_paths[0])) if file_paths else ""
    for (L, A, P), n in sorted(counts.items(), key=lambda kv: (-kv[1], kv[0])):
        Lm, Am, Pm = int(L) / 100.0, int(A) / 100.0, int(P) / 100.0
        desc = (f"Esquadria (janela/porta) {L}×{A}/{P} — {Lm:.2f} × {Am:.2f} m, "
                f"peitoril {Pm:.2f} m")
        items.append(BudgetItem(
            item_num="", description=desc, unit="un", quantity=float(n),
            observations=("Dimensão lida da cota escrita na prancha (estudo de layout — "
                          "confira o tipo e a quantidade; portas sem cota não entram)."),
            ref_sheet=_ref, confidence=Confidence.ESTIMADO,
            discipline="Esquadrias", origem="vision_pdf"))
    return items


_LAYOUT_COUNT_DISCIPLINES = {
    "Louças e Metais", "Portas e Ferragens", "Marcenaria", "Mobiliário",
    "Iluminação",
}
# Palavras que denunciam item de ÁREA/superfície (proibido no salvamento — regra
# dura nº1: sem m² fingido). Se a descrição bate, descarta o item.
_LAYOUT_COUNT_BLOCK = (
    "piso", "forro", "parede", "pintura", "revestiment", "rodap", "m²", "m2",
    "azulej", "porcelanat", "laje", "reboco", "contrapiso", "área", "area",
    "metro quadrado",
)

_LAYOUT_COUNT_SYSTEM = (
    "Você é um orçamentista de interiores experiente. A imagem é uma planta de "
    "LAYOUT (arrumação de móveis), sem quadro de áreas e sem cotas de dimensão "
    "além das esquadrias. Sua tarefa é LISTAR e CONTAR, por tipo, APENAS os itens "
    "CONTÁVEIS que você vê com CLAREZA no desenho.\n\n"
    "CATEGORIAS (use exatamente estes nomes em 'discipline'):\n"
    "- 'Louças e Metais': vasos sanitários/bacias, cubas/lavatórios, box/chuveiros, "
    "pia de cozinha, tanque, cooktop, torneiras.\n"
    "- 'Portas e Ferragens': portas internas (conte pelos arcos de abertura). NÃO "
    "conte janelas — elas já foram contadas à parte.\n"
    "- 'Marcenaria': roupeiros/armários embutidos, bancadas, ilha de cozinha, "
    "painel de TV/rack, aparador.\n"
    "- 'Mobiliário': camas, sofás, mesas, cadeiras, poltronas, banquetas.\n\n"
    "REGRAS DURAS (obrigatórias):\n"
    "1. NÃO invente. Se não tem certeza do tipo OU da quantidade, NÃO liste (é "
    "melhor faltar item do que errar número).\n"
    "2. É PROIBIDO produzir qualquer item de ÁREA em m² (piso, forro, parede, "
    "pintura, revestimento). Só contagem de peças (unidade 'un').\n"
    "3. Não repita esquadrias/janelas.\n\n"
    "Responda SOMENTE com JSON puro, sem texto fora dele, no formato:\n"
    '{"items":[{"description":"...","quantity":N,"discipline":"..."}]}\n'
    "quantity é um inteiro. Se não vê nada contável com clareza, retorne "
    '{"items":[]}.'
)


def _salvage_layout_ai_counts(client, file_paths, crops_dir):
    """Complemento do salvamento de layout: além das esquadrias (texto), pede pra
    IA CONTAR louças/metais, portas, marcenaria e mobiliário que ela vê na planta
    renderizada. ADITIVO e só roda no caminho de salvamento (0 itens da IA num
    layout vetorial) — nunca toca projeto que já funcionava. Tudo marcado
    'estimado' (o cliente confirma) e a IA é instruída a NÃO inventar nem gerar
    m². Qualquer falha (render, API, JSON) → [] (mantém só as esquadrias).

    ⚠ Sem verificação adversarial (ao contrário do estudo one-off): pode listar
    item a mais/menos. Aceitável porque só aparece onde HOJE daria erro seco, e
    todo item entra como 'estimado' pra conferência."""
    try:
        from processor import render_crops
        from analyzer import encode_image
        from llm_retry import call_with_retry_stream
        from models import BudgetItem, Confidence, SheetType
        import json as _json_lc, os as _os_lc
    except Exception as _imp_e:
        print(f"[salvage-ai] import falhou: {_imp_e}")
        return []
    _pdf = next((p for p in (file_paths or []) if str(p).lower().endswith(".pdf")), None)
    if not _pdf:
        return []
    try:
        crops = render_crops(_pdf, SheetType.ARQUITETURA, crops_dir, dpi=140,
                             page_index=0, out_stem="_salvage_layout")
    except Exception as _re_e:
        print(f"[salvage-ai] render falhou: {_re_e}")
        return []
    content = []
    for cp in (crops or [])[:3]:
        try:
            if _os_lc.path.exists(cp) and _os_lc.path.getsize(cp) <= 500_000:
                content.append({"type": "image", "source": {
                    "type": "base64",
                    "media_type": "image/jpeg" if cp.endswith(".jpg") else "image/png",
                    "data": encode_image(cp)}})
        except Exception:
            continue
    if not content:
        return []
    content.append({"type": "text", "text":
                    "Conte os itens desta planta de layout conforme as regras."})
    try:
        resp = call_with_retry_stream(
            client, tag="salvage-layout-counts", model="claude-sonnet-4-6",
            max_tokens=4000, temperature=0, system=_LAYOUT_COUNT_SYSTEM,
            messages=[{"role": "user", "content": content}])
        txt = resp.content[0].text
        if "```json" in txt:
            txt = txt.split("```json")[1].split("```")[0]
        elif "```" in txt:
            txt = txt.split("```")[1].split("```")[0]
        data = _json_lc.loads(txt.strip())
    except Exception as _ai_e:
        print(f"[salvage-ai] IA/JSON falhou: {_ai_e}")
        return []
    _ref = _os_lc.path.basename(str(_pdf))
    out = []
    for it in (data.get("items") or [])[:60]:
        try:
            desc = str(it.get("description", "")).strip()
            if len(desc) < 3:
                continue
            # Descarta item de ÁREA/superfície disfarçado de contagem (a IA às
            # vezes devolve "Piso porcelanato: 80 un" — proibido aqui).
            _dl = desc.lower()
            if any(b in _dl for b in _LAYOUT_COUNT_BLOCK):
                continue
            disc = str(it.get("discipline", "")).strip()
            # Fora das categorias de CONTAGEM → descarta (não remapeia pra genérico,
            # pra não virar depósito de item duvidoso).
            if disc not in _LAYOUT_COUNT_DISCIPLINES:
                continue
            try:
                qty = float(it.get("quantity", 0) or 0)
            except (TypeError, ValueError):
                qty = 0
            if qty <= 0 or qty > 500:
                continue
            out.append(BudgetItem(
                item_num="", description=desc, unit="un", quantity=qty,
                observations=("Contagem visual do layout pela IA — CONFIRA o tipo e a "
                              "quantidade (leitura de desenho, não medição)."),
                ref_sheet=_ref, confidence=Confidence.ESTIMADO,
                discipline=disc, origem="vision_pdf"))
        except Exception:
            continue
    return out


# ── Honestidade de ÁREA (regra dura nº1) ──────────────────────────────────────
# As unidades/keywords e is_floor_surface() vivem em engine_rules.py (regra
# determinística, testável sem IA — tests/test_engine_rules.py). Importados no topo
# como _AREA_UNITS_HONESTY / _FLOOR_M2_UNITS / _is_floor_surface.
def _derive_pintura_pe_direito(items, pe_direito: float) -> int:
    """Deriva PINTURA de parede quando só temos o COMPRIMENTO das paredes.

    Motivação (medido em 01/08/2026): em 22 dos 69 projetos com parede, a
    parede veio só em METRO LINEAR — e em TODOS os 22 a pintura ficou de fora
    da planilha. Reclamação real de cliente no chat ("não tem o quantitativo
    em m² da pintura"). Com o pé-direito INFORMADO no upload dá pra fechar:
    área = Σ(comprimento) × pé-direito × 2 faces.

    Regras duras respeitadas:
    - Só roda se o cliente informou o pé-direito (nunca assume 2,70).
    - Só cria pintura se NÃO existe nenhum item de pintura (nunca sobrescreve).
    - Sai como ESTIMADO com rótulo "pé-direito informado por você" (regra nº1).
    - Não desconta vãos (não os conhecemos aqui) — e diz isso na observação.
    Devolve 1 se criou o item, 0 caso contrário."""
    if not pe_direito or pe_direito <= 0:
        return 0
    from models import BudgetItem, Confidence
    tem_pintura = any(
        ("pintura" in (getattr(i, "description", "") or "").lower()
         or "látex" in (getattr(i, "description", "") or "").lower()
         or "latex" in (getattr(i, "description", "") or "").lower())
        for i in items)
    if tem_pintura:
        return 0
    total_m = 0.0
    ref = ""
    for i in items:
        d = (getattr(i, "description", "") or "").lower()
        u = (getattr(i, "unit", "") or "").strip().lower()
        if u in ("m", "ml") and ("parede" in d or "alvenaria" in d or "drywall" in d):
            try:
                q = float(getattr(i, "quantity", 0) or 0)
            except (TypeError, ValueError):
                continue
            if q > 0:
                total_m += q
                ref = ref or (getattr(i, "ref_sheet", "") or "")
    if total_m <= 0:
        return 0
    area = round(total_m * float(pe_direito) * 2.0, 1)   # 2 faces
    items.append(BudgetItem(
        item_num="PD.1",
        description=("Pintura látex sobre paredes internas — derivada do "
                     "comprimento de paredes × pé-direito informado"),
        unit="m²",
        quantity=area,
        observations=(f"⚠ ESTIMADO — {total_m:.1f} m de parede × pé-direito "
                      f"{pe_direito:.2f} m informado por você × 2 faces. Vãos de "
                      f"portas/janelas NÃO descontados — confira antes de orçar."),
        ref_sheet=ref,
        confidence=Confidence.ESTIMADO,
        discipline="Revestimentos",
    ))
    return 1


def _apply_area_honesty(items, total_area: float = 0, total_area_source: str = "") -> tuple[int, int]:
    """Aplica a regra dura nº1 aos itens de ÁREA que NÃO vieram da geometria do CAD:

    - Se o cliente INFORMOU a área (total_area_source='informado') e o item é uma
      SUPERFÍCIE HORIZONTAL (piso/forro/laje), PREENCHE com a área informada como
      ESTIMADO (laranja), rotulado 'informado por você — não medido'. É o que o
      cliente pediu ao informar a metragem: completar os itens de área com uma base
      honesta (não é medição nossa, segue como estimativa a conferir).
    - Caso contrário, ZERA a quantidade (Vision não mede geometria → evita m²
      inventado, caso Catarina 20/07) e anota que a área não foi medida.

    Contagens (un), verbas (vb) e itens medidos do CAD (origem 'dxf_geom') ficam
    intocados. Devolve (n_preenchidos, n_zerados)."""
    from models import Confidence
    informado = (total_area_source == "informado") and (total_area or 0) > 0
    filled = blanked = 0
    for it in items:
        if getattr(it, "origem", "") == "dxf_geom":
            continue
        u = (getattr(it, "unit", "") or "").strip().lower()
        if u not in _AREA_UNITS_HONESTY:
            continue
        try:
            q = float(getattr(it, "quantity", 0) or 0)
        except (TypeError, ValueError):
            q = 0.0
        if informado and u in _FLOOR_M2_UNITS and _is_floor_surface(getattr(it, "description", "")):
            it.quantity = round(float(total_area), 2)
            try:
                it.confidence = Confidence("estimado")
            except Exception:
                pass
            # tira aviso antigo "área NÃO medida" e põe o rótulo honesto da área informada
            _segs = [s.strip() for s in (it.observations or "").split("|")
                     if "não medida" not in s.lower() and "nao medida" not in s.lower()]
            _segs.append("Área informada por você (não medida): assumido = área total do "
                         "projeto. Confira antes de orçar.")
            it.observations = " | ".join(s for s in _segs if s)
            filled += 1
        elif q > 0:
            it.quantity = 0
            _obs = it.observations or ""
            if "não medida" not in _obs.lower():
                it.observations = (
                    _obs + " | Área NÃO medida (lida de PDF por IA, não da geometria) — "
                    "preencha a metragem, informe a área no upload ou envie o DXF pra medir."
                ).strip(" |")
            blanked += 1
    return filled, blanked


def _dedupe_revisoes(file_paths: list) -> tuple:
    """Quando a MESMA prancha vem em várias revisões (mesmo nome-base, só muda o
    sufixo -RNN), mantém só a revisão mais alta — senão o quantitativo conta a
    mesma prancha 2x (caso Rafael/Engefast 21/07: 008 R03+R04, 009 R02+R03).

    CONSERVADOR (regra nº1 — nunca apagar prancha real em silêncio):
    - só trata como REVISÃO quando veio "REV" explícito OU o número tem 2+ dígitos
      (R00, R01, R04...). Um "R" + 1 dígito só (R1, R2, R5) é ambíguo — em projeto
      BR pode ser Rua (viário/loteamento), modelo de casa (CUB R1/R8/R16), Raio de
      curva ou eixo — então NÃO deduplica (fica igual a antes: no máx conta 2x,
      nunca some prancha);
    - só funde quando o nome-base bate EXATAMENTE e a extensão é a MESMA;
    - arquivo sem sufixo de revisão NUNCA é descartado.
    Retorna (mantidos, descartados), descartados = lista de (path_descartado,
    path_mantido) — pra AVISAR o usuário (recuperável, nunca some em silêncio)."""
    import re
    # g1=base  g2="EV"/None (marca REV explícita)  g3=dígitos da revisão
    rev_rx = re.compile(r'^(.*?)[-_ ]R(EV\.?)?[-_ .]?(\d{1,3})$', re.IGNORECASE)
    grupos: dict = {}
    sem_rev: list = []
    for p in file_paths:
        stem, ext = os.path.splitext(os.path.basename(p))
        m = rev_rx.match(stem)
        if not m:
            sem_rev.append(p)
            continue
        base, rev_marker, digits = m.group(1), m.group(2), m.group(3)
        if not rev_marker and len(digits) < 2:
            # "R" + 1 dígito sem "REV": ambíguo (Rua/Raio/modelo) — não arrisca
            sem_rev.append(p)
            continue
        chave = (base.rstrip(" -_.").lower(), ext.lower())
        grupos.setdefault(chave, []).append((int(digits), p))
    mantidos_set = set(sem_rev)
    descartados: list = []
    for _chave, lst in grupos.items():
        if len(lst) == 1:
            mantidos_set.add(lst[0][1])
            continue
        lst.sort(key=lambda x: x[0])   # menor revisão primeiro
        vencedor = lst[-1][1]           # maior revisão vence
        mantidos_set.add(vencedor)
        for _rev, p in lst[:-1]:
            descartados.append((p, vencedor))
    mantidos = [p for p in file_paths if p in mantidos_set]  # preserva ordem
    return mantidos, descartados


def process_job(job_id: str, file_paths: list[str], work_dir: str,
                typology: str = "office",
                user_sheet_types: dict[str, str] | None = None,
                user_ambientes: dict[str, str] | None = None,
                project_type: str = "arquitetura",
                is_complement: bool = False,
                user_total_area: float = 0,
                user_pe_direito: float = 0):
    """Processa um job prancha por prancha. Aceita PDF, DWG e DXF.

    `typology` alimenta a camada de calibração por densidade — alertas
    comparam o projeto contra benchmarks da mesma categoria.

    `user_sheet_types` (opcional): mapa {caminho_arquivo: sheet_type}
    pra sobrescrever a detecção automática quando o usuário classificou
    a prancha manualmente no dashboard. Sem isso, o classificador por
    nome falha com numerações próprias de cada escritório.

    `user_ambientes` (opcional): mapa {caminho_arquivo: ambiente_canonico}
    quando sheet_type = detalhe_ambiente (banheiro_suite, cozinha, lavabo
    etc). Injeta contexto específico no prompt da IA."""
    user_sheet_types = user_sheet_types or {}
    user_ambientes = user_ambientes or {}
    # Projeto ESTRUTURAL (concreto armado) vs arquitetura — define qual
    # SYSTEM_PROMPT e quais unidades o motor usa. Vem do seletor do upload,
    # persistido em projects.project_type (sobrevive a restart/recovery).
    is_structural = (project_type or "").strip().lower() == "estrutura"
    import gc
    import anthropic
    from processor import identify_sheet_type, extract_text, render_crops
    from analyzer import analyze_sheet, SYSTEM_PROMPT
    from models import SheetInfo, SheetType, ProjectData, BudgetItem, Confidence

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        jobs.update_field(job_id, status="error")
        jobs.update_field(job_id, error_message="API key não configurada")
        return

    try:
        jobs.update_field(job_id, status="processing")
        # Espelha 'processing' no BANCO (fix disjuntor 2026-07-22): o disjuntor
        # anti-crash-loop da recuperação exige status='processing' no banco pra
        # achar o job venenoso — mas antes só o JobsStore LOCAL recebia esse status,
        # então a quarentena NUNCA disparava (o job ficava 'queued' no banco). Agora
        # o job ATIVO é marcado no banco → um OOM é quarentenado na 2ª tentativa em
        # vez de virar crash-loop de 6 boots.
        try:
            _supabase_update("projects", "job_id", job_id, {"status": "processing"})
        except Exception:
            pass
        jobs.update_field(job_id, progress=3)
        jobs.update_field(job_id, current_step="Iniciando processamento...")

        # Funções auxiliares (definir antes de usar)
        from analyzer import _normalize_br_number as _norm_br
        def sf(v):
            """Converte valor (string ou número) em float, respeitando notação
            PT-BR: vírgula pode ser decimal. Antes usávamos replace(',', '')
            que transformava '135,4' em 1354 — área de casa virava 10x maior."""
            if v is None: return 0
            s = str(v).replace('m²','').replace('m2','').replace('cm','').strip()
            s = _norm_br(s)
            try: return float(s)
            except: return 0

        project_data = ProjectData()
        # Acumulador de áreas: consenso ao final do loop (evita ordem de
        # processamento decidir qual valor sobrevive)
        _area_readings = {"total_area": [], "layout_area": [], "no_intervention_area": []}

        # Dedupe de REVISÃO: mesma prancha em várias revisões (só muda o -RNN) →
        # mantém a mais nova, senão o quantitativo conta a prancha 2x. Conservador
        # e TRANSPARENTE (avisa o usuário quais tratou como antigas). Vale pra todo
        # caminho (upload, reprocesso, add-file, recovery). Caso Rafael/Engefast 21/07.
        file_paths, _revs_desc = _dedupe_revisoes(file_paths)
        if _revs_desc:
            _pares = "; ".join(f"{os.path.basename(_v)} → {os.path.basename(_n)}"
                               for _v, _n in _revs_desc)
            print(f"[dedupe-revisao] {len(_revs_desc)} revisão(ões) antiga(s) ignorada(s): {_pares}")
            try:
                project_data.warnings = (project_data.warnings or []) + [
                    f"{len(_revs_desc)} prancha(s) vieram em mais de uma revisão — usei só a "
                    f"mais recente de cada pra não contar em dobro. Se alguma não for revisão "
                    f"da outra, me avise que eu incluo."
                ]
            except Exception:
                pass

        # Separar PDFs de DWG/DXF
        pdf_paths = [f for f in file_paths if f.lower().endswith('.pdf')]
        cad_paths = [f for f in file_paths if f.lower().endswith(('.dwg', '.dxf'))]

        # DXF supera DWG do MESMO nome-base: se o usuário re-exportou um DWG que
        # falhava (ex.: AEC/MEP) pra DXF, o DWG velho fica no Storage e re-falha a
        # cada processamento, gerando aviso "planilha INCOMPLETA" enganoso mesmo o
        # DXF tendo medido tudo. Vale pra TODO caminho (upload, reprocesso, add-file).
        # (Pedro 15/07, forro MEP.)
        _dxf_stems_cad = {os.path.splitext(os.path.basename(f))[0].lower()
                          for f in cad_paths if f.lower().endswith('.dxf')}
        if _dxf_stems_cad:
            cad_paths = [f for f in cad_paths
                         if not (f.lower().endswith('.dwg')
                                 and os.path.splitext(os.path.basename(f))[0].lower() in _dxf_stems_cad)]

        # Persistir TODOS os arquivos originais (PDF + DWG + DXF) no Storage.
        # Serve pra 3 coisas:
        # 1. Botão 👁 "Ver prancha" na revisão inline
        # 2. Reprocessar projeto com motor atualizado (baixa originais, roda de novo)
        # 3. RESILIÊNCIA (16/06): se o servidor reiniciar no meio do
        #    processamento (deploy/OOM/idle), o recovery encontra os arquivos
        #    aqui e RETOMA o job. Por isso o upload é SÍNCRONO e a PRIMEIRA
        #    coisa do processamento — antes era background e se perdia no
        #    crash (caso Adriano: 2 falhas, arquivo sumiu nas duas).
        _t0_upload = time.time()
        _n_ok = 0
        for _p in file_paths:
            try:
                _fname = os.path.basename(_p)
                _ok_up = _supabase_storage_upload_prancha(_p, job_id, _fname)
                if not _ok_up:
                    # 2ª tentativa: o timeout é de 60s e CAD grande passa perto
                    # disso. Falhar aqui não interrompe o processamento — mas deixa
                    # o projeto SEM o original, e aí um complemento futuro refaz sem
                    # o CAD e apaga a medição (caso Walter 30/07).
                    _ok_up = _supabase_storage_upload_prancha(_p, job_id, _fname)
                if _ok_up:
                    _n_ok += 1
                elif _p.lower().endswith((".dwg", ".dxf")):
                    # CAD que não subiu é grave e era INVISÍVEL: só ia pro log do
                    # Render. Agora aparece no painel — COM o motivo, que antes
                    # também se perdia (caso Ana 31/07: o DXF que mediu 72 itens
                    # não foi guardado e não deu pra saber se foi tamanho ou timeout).
                    _log_error("storage:cad-nao-guardado",
                               f"CAD não subiu pro Storage: {_fname} — complemento ou "
                               f"reprocesso deste projeto não vai conseguir medir de novo. "
                               f"Motivo: {_ULTIMA_FALHA_UPLOAD_PRANCHA or 'não registrado'}",
                               job_id, severity="warning")
            except Exception as _e:
                print(f"[upload-pranchas] erro {os.path.basename(_p)}: {_e}")
        print(f"[upload-pranchas] {_n_ok}/{len(file_paths)} no Storage "
              f"em {time.time()-_t0_upload:.1f}s (antes de processar)")

        # Render DWG/DXF → PNG pra preview na revisão (opção 1 acordada).
        # Só roda se tem CAD no projeto. Matplotlib é pesado, por isso
        # em thread separada.
        if cad_paths:
            def _render_cad_previews_bg():
                try:
                    from dxf_render import render_dxf_to_png_safe
                except Exception as _imp_e:
                    print(f"[cad-preview] import erro: {_imp_e}")
                    return
                for _cad in cad_paths:
                    _fname = os.path.basename(_cad)
                    # Procura o DXF convertido (se era DWG) ou usa direto
                    _dxf_path = _cad
                    if _cad.lower().endswith('.dwg'):
                        # Procurar DXF equivalente em work_dir
                        _base = os.path.splitext(_fname)[0]
                        for _cand in os.listdir(work_dir):
                            if _cand.lower().endswith('.dxf') and _base.lower() in _cand.lower():
                                _dxf_path = os.path.join(work_dir, _cand)
                                break
                        else:
                            continue  # DWG sem DXF convertido — pula
                    # Render
                    _png_path = os.path.join(work_dir, os.path.splitext(_fname)[0] + '.png')
                    if render_dxf_to_png_safe(_dxf_path, _png_path, timeout_s=60):
                        # Upload PNG pro Storage
                        _png_fname = os.path.splitext(_fname)[0] + '.png'
                        _supabase_storage_upload_prancha(_png_path, job_id, _png_fname)
            # NÃO inicia o preview aqui (fix OOM 2026-07-22): rodar o matplotlib
            # concorrente com a análise soma o pico de RAM do render EM CIMA do
            # pico da análise — justamente a janela de perigo em projeto com muitas
            # pranchas (6-9). O preview é cosmético (botão "Ver prancha"); iniciado
            # no FIM, depois da planilha, lendo os DXF que ficam no work_dir.

        # ── Pesos de progresso alinhados com percepção do usuário ──
        # A fase que o usuário percebe como "cada prancha processada" é a análise
        # IA (CoT + JSON), não a conversão DWG→DXF (rápida, parte invisível).
        # Alocação:
        #   0-5%:    upload/init
        #   5-20%:   conversão DWG→DXF (leve, 15% span pro total dos CADs)
        #   20-X%:   análise DXF (maior fatia — cada DXF = 1/N do span restante)
        #   X-95%:   análise PDFs (se houver)
        #   95-100%: consolidação + planilha
        has_cad = bool(cad_paths)
        has_pdf = bool(pdf_paths)
        # Partição principal do progresso
        if has_cad and has_pdf:
            # CAD conversão 5→15%, análise CAD 15→55%, PDF 55→92%
            conv_end_pct = 15
            cad_analysis_end_pct = 55
        elif has_cad:
            # Só CAD: conversão 5→20%, análise CAD 20→92% (bem maior pra DXF dominar)
            conv_end_pct = 20
            cad_analysis_end_pct = 92
        else:
            # Só PDF: 5→92%
            conv_end_pct = 5
            cad_analysis_end_pct = 92
        cad_end_pct = cad_analysis_end_pct  # compatibilidade com código abaixo

        # Converter DWG→DXF se necessário
        dxf_paths = []
        # (caminho, fator_para_metros) das pranchas DXF que a extração leu com
        # sucesso — alimenta a SOMBRA de montagem de cômodos no fim do job.
        _dxfrooms_units: list = []
        dwg_failed = []  # DWGs que não converteram — reportar mesmo quando outros deram certo (escopo garantido)
        dwg_via_libredwg = []  # convertidos pelo plano B — aviso pro cliente conferir (escopo garantido)
        _aec_failed = []  # DWGs que falharam E são AutoCAD MEP/Architecture (objetos AEC) → aviso preciso
        if cad_paths:
            jobs.update_field(job_id, progress=5)
            jobs.update_field(job_id, current_step="Processando arquivos DWG/DXF...")
            try:
                from dwg_extractor import extract_from_file, generate_budget_data, convert_dwg_to_dxf, dwg_has_aec_markers
                n_cad = len(cad_paths)
                conv_span = conv_end_pct - 5  # ex.: 15 ou 10 pts
                dwg_failed = []  # acumular DWGs que falharam pra reportar erro real depois
                for ci, cad_path in enumerate(cad_paths):
                    base = 5 + int((ci / max(n_cad, 1)) * conv_span)
                    ext = cad_path.lower().rsplit('.', 1)[-1]
                    if ext == 'dwg':
                        jobs.update_field(job_id, progress=base)
                        jobs.update_field(job_id, current_step=f"Convertendo DWG→DXF ({ci+1}/{n_cad}): {os.path.basename(cad_path)}")
                        dxf_path = convert_dwg_to_dxf(cad_path)
                        if dxf_path:
                            dxf_paths.append(dxf_path)
                            jobs.update_field(job_id, current_step=f"DWG convertido: {os.path.basename(dxf_path)}")
                            # Guarda-corpo do fallback (01/08): conversão que veio do
                            # libredwg (sufixo _libredwg.dxf) fica REGISTRADA e o
                            # projeto ganha um aviso pro cliente conferir medidas-chave.
                            # Validação disponível: 5/5 recusados convertem e abrem;
                            # 1/1 baseline com geometria idêntica (Δ 0%) — amostra de
                            # baseline não cresce porque o ODA falha em quase tudo.
                            if dxf_path.endswith("_libredwg.dxf"):
                                # 🪤 project_data ainda NÃO existe aqui (nasce depois da
                                # análise) — acumular na lista local, aplicar adiante.
                                dwg_via_libredwg.append(os.path.basename(cad_path))
                                _log_error("libredwg:usado-no-fluxo",
                                           f"{os.path.basename(cad_path)} convertido via "
                                           f"fallback libredwg (ODA recusou)",
                                           job_id, severity="info")
                        else:
                            dwg_failed.append(os.path.basename(cad_path))
                            # É arquivo AutoCAD MEP/Architecture (objetos AEC)? Aí a falha
                            # é esperada (conversor livre não abre proxy) e o aviso é preciso.
                            _is_aec = dwg_has_aec_markers(cad_path)
                            if _is_aec:
                                _aec_failed.append(os.path.basename(cad_path))
                            # Telemetria (20/07): a conversão DWG→DXF falha e some no disco
                            # efêmero do Render — foi a CAIXA-PRETA do caso estrutural Luciano
                            # (DWG não converteu → aço nunca foi lido → tudo estimado). Registra
                            # no error_log pra o Pedro VER por que o cliente não teve medição.
                            _log_error("dwg:convert-fail",
                                       f"DWG não converteu pra DXF: {os.path.basename(cad_path)}"
                                       + (" — objetos AEC/MEP (conversor livre não abre)" if _is_aec
                                          else " — cliente ganharia a medição mandando DXF direto"),
                                       job_id, severity="warning")
                            jobs.update_field(job_id, current_step=f"Falha ao converter DWG: {os.path.basename(cad_path)} (seguindo sem)")
                    else:
                        dxf_paths.append(cad_path)

                # Se TODOS os DWGs falharam E não tem PDFs, é fim de linha — marca como failed.
                # Mensagem instrutiva pro user resolver sozinho (95% dos casos).
                if dwg_failed and not dxf_paths and not pdf_paths:
                    arquivos = ', '.join(dwg_failed)
                    # Aviso PRECISO quando confirmamos objetos AEC (AutoCAD MEP/Architecture)
                    # no binário — em vez do genérico "costuma acontecer".
                    # 3 causas DIFERENTES, cada uma com uma saída diferente. Antes as
                    # três recebiam o mesmo texto de 12 linhas listando hipóteses — o
                    # cliente Thalison (29/07) leu "versão nova ou objetos especiais"
                    # quando o arquivo dele estava INCOMPLETO, reenviou o mesmo arquivo
                    # 2x e desistiu da prancha. Mensagem curta, em tópicos, 1 saída só.
                    try:
                        from dwg_extractor import dwg_failure_reason as _motivo_falha
                    except Exception:
                        _motivo_falha = lambda _p: ""
                    _truncados = [n for n in dwg_failed if _motivo_falha(n) == "truncado"]

                    if _aec_failed:
                        msg = (
                            f"Não conseguimos medir: {', '.join(_aec_failed)}\n\n"
                            f"É um arquivo do AutoCAD MEP/Architecture. Ele guarda paredes e "
                            f"móveis como \"objetos inteligentes\", que nenhum conversor abre "
                            f"direto — nem o \"Salvar como DXF\" comum. Não é defeito do seu "
                            f"arquivo.\n\n"
                            f"Como resolver, em 3 passos:\n"
                            f"1. No AutoCAD, com o arquivo aberto, digite EXPORTTOAUTOCAD e Enter\n"
                            f"2. Na janela que abrir, escolha a versão 2013 e confirme — ele cria "
                            f"um arquivo novo (o seu original não é alterado)\n"
                            f"3. Abra esse arquivo novo e salve como DXF. Suba o DXF aqui\n\n"
                            f"Não tem o AutoCAD Architecture? Mande o PDF vetorial da prancha "
                            f"(plotado com escala) — a gente mede pela geometria."
                        )
                    elif _truncados:
                        msg = (
                            f"Não conseguimos ler: {', '.join(_truncados)}\n\n"
                            f"Nosso conversor chegou ao fim do arquivo antes do esperado. Na "
                            f"maioria das vezes isso é limitação dele com esse DWG específico, "
                            f"não defeito do seu arquivo — se ele abre normal no seu CAD, o "
                            f"problema é do nosso lado.\n\n"
                            f"Como resolver (é o caminho que mais funciona):\n"
                            f"1. Abra o arquivo no seu CAD\n"
                            f"2. Salve como DXF (Salvar Como → DXF, versão 2013)\n"
                            f"3. Suba o DXF aqui — em geral ele mede sem problema\n\n"
                            f"Reenviar o mesmo DWG vai dar no mesmo resultado.\n\n"
                            f"Só se ele também NÃO abrir no seu CAD é que o arquivo está mesmo "
                            f"danificado — aí procure a última versão salva ou o backup "
                            f"automático (.bak)."
                        )
                    else:
                        msg = (
                            f"Não conseguimos abrir: {arquivos}\n\n"
                            f"O conversor não reconheceu o arquivo. Não quer dizer que ele esteja "
                            f"com defeito — pode ser uma versão de DWG que não lemos ainda.\n\n"
                            f"Como resolver:\n"
                            f"1. Abra o arquivo no seu CAD\n"
                            f"2. Salvar Como → DXF (versão 2013)\n"
                            f"3. Suba o DXF aqui\n\n"
                            f"O DXF é o formato que a gente mede melhor. Se preferir, o PDF "
                            f"vetorial da prancha também serve — mas aí sai como estimativa."
                        )
                    # COMPLEMENTO (/add-file) que falha NÃO pode derrubar o projeto que
                    # já funcionava. O add-file preserva os itens antigos (limpa só no
                    # sucesso). Se este job já tinha resultado, restaura 'done' + avisa —
                    # nunca marca erro. (Feedback Pedro 15/07: DWG-complemento que não abriu
                    # sumia a planilha da tela; ele achou que tinha perdido tudo.)
                    if is_complement and _complement_base_has_items(job_id):
                        _warn_txt = (
                            f"O arquivo CAD que você anexou ({arquivos}) não pôde ser aberto "
                            f"automaticamente (DWG de versão recente do AutoCAD ou com objetos "
                            f"de MEP/elétrica). Sua planilha anterior foi mantida — nada foi "
                            f"perdido. Pra medir pelo CAD: abra no AutoCAD ou BricsCAD, "
                            f"Salvar Como → DXF 2013, e anexe o DXF aqui.")
                        jobs.update_field(job_id, status="done", progress=100, error_message=None,
                                          current_step="Complemento não pôde ser lido — planilha anterior mantida")
                        try:
                            _supabase_update("projects", "job_id", job_id, {
                                "status": "done",
                                "error_message": None,
                                "warnings": [_warn_txt],
                                "completed_at": datetime.utcnow().isoformat(),
                            })
                        except Exception as _upe:
                            print(f"[add-file] restaurar done falhou: {_upe}")
                        print(f"[add-file] complemento CAD falhou, base preservada → done+aviso (sem erro)")
                        return
                    jobs.update_field(job_id, error_message=msg, current_step="❌ Arquivo CAD inválido — leia mensagem abaixo")
                    raise RuntimeError(msg)
            except Exception as e:
                jobs.update_field(job_id, error_message=f"Erro DWG→DXF: {e}")
                raise

        # Extrair dados de DXF e enviar pro Claude interpretar
        dxf_items = []
        # Coleta falha de IA por DXF (mesmo papel que sheet_errors no caminho
        # PDF). A interpretação Claude do DXF pode falhar por sobrecarga/timeout/
        # JSON inválido — sem registrar isso aqui, um projeto SÓ-DXF cuja única
        # chamada de IA falha cai no guard de 0-itens com a mensagem ERRADA
        # ("troque o arquivo / PDF escaneado") em vez de "IA sobrecarregada,
        # reprocesse grátis". Declarado ANTES do loop pra estar em escopo no
        # except da IA (armadilha #11 do CLAUDE.md reaparecendo no caminho DXF).
        dxf_errors: list[str] = []
        xref_warnings: list[str] = []  # xref não-resolvido / extração estéril → orienta o usuário (BIND / PDF)
        # Cross-check determinístico (promove estimado→confirmado por geometria).
        # 🚫 NÃO LIGAR EM PRODUÇÃO. Três rodadas de revisão adversarial (15/07) acharam,
        # cada uma, um jeito NOVO de gerar falso-medido (item de remoção pegando área do
        # piso novo; polígono do piso inteiro; abreviação de layer escapando da deny-list).
        # O ganho (poucas células a mais de branco) não paga o risco na regra nº1.
        # Fica OFF de propósito. A parte VALIOSA da caçada — rebaixar soma de layer
        # multi-hachura — virou rede de segurança SEMPRE ligada (ver _multi_hatch_sums,
        # independe deste flag). Se um dia for reativar, exige nova revisão + teste real.
        _XCHECK_ON = os.environ.get("DXF_CONFIRM_CROSSCHECK", "0") == "1"

        # ── CHECKPOINT: cache carregado UMA vez, serve DXF e PDF ───────────
        # Só na RETOMADA automática (auto_resume>0) e fora do reprocesso manual:
        # aí a prancha já analisada (DXF ou PDF) pula extração+IA. No 1º run o
        # cache fica vazio (nada a reaproveitar), mas os checkpoints são SALVOS
        # durante o run — pra existir se o servidor cair no meio. Antes só o
        # caminho PDF salvava/lia; o DXF não guardava nada, então um job de 43
        # DXF que caía refazia TUDO na retomada (caso perplan/Rafael 21/07).
        _ckpt_cache = {}
        try:
            import urllib.request as _ur_ck
            _qck = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
                    f"&select=auto_resume_count,reprocess_count")
            _rck = _ur_ck.Request(_qck, method="GET")
            _rck.add_header("apikey", SUPABASE_KEY)
            _rck.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            _rows_ck = _json.loads(_ur_ck.urlopen(_rck, timeout=8).read().decode("utf-8"))
            if _rows_ck and int(_rows_ck[0].get("auto_resume_count") or 0) > 0 \
                    and int(_rows_ck[0].get("reprocess_count") or 0) == 0:
                _ckpt_cache = _ckpt_load_all(job_id)
                if _ckpt_cache:
                    print(f"[ckpt] {job_id}: {len(_ckpt_cache)} prancha(s) com "
                          f"análise pronta — retomada vai pular a IA nelas")
        except Exception as _cke:
            print(f"[ckpt] cache indisponível (segue do zero): {_cke}")

        if dxf_paths:
            # Análise DXF começa onde a conversão termina
            extract_start = conv_end_pct
            jobs.update_field(job_id, progress=extract_start)
            jobs.update_field(job_id, current_step="Extraindo geometria dos DXF...")
            try:
                from dwg_extractor import extract_from_file, identify_architectural_elements, category_for_layer
                from analyzer import SYSTEM_PROMPT, SYSTEM_PROMPT_ESTRUTURA
                import json as _j

                # ── CONSENSO DE UNIDADE DO PROJETO (pré-passe barato) ──────────
                # Prancha SEM cota chuta a escala — num projeto BR sem cota o motor
                # cai em "pés" (0,3048) e o comprimento sai errado/zerado. Aqui
                # sondamos algumas pranchas (planta primeiro, menor primeiro) até
                # UMA provar a escala por cota, e usamos essa nas pranchas sem cota.
                # Cota PRÓPRIA da prancha sempre vence (não sobrescreve quem tem
                # prova). Caso Rafael 21/07: 501 provou metros, 004/201 chutavam
                # pés — achado na 6ª sondagem em ~10s. Kill: DXF_UNIT_CONSENSUS=0.
                _unit_consensus = None
                if os.getenv("DXF_UNIT_CONSENSUS", "1") != "0" and len(dxf_paths) >= 2:
                    try:
                        from dwg_extractor import probe_unit as _probe_unit
                        def _planta_score(_p):
                            _up = os.path.basename(_p).upper()
                            _s = 0
                            if ".PL." in _up or "-PL-" in _up or "PLANTA" in _up:
                                _s -= 100          # planta: mais provável ter cota
                            if any(_k in _up for _k in ("DET", "DIAG", "-DG", "DT.DC",
                                                        "DT.DP", "ESQ.", "VS.", "UNIFIL")):
                                _s += 50           # detalhe/diagrama: raramente cotado
                            try:
                                _sz = os.path.getsize(_p)
                            except OSError:
                                _sz = 0
                            return (_s, _sz)       # depois menor arquivo (readfile + barato)
                        for _i_pb, _pp in enumerate(sorted(dxf_paths, key=_planta_score)):
                            if _i_pb >= 12:        # cap: não lê o projeto inteiro só pra unidade
                                break
                            # a sondagem roda FORA do timeout de 900s — não lê prancha
                            # gigante (a que prova a escala é planta normal, <20MB).
                            try:
                                if os.path.getsize(_pp) > 60 * 1024 * 1024:
                                    continue
                            except OSError:
                                pass
                            _f = _probe_unit(_pp)
                            if _f and _f > 0:
                                _unit_consensus = _f
                                print(f"[unit-consenso] escala do projeto provada por cota em "
                                      f"{os.path.basename(_pp)}: fator {_f} (sondagem {_i_pb+1})")
                                break
                        if _unit_consensus is None:
                            print("[unit-consenso] nenhuma prancha provou a escala por cota — "
                                  "cada prancha segue com a própria detecção")
                    except Exception as _uce:
                        print(f"[unit-consenso] pré-passe falhou (segue sem): {_uce}")

                n_dxf = len(dxf_paths)
                dxf_span = cad_end_pct - extract_start
                for idx, dxf_path in enumerate(dxf_paths):
                    # 🛡️ Freio de MEMÓRIA: se o container está chegando perto do limite
                    # de RAM, aborta ANTES do OOM matar o servidor inteiro. Um projeto
                    # pesado falha sozinho (com orientação de dividir em lotes) e o
                    # servidor fica de pé pra todos os outros. Só depois da 1ª prancha
                    # (idx>0) — se estourar já na 1ª, é caso de prancha única densa.
                    if idx > 0 and _mem_pressure(0.85):
                        _abort_job_mem(job_id, idx, n_dxf)
                        return
                    # Cada DXF ocupa 1/N da faixa. Extração 30% + IA 70% dentro da faixa.
                    dxf_base = extract_start + int((idx / max(n_dxf, 1)) * dxf_span)
                    dxf_next = extract_start + int(((idx + 1) / max(n_dxf, 1)) * dxf_span)
                    dxf_mid = dxf_base + int((dxf_next - dxf_base) * 0.3)

                    jobs.update_field(job_id, progress=dxf_base)
                    jobs.update_field(job_id, current_step=f"DXF {idx+1}/{n_dxf}: Extraindo {os.path.basename(dxf_path)}...")

                    # 🪤 Nome REAL da prancha, pra cada item saber de onde veio.
                    # Antes os itens de CAD nasciam com ref_sheet = a palavra
                    # "DXF" — 1.798 dos 3.851 itens do banco (47%) ficaram assim.
                    # Consequência visível: o card "Arquivos que alimentaram o
                    # projeto" ficava VAZIO em todo projeto medido por CAD, porque
                    # ele monta a lista a partir desse campo (Pedro, 31/07). Os
                    # itens vindos de PDF sempre guardaram o nome — só o CAD perdia.
                    _dxf_nome = os.path.basename(dxf_path)
                    # ── CHECKPOINT DXF (buraco tapado 21/07): se esta prancha já
                    # foi analisada num run anterior (retomada), restaura os itens
                    # + contribuições salvos e PULA extração/IA (que é o caro). A
                    # chave leva prefixo "dxfck_" pra nunca colidir com as do PDF
                    # (que são "<stem>_p<n>"). Passa o stem CRU pro _ckpt_save (que
                    # sanitiza 1x, igual o PDF) e sanitiza só no lookup — senão a
                    # dupla-sanitização (o "__" colapsa pra "_") faz a chave não bater.
                    _dxf_raw_stem = "dxfck_" + os.path.splitext(os.path.basename(dxf_path))[0]
                    _dxf_stem = _sanitize_filename_for_storage(_dxf_raw_stem)
                    if _dxf_stem in _ckpt_cache:
                        _cp = _ckpt_cache[_dxf_stem]
                        from models import BudgetItem as _BI, Confidence as _Conf
                        _rest = 0
                        for _r in (_cp.get("items") or []):
                            try:
                                dxf_items.append(_BI(
                                    item_num=_r.get("item_num", "") or "",
                                    description=_r.get("description", "") or "",
                                    unit=_r.get("unit", "vb") or "vb",
                                    quantity=float(_r.get("quantity") or 0),
                                    observations=_r.get("observations", "") or "",
                                    ref_sheet=_r.get("ref_sheet") or _dxf_nome,
                                    confidence=_Conf(_r.get("confidence", "estimado") or "estimado"),
                                    discipline=_r.get("discipline", "Complementares") or "Complementares",
                                    origem=_r.get("origem", "dxf_geom") or "dxf_geom",
                                ))
                                _rest += 1
                            except Exception:
                                continue
                        for _w in (_cp.get("xref_warnings") or []):
                            xref_warnings.append(_w)
                        project_data.warnings = (project_data.warnings or []) + list(_cp.get("pd_warnings") or [])
                        if _cp.get("new_rooms"): project_data.new_rooms.extend(_cp["new_rooms"])
                        if _cp.get("kept_elements"): project_data.kept_elements.extend(_cp["kept_elements"])
                        if _cp.get("demolition_notes"): project_data.demolition_notes.extend(_cp["demolition_notes"])
                        for _k, _vs in (_cp.get("area_readings") or {}).items():
                            if _k in _area_readings:
                                _area_readings[_k].extend(_vs)
                        if _cp.get("name") and not project_data.name:
                            project_data.name = _cp["name"]
                        jobs.update_field(job_id, current_step=f"DXF {idx+1}/{n_dxf}: análise já concluída, retomando ✓")
                        print(f"[ckpt] {os.path.basename(dxf_path)}: {_rest} itens reaproveitados do checkpoint")
                        continue
                    # snapshot dos acumuladores ANTES de processar (pro checkpoint no fim)
                    _snap = {
                        "items": len(dxf_items),
                        "xref": len(xref_warnings),
                        "warn": len(project_data.warnings or []),
                        "rooms": len(project_data.new_rooms or []),
                        "kept": len(project_data.kept_elements or []),
                        "demo": len(project_data.demolition_notes or []),
                        "area": {_k: len(_v) for _k, _v in _area_readings.items()},
                        "name": project_data.name or "",
                    }

                    # 1. Extrair dados estruturados do DXF
                    # ISOLADO por-arquivo: se o ezdxf não parseia ESTA prancha (objeto
                    # proxy, DXF truncado, entidade não suportada), registra em dxf_errors
                    # e PULA pro próximo — não derruba os outros DXF/PDF do job. Antes a
                    # exceção subia pro except externo que fazia raise e matava a planilha
                    # inteira por causa de uma prancha só (caso Thamiry: 22 DWGs).
                    try:
                        # ARQUIVO GRANDE (19/07): DXF acima do teto de leitura
                        # segura passa por emagrecimento streaming (iterdxf) que
                        # mantém só o que o motor mede — em vez de recusar ou
                        # estourar a memória. Falhou/não rendeu → segue original.
                        #
                        # TIMEOUT DE EXTRAÇÃO (21/07, incidente Everton): um DXF gigante
                        # de hospital pendurava extract_from_file por ~20min — o job travava
                        # e segurava o _JOB_SEMAPHORE (próximo upload ficava preso atrás).
                        # Agora emagrecimento + parse + prompt rodam num worker com teto de
                        # 900s (15min); estourou → pula ESTA prancha (job segue) em vez de
                        # pendurar. Teto GENEROSO de propósito (servidor 25GB, 21/07): só
                        # pega hang de verdade — prancha grande legítima leva minutos e conclui.
                        import concurrent.futures as _cf
                        # Extração ISOLADA num subprocesso matável com teto de RAM
                        # (fix confiabilidade 2026-07-22): uma prancha densa que
                        # estouraria a memória mata só o filho — o servidor fica de pé.
                        # timeout → _cf.TimeoutError (pula prancha); filho morto/erro →
                        # RuntimeError (pula prancha); sem subprocesso → fallback in-process.
                        extraction, structured_text, dxf_path = _extract_dxf_isolated(
                            dxf_path, _unit_consensus, timeout_s=900)
                    except _cf.TimeoutError:
                        _bn_dxf = os.path.basename(dxf_path)
                        print(f"[dxf] extração ESTOUROU 900s em {_bn_dxf} — pulando (job segue)")
                        try:
                            _log_error("dxf:extract-timeout",
                                       f"{_bn_dxf}: extração de geometria > 900s (prancha grande/complexa)", job_id)
                        except Exception:
                            pass
                        dxf_errors.append(f"{_bn_dxf}: a leitura da geometria demorou demais (prancha muito grande/pesada) — mande essa prancha isolada, ou reexporte só o pavimento que você precisa")
                        continue
                    except Exception as _ex_dxf:
                        _bn_dxf = os.path.basename(dxf_path)
                        print(f"[dxf] extração falhou em {_bn_dxf}: {_ex_dxf}")
                        dxf_errors.append(f"{_bn_dxf}: não consegui ler a geometria desse arquivo (pode estar corrompido ou ter objetos não suportados)")
                        continue
                    # Cap de segurança (auditoria 06/07): projeto gigante pode gerar
                    # prompt enorme e estourar RAM/contexto do modelo. 300k chars é
                    # folgado pra uma prancha real.
                    if len(structured_text) > 300_000:
                        structured_text = structured_text[:300_000] + "\n[... texto truncado por tamanho ...]"

                    # ── TRAVA DE PROCEDÊNCIA (regra nº1) ──────────────────────
                    # Se a extração geométrica veio com RESSALVA (estéril / unidade
                    # suspeita / xref não resolvido), o CÓDIGO não deixa nenhum item
                    # deste DXF ser 'confirmado' (branco/medido), independente do que
                    # a IA marcar. Só REBAIXA pra estimado — nunca promove. Fecha o
                    # furo de "a IA carimba medido num número que não dá pra confiar".
                    _dxf_sem_procedencia = _extraction_has_quality_caveat(extraction.metadata)
                    # Guarda (caminho, fator) pra SOMBRA de montagem de cômodos —
                    # roda depois do done, fora do caminho do cliente. O fator tem
                    # que ser o MESMO que a extração real usou (já validado pelas
                    # cotas), senão a sombra mede noutra escala e não serve de nada.
                    try:
                        _dxfrooms_units.append(
                            (dxf_path, float(extraction.metadata.get("fator_para_metros") or 1.0)))
                    except (TypeError, ValueError):
                        pass
                    # Aviso ao usuário (não só rebaixar a cor): xref não-resolvido é a
                    # causa nº1 de "planilha só laranja" em CAD — o sinal existe na
                    # metadata mas antes morria no downgrade. Agora orienta a dar BIND
                    # ou mandar o PDF plotado.
                    _bn_w = os.path.basename(dxf_path)
                    if extraction.metadata.get("xref_nao_resolvido"):
                        xref_warnings.append(
                            f"{_bn_w}: parece que falta o desenho de referência (xref) — o arquitetônico "
                            f"pode estar num arquivo externo que não veio junto. No AutoCAD, use BIND pra "
                            f"incorporar os xrefs e reexporte o DXF, ou mande o PDF plotado da prancha.")
                    elif extraction.metadata.get("extracao_esteril"):
                        xref_warnings.append(
                            f"{_bn_w}: não consegui ler geometria mensurável nesse arquivo. Reexporte "
                            f"do CAD com tudo incorporado (BIND, sem xref solto) ou mande o PDF plotado da prancha.")

                    # Medidas DURAS desta prancha (pro cross-check), POR CATEGORIA e com
                    # 3 guards da 2ª revisão adversarial (15/07):
                    #  • furo C: medida fica atada à categoria física do PRÓPRIO item
                    #    (nunca confirma forro com área de piso).
                    #  • furo E (CRÍTICO): área só entra se o layer tem UMA hachura só.
                    #    Layer com várias hachuras = acabamentos mistos (porcelanato+cerâmica
                    #    no mesmo "ARQ-PISO") → a soma não mede NENHUM acabamento sozinho.
                    #  • furo F: layers que casam o token mas não são superfície instalável
                    #    (PISCINA→"PIS", PAVIMENTO→"PAV", FÔRMA/FORNO→"FOR") ficam FORA.
                    #  • furo G: layer de EXISTENTE / A-DEMOLIR não é superfície nova a instalar.
                    # Comprimento (parede/demolição) segue somado por layer — muitos segmentos
                    # = uma parede é o agregado correto (não tem o problema de acabamento misto).
                    _cats_geo = identify_architectural_elements(extraction)
                    _AREA_CATS = {"piso", "forro", "pintura"}   # m² físico
                    _LEN_CATS = {"paredes", "demolicao"}          # m físico (linear)
                    # tokens que casam categoria por prefixo mas NÃO são superfície nova
                    _XC_AREA_DENY = ("piscina", "pavimento", "forma", "fôrma", "forno",
                                     "parapeito", "peitoril", "soleira", "escada", "rampa",
                                     "diversos", "generico", "genérico")

                    def _xc_layer_ok(_lyr, _cat):
                        """True se o layer pode virar medida DURA pro cross-check.
                        • Superfície/parede NOVA não pode ser existente nem a-remover.
                        • ÁREA também barra 'demolir/demol' (piso a demolir) e nomes que
                          poluem (piscina/pavimento/fôrma). ATENÇÃO: a categoria LINEAR
                          'demolicao' é legítima (queremos medir demolição) — por isso o
                          barramento de 'demolir' é SÓ na área, nunca em demolicao."""
                        _low = (_lyr or "").lower()
                        if _cat in _AREA_CATS or _cat == "paredes":
                            if any(_t in _low for _t in ("exist", "remov", "retir", "a-demol", "a demol")):
                                return False
                        if _cat in _AREA_CATS:
                            if "demolir" in _low or "demol" in _low:
                                return False
                            if any(_t in _low for _t in _XC_AREA_DENY):
                                return False
                        return True

                    _hard_area_by_cat: dict[str, set] = {}   # cat -> {áreas m² inequívocas}
                    _hard_len_by_cat: dict[str, set] = {}    # cat -> {comprimentos m por layer}
                    for _ck, _cd in _cats_geo.items():
                        if _ck in _AREA_CATS:
                            # agrupa hachuras por layer; SÓ layer com EXATAMENTE 1 hachura
                            # (superfície única, sem ambiguidade de acabamento) entra.
                            _byl: dict[str, list] = {}
                            for _h in _cd.get("hatches", []):
                                if not _xc_layer_ok(_h.layer, _ck):
                                    continue
                                _byl.setdefault(_h.layer, []).append(getattr(_h, "area", 0) or 0)
                            for _areas in _byl.values():
                                if len(_areas) == 1 and _areas[0] > 0:
                                    _hard_area_by_cat.setdefault(_ck, set()).add(round(_areas[0], 2))
                        if _ck in _LEN_CATS:
                            _pl2: dict[str, float] = {}
                            for _w in _cd.get("walls", []):
                                if not _xc_layer_ok(_w.layer, _ck):
                                    continue
                                _pl2[_w.layer] = _pl2.get(_w.layer, 0.0) + (getattr(_w, "length", 0) or 0)
                            _s2 = {round(_v, 2) for _v in _pl2.values() if _v > 0}
                            if _s2:
                                _hard_len_by_cat.setdefault(_ck, set()).update(_s2)
                    # Polígonos fechados: individuais (não somados), SÓ piso/forro e passando
                    # os mesmos guards (furo D + F + G). Cada contorno já é 1 superfície.
                    for _pa in (getattr(extraction, "polygon_areas", None) or []):
                        try:
                            _pav = getattr(_pa, "area", None)
                            _ply = getattr(_pa, "layer", "") or ""
                            if not (_pav and float(_pav) > 0):
                                continue
                            _pcat = category_for_layer(_ply)
                            if _pcat in ("piso", "forro") and _xc_layer_ok(_ply, _pcat):
                                _hard_area_by_cat.setdefault(_pcat, set()).add(round(float(_pav), 2))
                        except Exception:
                            pass

                    # REDE DE SEGURANÇA (regra nº1, SEMPRE ligada — independe do flag):
                    # o prompt lista a ÁREA HACHURADA somada POR LAYER. Um layer GENÉRICO com
                    # VÁRIAS hachuras pode misturar acabamentos (porcelanato + cerâmica no
                    # mesmo "ARQ-PISO") → a soma não mede um acabamento sozinho. Se a IA marcar
                    # 'confirmado' num m² igual a essa soma, o código rebaixa pra estimado.
                    # Só REBAIXA — nunca cria medido (Finding 1 da revisão 15/07).
                    # EXCEÇÃO: layer que já NOMEIA o acabamento (PISO-PORCELANATO com 5 salas)
                    # é mono-acabamento por construção — a soma É legítima, não rebaixa.
                    _FINISH_TOKENS = ("porcelanato", "ceramic", "cerâmic", "vinil", "viníl",
                                      "laminad", "madeira", "granito", "marmore", "mármore",
                                      "ardosia", "ardósia", "cimenticio", "cimentício", "epoxi",
                                      "epóxi", "carpete", "taco", "pedra", "deck", "granilite",
                                      "paviflex", "borracha")
                    _multi_hatch_sums = set()
                    _mh_by_layer: dict[str, list] = {}
                    for _h in (extraction.hatches or []):
                        _mh_by_layer.setdefault(_h.layer, []).append(getattr(_h, "area", 0) or 0)
                    for _lyr, _areas in _mh_by_layer.items():
                        if len(_areas) <= 1:
                            continue
                        if any(_ft in (_lyr or "").lower() for _ft in _FINISH_TOKENS):
                            continue  # layer nomeia o acabamento → mono-acabamento, soma legítima
                        _tot = round(sum(_a for _a in _areas if _a and _a > 0), 2)
                        if _tot > 0:
                            _multi_hatch_sums.add(_tot)

                    # 2. Enviar pro Claude interpretar
                    jobs.update_field(job_id, progress=dxf_mid)
                    jobs.update_field(job_id, current_step=f"DXF {idx+1}/{n_dxf}: Nossa IA está analisando os dados extraídos...")
                    dxf_client = anthropic.Anthropic(api_key=api_key, timeout=300.0)

                    _proj_kind = "ESTRUTURA (concreto armado)" if is_structural else "arquitetura"
                    _estrutura_directive = (
                        "\n⚠ PROJETO ESTRUTURAL: gere quantitativo de ESTRUTURA — CONCRETO em m³, "
                        "FÔRMA em m², AÇO/ARMADURA/ESTRIBO em kg (nunca m²/m/un), discipline='Estrutura'. "
                        "Se houver quadro/resumo de aço nos dados, use o peso de lá (medido). Volume/área/"
                        "peso que você calcular é 'estimado'. Não invente bitola/fck.\n"
                    ) if is_structural else ""
                    dxf_prompt = f"""Analise os dados extraídos de um arquivo DXF de projeto de {_proj_kind}.
Os dados abaixo foram extraídos automaticamente do arquivo CAD (blocos, textos, layers, comprimentos, áreas).
Gere itens quantitativos (descrição + unidade + quantidade, SEM preço) com base nesses dados.
{_estrutura_directive}
{structured_text}

════════════════════════════════════════════════════════
REGRA CRÍTICA DE CONFIANÇA — NUNCA ESTIMAR, SÓ MEDIR OU SUGERIR
════════════════════════════════════════════════════════

O campo "confidence" TEM apenas duas categorias possíveis:

1. "confirmado" — SÓ quando a quantidade corresponde EXATAMENTE a uma medição objetiva do DXF:
   - Contagem literal de blocos (INSERT) que aparece em "CONTAGEM DE BLOCOS"
   - Contagem literal de esquadrias na seção "ESQUADRIAS" (com dimensão W×H)
   - Comprimento calculado em "COMPRIMENTOS POR LAYER" (valor em metros)
   - Área calculada em "ÁREAS HACHURADAS POR LAYER" (valor em m²)
   - Cota numérica que aparece em "COTAS/DIMENSÕES"
   A quantidade do item TEM que bater com o número extraído. Se você multiplicou, somou
   ou fez qualquer cálculo além de copiar o valor, NÃO é confirmado.

IMPORTANTE — SEÇÃO ESQUADRIAS (quando presente nos dados):
Cada linha tem o formato "NOME: N un | ~Wm × Hm = Xm²". Isso é DADO ESTRUTURADO
de portas/janelas com dimensão REAL extraída do CAD. Use para:
   - Gerar item de portas/janelas com a quantidade e dimensão exatas
   - Aplicar regra TCPO: se área ≤ 2m², NÃO descontar esse vão da pintura;
     se > 2m², descontar o excedente (área_vão - 2m²) da pintura adjacente
   - Toda esquadria extraída aqui pode ser marcada "confirmado" (veio de medição)

2. "estimado" — para todo o resto, SEM EXCEÇÃO:
   - Quantidades derivadas de texto/legenda ("demolir X" → qtd=1)
   - Itens sugeridos de práxis (administração local, limpeza final, instalação de placa)
   - Qualquer item cuja quantidade você não conseguiu ler DIRETO dos dados extraídos
   - Itens "vb" (verba) de valor único
   - Composições inferidas ("se tem drywall, precisa de montante" sem count no CAD)

REGRA DE OURO: **NA DÚVIDA, MARQUE "estimado".** É preferível 100 itens laranja que o
usuário confirma um a um, do que 1 item branco com número inventado. O usuário quer
poder confiar que "branco = aprovado direto", então só marque branco quando não houver
NENHUMA dúvida.

Não existe "verificar" nesta fase — use "estimado" pra qualquer incerteza.

════════════════════════════════════════════════════════
REGRA DE DETERMINISMO — UM ITEM POR BLOCO ÚNICO
════════════════════════════════════════════════════════

Ao gerar os itens a partir de "CONTAGEM DE BLOCOS":
- Cada nome de bloco único = **um item só** na planilha, com a quantidade literal da contagem. Não reagrupar, não dividir, não combinar blocos diferentes.
- Se a contagem listou "lum R4 remanejada: 20 un" e "lum R4 nova: 2 un", gerar DOIS itens separados com essas quantidades exatas. NÃO inferir que "são ambos R4" e somar, nem dividir um único em sub-itens por intuição.
- Se um bloco tem nome genérico/estranho ("BLOCO1", "INSERT_0"), mantenha — marcar como estimado pra o usuário identificar.
- A descrição do item pode ser enriquecida (modelo, fabricante) mas o IDENTIFICADOR e a QUANTIDADE são literais do DXF.

Essa regra garante que subir o mesmo arquivo duas vezes retorne o MESMO resultado.

════════════════════════════════════════════════════════
UNIDADES CORRETAS — ml VS un VS m² VS vb
════════════════════════════════════════════════════════

**A unidade vem do TIPO DE DADO no DXF, não do que parece intuitivo:**
- Valor vindo de "COMPRIMENTOS POR LAYER" → **ml** (metro linear). Use o número LITERAL.
- Valor vindo de "ÁREAS HACHURADAS POR LAYER" → **m²**. Use o número LITERAL.
- Valor vindo de "CONTAGEM DE BLOCOS" → **un**. Use o número LITERAL.
- Valor vindo de "COTAS/DIMENSÕES" → **m** (metro simples). Use o número LITERAL.
- Verba sem medida clara → **vb** com quantity=0 (laranja, usuário preenche).

REGRA GERAL: a unidade vai PARTE COM O DADO EXTRAÍDO. Se você usa um valor de
"ÁREAS HACHURADAS" e coloca unidade "ml", está errado — o valor é m² por definição.

CASO ESPECIAL — lineares aparecem 2× no DXF:
- "LED LINE 45°" ou "perfil linear" podem aparecer em **CONTAGEM DE BLOCOS** (como
  "un") E em **COMPRIMENTOS POR LAYER** (como "ml"). São o mesmo item físico.
  Escolha o COMPRIMENTO (ml), não a contagem, pois é a medida útil pra orçar.
  Exemplo: "LUMINI LED LINE: 2 un" + "layer LUM-LINE: 23.24 m" → gera item com
  ml=23.24 (não un=2).
- Rodapés, tabicas, eletrocalhas, perfis: sempre ml.

CASO ESPECIAL — INFRA LINEAR (eletroduto, eletrocalha, perfilado, leito de
cabos, conduto, tubulação): a ESPECIFICAÇÃO quase sempre está num layer de
TEXTO (ex.: ELE-TEXTOS, ELE-CHAMADA, DI-Textos) e o DESENHO num layer próprio
de condutos com OUTRO nome (ex.: EL-Condutos, ELE-T-EMBUTIDO, LO-Condutos,
TV-Condutos, HID-TUB, nomes com CONDU/ELETR/PERFIL/TUB/CALHA).
- Antes de zerar a quantidade de um item desses, PROCURE em "COMPRIMENTOS POR
  LAYER" um layer de condutos compatível e USE o comprimento como quantity em
  **ml**, marcando **estimado** e citando o layer e o valor na observação.
  Deixar quantity=0 com um layer de condutos MEDIDO na lista é jogar medição fora.
- Se VÁRIOS diâmetros/bitolas compartilham o mesmo layer de condutos, NÃO
  repita o comprimento em cada item (dupla contagem!): ponha o TOTAL no item
  mais genérico com observação "dividir por bitola na revisão" e deixe os
  demais com quantity=0 apontando pra esse item.
- Layers de CHAMADA/TEXTO/COTA/LEGENDA/TITULO **não são condutos** — o
  comprimento deles é de setas e letras; NUNCA use como quantidade.

PRIORIDADE SEMÂNTICA — a unidade do ITEM é definida pelo TIPO DE SERVIÇO, não só
pelo dado extraído:
- **Pisos** (carpete, cerâmica, porcelanato, vinílico, laminado, madeira) → **SEMPRE m²**
  mesmo que o dado do DXF venha como comprimento de polyline. Pisos se orçam por área.
- **Forros** (modular, gesso, ripado) → **SEMPRE m²**.
- **Pinturas e revestimentos verticais** (parede, azulejo, tijolinho, papel) → **m²**.
- **Rodapés, tabicas, soleiras, perfis, eletrocalhas, molduras** → **ml** (linear).
- **Luminárias, portas, difusores, interruptores, tomadas** → **un** (contagem).

Se o DXF te der um comprimento (em COMPRIMENTOS POR LAYER) pra uma SUPERFÍCIE
como piso/forro/pintura, esse comprimento provavelmente é o perímetro da área,
não uma medida linear pra orçar. Nesse caso, procure a área correspondente em
"ÁREAS HACHURADAS" e use m². Se não houver área hachurada, marque o item como
"estimado" com `quantity=0` e pede confirmação na observação.

════════════════════════════════════════════════════════
QUANDO MARCAR "estimado" (LARANJA)
════════════════════════════════════════════════════════

**NÃO seja tímido com "confirmado" quando o DADO EXISTE.** Se o DXF tem um
comprimento somado de polylines de uma layer (ex.: layer de divisória com um
valor em metros), use esse valor como "confirmado" — é uma medição objetiva
direta do arquivo.

Você só deve marcar "estimado" (laranja) nos casos abaixo:

(a) Quantidade você INFERIU de texto/contexto, não leu direto:
    - "demolir X" → qtd=1 (texto não numérico)
    - "administração local 2 meses" (não vem do arquivo)
    - "retrofit de lâmpadas existentes" (sem contagem)
    → **quantity=0** (vazio na planilha, usuário preenche)

(b) Clara dupla contagem de áreas LEV vs FOR:
    Se "ÁREAS HACHURADAS" tem AMBOS "LEV-X: A m²" E "FOR-X: B m²" e você
    está tentando orçar o mesmo tipo, escolha apenas UM (geralmente o
    "FOR-*" / "NOV-*" / "ARQ-*" = novo projeto) e coloque confirmado.
    Se ficar em dúvida entre os dois, marque estimado.

(c) Área total > área da laje (impossível):
    Se a soma de áreas de um tipo fica > "Área construída" das PREMISSAS,
    suspeite de dupla contagem e marque estimado.

**NÃO marque estimado só por precaução em valores medidos.** Se COMPRIMENTOS
POR LAYER te dá um valor em metros pra uma layer de projeto, use esse valor
como confirmado. Não desconfie do número só porque vem de soma de linhas —
medir soma de linhas É a medição.

════════════════════════════════════════════════════════
LEV- vs FOR- — convenção de layers em reforma (regra B acima)
════════════════════════════════════════════════════════

Em projetos BR:
- "LEV-" = LEVANTAMENTO (existente no imóvel). Não orçar, exceto como demolição.
- "FOR-" / "NOV-" / "ARQ-" = PROJETO NOVO (a construir).
- "DEM-" = DEMOLIÇÃO.

**Se aparecerem AMBOS (LEV e FOR do mesmo tipo)**, use só o FOR/NOV.
**Se aparecer SÓ UM**, use-o (pode ser o único dado disponível).
**Nunca some LEV + FOR** — são momentos distintos (antes/depois da reforma).

════════════════════════════════════════════════════════
FORMATO DE RESPOSTA — RACIOCÍNIO EXPLÍCITO ANTES DO JSON
════════════════════════════════════════════════════════

Antes de retornar o JSON, PENSE em voz alta em 4 passos. O texto do raciocínio
é obrigatório — ele ajuda você a errar menos e ajuda o revisor humano a
confiar no resultado. Formato:

```
RACIOCÍNIO:

Passo 1 — Inventário de layers:
  Para cada LAYER relevante, uma linha:
    "<nome do layer>: <tipo de dado> — <quantidade extraída> — representa <item>"
  Use os nomes de layer DESTE arquivo. Ignore layers de anotação, xrefs e aux.

Passo 2 — Checagem de LEV vs FOR:
  Liste pares conflitantes (mesmo tipo em layer LEV/existente e FOR/novo).
  Para cada par, diga qual escolheu e por quê (geralmente o layer de projeto
  NOVO vence sobre o de levantamento do existente).

Passo 3 — Plausibilidade:
  Para cada grupo de itens, verifique se a soma faz sentido:
  - Áreas de piso/forro somadas ≤ área da laje construída
  - Contagens (un) são números inteiros
  - Comprimentos e áreas são positivos
  Se algo parece absurdo, marque estimado.

Passo 4 — Geração dos itens:
  Para cada item que vai no JSON, uma linha:
    "<descrição>: <qtd> <un> — fonte: <layer/bloco exato>"
```

Depois do raciocínio, retorne o JSON em bloco de código (```json...```):

{{
  "project_data": {{
    "name": "",
    "total_area": 0,
    "layout_area": 0,
    "workstations": 0,  // apenas pra ESCRITÓRIOS; em residencial/clínica etc deixe 0
    "departments": [],  // apenas pra escritório/escola; em residencial use new_rooms
    "demolition_notes": [],
    "new_rooms": [],
    "kept_elements": []
  }},
  "items": [
    {{
      "item_num": "1",
      "description": "Descrição completa",
      "unit": "m²",
      "quantity": 100,
      "observations": "Fonte: <layer/bloco/texto exato da extração>",
      "ref_sheet": "deixe VAZIO — o código preenche com o nome real do arquivo",
      "confidence": "confirmado ou estimado — nunca inventar",
      "discipline": "Categoria"
    }}
  ]
}}

REGRA DA OBSERVATION: o campo "observations" deve SEMPRE citar a fonte exata
do número no DXF, usando os nomes de layer/bloco desta prancha — ex.:
"Fonte: <N> INSERTs do bloco '<nome_bloco_real>'" ou "Fonte: área hachurada
do layer '<nome_layer_real>' = <valor> m²". Nunca invente nomes de layer ou
bloco — só cite os que estão no inventário deste arquivo."""

                    try:
                        # STREAMING: resposta longa (CoT + JSON grande de planta
                        # grande) estourava o timeout no modo create() não-streaming
                        # e falhava como "IA sobrecarregada". Stream não estoura.
                        from llm_retry import call_with_retry_stream as _llm_retry
                        # Modelo da extração configurável via env (A/B sem deploy).
                        # Default Opus 4.8 (teste de precisão vs Sonnet 4.6). Atenção:
                        # Opus 4.7/4.8 e Fable NÃO aceitam `temperature` (dá 400) —
                        # só mando pros modelos que aceitam (Sonnet/Haiku).
                        _dxf_model = os.environ.get("DXF_EXTRACT_MODEL", "claude-sonnet-4-6")
                        _dxf_kwargs = dict(
                            tag=f"dxf:{os.path.basename(dxf_path)}",
                            model=_dxf_model,
                            max_tokens=32000,  # CoT + JSON de planta GRANDE (16k truncava o JSON -> 0 itens)
                            # SYSTEM_PROMPT (~4,4k tok) é o mesmo em toda prancha →
                            # cacheado, custa ~90% menos na leitura (23/07).
                            cache_system=True,
                            system=(SYSTEM_PROMPT_ESTRUTURA if is_structural else SYSTEM_PROMPT),
                            messages=[{"role": "user", "content": dxf_prompt}],
                        )
                        if not any(_t in _dxf_model for _t in ("opus-4-8", "opus-4-7", "fable")):
                            _dxf_kwargs["temperature"] = 0  # determinismo (Sonnet/Haiku aceitam)
                        response = _llm_retry(dxf_client, **_dxf_kwargs)

                        text = response.content[0].text
                        # #7 sinal de 1ª classe: resposta cortada no teto (max_tokens) =
                        # leitura possivelmente INCOMPLETA, mesmo que o JSON ainda parseie ok.
                        _dxf_truncado = _response_truncated(getattr(response, "stop_reason", ""))
                        # Parser robusto: agora o Claude pode retornar raciocínio ANTES do JSON (CoT).
                        # Tentar em ordem: bloco ```json, bloco ```, último objeto {...} do texto.
                        json_str = None
                        if "```json" in text:
                            json_str = text.split("```json")[-1].split("```")[0].strip()
                        elif "```" in text:
                            # Pegar o último bloco de código (caso tenha múltiplos)
                            parts = text.split("```")
                            if len(parts) >= 3:
                                json_str = parts[-2].strip()
                        if not json_str or not json_str.startswith("{"):
                            # Fallback: regex pra achar último JSON object "compatível" no texto
                            import re as _re_parse
                            # Matches { ... } que contém "items" ou "project_data"
                            candidates = _re_parse.findall(
                                r"\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}",
                                text, flags=_re_parse.DOTALL
                            )
                            for cand in reversed(candidates):
                                if '"items"' in cand or '"project_data"' in cand:
                                    json_str = cand
                                    break
                        if not json_str:
                            json_str = text.strip()

                        try:
                            result = _j.loads(json_str)
                        except Exception as _je:
                            # JSON truncado (resposta cortada no teto de tokens numa
                            # planta grande): em vez de perder TUDO, recupera os itens
                            # completos. Antes isso virava 0 itens + erro "sobrecarregada"
                            # enganoso (caso Ademir, DWG de prédio público).
                            result = _salvage_truncated_json(json_str)
                            print(f"DXF JSON truncado ({type(_je).__name__}: {_je}); "
                                  f"salvados {len(result.get('items', []))} itens")
                            _dxf_truncado = True  # JSON quebrou = quase sempre corte no teto

                        # Robustez: array cru [...] vira {"items":[...]} (engine_rules,
                        # testado). Evita 'list object has no attribute get'.
                        result = _normalize_items_payload(result)

                        # #7 leitura possivelmente INCOMPLETA (corte no teto via stop_reason
                        # OU JSON truncado): avisa — não entrega parcial calado (caso Ademir).
                        # Aviso único, dedup por arquivo.
                        if _dxf_truncado:
                            try:
                                _n_salv = len(result.get("items", []))
                                project_data.warnings = (project_data.warnings or []) + [
                                    f"A leitura de '{os.path.basename(dxf_path)}' pode estar INCOMPLETA "
                                    f"(a resposta da IA foi cortada por tamanho — li {_n_salv} itens, mas "
                                    f"pode faltar algum). Reprocessar pode completar a planilha."
                                ]
                            except Exception:
                                pass

                        # Extrair project_data
                        if "project_data" in result:
                            pd = result["project_data"]
                            # Acumula em vez de sobrescrever — resolvido via consenso depois
                            for _fld in ("total_area", "layout_area", "no_intervention_area"):
                                _v = pd.get(_fld)
                                if _v:
                                    _vf = sf(_v)
                                    if _vf > 0:
                                        _area_readings[_fld].append(_vf)
                            if pd.get("name") and not project_data.name: project_data.name = pd["name"]
                            if pd.get("demolition_notes"): project_data.demolition_notes.extend(pd["demolition_notes"])
                            if pd.get("new_rooms"): project_data.new_rooms.extend(pd["new_rooms"])
                            if pd.get("kept_elements"): project_data.kept_elements.extend(pd["kept_elements"])

                        # Extrair itens
                        for item_data in result.get("items", []):
                            try:
                                desc = item_data.get("description", "")
                                if not desc or len(desc) < 3: continue
                                discipline = item_data.get("discipline", "Complementares")
                                conf = item_data.get("confidence", "estimado")
                                if conf not in ["confirmado", "estimado", "verificar"]: conf = "estimado"
                                # trava de procedência: extração com ressalva → nunca confirmado
                                if _dxf_sem_procedencia and conf == "confirmado":
                                    conf = "estimado"
                                    item_data["_procedencia_rebaixada"] = True
                                qty = sf(item_data.get("quantity", 0))
                                # qty=0 é permitido para itens estimados sem número concreto
                                # (virará vazio na planilha pro usuário preencher).
                                # Só forçamos qty=1 em CONFIRMADO (deveria ter número real).
                                if qty < 0:
                                    qty = 0
                                if qty == 0 and conf == "confirmado":
                                    qty = 1  # defensivo: confirmado sem número cai em vb=1

                                # Normalização pós-IA: força unidade correta pra descrição
                                # (ex.: "piso vinílico" sempre m², nunca ml).
                                original_unit = item_data.get("unit", "vb")
                                normalized_unit, unit_corrected = _normalize_unit_for_item(desc, original_unit)
                                obs_raw = item_data.get("observations", "Fonte: DXF")
                                if unit_corrected:
                                    # IA escolheu unidade inconsistente com o tipo — marcar estimado
                                    # pra usuário conferir o número
                                    conf = "estimado"
                                    obs_raw = (f"{obs_raw} | Unidade ajustada de {original_unit} "
                                               f"para {normalized_unit} (revisar quantidade)")
                                if item_data.get("_procedencia_rebaixada"):
                                    obs_raw = (f"{obs_raw} | Procedência: extração com ressalva "
                                               f"(estéril/unidade/xref) — quantidade não confirmada, revisar").strip(" |")

                                # CROSS-CHECK determinístico (opt-in via env DXF_CONFIRM_CROSSCHECK):
                                # promove 'estimado' → 'confirmado' SÓ quando TODAS batem:
                                #  (a) o item cai numa categoria física clara (piso/forro/parede/...);
                                #  (b) a unidade corresponde ao tipo da categoria (m²=área, m=linear);
                                #  (c) a quantidade bate byte-a-byte com uma medida dura DAQUELA
                                #      categoria (não de outro elemento — furo C);
                                #  (d) a medida é INEQUÍVOCA — não colide com outra categoria
                                #      (forro≈piso não promove nenhum dos dois).
                                # Sem ressalva, sem unidade ajustada, qty>0. Conserta a "timidez"
                                # da IA sem falso-medido. Revisão adversarial 15/07.
                                if (_XCHECK_ON and conf == "estimado" and not _dxf_sem_procedencia
                                        and not unit_corrected and qty > 0):
                                    _q2 = round(qty, 2)
                                    _icat = _item_geo_category(desc, discipline)
                                    _promo = False
                                    if _icat in _AREA_CATS and normalized_unit in ("m²", "m2"):
                                        if (_q2 in _hard_area_by_cat.get(_icat, set())
                                                and _measure_unambiguous(_q2, _icat, _hard_area_by_cat)):
                                            _promo = True
                                    elif _icat in _LEN_CATS and normalized_unit in ("m", "ml"):
                                        if (_q2 in _hard_len_by_cat.get(_icat, set())
                                                and _measure_unambiguous(_q2, _icat, _hard_len_by_cat)):
                                            _promo = True
                                    if _promo:
                                        conf = "confirmado"
                                        obs_raw = (f"{obs_raw} | Medido: confere com a geometria "
                                                   f"de {_icat} extraída ({_q2} {normalized_unit})").strip(" |")

                                # REDE DE SEGURANÇA (Finding 1, SEMPRE ligada): se a IA marcou
                                # 'confirmado' num m² que é a SOMA de um layer com várias
                                # hachuras (acabamento possivelmente misto), rebaixa pra
                                # estimado. Roda depois de tudo — só rebaixa, nunca cria medido.
                                if (conf == "confirmado" and normalized_unit in ("m²", "m2")
                                        and round(qty, 2) in _multi_hatch_sums):
                                    conf = "estimado"
                                    obs_raw = (f"{obs_raw} | Revisar: área = soma de várias hachuras "
                                               f"no mesmo layer (possível acabamento misto) — confira "
                                               f"o valor por ambiente").strip(" |")

                                item = BudgetItem(
                                    item_num=str(item_data.get("item_num", "")),
                                    description=desc,
                                    unit=normalized_unit,
                                    quantity=qty,
                                    observations=obs_raw,
                                    ref_sheet=_dxf_nome,
                                    confidence=Confidence(conf),
                                    discipline=discipline,
                                    origem="dxf_geom",  # medido por geometria ezdxf
                                )
                                dxf_items.append(item)
                            except: continue

                        print(f"DXF {os.path.basename(dxf_path)}: {len(result.get('items', []))} itens extraídos via Claude")

                        # ── CHECKPOINT DXF: guarda os itens + contribuições desta
                        # prancha pra retomada não recomeçar do zero. Só chega aqui
                        # se a IA e o parse deram certo (falha cai no except abaixo,
                        # sem checkpoint → re-tenta na retomada, como o PDF já fazia).
                        try:
                            def _ser_ck(_it):
                                _cf = getattr(_it, "confidence", None)
                                return {
                                    "item_num": getattr(_it, "item_num", "") or "",
                                    "description": getattr(_it, "description", "") or "",
                                    "unit": getattr(_it, "unit", "") or "vb",
                                    "quantity": float(getattr(_it, "quantity", 0) or 0),
                                    "observations": getattr(_it, "observations", "") or "",
                                    "ref_sheet": getattr(_it, "ref_sheet", "DXF") or "DXF",
                                    "confidence": getattr(_cf, "value", None) or str(_cf or "estimado"),
                                    "discipline": getattr(_it, "discipline", "Complementares") or "Complementares",
                                    "origem": getattr(_it, "origem", "dxf_geom") or "dxf_geom",
                                }
                            _cp_payload = {
                                "_kind": "dxf",
                                "items": [_ser_ck(it) for it in dxf_items[_snap["items"]:]],
                                "xref_warnings": list(xref_warnings[_snap["xref"]:]),
                                "pd_warnings": list((project_data.warnings or [])[_snap["warn"]:]),
                                "new_rooms": list((project_data.new_rooms or [])[_snap["rooms"]:]),
                                "kept_elements": list((project_data.kept_elements or [])[_snap["kept"]:]),
                                "demolition_notes": list((project_data.demolition_notes or [])[_snap["demo"]:]),
                                "area_readings": {_k: list(_area_readings[_k][_snap["area"].get(_k, 0):]) for _k in _area_readings},
                                "name": (project_data.name if (not _snap["name"] and project_data.name) else ""),
                            }
                            _ckpt_save(job_id, _dxf_raw_stem, _cp_payload)
                        except Exception as _cpe:
                            print(f"[ckpt] save DXF {os.path.basename(dxf_path)} falhou (segue): {_cpe}")

                    except Exception as e:
                        # Falha de IA neste DXF (sobrecarga/timeout/JSON inválido).
                        # Registra em dxf_errors pra o guard de 0-itens distinguir
                        # "IA falhou (reprocessável grátis)" de "arquivo sem
                        # conteúdo medível" — espelha o tratamento do caminho PDF.
                        # Preserva status_code E nome da classe do erro-raiz pra
                        # classificar por tipo (não re-adivinhar por substring). O
                        # type= pega erro de rede cru do streaming (BrokenPipe/
                        # RemoteProtocol) que vem sem status mas é transitório.
                        _status = getattr(e, "status_code", None)
                        _bits = ([f"status={_status}"] if _status is not None else []) + [f"type={type(e).__name__}"]
                        err_msg = f"[{' '.join(_bits)}] {str(e)[:200]}"
                        dxf_errors.append(f"{os.path.basename(dxf_path)}: {err_msg}")
                        jobs.update_field(job_id, current_step=f"Erro IA (DXF): {err_msg}")
                        print(f"Erro Claude DXF: {e}")
                        _log_error("dxf:claude", f"{os.path.basename(dxf_path)}: {e}", job_id)

                    del structured_text
                    try:
                        del extraction
                    except Exception:
                        pass
                    # Libera o DXF convertido (temp dir arq_dxf_*) desta prancha —
                    # senão os N DXFs convertidos acumulam no /tmp do dyno durante o
                    # job (leak achado na auditoria 06/07; caso Thamiry, 22 DWGs). O
                    # preview lê de work_dir, não desses temp dirs, então é seguro.
                    try:
                        _dxf_dir = os.path.dirname(dxf_path)
                        if os.path.basename(_dxf_dir).startswith("arq_dxf_"):
                            shutil.rmtree(_dxf_dir, ignore_errors=True)
                    except Exception:
                        pass
                    gc.collect()

            except Exception as e:
                jobs.update_field(job_id, error_message=f"Erro extração DXF: {str(e)[:500]}")
                jobs.update_field(job_id, current_step=f"ERRO DXF: {str(e)[:200]}")
                import traceback
                traceback.print_exc()
                raise  # Deixar o erro aparecer

        total = len(pdf_paths)
        client = anthropic.Anthropic(api_key=api_key, timeout=300.0)
        all_items = list(dxf_items)  # Começar com itens DXF
        crops_dir = os.path.join(work_dir, "crops")
        os.makedirs(crops_dir, exist_ok=True)

        # Coleta erros de IA por prancha. analyze_sheet NÃO lança exceção quando
        # a chamada Claude falha (timeout/sobrecarga/JSON inválido) — retorna
        # {"items": [], "error": "..."}. Sem coletar isso aqui, o job terminava
        # "done" com planilha vazia e o usuário não sabia que houve falha.
        # Bug Vinícius (2026-05-21): PDF processou em 17s, 0 itens, status done.
        sheet_errors: list[str] = []

        # Ordenar PDFs por prioridade (layout primeiro)
        priority = {"layout_novo": 0, "layout_atual": 1, "demolir": 2, "arquitetura": 3,
                     "forro": 4, "piso": 5, "pontos": 6, "mobiliario": 7, "marcenaria": 8,
                     "det_forro": 9, "detalhe_ambiente": 10}

        pdf_infos = []
        for pdf_path in pdf_paths:
            filename = os.path.basename(pdf_path)
            # Prioridade: escolha manual do usuário no dropdown > detecção por nome.
            # Sem isso, nomes arbitrários tipo "225.AFS.700.DET BANH SUITE" caíam
            # em DESCONHECIDO ou eram misclassificados.
            user_st_str = user_sheet_types.get(pdf_path, "")
            if user_st_str:
                try:
                    sheet_type = SheetType(user_st_str)
                except ValueError:
                    sheet_type = identify_sheet_type(filename)
            else:
                sheet_type = identify_sheet_type(filename)
            pdf_infos.append((pdf_path, filename, sheet_type))

        pdf_infos.sort(key=lambda x: priority.get(x[2].value, 99))

        # Agrupar pranchas IRMÃS: pranchas do mesmo ambiente (ex.: DET LAVABO
        # pode ter planta baixa + elevações em PDFs separados). Saber que
        # existem irmãs evita warning de "prancha órfã" e permite o prompt
        # cruzar informação entre elas (ex.: legenda de códigos 01-14 que
        # está no PDF da planta mas não nas elevações).
        from processor import identify_ambiente as _id_amb_sib
        _siblings_map: dict[str, list[str]] = {}
        for _p, _fn, _st in pdf_infos:
            if _st == SheetType.DETALHE_AMBIENTE:
                _amb = user_ambientes.get(_p, "") or _id_amb_sib(_fn)
                if _amb:
                    _siblings_map.setdefault(_amb, []).append(_fn)

        # Faixa de progresso reservada para PDFs: após cad_end_pct (se houver CAD) até 90%
        # 🪤 `has_cad` diz "veio DWG/DXF no envio", NÃO "sobrou CAD pra analisar".
        # Quando o DWG não converte (caso AEC/MEP), a faixa de análise CAD não roda —
        # e reservá-la mesmo assim dava 40 pontos de barra por trabalho que não
        # aconteceu: o cliente via "55%" ainda na prancha 1/7 (caso Walter 29/07).
        # Só reserva a faixa se de fato houve DXF pra analisar.
        _cad_analisou = bool(dxf_paths)
        pdf_start_pct = (cad_end_pct if _cad_analisou else conv_end_pct) if has_cad else 5
        pdf_end_pct = 90
        pdf_span = pdf_end_pct - pdf_start_pct

        # ── Explode PDF multi-página em unidades de página ──
        # Um "PROJETO EXECUTIVO.pdf" traz VÁRIAS pranchas num arquivo só. Antes o
        # motor lia só a página 1 (às vezes a capa!) e ainda estourava a memória
        # varrendo o texto de todas as páginas de uma vez. Agora cada página vira
        # uma prancha, processada uma de cada vez com memória limitada. Caso real
        # sumi/lia (06/07): PDF de 13 MB derrubava o Render de 2 GB toda tentativa.
        from processor import pdf_page_count as _pdf_pages
        MAX_PAGES_PER_PDF = 40
        page_units = []  # (pdf_path, filename, sheet_type, page_index, page_count)
        for _p, _fn, _st in pdf_infos:
            _npg = _pdf_pages(_p)
            for _pi in range(min(_npg, MAX_PAGES_PER_PDF)):
                page_units.append((_p, _fn, _st, _pi, _npg))
            if _npg > MAX_PAGES_PER_PDF:
                try:
                    project_data.warnings = (project_data.warnings or []) + [
                        f"'{_fn}' tem {_npg} páginas — lemos as primeiras {MAX_PAGES_PER_PDF}. "
                        f"Pro restante, divida o arquivo em partes ou envie em DXF."
                    ]
                except Exception:
                    pass
        total = len(page_units)

        # CHECKPOINT: o cache (_ckpt_cache) já foi carregado UMA vez lá em cima,
        # antes do bloco DXF — serve os dois loops (DXF e PDF). Em RETOMADA
        # automática, prancha já analisada pula extração+IA; reprocesso MANUAL
        # ignora de propósito (roda o motor novo na prancha inteira).
        for i, (pdf_path, filename, sheet_type, page_index, page_count) in enumerate(page_units):
            # 🛡️ Freio de MEMÓRIA (idem loop DXF): aborta limpo antes do OOM,
            # mantendo o servidor de pé pros outros clientes.
            if i > 0 and _mem_pressure(0.85):
                _abort_job_mem(job_id, i, total)
                return
            _disp = filename if page_count <= 1 else f"{filename} · pág {page_index+1}/{page_count}"
            step_pct = pdf_start_pct + int((i / max(total, 1)) * pdf_span)
            jobs.update_field(job_id, progress=step_pct)
            jobs.update_field(job_id, current_step=f"Prancha {i+1}/{total}: {_disp}")

            # Se o auto-detect falhou E o usuário não classificou manualmente:
            # - se o nome contém palavra-chave de AMBIENTE (banh, cozinha, lavabo,
            #   suíte, etc), trata como DETALHE_AMBIENTE e extrai o ambiente do nome.
            # - senão, usa ARQUITETURA como prompt-genérico de fallback.
            if sheet_type == SheetType.DESCONHECIDO:
                from processor import identify_ambiente as _id_amb2
                _auto_amb = _id_amb2(filename)
                if _auto_amb:
                    print(f"[fallback] {filename}: detectado ambiente '{_auto_amb}' — usando PROMPT_DETALHE_AMBIENTE")
                    sheet_type = SheetType.DETALHE_AMBIENTE
                    # Guarda auto-detect no mapa (vai ser lido por user_ambientes.get)
                    user_ambientes.setdefault(pdf_path, _auto_amb)
                    jobs.update_field(job_id, current_step=f"Prancha {i+1}/{total}: {filename} (detalhe de {_auto_amb})")
                else:
                    print(f"[fallback] {filename}: tipo não detectado, usando PROMPT_ARQUITETURA como genérico")
                    sheet_type = SheetType.ARQUITETURA
                    jobs.update_field(job_id, current_step=f"Prancha {i+1}/{total}: {filename} (tipo não identificado — extração genérica)")

            _stem = f"{os.path.splitext(os.path.basename(pdf_path))[0]}_p{page_index}"
            _ck_key = _sanitize_filename_for_storage(_stem)
            if _ck_key in _ckpt_cache:
                # CHECKPOINT: prancha já analisada antes do restart — reusa o
                # resultado salvo (pula texto, crops e IA).
                result = _ckpt_cache[_ck_key]
                jobs.update_field(job_id, current_step=f"Prancha {i+1}/{total}: {_disp} — análise já concluída, retomando ✓")
                print(f"[ckpt] {_stem}: análise reaproveitada do checkpoint")
            else:
                # 1. Extrair texto (só da página desta unidade — leve, bounded)
                text = extract_text(pdf_path, page_index)

                # 2. Renderizar crops (1 página de cada vez; stem único por página)
                crop_paths = render_crops(pdf_path, sheet_type, crops_dir,
                                          page_index=page_index, out_stem=_stem)

                # 2.5 PROMOÇÃO DO VETORIAL (19/07): mede a página com o motor
                # vetorial ANTES da IA. SÓ injeta quando a escala foi VALIDADA
                # POR COTA da própria prancha (≥2 cotas independentes batendo
                # ±2% com elementos medidos) — sem prova, comportamento antigo.
                # Subprocess com timeout: página pesada não trava o job.
                _vet_secao = ""
                try:
                    import subprocess as _sp, json as _jv, sys as _sysv
                    _cmd = [_sysv.executable, "-c", (
                        "import sys, json; sys.path.insert(0, r'" +
                        os.path.dirname(os.path.abspath(__file__)) + "'); "
                        "from pdf_vector import _measure_page; "
                        f"print(json.dumps(_measure_page(r'{pdf_path}', {page_index}, '')))"
                    )]
                    _pr = _sp.run(_cmd, capture_output=True, text=True, timeout=75)
                    _vm = _jv.loads(_pr.stdout.strip().splitlines()[-1]) if _pr.returncode == 0 and _pr.stdout.strip() else {}
                    if _vm.get("escala_validada") and (_vm.get("n_rooms") or _vm.get("walls_m")):
                        _linhas = [
                            "",
                            "=== MEDIÇÕES VETORIAIS DA PRANCHA (validadas por cota) ===",
                            f"Escala 1:{_vm.get('scale')} confirmada por {_vm.get('cotas_batem')} cota(s) da própria prancha.",
                        ]
                        if _vm.get("n_rooms"):
                            _linhas.append(f"Ambientes MEDIDOS geometricamente: {_vm['n_rooms']} somando {_vm.get('rooms_m2', 0)} m² (maiores: {_vm.get('top_rooms')} m²).")
                        if _vm.get("walls_m"):
                            _linhas.append(f"Paredes/divisórias MEDIDAS: {_vm['walls_m']} m ({_vm.get('n_walls')} segmentos).")
                        _linhas.append(
                            "REGRA: itens de ÁREA (piso/forro/pintura de área) e COMPRIMENTO de parede "
                            "baseados EXATAMENTE nestes valores podem sair com confidence 'confirmado' "
                            "(medição geométrica validada). Qualquer outro valor segue 'estimado'.")
                        _vet_secao = "\n".join(_linhas)
                        print(f"[pdfvec-promo] {_stem}: escala validada por cota — seção injetada")
                except Exception as _ve:
                    print(f"[pdfvec-promo] {_stem}: sem promoção ({_ve})")

                # 3. Analisar com IA
                jobs.update_field(job_id, current_step=f"Prancha {i+1}/{total}: Nossa IA está analisando {_disp}...")
                sheet = SheetInfo(
                    filename=filename,
                    sheet_type=sheet_type,
                    text_content=text[:5000] + _vet_secao,
                    crops=crop_paths,
                )
                # Ambiente: user escolheu manualmente > auto-detect por keyword
                from processor import identify_ambiente as _id_amb
                amb_for_sheet = user_ambientes.get(pdf_path, "") or (
                    _id_amb(filename) if sheet_type == SheetType.DETALHE_AMBIENTE else ""
                )
                # Irmãs: outras pranchas do mesmo ambiente (ex: LAVABO em 704+705)
                siblings = []
                if amb_for_sheet and amb_for_sheet in _siblings_map:
                    siblings = [s for s in _siblings_map[amb_for_sheet] if s != filename]
                result = analyze_sheet(client, sheet, typology=typology,
                                       ambiente=amb_for_sheet, siblings=siblings,
                                       is_structural=is_structural)
                # Salva o checkpoint da prancha (só análise SEM erro — erro
                # deve re-tentar na retomada, não ser reaproveitado)
                if not result.get("error"):
                    _ckpt_save(job_id, _stem, result)

            # 3b. Capturar falha de IA nesta prancha (não interrompe o loop —
            # outras pranchas podem ter sucesso — mas registra pra decidir o
            # status final do job se nada for extraído).
            if result.get("error"):
                err_msg = str(result.get("error"))[:200]
                sheet_errors.append(f"{filename}: {err_msg}")
                print(f"[analyze-erro] {filename}: {err_msg}")
                _log_error("pdf:analyze", f"{filename}: {err_msg}", job_id)

            # #7 leitura possivelmente INCOMPLETA nesta prancha (resposta da IA
            # cortada no teto). Avisa — não entrega parcial calado (caso Ademir).
            if result.get("_truncated"):
                try:
                    project_data.warnings = (project_data.warnings or []) + [
                        f"A leitura de '{filename}' pode estar INCOMPLETA (a resposta da IA "
                        f"foi cortada por tamanho; pode faltar item). Reprocessar pode completar."
                    ]
                except Exception:
                    pass

            # 4. Extrair dados do projeto
            if "project_data" in result:
                pd = result["project_data"]
                # Acumula em vez de sobrescrever — consenso ao final do loop
                for _fld in ("total_area", "layout_area", "no_intervention_area"):
                    _v = pd.get(_fld)
                    if _v:
                        _vf = sf(_v)
                        if _vf > 0:
                            _area_readings[_fld].append(_vf)
                if pd.get("workstations"):
                    try: project_data.workstations = int(float(str(pd["workstations"]).replace('un','').strip()))
                    except: pass
                if pd.get("departments"): project_data.departments = pd["departments"]
                if pd.get("demolition_notes"): project_data.demolition_notes.extend(pd["demolition_notes"])
                if pd.get("new_rooms"): project_data.new_rooms.extend(pd["new_rooms"])
                if pd.get("kept_elements"): project_data.kept_elements.extend(pd["kept_elements"])
                if pd.get("name") and not project_data.name: project_data.name = pd["name"]
                if pd.get("address") and not project_data.address: project_data.address = pd["address"]
                if pd.get("architect") and not project_data.architect: project_data.architect = pd["architect"]

            # 5. Extrair itens
            valid_disciplines = [
                "Serviços Preliminares", "Demolição e Remoção", "Fechamentos Verticais",
                "Revestimentos", "Pisos e Rodapés", "Forros", "Portas e Ferragens",
                "Divisórias e Vidros", "Persianas e Cortinas", "Iluminação",
                "Instalações Elétricas e Dados",
                "Instalações Hidráulicas", "Instalações de Gás",
                "Ar-Condicionado", "Incêndio e Segurança",
                "Marcenaria", "Mobiliário", "Estrutura", "Complementares"
            ]
            for item_data in result.get("items", []):
                try:
                    desc = item_data.get("description", "")
                    if not desc or len(desc) < 3: continue
                    discipline = item_data.get("discipline", "Complementares")
                    if discipline not in valid_disciplines: discipline = "Complementares"
                    conf = item_data.get("confidence", "estimado")
                    if conf not in ["confirmado", "estimado", "verificar"]: conf = "estimado"
                    # PDF/Vision NÃO mede geometria — só lê número numa imagem.
                    # Pela regra dura, "medido do CAD" (branco) vem só de geometria
                    # (ezdxf). "confirmado" vindo de PDF é rebaixado pra estimado.
                    _pdf_downgrade = (conf == "confirmado")
                    if _pdf_downgrade:
                        conf = "estimado"
                    qty_raw = item_data.get("quantity", 0)
                    qty = sf(qty_raw) if qty_raw else 0
                    # qty=0 permitido em "estimado" (usuário preenche); forçar 1 só em confirmado
                    if qty < 0:
                        qty = 0
                    if qty == 0 and conf == "confirmado":
                        qty = 1

                    # Normalização pós-IA: força unidade consistente com descrição
                    original_unit = item_data.get("unit", "vb")
                    normalized_unit, unit_corrected = _normalize_unit_for_item(desc, original_unit)
                    obs_raw = item_data.get("observations", "")
                    if unit_corrected:
                        conf = "estimado"
                        obs_raw = (f"{obs_raw} | Unidade ajustada de {original_unit} "
                                   f"para {normalized_unit}").strip(" |")
                    if _pdf_downgrade:
                        obs_raw = (f"{obs_raw} | Estimativa: lido de PDF, não medido em "
                                   f"geometria — envie DWG/DXF pra medir").strip(" |")

                    # ref_sheet SEMPRE tem o filename real pra que o botão
                    # "Ver prancha" na revisão inline funcione. Hint da IA
                    # vai entre parênteses se for diferente do nome.
                    ia_hint = (item_data.get("ref_sheet") or "").strip()
                    if ia_hint and ia_hint.lower() not in filename.lower():
                        _ref = f"{filename} ({ia_hint[:60]})"
                    else:
                        _ref = filename
                    item = BudgetItem(
                        item_num=str(item_data.get("item_num", "")),
                        description=desc,
                        unit=normalized_unit,
                        quantity=qty,
                        observations=obs_raw,
                        ref_sheet=_ref,
                        confidence=Confidence(conf),
                        discipline=discipline,
                        origem="vision_pdf",  # lido por Vision, não medido em geometria
                    )
                    all_items.append(item)
                except: continue

            # 6. Liberar memória desta prancha
            del text, crop_paths, sheet, result
            gc.collect()

        # ── Aço SEMPRE em kg (projeto estrutural) ──
        # Override determinístico: por mais que o SYSTEM_PROMPT_ESTRUTURA mande,
        # a IA às vezes devolve item de aço/armadura/estribo em m² (caso Luciano).
        # Em estrutura, aço é SEMPRE kg (regra de norma, universal — não é
        # benchmark de projeto, então não fere isolamento). Não toca em fôrma (m²)
        # nem concreto (m³): só força quando a descrição é claramente de aço.
        if is_structural:
            # Aço SEMPRE em kg (regra de norma) + guardrail de tipo errado.
            # Regras em engine_rules.py (testadas em tests/test_engine_rules.py).
            _fixed_kg = 0
            for _it in all_items:
                if _should_force_steel_kg(getattr(_it, "description", "")) and getattr(_it, "unit", "") != "kg":
                    _it.unit = "kg"
                    _fixed_kg += 1
            if _fixed_kg:
                print(f"[estrutural] forcei unit=kg em {_fixed_kg} item(ns) de aço")

            # Guardrail (caso Magno): estrutural com quase tudo zerado provavelmente
            # é arquitetura marcada errada no upload. Avisa em vez de entregar lixo.
            if _is_likely_wrong_type([getattr(_it, "quantity", 0) for _it in all_items]):
                _warn_tipo = ("⚠ Este arquivo parece ser de ARQUITETURA, não de estrutura — "
                              "quase nenhum item estrutural pôde ser medido. Reenvie marcando "
                              "\"Arquitetura\" no tipo de projeto (ou responda o e-mail do "
                              "AI.arq que a gente reprocessa pra você).")
                try:
                    if _warn_tipo not in (project_data.warnings or []):
                        project_data.warnings.insert(0, _warn_tipo)
                except Exception:
                    pass
                print("[estrutural] possivel tipo incorreto: muitos itens zerados")

        # ── Consolidação pós-IA ──
        jobs.update_field(job_id, progress=91)
        # Remove duplicatas similares (mesmo item com qty idêntica repetido em
        # múltiplas pranchas), consolida réplicas por departamento/zona e valida
        # un=inteiro (corrige quando a IA devolve un com decimais suspeitos).
        jobs.update_field(job_id, current_step="Consolidando itens duplicados...")
        n_before = len(all_items)
        all_items = _consolidate_items(all_items)
        all_items = _dedupe_by_block(all_items)  # funde itens do mesmo bloco CAD (anti-duplicação)
        all_items = _drop_nonsense_items(all_items)       # tira "seção transversal" e afins
        all_items = _consolidate_by_type_code(all_items)  # funde mesmo tipo (DRY 07) entre pranchas
        # Validar qty/unit após consolidação
        for it in all_items:
            new_qty, adjusted = _validate_quantity_for_unit(it)
            if adjusted:
                it.quantity = new_qty
                try:
                    from models import Confidence
                    it.confidence = Confidence("estimado")
                except Exception:
                    pass
                it.observations = (
                    (it.observations or "") +
                    " | Qty de un ajustada: valor original não-inteiro"
                ).strip(" |")
        n_after = len(all_items)
        if n_before != n_after:
            print(f"[consolidação] {n_before} → {n_after} itens ({n_before - n_after} consolidados)")

        # ── Regras pós-consolidação (auditoria 2026-05-17 — projeto Yuri) ──
        # 🅑 Dedup m² por layer (evita pisos duplicados somando > área do projeto)
        # 🅒 "Conforme projeto" / "a definir" → estimado obrigatório (regra dura #1)
        # 🅓 "Fundido"/"Consolidado de N entradas" → estimado obrigatório
        all_items, n_post_changed = _apply_post_consolidation_rules(all_items)
        if n_post_changed > 0:
            print(f"[pós-consolidação] {n_post_changed} itens marcados como estimado por regras 🅑+🅒+🅓")

        # 🅐 Detector multifamiliar — sinaliza warning visível pro cliente
        is_multi, evidencias = _detect_multifamiliar_signal(all_items, typology)
        if is_multi:
            multi_warning = (
                "Tipologia salva como 'residential' mas as plantas indicam projeto MULTIFAMILIAR. "
                f"Detectado: {', '.join(evidencias)}. "
                "Recomendo trocar a tipologia em 'Dados do projeto' pra que a calibração compare "
                "com benchmarks corretos. Sem isso, alertas de quantitativo podem estar errados."
            )
            project_data.warnings = (project_data.warnings or []) + [multi_warning]
            print(f"[multifamiliar] sinal detectado: {evidencias}")

        # ── Validação de plausibilidade ──
        # Detecta disciplina×unidade mismatch, range absurdo, área > laje×1.5.
        # Marca estimado (laranja) e anota o motivo pra usuário revisar.
        jobs.update_field(job_id, current_step="Validando plausibilidade dos itens...")
        flagged_count = 0
        laje_area = project_data.total_area or 0
        for it in all_items:
            plausible, reason = _check_plausibility(it, laje_area)
            if not plausible:
                try:
                    from models import Confidence
                    it.confidence = Confidence("estimado")
                except Exception:
                    pass
                it.observations = (
                    (it.observations or "") + f" | ⚠ Revisar: {reason}"
                ).strip(" |")
                flagged_count += 1
        if flagged_count > 0:
            print(f"[plausibilidade] {flagged_count} itens flagados pra revisão")

        # ── Calibração por DENSIDADE (ratios qty/área) ──
        # Compara a densidade (qty/área) de cada item contra benchmarks
        # agregados de projetos históricos (mesma tipologia). Desvio > ±2σ
        # vira observação laranja. NUNCA promove pra confirmado.
        # Área de referência: layout_area se disponível, senão total_area.
        ref_area = project_data.layout_area or project_data.total_area or 0
        if HAS_DENSITY_CAL and ref_area > 0:
            try:
                from density_calibration import check_density_anomaly
                benchmarks = density_get_benchmarks(typology=typology)
                density_flagged = 0
                for it in all_items:
                    is_anom, reason = check_density_anomaly(
                        it, ref_area, benchmarks=benchmarks, typology=typology,
                    )
                    if is_anom:
                        try:
                            from models import Confidence
                            it.confidence = Confidence("estimado")
                        except Exception:
                            pass
                        it.observations = (
                            (it.observations or "") + f" | ⚠ Calibração: {reason}"
                        ).strip(" |")
                        density_flagged += 1
                if density_flagged > 0:
                    print(f"[densidade] {density_flagged} itens fora do padrão histórico")
            except Exception as e:
                print(f"[densidade] Erro no check de anomalia: {e}")

        # ── SALVAMENTO DE LAYOUT (20/07, caso Catarina) ──
        # Antes de declarar "0 itens = falha": uma planta de ESTUDO DE LAYOUT
        # (interiores, sem quadro de áreas) é legível mas não rende item clássico
        # — a IA rodava de boa e devolvia vazio → erro. Só que a prancha tem
        # ESQUADRIAS cotadas no texto vetorial (medida escrita: "160x150/86").
        # Extrai essas esquadrias de forma DETERMINÍSTICA (sem IA, sem inventar)
        # pra planilha ter conteúdo real em vez de erro. Só quando a IA NÃO deu
        # erro (senão é falha real do provedor, não layout) — aí mantém o erro
        # honesto. Não roda em complemento (a base já tem itens).
        if len(all_items) == 0 and not is_complement and not (sheet_errors or dxf_errors):
            _salv = _salvage_layout_esquadrias(file_paths)
            if _salv:
                all_items.extend(_salv)
                _n_esq = sum(int(i.quantity) for i in _salv)
                # Esquadrias confirmam que é um layout vetorial legível → vale o
                # passo de contagem visual (louças/móveis/portas). Aditivo e robusto:
                # falha → só as esquadrias. Não roda se esquadrias não acharam nada
                # (provável prancha ilegível/escaneada — não gasta chamada de IA).
                _n_ai = 0
                try:
                    _salv_ai = _salvage_layout_ai_counts(client, file_paths, crops_dir)
                    if _salv_ai:
                        all_items.extend(_salv_ai)
                        _n_ai = len(_salv_ai)
                except Exception as _sai_e:
                    print(f"[salvage-ai] falhou (mantendo só esquadrias): {_sai_e}")
                project_data.warnings = (project_data.warnings or []) + [
                    f"Esta prancha é um ESTUDO DE LAYOUT (sem quadro de áreas/cotas de "
                    f"dimensão). Não dá pra medir áreas (piso/forro) com honestidade daqui. "
                    f"Extraímos {_n_esq} esquadrias cotadas na prancha"
                    + (f" e {_n_ai} itens de contagem (louças/móveis, ⚠ confira)" if _n_ai else "")
                    + f". Pra o quantitativo COM áreas, informe a área total no upload ou envie o DXF."
                ]
                print(f"[salvage-layout] job={job_id}: 0 itens da IA → salvei {_n_esq} "
                      f"esquadrias ({len(_salv)} tipos) + {_n_ai} itens de contagem visual")

        # ── Resultado vazio: NÃO marcar "done" silencioso ──
        # Bug Vinícius (2026-05-21): 1 PDF processou em 17s, 0 itens, status
        # done. Usuário (1º projeto, grátis) recebeu planilha vazia achando
        # que "concluiu". Planilha com 0 itens é SEMPRE falha — qualquer
        # prancha de arquitetura real tem ao menos paredes/piso/forro.
        # Distingue 2 causas pra orientar o usuário com mensagem certa:
        if len(all_items) == 0:
            # EXCEÇÃO — COMPLEMENTO (/add-file) sobre projeto que já tem planilha:
            # a regra "0 itens = falha" vale pro projeto NORMAL (bug Vinícius). Num
            # complemento, o arquivo anexado não rendeu nada mas a planilha anterior
            # segue intacta (_persist_items_to_supabase só troca os itens no sucesso).
            # Marcar erro aqui sumia a planilha da tela e mandava email de falha pra
            # quem não perdeu nada — mesmo bug que o guard do DWG inválido conserta,
            # só que por outra porta. Restaura 'done' + aviso; nunca marca erro.
            if is_complement and _complement_base_has_items(job_id):
                _anexados = ', '.join(os.path.basename(p) for p in (file_paths or [])) or 'o arquivo'
                _warn_zero = (
                    f"O arquivo que você anexou ({_anexados}) foi lido, mas não rendeu "
                    f"nenhum item quantificável — pode ser prancha só de layout, sem "
                    f"quadros/legendas, ou um PDF escaneado. Sua planilha anterior foi "
                    f"mantida — nada foi perdido. Pra medir pelo CAD, anexe o DWG/DXF "
                    f"da planta arquitetônica.")
                jobs.update_field(job_id, status="done", progress=100, error_message=None,
                                  current_step="Complemento sem itens — planilha anterior mantida")
                try:
                    _supabase_update("projects", "job_id", job_id, {
                        "status": "done",
                        "error_message": None,
                        "warnings": [_warn_zero],
                        "completed_at": datetime.utcnow().isoformat(),
                    })
                except Exception as _upe:
                    print(f"[add-file] restaurar done (0 itens) falhou: {_upe}")
                print(f"[add-file] complemento rendeu 0 itens, base preservada → done+aviso (sem erro)")
                return

            # Considera falhas dos DOIS caminhos: PDF (sheet_errors) e DXF
            # (dxf_errors). Um projeto só-DXF cuja IA falhou tem dxf_errors mas
            # sheet_errors vazio — sem somar aqui, cairia na mensagem errada
            # de "troque o arquivo".
            ai_errors = sheet_errors + dxf_errors
            if ai_errors:
                # Classificação HONESTA (QW2, 20/07): só chamar de "provedor
                # sobrecarregado" com PROVA de transitório (429/529/timeout). O
                # DEFAULT deixou de ser "sobrecarga" — erro desconhecido/permanente
                # (404 de model-id errado, 400 invalid_request, surrogate) vira
                # mensagem honesta e NÃO casa o _TRANSIENT_ERR_RX, então NÃO entra
                # em auto-retry infinito — era exatamente esse loop que prendeu o
                # Rodrigo (19/07). Classificador único em llm_retry.classify_error_text.
                from llm_retry import classify_error_text
                _errblob = " ".join(str(e) for e in ai_errors)
                _verdict = classify_error_text(_errblob)
                if _verdict == "transient":
                    # PROVA de sobrecarga/timeout — é do provedor, reprocessar resolve.
                    raise RuntimeError(
                        "⚠ Os servidores de IA estavam sobrecarregados neste "
                        "momento — é um problema temporário do provedor, NÃO do "
                        "seu arquivo. O sistema já tentou sozinho várias vezes "
                        "(alguns minutos) antes de desistir. É só reprocessar daqui "
                        "a alguns minutos — é grátis e não conta no seu limite."
                    )
                # permanent OU unknown: NUNCA culpar o provedor sem prova. Mensagem
                # honesta — problema técnico do nosso lado, reprocessável, com suporte.
                _low = _errblob.lower()
                _detalhe = ("um caractere inválido no arquivo do CAD"
                            if ("surrogate" in _low or "invalid high surrogate" in _low)
                            else "um problema técnico do nosso lado")
                raise RuntimeError(
                    f"⚠ Tivemos um problema técnico ao processar este projeto "
                    f"({_detalhe}). Já estamos de olho nisso do nosso lado. "
                    f"Reprocesse — se persistir, fale com o suporte pelo botão "
                    f"'Reportar problema' que a gente resolve rápido."
                )
            else:
                # A IA rodou sem erro mas não achou nada quantificável.
                # Reprocessar o MESMO arquivo daria o mesmo resultado —
                # a mensagem orienta a trocar o arquivo de entrada.
                raise RuntimeError(
                    "Nenhum item quantificável foi identificado neste "
                    "arquivo. Causas mais comuns: (1) o PDF é uma imagem "
                    "escaneada ou fotografada — o motor lê PDF vetorial "
                    "exportado direto do CAD (AutoCAD/Revit); (2) a prancha "
                    "tem só o desenho de layout, sem quadros de áreas, "
                    "legendas ou especificações. Reenvie a planta "
                    "arquitetônica completa exportada do CAD, ou fale com "
                    "o suporte pelo botão 'Reportar problema'."
                )

        # ── Falha PARCIAL: vieram itens, mas pranchas/DXF falharam ──
        # O guard acima só pega o caso de ZERO itens. Se sobram itens mas uma
        # disciplina inteira caiu por um pico passageiro da IA, o usuário recebia
        # planilha que PARECE completa (bug Vinícius ainda meio aberto).
        # DWG que falhou mas tem PDF IRMÃO no mesmo job (mesmo nome-base):
        # a prancha NÃO está perdida — entrou pela leitura do PDF. Caso real
        # 19/07 (validação estrutural): FORMA.dwg falhou na conversão, mas
        # FORMA.pdf estava no job e foi processado — e o aviso dizia "faltou",
        # assustando à toa. Só o DWG SEM irmão conta como prancha perdida.
        def _stem_norm(p):
            import unicodedata as _ud
            s = os.path.splitext(os.path.basename(p))[0].strip().lower()
            s = _ud.normalize("NFKD", s)
            return "".join(c for c in s if not _ud.combining(c))
        _pdf_stems = {_stem_norm(p) for p in (pdf_paths or [])}
        # Aviso do plano B (01/08): DWG convertido via libredwg — medições saem
        # normais, mas o cliente confere medidas-chave (honestidade sobre origem).
        for _n_lw in (dwg_via_libredwg or []):
            project_data.warnings = (project_data.warnings or []) + [
                f"O arquivo {_n_lw} foi convertido pelo leitor alternativo (plano B). "
                f"As medições saíram normalmente, mas vale conferir 2-3 medidas-chave "
                f"da planilha contra o projeto antes de fechar orçamento."]
        _dwg_com_irmao = [n for n in (dwg_failed or []) if _stem_norm(n) in _pdf_stems]
        _dwg_sem_irmao = [n for n in (dwg_failed or []) if _stem_norm(n) not in _pdf_stems]
        if _dwg_com_irmao:
            project_data.warnings = (getattr(project_data, 'warnings', None) or []) + [
                f"ℹ {len(_dwg_com_irmao)} DWG(s) não converteram, mas a(s) prancha(s) "
                f"entraram pela versão em PDF do mesmo nome: "
                f"{', '.join(_dwg_com_irmao)[:200]}. Pra medir direto da geometria, "
                f"rode EXPORTTOAUTOCAD no CAD e reenvie o DWG."
            ]
        _dwg_failed_msgs = [f"{n}: não consegui converter esse DWG (talvez versão nova do AutoCAD ou objetos especiais)"
                            for n in _dwg_sem_irmao]
        partial_errors = (sheet_errors or []) + (dxf_errors or []) + _dwg_failed_msgs
        partial_failure = bool(partial_errors)
        if partial_failure:
            _falhos = "; ".join(e.split(":")[0] for e in partial_errors)
            # 🪤 "Reprocesse (grátis)" é conselho ERRADO quando o que falhou foi um
            # DWG que não converte: reprocessar roda o MESMO arquivo no MESMO
            # conversor e falha igual. Foi o que fez o cliente Thalison (29/07)
            # reenviar 2x e desistir da prancha. Só sugere reprocessar quando a
            # falha for de prancha (pode ter sido soluço da IA); quando for DWG,
            # o caminho é mandar o arquivo em DXF.
            if _dwg_sem_irmao and not (sheet_errors or dxf_errors):
                _saida = ("Reprocessar não resolve — o conversor vai falhar igual. "
                          "Abra no seu CAD e salve como DXF, e suba o DXF aqui.")
            elif _dwg_sem_irmao:
                _saida = ("Pras pranchas, reprocessar (grátis) pode completar. "
                          "Pro DWG que não abriu, reprocessar não resolve: salve como "
                          "DXF no seu CAD e suba o DXF.")
            else:
                _saida = "Reprocessar é grátis e pode completar."
            _aviso_cob = (f"⚠ {len(partial_errors)} prancha(s)/arquivo(s) não entraram nesta "
                          f"planilha — ela pode estar INCOMPLETA. {_saida} "
                          f"Faltaram: {_falhos[:280]}")
            project_data.warnings = (getattr(project_data, 'warnings', None) or []) + [_aviso_cob]

        # Aviso de xref/estéril por-arquivo (independe de falha parcial): orienta o
        # usuário a incorporar o desenho externo, em vez de só ver tudo laranja.
        if xref_warnings:
            project_data.warnings = (getattr(project_data, 'warnings', None) or []) + xref_warnings

        # Gerar planilha
        jobs.update_field(job_id, progress=92)
        jobs.update_field(job_id, current_step=f"Gerando planilha com {len(all_items)} itens...")

        # Normalizar nome do projeto: quando há múltiplos arquivos, evitar
        # que o project.name inferido da IA (que geralmente pega o nome do
        # primeiro DXF) dê uma impressão errada de "projeto de só uma coisa".
        # Sempre sobrescreve quando >1 arquivo, pois a IA processa um por vez
        # e escolhe o nome do que viu primeiro.
        if len(file_paths) > 1:
            project_data.name = f"Quantitativos — {len(file_paths)} arquivos processados"

        # Dedup de ambientes (kept_elements/new_rooms) case-insensitive.
        # A IA extrai de várias pranchas e acaba gerando "BANHEIRO" + "Banheiro"
        # + "Banheiro suíte" como itens separados. Agrupa por forma canônica.
        project_data.kept_elements = _dedupe_rooms(project_data.kept_elements)
        project_data.new_rooms = _dedupe_rooms(project_data.new_rooms)

        # Consenso de área: pega a MODA entre leituras de várias pranchas.
        # Sem isso, a última prancha processada sobrescrevia — uma leitura
        # errada da IA (135 vs 270 m²) dependia da ordem de análise.
        project_data.total_area = _pick_area_consensus(_area_readings["total_area"])
        project_data.layout_area = _pick_area_consensus(_area_readings["layout_area"])
        project_data.no_intervention_area = _pick_area_consensus(_area_readings["no_intervention_area"])
        # ÁREA INFORMADA PELO CLIENTE (campo no upload): só usa como base quando a
        # planta NÃO deu área nenhuma (típico de estudo de layout sem cota, ex.
        # Catarina). Geometria/leitura da própria planta SEMPRE tem prioridade —
        # a área do cliente nunca sobrescreve o que a planta forneceu. Rótulo
        # 'informado' faz a planilha dizer "informada por você, não medida"
        # (regra dura nº1: não é medição nossa, segue tratada como base a conferir).
        # Pé-direito informado: independente da área — grava sempre que veio.
        try:
            _upd = float(user_pe_direito or 0)
        except (TypeError, ValueError):
            _upd = 0
        if _upd > 0:
            project_data.user_pe_direito = round(_upd, 2)
            print(f"[pe-direito-informado] job={job_id}: {_upd} m")
        try:
            _uta = float(user_total_area or 0)
        except (TypeError, ValueError):
            _uta = 0
        if _uta > 0 and not (project_data.total_area or 0):
            project_data.total_area = round(_uta, 2)
            project_data.total_area_source = "informado"
            project_data.warnings = (project_data.warnings or []) + [
                f"Área total de {_uta:.0f} m² foi INFORMADA POR VOCÊ no upload (a planta não trazia "
                f"cota/quadro pra medir). Ela entra como BASE pros itens de área — confira antes de orçar."
            ]
            print(f"[area-informada] job={job_id}: usando área do cliente {_uta} m² (planta sem medição)")
        elif not (project_data.total_area or 0):
            # 🚨 63% dos projetos concluídos terminam SEM área total (47% dos com
            # CAD, 76% dos só-PDF — medido em 30/07). A área sai UNICAMENTE da IA
            # lendo o quadro de áreas da prancha; não há cálculo geométrico por
            # trás. Quando não acha, o cliente recebia a planilha em SILÊNCIO: sem
            # saber que faltou a base dos itens de área, e sem saber que podia
            # informar a metragem. Agora ele sabe e tem o que fazer.
            project_data.warnings = (project_data.warnings or []) + [
                "⚠ Não encontramos a área total do projeto — a prancha não trazia um quadro "
                "de áreas legível. Os itens medidos em m² saíram sem essa base de conferência, "
                "então confira com atenção. Pra resolver: reenvie informando a área total no "
                "campo do envio, ou mande também a prancha que tem o quadro de áreas."
            ]
            print(f"[area-ausente] job={job_id}: sem área total e sem área informada")

        # ── HONESTIDADE DE ÁREA (regra dura nº1) — helper _apply_area_honesty ──
        # Vision (PDF) às vezes CHUTA metragem numa planta sem cota ("Forro Sala 52 m²")
        # — número que parece medido mas é inventado (caso Catarina 20/07): ZERA.
        # PORÉM, se o cliente INFORMOU a área (no upload ou depois), os itens de
        # piso/forro/laje recebem essa área como ESTIMADO rotulado "informado por
        # você" — completar os itens com base honesta, sem fingir medição. O m²
        # medido de verdade (origem 'dxf_geom') nunca é tocado.
        _n_fill, _blanked = _apply_area_honesty(
            all_items, project_data.total_area,
            getattr(project_data, "total_area_source", ""))
        if _n_fill:
            print(f"[honestidade-m2] job={job_id}: preenchi {_n_fill} itens de piso/forro/laje "
                  f"com a área INFORMADA {project_data.total_area} m² (estimado, a conferir)")
        if _blanked:
            print(f"[honestidade-m2] job={job_id}: zerei a quantidade de {_blanked} itens de área "
                  f"não-medidos (Vision) — evita m² inventado (regra nº1)")
        # Pintura derivada do pé-direito informado (01/08/2026) — só quando a
        # planilha tem parede em metro linear e nenhum item de pintura.
        try:
            if getattr(project_data, "user_pe_direito", 0):
                if _derive_pintura_pe_direito(all_items, project_data.user_pe_direito):
                    print(f"[pe-direito] job={job_id}: pintura derivada de "
                          f"{project_data.user_pe_direito} m informado (estimado)")
        except Exception as _epd:
            print(f"[pe-direito] job={job_id}: derivação falhou: {_epd}")
        # Coerência de unidade (caso Rafael 01/08): item CONTÁVEL (condulete,
        # tomada, ponto...) com unidade linear/área não pode ficar CONFIRMADO —
        # a quantidade veio de outra medição pendurada na linha errada.
        try:
            from models import Confidence as _Conf
            _n_rebaixados = 0
            for _it in all_items:
                if (_it.confidence == _Conf.CONFIRMADO
                        and _is_unit_mismatch_countable(_it.description, _it.unit)):
                    _it.confidence = _Conf.ESTIMADO
                    _it.observations = ((_it.observations + " | ") if _it.observations else "") + (
                        "⚠ REBAIXADO: item contável (un) veio com unidade "
                        f"{_it.unit} — a quantidade pode pertencer a outra linha "
                        "(ex.: metros de eletroduto). Confira antes de orçar.")
                    _n_rebaixados += 1
            if _n_rebaixados:
                print(f"[unidade-contavel] job={job_id}: rebaixei {_n_rebaixados} "
                      f"item(ns) contável(is) com unidade linear (falso-medido)")
        except Exception as _euc:
            print(f"[unidade-contavel] job={job_id}: checagem falhou: {_euc}")

        for _fld, _reads in _area_readings.items():
            if len(set(_reads)) > 1:
                print(f"[area-consensus] {_fld}: leituras={_reads} → "
                      f"escolhido={getattr(project_data, _fld)}")

        # Enriquece itens com matches SINAPI (Caixa) + TCPO BIM (Pini).
        # SINAPI = referência principal (preço oficial gov BR, atualizado mensal).
        # TCPO = referência técnica complementar (insumos detalhados).
        # Ambos best-effort, nunca bloqueiam planilha.
        # BUSCA JUNTA + IA ESCOLHE (17/07). Antes: uma busca por texto escolhia
        # sozinha e o que ela achava virava "referência SINAPI". Ela não sabe que
        # spot é luz e cuba é pia, então "Piso porcelanato 60x60" saía como PISO DE
        # BORRACHA com 100% de confiança. Agora a busca só JUNTA candidatos (é boa
        # nisso) e a Haiku ESCOLHE — e responde "nenhum" quando o SINAPI não tem o
        # serviço, em vez de empurrar o vizinho mais parecido.
        try:
            from sinapi_matcher import candidates_for, apply_llm_pick
            from concurrent.futures import ThreadPoolExecutor

            def _cands(it):
                try:
                    return {"description": it.description,
                            "unit": getattr(it, "unit", "") or "",
                            "candidates": candidates_for(it.description, limit=60),
                            "_item": it}
                except Exception:
                    return None

            # Em paralelo: cada item faz ~5 consultas e o tempo é de ESPERA de rede,
            # não de conta — em série dava ~1,4s/item (2min+ numa planilha de 100).
            # Best-effort: item que falhar sai fora sem derrubar a planilha.
            with ThreadPoolExecutor(max_workers=5) as _ex:   # 8→5: gentileza com a CPU da instância Supabase (fix SINAPI timeout 2026-07-22)
                lote = [r for r in _ex.map(_cands, all_items) if r]
            n_conf = apply_llm_pick(lote)   # sem ANTHROPIC_API_KEY → 0, cai no corte por nota
            for e in lote:
                if e["candidates"]:
                    e["_item"].sinapi_matches = e["candidates"][:3]
            print(f"[sinapi] {n_conf}/{len(lote)} itens com código conferido pela IA")
            # Libera o pool 60-wide de candidatos SINAPI (de TODAS as pranchas)
            # antes da geração do XLSX + upload — não é mais usado. (higiene de RAM)
            for e in lote:
                e["candidates"] = None
            del lote
            gc.collect()
        except ImportError:
            pass

        # TCPO REMOVIDO da planilha (16/07): decisão de produto = catálogo é SINAPI,
        # TCPO fora da fase atual ([[project_ai_arq_catalogo]]). Parar de rodar o
        # matcher em TODO item economiza tempo/consulta e não popula tcpo_matches, então
        # a planilha não mostra mais TCPO. Endpoints /api/tcpo/* seguem (mortos, sem UI).

        # Checa heurísticas de mercado (dispersão, cobertura, share MAT/MO)
        # contra as categorias de cada item. Adiciona alertas nas observations
        # pra ajudar o cliente a saber onde pedir 3 orçamentos / revisar escopo.
        # REGRA DURA: heurísticas são agregadas e anônimas — nunca valor
        # específico de projeto. Vem da tabela market_heuristics.
        try:
            from market_heuristics import check_item_anomaly
            for it in all_items:
                try:
                    alertas = check_item_anomaly(it, typology=typology)
                    if alertas:
                        sep = " | " if it.observations else ""
                        it.observations = (it.observations or "") + sep + " ".join(alertas)
                except Exception:
                    pass  # alertas são best-effort
        except ImportError:
            pass  # market_heuristics opcional

        output_path = os.path.join(work_dir, f"orcamento_{job_id}.xlsx")
        generate_spreadsheet(project_data, all_items, output_path, typology=typology)

        # Preview das pranchas (cosmético): inicia SÓ AGORA, depois do pico de RAM
        # da análise+planilha, pra o matplotlib não coincidir com o momento crítico
        # (fix OOM 2026-07-22). Lê os DXF do work_dir, que sobrevive ao fim do job.
        if cad_paths:
            try:
                import threading as _t2
                _t2.Thread(target=_render_cad_previews_bg, daemon=True).start()
            except Exception:
                pass

        # Persistir itens individuais no Supabase pra permitir revisão inline
        # no navegador (endpoint /api/items/{job_id}). Sem isso, os itens só
        # existem no xlsx — a revisão só poderia ser feita no Excel offline.
        _persist_items_to_supabase(job_id, all_items)

        # Planilha e itens acabaram de nascer juntos: carimba a origem pra o
        # site saber, depois, se a revisão do cliente deixou o .xlsx pra trás.
        _carimbar_planilha(job_id)

        # Persistir warnings do motor (prancha órfã, legenda ausente) no
        # campo `warnings` da tabela projects — exibido em Meus Projetos
        # como alerta "precisa de complemento".
        if getattr(project_data, 'warnings', None):
            _supabase_update("projects", "job_id", job_id, {
                "warnings": project_data.warnings,
            })

        # Persistir no Supabase Storage pra sobreviver redeploy do Render
        # (o /tmp do dyno é volátil — sem isso, agente e download quebram).
        _storage_ok = _supabase_storage_upload(output_path, f"{job_id}.xlsx")
        print(f"[storage] upload {job_id}.xlsx ok={_storage_ok}")

        jobs.update_field(job_id, progress=100)
        jobs.update_field(job_id, status="done")
        jobs.update_field(job_id, current_step="Concluído!")
        jobs.update_field(job_id, download_url=f"/api/download/{job_id}")

        # Atualizar projeto no Supabase (log explícito do resultado pra rastrear
        # falhas que antes passavam silenciosas)
        _supa_ok = _supabase_update("projects", "job_id", job_id, {
            "status": "done",
            "items_count": len(all_items),
            "total_area": project_data.total_area if project_data.total_area else None,
            "layout_area": project_data.layout_area if project_data.layout_area else None,
            "completed_at": datetime.utcnow().isoformat(),
        })
        print(f"[supabase] update job={job_id} status=done items={len(all_items)} "
              f"total_area={project_data.total_area} layout_area={project_data.layout_area} "
              f"ok={_supa_ok}")

        # Checkpoints por prancha não servem mais depois do done — limpa em
        # thread (best-effort, não atrasa a conclusão).
        try:
            import threading as _thck
            _thck.Thread(target=_ckpt_limpar, args=(job_id,), daemon=True).start()
        except Exception:
            pass

        # ── SHADOW: Medição Vetorial de PDF v1 (pdf_vector.py) ──
        # Mede a geometria vetorial das primeiras páginas PDF DEPOIS do done,
        # em thread daemon, e loga em error_log (stage pdfvec:shadow) pra
        # comparação com o Vision. Zero impacto no cliente; PDFVEC_SHADOW=0
        # desliga. work_dir sobrevive ao fim do job (sem rmtree), então a
        # thread lê os arquivos numa boa; se sumirem (restart), ela só pula.
        try:
            if page_units and os.environ.get("PDFVEC_SHADOW", "1") != "0":
                from pdf_vector import shadow_measure_async
                shadow_measure_async(page_units, job_id, api_key, _log_error)
        except Exception as _sve:
            print(f"[pdfvec] shadow não iniciado: {_sve}")

        # ── SHADOW: montagem de cômodos a partir do DXF (dxf_rooms_shadow.py) ──
        # A área total falta em 63% dos projetos porque o motor só reconhece
        # ambiente desenhado como polígono FECHADO — e quase ninguém desenha
        # assim. O spike de 30/07 montou as faces a partir dos traços do DXF e
        # chegou a 99% da área declarada na prancha de arquitetura (contra 33-43%
        # do caminho do PDF). Aqui ele CALCULA e REGISTRA, sem entregar nada ao
        # cliente — regra dura nº1: vira medição só com prova de vários projetos.
        # DXFROOMS_SHADOW=0 desliga.
        try:
            if _dxfrooms_units:
                from dxf_rooms_shadow import shadow_rooms_async
                # Passa a área declarada pra sombra já gravar o veredito
                # (maior grupo ÷ declarada) em vez de exigir conferência
                # manual job a job — é assim que a regra se prova sozinha.
                shadow_rooms_async(_dxfrooms_units, job_id, _log_error,
                                   area_declarada=(project_data.total_area or None))
        except Exception as _sre:
            print(f"[dxfrooms] shadow não iniciado: {_sre}")

        # Email "planilha pronta" pro usuário (best-effort; falha não derruba o job)
        try:
            import html as _html, urllib.request as _ur2
            _q = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
                  f"&select=user_email,user_name,project_name,reprocess_count")
            _rq = _ur2.Request(_q, method="GET")
            _rq.add_header("apikey", SUPABASE_KEY)
            _rq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            _rows = _json.loads(_ur2.urlopen(_rq, timeout=10).read().decode("utf-8"))
            _pe = (_rows[0].get("user_email") if _rows else "") or ""
            # REPROCESSO (reprocess_count>0): o cliente já clicou reprocessar e está
            # acompanhando, OU é revisão interna nossa — NÃO re-emailar "planilha
            # pronta" (evita spam de email a cada reprocesso). Só o 1º envio notifica.
            _is_reproc = bool(_rows and (_rows[0].get("reprocess_count") or 0) > 0)
            # Nome do cliente pra saudação personalizada (regra do Pedro 19/07):
            # se o projeto não trouxe user_name, busca em profiles/auth.
            _nm = _resolve_client_name(_pe, hint=(_rows[0].get("user_name") if _rows else "") or "")
            if _pe and is_complement:
                # add-file: refizemos o projeto medindo pelo CAD que o cliente anexou.
                # Email PRÓPRIO (não cai no dedup dos outros), 1x por job — garante que
                # a conclusão do "complementar" SEMPRE notifica, independente de
                # reprocess_count. Se um resume pós-restart perder o is_complement, cai
                # no email "planilha pronta" normal (sem dedup) → cliente notificado mesmo.
                if _email_auto_ja_enviado(_pe, "complemento_pronto", ref=job_id):
                    print(f"[email] complemento-pronto já enviado pra este job — pulando")
                else:
                    _pn_c = _html.escape(_rows[0].get("project_name") or "seu projeto")
                    _greet_c = _greeting_line(_html.escape(_nm))
                    _exts_c = [os.path.splitext(p)[1].lower() for p in file_paths]
                    _n_pdf_c = sum(1 for e in _exts_c if e == ".pdf")
                    _n_cad_c = sum(1 for e in _exts_c if e in (".dwg", ".dxf"))
                    _diag_c = _build_reading_diagnostic(all_items, _n_pdf_c, _n_cad_c, project_type, project_data)
                    _n_med_c = sum(1 for it in all_items
                                   if str(getattr(getattr(it, "confidence", None), "value",
                                                  getattr(it, "confidence", "")) or "") == "confirmado")
                    _proximos_c = _next_steps_html(job_id, _n_med_c, len(all_items), _n_cad_c == 0 and _n_pdf_c > 0)
                    _body_c = (f"{_greet_c}<br><br>Refizemos o projeto <b>{_pn_c}</b> medindo pelo <b>CAD</b> "
                               f"que você anexou — a planilha foi atualizada, agora com <b>{len(all_items)} itens</b>."
                               f"{_diag_c}{_proximos_c}")
                    _pn_c_raw = (_rows[0].get("project_name") or "").strip()
                    _subj_c = (f"{_pn_c_raw} — medimos com o CAD, planilha atualizada"
                               if _pn_c_raw else "Medimos seu projeto com o CAD — planilha atualizada")
                    _ok_c = _send_email_smtp(
                        _pe, _subj_c,
                        _email_wrap("Planilha atualizada com o CAD", _body_c,
                                    "Abrir meu projeto", f"https://ai.arq.br/projeto.html?job_id={job_id}",
                                    badge="&#10003; Medido"))
                    if _ok_c:
                        _email_auto_registrar(_pe, "complemento_pronto", ref=job_id)
                    print(f"[email] complemento-pronto -> enviado={_ok_c}")
            if _pe and not is_complement and _is_reproc:
                # Antes: mudo (anti-spam). Agora: email PRÓPRIO de reprocesso, 1x por
                # job (dedup em email_auto_log). Fecha o buraco onde o cliente
                # reprocessava (ou a gente resgatava) e ninguém avisava — caso
                # Thamiry 06/07. Automação decidida com o Pedro em 07/07.
                if _email_auto_ja_enviado(_pe, "reprocesso_pronto", ref=job_id):
                    print(f"[email] reprocesso-pronto já enviado pra este job — pulando")
                else:
                    _pn_r = _html.escape(_rows[0].get("project_name") or "seu projeto")
                    _greet_r = _greeting_line(_html.escape(_nm))
                    _body_r = (f"{_greet_r}<br><br>Reprocessamos o projeto <b>{_pn_r}</b> com o motor "
                               f"mais recente e a planilha atualizada ficou pronta, com "
                               f"<b>{len(all_items)} itens</b> de quantitativo.<br><br>"
                               f"Ela substitui a versão anterior — abra pra revisar e baixar. "
                               f"Cada item vem marcado como <b>medido</b> (direto do CAD) ou "
                               f"<b>estimativa</b> (pra você conferir).")
                    _pn_r_raw = (_rows[0].get("project_name") or "").strip()
                    _subj_r = (f"{_pn_r_raw} — reprocessamos, planilha atualizada"
                               if _pn_r_raw else "Reprocessamos seu projeto no AI.arq — planilha atualizada")
                    _ok_r = _send_email_smtp(
                        _pe, _subj_r,
                        _email_wrap("Planilha atualizada", _body_r,
                                    "Ver minha planilha", "https://ai.arq.br/dashboard.html",
                                    badge="&#10003; Atualizado"))
                    if _ok_r:
                        _email_auto_registrar(_pe, "reprocesso_pronto", ref=job_id)
                    print(f"[email] reprocesso-pronto -> enviado={_ok_r}")
            if _pe and not is_complement and not _is_reproc:
                _aviso_html = ""
                if partial_failure:
                    # 🪤 Mesma correção do aviso na tela: não sugerir reprocesso quando
                    # o que falhou foi DWG que não converte — reprocessar falha igual.
                    _rep = ("" if _dwg_sem_irmao
                            else " Reprocessar é grátis e tenta completar.")
                    _aviso_html = (f"<br><br><b>&#9888; Atenção:</b> {len(partial_errors)} prancha(s) "
                                   f"não entraram — a planilha pode estar incompleta.{_rep}")
                # DWG que não abriu: o cliente PRECISA saber, senão ele acha que o
                # resultado fraco é o normal do produto. Caso Walter 29/07 — mandou
                # 1 DWG + 7 PDF, o DWG não converteu, e o email não dizia uma palavra.
                if dwg_failed:
                    import html as _hesc
                    _q = len(dwg_failed)
                    _nomes = ", ".join(_hesc.escape(os.path.basename(str(p))) for p in dwg_failed[:3])
                    _e_aec = " (é um arquivo do AutoCAD Architecture/MEP)" if _aec_failed else ""
                    _aviso_html += (
                        f"<br><br><b>&#9888; {'Seu arquivo DWG não abriu' if _q == 1 else f'{_q} arquivos DWG não abriram'}"
                        f"</b>{_e_aec}: <i>{_nomes}</i>. Era o arquivo que mediria de verdade &mdash; "
                        f"sem ele, o que veio de PDF sai como <b>estimativa</b>. "
                        f"Pra resolver: abra no AutoCAD ou BricsCAD, <b>Salvar Como &rarr; DXF 2013</b>, "
                        f"e suba o DXF no mesmo projeto. A gente refaz <b>medindo</b> &mdash; de graça.")
                # Diagnóstico de leitura personalizado: explica COMO lemos e POR QUE
                # a planilha ficou assim (tipo de arquivo + medido vs estimado + avisos).
                _exts = [os.path.splitext(p)[1].lower() for p in file_paths]
                _n_pdf = sum(1 for e in _exts if e == ".pdf")
                # 🪤 Conta o CAD que REALMENTE virou análise, não o que foi enviado.
                # DWG que não converteu não é CAD lido — tratar como tal fazia o email
                # culpar o desenho do cliente ("elementos desenhados como linhas soltas")
                # e sumir com o passo "complemente com o CAD", que era a saída dele.
                _n_cad = len(dxf_paths)
                _diag = _build_reading_diagnostic(all_items, _n_pdf, _n_cad, project_type, project_data)
                # Próximos passos PERSONALIZADOS: PDF sem nada medido → puxa o CAD;
                # senão lidera com revisão (citando quantos ficaram em laranja).
                _n_med = sum(1 for it in all_items
                             if str(getattr(getattr(it, "confidence", None), "value",
                                            getattr(it, "confidence", "")) or "") == "confirmado")
                _veio_pdf = (_n_cad == 0 and _n_pdf > 0)
                _proximos = _next_steps_html(job_id, _n_med, len(all_items), _veio_pdf)
                _subj_pp, _html_pp = _build_planilha_pronta_email(
                    _nm,
                    _rows[0].get("project_name") or "seu projeto",
                    job_id, len(all_items), f"{_aviso_html}{_diag}{_proximos}")
                _send_email_smtp(_pe, _subj_pp, _html_pp, log_kind="planilha_pronta")
        except Exception as _ee:
            print(f"[email] planilha-pronta nao enviada (nao-fatal): {_ee}")

    except Exception as e:
        jobs.update_field(job_id, status="error")
        jobs.update_field(job_id, error_message=str(e))
        jobs.update_field(job_id, current_step=f"Erro: {str(e)[:200]}")
        import traceback as _tb_err
        _log_error("process_job", f"{type(e).__name__}: {e}\n{_tb_err.format_exc()[:1500]}", job_id)

        # Atualizar erro no Supabase
        _supabase_update("projects", "job_id", job_id, {
            "status": "error",
            "error_message": str(e)[:500],
        })

        # Email pro cliente (best-effort; sem jargão técnico). Distingue falha
        # passageira (reprocessar resolve) de arquivo não-quantificável (trocar
        # o arquivo) pela mensagem da exceção — a mesma distinção feita lá no
        # raise de "0 itens". O helper faz dedup pra não repetir por job.
        try:
            # Reprocessável SÓ se o erro for passageiro (infra/IA) — mesmo detector
            # do alerta interno (_TRANSIENT_ERR_RX). DWG que não abre, DXF grande
            # demais, 0 itens = problema de arquivo: reprocessar o mesmo NÃO resolve,
            # o email orienta a trocar/corrigir o arquivo. (bug eletrivan/Luciano 14/07)
            _reproc = bool(_TRANSIENT_ERR_RX.search(str(e)))
            _email_falha_cliente(job_id, reprocessavel=_reproc)
        except Exception as _ee3:
            print(f"[email] erro-cliente nao enviado (nao-fatal): {_ee3}")


@app.get("/")
async def root():
    # `commit` existe pro smoke test saber SE O DEPLOY NOVO JÁ SUBIU.
    # 🪤 28/07/2026: o Render mantém a versão antiga no ar até a nova ficar
    # pronta (deploy sem downtime). O poll do workflow só checava se alguém
    # respondia — e respondia: o código VELHO. Resultado: o teste rodava contra
    # a versão anterior e acusava falha falsa a cada mudança de backend.
    # Agora o workflow espera este commit bater com o do push.
    return {
        "service": "AI.arq API",
        "version": "1.0.0",
        "status": "online",
        "commit": (os.getenv("RENDER_GIT_COMMIT") or "")[:7],
    }


_VALID_TYPOLOGIES = {"office", "residential", "retail", "hospital", "educational"}


async def _stream_upload_to_disk(upload_file, file_path, *, head_bytes=16, chunk=1024 * 1024):
    """Grava o upload em disco em pedaços, SEM bufferizar o arquivo inteiro na RAM.

    Motivo (anti-OOM): `await upload_file.read()` (sem argumento) puxa o arquivo
    todo pra memória de uma vez — um DXF de 300 MB dava um pico de 300 MB, e em
    cima do resto isso estourava a RAM do Render e reiniciava o serviço pra todos.
    O parser multipart do Starlette já derrama uploads grandes num arquivo temp em
    disco; o vilão era só o `.read()`. Lendo em pedaços de 1 MB, o pico de RAM do
    upload fica em ~1 MB, independente do tamanho do arquivo.

    Retorna (bytes_gravados, head): `head` são os primeiros bytes, pra checagem de
    assinatura (ex.: 'AC' do DWG) sem precisar reabrir o arquivo.
    """
    try:
        await upload_file.seek(0)
    except Exception:
        pass
    total = 0
    head = b""
    try:
        with open(file_path, "wb") as out:
            while True:
                buf = await upload_file.read(chunk)
                if not buf:
                    break
                if len(head) < head_bytes:
                    head += buf[: head_bytes - len(head)]
                out.write(buf)
                total += len(buf)
    except OSError as e:
        # Disco cheio / IOError no meio da gravação: apaga o parcial e devolve
        # erro claro (senão o FastAPI daria 500 mudo e o arquivo meio-escrito
        # ficaria no work_dir). Ambos os chamadores herdam esse tratamento.
        try: os.remove(file_path)
        except OSError: pass
        raise HTTPException(507, "Não consegui salvar o arquivo no servidor "
                            "(sem espaço em disco no momento). Tente de novo em instantes.") from e
    return total, head


@app.post("/api/process")
async def process_files(
    request: Request,
    background_tasks: BackgroundTasks,
    files: list[UploadFile] = File(...),
    sheet_types: list[str] = Form(default=[]),
    sheet_ambientes: list[str] = Form(default=[]),
    typology: str = "office",
    project_type: str = "arquitetura",
    project_name: str = "",
    user_id: str = "",
    user_email: str = "",
    user_name: str = "",
    credits_to_consume_cents: int = 0,
    user_total_area: float = 0,
    user_pe_direito: float = 0,
):
    """Recebe PDF, DWG ou DXF e inicia processamento em background.

    - `files`: arquivos de prancha.
    - `sheet_types` (opcional): lista paralela a `files`. String vazia
      em uma posição significa "detectar automaticamente". Valor não-vazio
      sobrescreve a detecção e força o prompt do tipo escolhido pelo user.
      Valores aceitos: demolir, layout_novo, layout_atual, arquitetura,
      pontos, piso, forro, det_forro, mobiliario, marcenaria.
    - `typology` (opcional, default `office`): usado pela calibração por
      densidade pra comparar o projeto com padrões da mesma categoria.
    - `project_name` (opcional): apelido amigável dado pelo cliente.
    - `user_id` (opcional): se informado, vincula o projeto ao usuário e
      permite consumo/crédito. Validado contra JWT — não dá pra criar
      projeto em nome de outro usuário.
    - `credits_to_consume_cents` (opcional): se > 0, consome esse valor
      de créditos do user (usado quando checkout retornou is_free=true
      por saldo suficiente).
    """
    # AUTORIZAÇÃO: se o Form passou user_id (não-anônimo), tem que ter JWT
    # correspondente. Sem isso, atacante manda user_id de outro user e
    # cria projeto+consome créditos no nome dele.
    # Processar EXIGE login sempre — fecha a "porta dos fundos": impede chamada
    # anônima direta na API (por fora do site) que dispararia IA cara sem dono.
    # O painel já loga e manda o Bearer, então não muda nada pro usuário real.
    jwt_user = _get_user_from_request(request)
    if not jwt_user:
        raise HTTPException(401, "Faça login para enviar um projeto.")
    if user_id and user_id != "anonymous" and jwt_user.get("id") != user_id and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
        raise HTTPException(403, "user_id não corresponde ao token de autenticação")
    # Dono autoritativo vem do token (não confia só no parâmetro).
    if not user_id or user_id == "anonymous":
        user_id = jwt_user.get("id") or user_id

    # Freio anti-abuso de CUSTO (achado auditoria 27/07): processar dispara IA
    # cara (Anthropic/Replicate) + RAM. No beta grátis/ilimitado, um usuário
    # logado num loop queimaria gasto sem teto. Limite por usuário+IP, generoso —
    # um humano sobe poucos projetos por sessão; robô em loop é barrado.
    if not _rate_limit_ok(f"process:{user_id}", request, limit=12, window_s=600):
        raise HTTPException(429, "Muitos projetos enviados em pouco tempo. Espere "
                            "alguns minutos e tente de novo.")

    # Teto de tamanho do REQUEST (anti-OOM). Desde 21/07 o upload é gravado em
    # disco em pedaços (_stream_upload_to_disk) — NÃO bufferiza mais o arquivo
    # inteiro na RAM, então o pico de memória do upload é ~1 MB por arquivo,
    # independente do tamanho. Render agora em 4 GB de RAM (o "25 GB" é DISCO).
    # O que ainda consome RAM é o PROCESSAMENTO (ezdxf carrega o DXF na memória,
    # prancha a prancha) — por isso o cap por-arquivo (DXF 150MB) no motor segue
    # sendo a trava real; este teto de request só limita quantas pranchas vêm de
    # uma vez. 320→450 MB após o fix de streaming + upgrade pra 4 GB (incidente
    # 21/07: em 2 GB, cap de 700MB + upload de 506MB deu OOM; agora é outro cenário).
    _clen = request.headers.get("content-length") or request.headers.get("Content-Length")
    if _clen and _clen.isdigit() and int(_clen) > 450 * 1024 * 1024:
        raise HTTPException(413, "Arquivos muito grandes (máx. ~450 MB no total). Envie as pranchas do projeto — se for um projeto enorme, mande em 2 lotes.")

    if typology not in _VALID_TYPOLOGIES:
        typology = "office"
    # Tipo de projeto escolhido no upload: roteia o motor (arquitetura vs estrutura)
    project_type = (project_type or "arquitetura").strip().lower()
    if project_type not in ("arquitetura", "estrutura"):
        project_type = "arquitetura"
    # Área informada pelo cliente (campo opcional no upload): sanitiza contra
    # número absurdo/negativo. 0 = não informou (comportamento antigo).
    try:
        user_pe_direito = float(user_pe_direito or 0)
    except (TypeError, ValueError):
        user_pe_direito = 0
    if not (1.8 <= user_pe_direito <= 8.0):   # fora disso não é pé-direito plausível
        user_pe_direito = 0
    try:
        user_total_area = float(user_total_area or 0)
    except (TypeError, ValueError):
        user_total_area = 0
    if user_total_area < 0 or user_total_area > 1_000_000:
        user_total_area = 0
    if not files:
        raise HTTPException(400, "Nenhum arquivo enviado")

    # Validar arquivos (aceitar PDF, DWG e DXF)
    valid_extensions = ('.pdf', '.dwg', '.dxf')
    # Guarda triplas (arquivo, sheet_type, ambiente) pra manter o mapeamento
    # após filtragem. sheet_types/sheet_ambientes podem estar vazias.
    valid_pairs = []
    for idx, f in enumerate(files):
        if f.filename and f.filename.lower().endswith(valid_extensions):
            st = sheet_types[idx] if idx < len(sheet_types) else ""
            amb = sheet_ambientes[idx] if idx < len(sheet_ambientes) else ""
            valid_pairs.append((f, st, amb))
    if not valid_pairs:
        raise HTTPException(400, "Nenhum arquivo válido encontrado. Aceito: PDF, DWG ou DXF.")

    if len(valid_pairs) > 50:
        raise HTTPException(400, "Máximo de 50 arquivos por projeto")

    # Criar job
    job_id = str(uuid.uuid4())[:8]
    work_dir = os.path.join(WORK_DIR, job_id)
    os.makedirs(work_dir, exist_ok=True)

    # Salvar arquivos + mapas {caminho_local: sheet_type / ambiente}
    file_paths = []
    user_sheet_types: dict[str, str] = {}
    user_ambientes: dict[str, str] = {}
    file_types = {'pdf': 0, 'dwg': 0, 'dxf': 0}
    avisos_aec: list[str] = []   # DWGs com objetos AEC — aviso imediato (ver abaixo)
    # Arquivos com cara de projeto ESTRUTURAL enviados em modo arquitetura
    # (01/08/2026: só 3 de 107 projetos escolheram o modo estrutura no dropdown,
    # enquanto FORMA.pdf/fundacao.pdf passavam como arquitetura — onde laje e
    # aço nunca são medidos). Detecção por nome do arquivo, aviso não-bloqueante.
    avisos_estrutural: list[str] = []
    import re as _re_estrut
    # 🪤 "forma" precisa de fronteira: sem ela, "inFORMAtivo" e "plataFORMA"
    # disparavam o aviso (pego no teste antes do deploy).
    _RX_ESTRUT = _re_estrut.compile(
        r"estrut|(?<![a-zà-ü])f[oô]rmas?(?![a-zà-ü])|funda[cç]|arma[cç]"
        r"|pilar|viga|laje|baldrame|sapata",
        _re_estrut.IGNORECASE)
    # 🪤 Importar AQUI e reclamar alto se falhar. A detecção só valia a pena se
    # rodasse mesmo; escondida atrás de um try/except mudo, um ImportError
    # deixaria o aviso desligado pra sempre sem ninguém perceber.
    try:
        from dwg_extractor import dwg_has_aec_markers as _detecta_aec
    except Exception as _e_imp:
        _detecta_aec = None
        print(f"[upload] AVISO: detecção de AEC indisponível ({type(_e_imp).__name__}: {_e_imp})")
        try:
            _log_error("dwg:aec-detector-indisponivel",
                       f"Não consegui importar dwg_has_aec_markers no upload: "
                       f"{type(_e_imp).__name__}: {_e_imp}", None, severity="warning")
        except Exception:
            pass
    for upload_file, user_st, user_amb in valid_pairs:
        # Anti path-traversal: nunca confiar em upload_file.filename
        safe_name = _safe_local_filename(upload_file.filename)
        file_path = os.path.join(work_dir, safe_name)
        # Grava em disco em pedaços (não bufferiza o arquivo inteiro na RAM).
        n_written, head = await _stream_upload_to_disk(upload_file, file_path)

        # Validação de integridade (Bug Rafael 2026-05-04: DWG chegou
        # truncado e backend processou sem detectar, gerando planilha vazia)
        if upload_file.size and n_written != upload_file.size:
            try: os.remove(file_path)
            except OSError: pass
            raise HTTPException(
                400,
                f"Arquivo '{upload_file.filename}' chegou incompleto: "
                f"recebido {n_written} de {upload_file.size} bytes. "
                f"Provável conexão instável durante upload — tente de novo."
            )
        ext = upload_file.filename.lower().rsplit('.', 1)[-1]
        if ext == "dwg":
            if n_written < 100:
                try: os.remove(file_path)
                except OSError: pass
                raise HTTPException(
                    400,
                    f"DWG '{upload_file.filename}' muito pequeno ({n_written} bytes) — "
                    f"provavelmente corrompido. Verifique se o arquivo abre no AutoCAD."
                )
            if head[:2] != b"AC":
                try: os.remove(file_path)
                except OSError: pass
                raise HTTPException(
                    400,
                    f"DWG '{upload_file.filename}' não tem assinatura válida "
                    f"(esperado iniciar com 'AC10xx'). Arquivo corrompido — "
                    f"verifique no AutoCAD ou exporte como PDF e suba o PDF."
                )
            # AEC/MEP: avisar AGORA, não depois de 5 min de processamento.
            # Medido em 01/08/2026: 13 das 29 falhas de DWG são objetos AEC
            # (AutoCAD Architecture/MEP), que nenhum conversor livre abre. A
            # detecção já existia — mas só rodava DEPOIS da conversão falhar,
            # então o cliente esperava o processamento inteiro pra descobrir.
            # NÃO bloqueia: o mesmo envio pode trazer um DXF ou PDF que salve
            # o projeto. Só entrega a informação enquanto ele ainda está na tela.
            if _detecta_aec is not None:
                try:
                    if _detecta_aec(file_path):
                        avisos_aec.append(upload_file.filename)
                except Exception as _e_aec:
                    print(f"[upload] detecção AEC falhou em {safe_name}: {_e_aec}")

        if project_type != "estrutura" and _RX_ESTRUT.search(upload_file.filename or ""):
            avisos_estrutural.append(upload_file.filename)
        file_paths.append(file_path)
        if user_st:
            user_sheet_types[file_path] = user_st
        if user_amb:
            user_ambientes[file_path] = user_amb
        file_types[ext] = file_types.get(ext, 0) + 1

    # Resumo de tipos recebidos
    types_summary = ", ".join(f"{v} {k.upper()}" for k, v in file_types.items() if v > 0)

    # Criar status
    jobs[job_id] = ProcessingStatus(
        job_id=job_id,
        status="queued",
        progress=0,
        current_step=f"Recebidos {len(file_paths)} arquivos ({types_summary}). Iniciando processamento...",
        total_steps=3,
    )

    # Salvar projeto no Supabase
    _supabase_insert("projects", {
        "job_id": job_id,
        "user_id": user_id or "anonymous",
        "user_email": user_email or "",
        "user_name": user_name or "",
        "project_name": project_name or "Sem nome",
        "typology": typology,
        "project_type": project_type,
        "files_count": len(file_paths),
        "file_types": file_types,
        "status": "queued",
        "user_total_area": user_total_area if user_total_area > 0 else None,
    })

    # Consome créditos de cashback/cupom se o frontend declarou uso
    # (quando crédito cobriu 100% do preço e pulou o Stripe).
    if credits_to_consume_cents > 0 and user_id and user_id != "anonymous":
        consumed = _consume_credits(user_id, credits_to_consume_cents, job_id)
        print(f"[credits] job={job_id} user={user_id} consumed={consumed}/{credits_to_consume_cents}")

    # Iniciar processamento em thread separada (não bloqueia HTTP).
    # Usa _process_job_throttled (semáforo 1-por-vez) pra não estourar RAM
    # com jobs concorrentes — ver _JOB_SEMAPHORE acima.
    import threading
    t = threading.Thread(
        target=_process_job_throttled,
        args=(job_id, file_paths, work_dir),
        kwargs={"typology": typology,
                "user_sheet_types": user_sheet_types,
                "user_ambientes": user_ambientes,
                "project_type": project_type,
                "user_total_area": user_total_area,
                "user_pe_direito": user_pe_direito},
        daemon=True,
    )
    t.start()

    # Alerta interno pro Pedro: novo projeto entrou (em thread — não atrasa a resposta)
    try:
        import html as _h4, threading as _th4
        _types_str = ", ".join(f"{v} {k.upper()}" for k, v in file_types.items() if v > 0)
        _alert_body = (f"<b>Usuário:</b> {_h4.escape(user_name or '—')} &lt;{_h4.escape(user_email or 'anônimo')}&gt;<br>"
                       f"<b>Projeto:</b> {_h4.escape(project_name or '(sem nome)')}<br>"
                       f"<b>Arquivos:</b> {len(file_paths)} ({_types_str})<br>"
                       f"<b>Código:</b> {job_id}")
        _th4.Thread(target=_notify_admin, args=("Novo projeto recebido", _alert_body), daemon=True).start()
    except Exception as _na:
        print(f"[notify] alerta novo-projeto falhou: {_na}")

    resp = {"job_id": job_id, "files_received": len(file_paths),
            "file_types": file_types, "status": "queued", "typology": typology,
            "project_type": project_type}
    if avisos_estrutural and project_type != "estrutura":
        resp["aviso_estrutural"] = {
            "arquivos": avisos_estrutural,
            "titulo": "Esses arquivos parecem de projeto ESTRUTURAL",
            "texto": ("Pelo nome, são pranchas de estrutura (fôrma, fundação, armação...). "
                      "Este envio está no modo Arquitetura, que não mede aço, pilares nem lajes.\n\n"
                      "Se o objetivo é o quantitativo estrutural: envie de novo escolhendo "
                      "'Estrutura (concreto armado)' no tipo de projeto — lá o motor mede aço "
                      "por tabela, pilares e lajes.\n\n"
                      "Se essas pranchas são só apoio do projeto de arquitetura, ignore este aviso."),
        }
        try:
            _log_error("upload:estrutural-em-modo-arquitetura",
                       f"Arquivos com cara de estrutural em modo arquitetura: "
                       f"{', '.join(avisos_estrutural[:5])} — cliente avisado no envio",
                       job_id, severity="info")
        except Exception:
            pass
    if avisos_aec:
        _tem_alternativa = (file_types.get('dxf', 0) + file_types.get('pdf', 0)) > 0
        resp["aviso_aec"] = {
            "arquivos": avisos_aec,
            "tem_alternativa": _tem_alternativa,
            "titulo": ("Um dos seus arquivos é do AutoCAD Architecture/MEP"
                       if len(avisos_aec) == 1 else
                       "Alguns arquivos são do AutoCAD Architecture/MEP"),
            "texto": (
                "Esse tipo de DWG guarda paredes e móveis como \"objetos inteligentes\", "
                "que nenhum conversor abre direto — nem o \"Salvar como DXF\" comum. "
                "Não é defeito do seu arquivo.\n\n"
                "Como resolver, em 3 passos:\n"
                "1. No AutoCAD, com o arquivo aberto, digite EXPORTTOAUTOCAD e Enter\n"
                "2. Escolha a versão 2013 e confirme — ele cria um arquivo novo "
                "(o seu original não é alterado)\n"
                "3. Abra esse arquivo novo, salve como DXF e anexe aqui\n\n"
                + ("Seu projeto vai seguir processando com os outros arquivos, "
                   "mas o que estiver só nesse DWG não vai ser medido."
                   if _tem_alternativa else
                   "Como esse é o único arquivo do envio, o processamento "
                   "provavelmente não vai conseguir medir nada.")),
        }
        try:
            _log_error("dwg:aec-detectado-no-upload",
                       f"AEC detectado ANTES de processar: {', '.join(avisos_aec)} — "
                       f"cliente avisado na hora (alternativa no envio: {_tem_alternativa})",
                       job_id, severity="info")
        except Exception:
            pass
    return resp


@app.get("/api/debug/supa-log")
async def debug_supa_log(request: Request, tail: int = 50):
    """Últimas N linhas do log de operações Supabase — pra investigar por que
    updates silenciosos falham sem ter acesso direto ao log do Render.

    Restrito a admin (vaza queries internas com payloads e IDs).
    """
    _require_admin(request)
    try:
        if not os.path.exists(_SUPA_LOG_PATH):
            return {"status": "ok", "lines": [], "note": "log vazio ou ainda não criado"}
        with open(_SUPA_LOG_PATH, "r", encoding="utf-8") as f:
            all_lines = f.readlines()
        last = all_lines[-tail:] if tail > 0 else all_lines
        return {"status": "ok", "total_lines": len(all_lines),
                "returned": len(last), "lines": [ln.rstrip("\n") for ln in last]}
    except Exception as e:
        return {"status": "error", "error": f"{type(e).__name__}: {e}"}


# Estado do teste de libredwg (roda em thread — o Render corta HTTP em ~100 s,
# e o lote leva ate 4 min; "Load failed" no admin em 01/08 foi isso).
_LIBREDWG_TESTE = {"rodando": False, "resultados": [], "testados": 0,
                   "converteram_e_abrem": 0, "jobs_na_fila": 0, "fim": None}


def _medir_dxf_geometria(dxf_path: str) -> dict:
    """Métricas de geometria pra comparar conversores (qualidade, não estética):
    nº de entidades, extensão do desenho e soma dos comprimentos de linha.
    Se dois DXF do MESMO DWG batem nisso, medem igual no motor."""
    import ezdxf, math
    doc = ezdxf.readfile(dxf_path)
    msp = doc.modelspace()
    n = 0
    total_len = 0.0
    xs, ys = [], []
    for e in msp:
        n += 1
        try:
            t = e.dxftype()
            if t == "LINE":
                dx = e.dxf.end.x - e.dxf.start.x
                dy = e.dxf.end.y - e.dxf.start.y
                total_len += math.hypot(dx, dy)
                xs += [e.dxf.start.x, e.dxf.end.x]; ys += [e.dxf.start.y, e.dxf.end.y]
            elif t == "LWPOLYLINE":
                pts = [(p[0], p[1]) for p in e.get_points()]
                for a, b in zip(pts, pts[1:]):
                    total_len += math.hypot(b[0]-a[0], b[1]-a[1])
                xs += [p[0] for p in pts]; ys += [p[1] for p in pts]
        except Exception:
            continue
    ext = ((max(xs)-min(xs)) if xs else 0.0, (max(ys)-min(ys)) if ys else 0.0)
    return {"entidades": n, "soma_linhas": round(total_len, 1),
            "extensao": (round(ext[0], 1), round(ext[1], 1))}


@app.get("/api/debug/libredwg-batch")
async def debug_libredwg_batch(request: Request, limite: int = 5,
                               iniciar: int = 0, modo: str = "recusados"):
    """Mede o fallback libredwg contra os DWG que o ODA recusou. Admin-only.

    Contexto (01/08/2026): 6+ envios falharam com 'não conseguimos abrir' e o
    caso Ana provou que às vezes o arquivo é bom — o conversor é que não lê.
    O dwg2dxf (libredwg) está compilado no Docker desde 29/07 e NUNCA foi
    medido contra um arquivo real recusado. Este endpoint baixa do Storage os
    DWG de jobs com esse erro, roda o dwg2dxf em cada um e reporta:
    converteu? o DXF resultante abre no ezdxf? quantas entidades tem?

    Só LÊ e mede — não reprocessa nada, não mexe em projeto nenhum.

    🪤 Roda em THREAD e o navegador faz POLLING: o Render corta respostas HTTP
    em ~100 s e o lote leva até 4 min — a 1ª versão síncrona morreu com
    "Load failed" no admin (01/08). `?iniciar=1` dispara; sem parâmetro,
    devolve o andamento. `limite` = arquivos por lote (1-10).
    """
    _require_admin(request)
    import shutil as _sh
    dwg2dxf = _sh.which("dwg2dxf")
    if not dwg2dxf:
        return {"erro": "dwg2dxf não está no PATH deste servidor",
                "dica": "confirmar no /api/debug/dwg se libredwg_which é null"}

    if not iniciar or _LIBREDWG_TESTE["rodando"]:
        return dict(_LIBREDWG_TESTE)   # snapshot do andamento (ou do último lote)

    limite = max(1, min(10, int(limite)))
    _LIBREDWG_TESTE.update({"rodando": True, "resultados": [], "testados": 0,
                            "converteram_e_abrem": 0, "jobs_na_fila": 0, "fim": None})

    def _rodar_qualidade(_sp, _tf, t0):
        """Converte o MESMO DWG pelo ODA e pelo libredwg e compara a geometria.

        É a validação que a trava de qualidade do fallback pede: DXF que abre
        mas mede diferente geraria número branco falso (regra nº1). Roda nos
        DWGs de projetos DONE (onde o ODA funcionou e dá baseline)."""
        from dwg_extractor import convert_dwg_to_dxf
        _st_q, rows = _supa_rest_service(
            "GET",
            "projects?status=eq.done&select=job_id,file_types,created_at"
            "&order=created_at.desc&limit=80")
        rows = rows if isinstance(rows, list) else []
        jobs = [r["job_id"] for r in rows
                if (r.get("file_types") or {}).get("dwg", 0) > 0]
        _LIBREDWG_TESTE["jobs_na_fila"] = len(jobs)
        dwg2dxf_bin = __import__("shutil").which("dwg2dxf")
        # 🪤 1ª rodada (print do Pedro, 01/08): 4 dos 5 sorteados eram DWGs que
        # o ODA nem converte — projetos DONE via PDF/DXF com o DWG falhando em
        # silêncio. Sem baseline não há comparação, então esses não podem
        # consumir o limite; só contam num teto de tentativas separado.
        _tentativas = 0
        for job in jobs:
            if (_LIBREDWG_TESTE["testados"] >= limite or _tentativas >= 25
                    or time.time() - t0 > 600):
                break
            try:
                nomes = [n for n in _supabase_storage_list(PRANCHAS_BUCKET, f"{job}/")
                         if n.lower().endswith(".dwg")]
            except Exception:
                continue
            for nome in nomes[:1]:   # 1 DWG por job basta pra amostra
                if (_LIBREDWG_TESTE["testados"] >= limite or _tentativas >= 25
                        or time.time() - t0 > 600):
                    break
                _tentativas += 1
                item = {"job_id": job, "arquivo": nome}
                try:
                    dados = _supabase_storage_download_prancha(job, nome)
                    if not dados:
                        item["resultado"] = "não consegui baixar"
                    else:
                        with _tf.TemporaryDirectory() as tmp:
                            src = os.path.join(tmp, "in.dwg")
                            open(src, "wb").write(dados)
                            item["mb"] = round(len(dados) / 1048576, 1)
                            dxf_oda = convert_dwg_to_dxf(src)   # gate off ⇒ só ODA
                            dst_lw = os.path.join(tmp, "lw.dxf")
                            try:
                                _sp.run([dwg2dxf_bin, "-y", "-o", dst_lw, src],
                                        capture_output=True, timeout=90)
                            except _sp.TimeoutExpired:
                                pass
                            if not dxf_oda:
                                item["resultado"] = "ODA não converteu (sem baseline)"
                            elif not os.path.exists(dst_lw) or os.path.getsize(dst_lw) == 0:
                                item["resultado"] = "libredwg não converteu este"
                            else:
                                m_oda = _medir_dxf_geometria(dxf_oda)
                                m_lw = _medir_dxf_geometria(dst_lw)
                                item["oda"] = m_oda
                                item["libredwg"] = m_lw
                                base = m_oda["soma_linhas"] or 1e-9
                                delta = abs(m_lw["soma_linhas"] - m_oda["soma_linhas"]) / base
                                item["delta_linhas_pct"] = round(delta * 100, 2)
                                iguais = delta <= 0.01
                                item["resultado"] = ("GEOMETRIA BATE (≤1%)" if iguais
                                                     else "GEOMETRIA DIVERGE")
                                if iguais:
                                    _LIBREDWG_TESTE["converteram_e_abrem"] += 1
                except Exception as e:
                    item["resultado"] = f"erro: {type(e).__name__}: {e}"[:150]
                _LIBREDWG_TESTE["resultados"].append(item)
                if "delta_linhas_pct" in item:
                    _LIBREDWG_TESTE["testados"] += 1   # só comparações reais contam
        try:
            _log_error("libredwg:qualidade",
                       f"{_LIBREDWG_TESTE['converteram_e_abrem']} de "
                       f"{_LIBREDWG_TESTE['testados']} DWGs com geometria batendo "
                       f"ODA×libredwg (±1%, soma de linhas)", None, severity="info")
        except Exception:
            pass

    def _rodar():
        import subprocess as _sp, tempfile as _tf
        t0 = time.time()
        try:
            # 🪤 _supa_rest_service devolve (status, dados) — iterar a tupla
            # direto quebrava a thread em 0,1 s ("0 de 0" no admin, 01/08).
            if modo == "qualidade":
                _rodar_qualidade(_sp, _tf, t0)
                return
            _st_lw, rows = _supa_rest_service(
                "GET",
                "projects?status=eq.error&select=job_id,error_message,created_at"
                "&order=created_at.desc&limit=60")
            rows = rows if isinstance(rows, list) else []
            jobs = [r["job_id"] for r in rows
                    if "conseguimos abrir" in (r.get("error_message") or "").lower()]
            _LIBREDWG_TESTE["jobs_na_fila"] = len(jobs)
            for job in jobs:
                if _LIBREDWG_TESTE["testados"] >= limite or time.time() - t0 > 600:
                    break
                try:
                    nomes = [n for n in _supabase_storage_list(PRANCHAS_BUCKET, f"{job}/")
                             if n.lower().endswith(".dwg")]
                except Exception:
                    continue
                for nome in nomes:
                    if _LIBREDWG_TESTE["testados"] >= limite or time.time() - t0 > 600:
                        break
                    item = {"job_id": job, "arquivo": nome}
                    try:
                        dados = _supabase_storage_download_prancha(job, nome)
                        if not dados:
                            item["resultado"] = "não consegui baixar do Storage"
                        else:
                            with _tf.TemporaryDirectory() as tmp:
                                src = os.path.join(tmp, "in.dwg")
                                dst = os.path.join(tmp, "out.dxf")
                                open(src, "wb").write(dados)
                                item["mb"] = round(len(dados) / 1048576, 1)
                                try:
                                    proc = _sp.run([dwg2dxf, "-o", dst, src],
                                                   capture_output=True, timeout=90)
                                    item["exit"] = proc.returncode
                                    if not os.path.exists(dst) or os.path.getsize(dst) == 0:
                                        item["resultado"] = "dwg2dxf não gerou DXF"
                                        item["stderr"] = (proc.stderr or b"")[-200:].decode(
                                            "utf-8", "replace")
                                    else:
                                        item["dxf_mb"] = round(os.path.getsize(dst) / 1048576, 1)
                                        try:
                                            import ezdxf
                                            doc = ezdxf.readfile(dst)
                                            item["resultado"] = "CONVERTEU e o DXF abre"
                                            item["entidades"] = len(doc.modelspace())
                                            _LIBREDWG_TESTE["converteram_e_abrem"] += 1
                                        except Exception as e_dxf:
                                            item["resultado"] = "converteu mas o DXF não abre no ezdxf"
                                            item["erro_leitura"] = (
                                                f"{type(e_dxf).__name__}: {e_dxf}"[:120])
                                except _sp.TimeoutExpired:
                                    item["resultado"] = "timeout (90s)"
                    except Exception as e:
                        item["resultado"] = f"erro no teste: {type(e).__name__}: {e}"[:150]
                    _LIBREDWG_TESTE["resultados"].append(item)
                    _LIBREDWG_TESTE["testados"] += 1
            # registra o veredito no error_log — sobrevive a restart e aparece no painel
            try:
                _log_error("libredwg:batch",
                           f"{_LIBREDWG_TESTE['converteram_e_abrem']} de "
                           f"{_LIBREDWG_TESTE['testados']} DWG recusados converteram e abrem "
                           f"({round(time.time()-t0)}s)", None, severity="info")
            except Exception:
                pass
        finally:
            _LIBREDWG_TESTE["rodando"] = False
            _LIBREDWG_TESTE["fim"] = round(time.time() - t0, 1)

    threading.Thread(target=_rodar, daemon=True).start()
    return {"iniciado": True, **{k: v for k, v in _LIBREDWG_TESTE.items()}}


@app.get("/api/debug/storage-limit")
async def debug_storage_limit(request: Request, mb: int = 60):
    """Descobre o limite REAL de upload do Storage sem usar cliente como cobaia.

    Sobe um arquivo sintético de `mb` MB no bucket das pranchas e apaga em
    seguida. Devolve OK ou o motivo exato da recusa (413 = passou do limite
    global do projeto; timeout = lentidão, não limite).

    Existe porque em 31/07 (caso Ana) 4 de 37 projetos com CAD estavam sem o
    arquivo original e não dava pra saber se era o teto de 50 MB do plano Free
    ou o timeout de 60s — os dois suspeitos empilhados. Restrito a admin.
    """
    _require_admin(request)
    import tempfile, urllib.request, urllib.parse as _up
    PASTA_DIAG = "diagnostico"
    mb = max(1, min(600, int(mb)))
    _t0 = time.time()
    caminho = None
    try:
        # Conteúdo tipo DXF (texto repetido) pra o teste parecer o caso real:
        # bytes aleatórios não comprimem, texto sim, e é texto que a gente sobe.
        linha = b"  0\nLINE\n  8\nTESTE-DIAGNOSTICO\n 10\n0.0\n 20\n0.0\n"
        with tempfile.NamedTemporaryFile(suffix=".dxf", delete=False) as fh:
            caminho = fh.name
            escrito = 0
            alvo = mb * 1048576
            bloco = linha * 2048
            while escrito < alvo:
                fh.write(bloco)
                escrito += len(bloco)
        tamanho_real = os.path.getsize(caminho) / 1048576.0
        nome = "diagnostico-limite.dxf"
        ok = _supabase_storage_upload_prancha(caminho, PASTA_DIAG, nome)
        resp = {
            "tamanho_testado_mb": round(tamanho_real, 1),
            "subiu": ok,
            "motivo": None if ok else (_ULTIMA_FALHA_UPLOAD_PRANCHA or "não registrado"),
            "segundos": round(time.time() - _t0, 1),
        }
        # Limpa TUDO que existir na pasta de diagnóstico, não só o desta rodada.
        # 🪤 Antes eu montava a chave com o nome ORIGINAL, mas o upload sanitiza
        # o nome antes de gravar (`_sanitize_filename_for_storage`) — o delete
        # batia num caminho que não existia e deixava 250 MB parados no bucket.
        # Listar e apagar o que está lá de verdade não depende de adivinhar o nome.
        # 🪤 Varre a pasta atual E a antiga "_diagnostico": a 1ª versão gravava com
        # underscore (o job_id não passa pelo sanitizador, só o nome do arquivo).
        # Ao renomear a pasta eu quase deixei 250 MB órfãos pra sempre, porque a
        # limpeza nova só olhava o nome novo.
        apagados, erros = [], []
        for pasta in (PASTA_DIAG, "_diagnostico"):
            try:
                for n in _supabase_storage_list(PRANCHAS_BUCKET, f"{pasta}/"):
                    alvo = f"{pasta}/{_up.quote(n)}"
                    if _supabase_storage_delete(PRANCHAS_BUCKET, alvo):
                        apagados.append(f"{pasta}/{n}")
                    else:
                        erros.append(f"{pasta}/{n}")
            except Exception as e:
                erros.append(f"{pasta}: {type(e).__name__}: {e}")
        resp["apagado"] = (not erros)
        resp["apagados"] = apagados
        if erros:
            resp["erro_ao_apagar"] = "; ".join(str(x) for x in erros)
        return resp
    finally:
        if caminho:
            try:
                os.unlink(caminho)
            except Exception:
                pass


@app.get("/api/debug/dwg")
async def debug_dwg(request: Request):
    """Diagnóstico do suporte DWG. Restrito a admin (revela paths do FS)."""
    _require_admin(request)
    import shutil
    result = {
        "oda_which": shutil.which("ODAFileConverter"),
        # O Dockerfile instala libredwg-tools com "|| echo" — se o pacote não vier,
        # o build passa em silêncio e o fallback simplesmente não existe, sem ninguém
        # saber. Aqui a resposta fica explícita (null = fallback indisponível).
        "libredwg_which": shutil.which("dwg2dxf"),
        "oda_paths_checked": [],
    }
    # Verificar caminhos
    for p in ["/usr/bin/ODAFileConverter", "/usr/local/bin/ODAFileConverter",
              "/opt/ODAFileConverter/ODAFileConverter"]:
        result["oda_paths_checked"].append({"path": p, "exists": os.path.exists(p)})

    # Tentar importar dwg_extractor
    try:
        from dwg_extractor import _find_oda_converter, extract_from_file
        result["dwg_extractor_import"] = True
        oda = _find_oda_converter()
        result["oda_found_by_extractor"] = oda
    except Exception as e:
        result["dwg_extractor_import"] = False
        result["dwg_extractor_error"] = str(e)

    # libredwg (fallback open-source — opcional)
    result["libredwg_dwg2dxf"] = shutil.which("dwg2dxf")
    # Removido o find recursivo: caro (timeout) e revelador de layout do FS.

    return result


@app.get("/api/debug/oda-log/{job_id}")
async def get_oda_log(request: Request, job_id: str):
    """Retorna o log do ODA File Converter para um job.

    Restrito ao dono do projeto (ou admin) — log pode revelar nomes de
    arquivos enviados pelo usuário.
    """
    _require_project_owner(request, job_id)
    log_path = os.path.join(WORK_DIR, job_id, "_oda_log.txt")
    if os.path.exists(log_path):
        with open(log_path, 'r') as f:
            return {"log": f.read()}
    return {"log": "Log não encontrado", "path_checked": log_path}


@app.get("/api/status/{job_id}")
async def get_status(job_id: str, request: Request):
    """Retorna o status de processamento de um job.

    Restrito ao dono (fix IDOR 2026-07-22): current_step/error_message carregam
    o nome do arquivo enviado. A linha do projeto é inserida de forma síncrona no
    /api/process antes de devolver o job_id, então o owner já existe quando o
    frontend começa a pollar. Projetos anônimos ficam livres via _require_project_owner."""
    _require_project_owner(request, job_id)
    if job_id not in jobs:
        raise HTTPException(404, "Job não encontrado")
    return jobs[job_id]


@app.get("/api/download/{job_id}")
async def download_file(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Baixa a planilha gerada. Tenta cache local primeiro; se sumiu
    (Render redeploy), busca no Supabase Storage."""
    # Suaviza a checagem de job — se o JSON foi limpo no restart mas o
    # arquivo está no Storage, ainda servimos pra não perder o cliente.
    output_path = get_planilha_path(job_id)
    if not output_path:
        raise HTTPException(404, "Planilha não encontrada (nem em cache nem no Storage)")

    return FileResponse(
        output_path,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=f"quantitativos_aiarq_{job_id}.xlsx",
    )


@app.post("/api/notify/welcome")
async def notify_welcome(request: Request):
    """Boas-vindas no cadastro. Manda o email SÓ pro email do próprio JWT
    (não dá pra spammar terceiros) e avisa o Pedro de novo cliente.
    Chamado pelo frontend logo após o signUp."""
    user = _get_user_from_request(request)
    if not user or not user.get("email"):
        raise HTTPException(401, "Autenticação requerida")
    email = user["email"]
    # Nome + idade do profile. O disparo vem do 1º load do dashboard, mas o
    # guard no frontend é localStorage (por-browser): sem este gate, os
    # usuários ANTIGOS receberiam boas-vindas + gerariam alerta FALSO de "novo
    # cliente" no próximo acesso (poluindo a métrica de crescimento). Só trata
    # como NOVO quem criou o perfil agora há pouco — o perfil nasce no cadastro,
    # momentos antes deste 1º load; perfis antigos têm created_at de dias atrás.
    name = ""
    is_new = False
    try:
        import urllib.request as _urw
        _qw = f"{SUPABASE_URL}/rest/v1/profiles?user_id=eq.{user['id']}&select=full_name,created_at"
        _rw = _urw.Request(_qw, method="GET")
        _rw.add_header("apikey", SUPABASE_KEY)
        _rw.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _rows = _json.loads(_urw.urlopen(_rw, timeout=8).read().decode("utf-8"))
        if _rows:
            name = _rows[0].get("full_name") or ""
            try:
                from datetime import datetime as _dtw, timezone as _tzw, timedelta as _tdw
                _cad = _dtw.fromisoformat((_rows[0].get("created_at") or "").replace("Z", "+00:00"))
                is_new = (_dtw.now(_tzw.utc) - _cad) <= _tdw(hours=1)
            except Exception:
                is_new = False
    except Exception:
        pass
    # Cliente antigo (ou perfil sem data confiável): não manda nada.
    if not is_new:
        return {"status": "ok", "sent": False, "reason": "not_new"}
    import html as _hw
    sent = _send_welcome_email(email, name)
    # Alerta interno pro Pedro: novo cliente (em thread)
    try:
        import threading as _thw
        _ab = (f"<b>Novo cadastro! 🎉</b><br><b>Email:</b> {_hw.escape(email)}<br>"
               f"<b>Nome:</b> {_hw.escape(name or '(ainda não preencheu)')}")

        def _alerta_novo_cliente():
            # 🪤 Registra no email_auto_log com o MESMO kind/ref que o tick horário
            # usa pra deduplicar ("alerta_novo_cadastro"). Sem isso o Pedro recebe
            # DOIS alertas da mesma pessoa: este (imediato) e o do tick, até 1h
            # depois — foi o que aconteceu com o cliente Walter em 29/07.
            # Os dois existem de propósito e se completam: este é imediato mas só
            # dispara quando a pessoa ABRE o painel; o do tick pega também quem
            # criou conta e nunca voltou. Quem chegar primeiro cala o outro.
            if _notify_admin("Novo cliente cadastrado", _ab):
                try:
                    _email_auto_registrar(NOTIFY_EMAIL, "alerta_novo_cadastro",
                                          ref=(email or "").strip().lower())
                except Exception:
                    pass

        _thw.Thread(target=_alerta_novo_cliente, daemon=True).start()
    except Exception:
        pass
    return {"status": "ok", "sent": sent}


@app.post("/api/admin/send-welcome")
async def admin_send_welcome(request: Request):
    """Reenvia o email de boas-vindas pra um usuário específico (admin-only).
    Para quem se cadastrou ANTES do welcome existir. NÃO dispara o alerta de
    'novo cliente' (não é cadastro novo) e ignora o gate de created_at."""
    _require_admin(request)
    try:
        _data = await request.json()
    except Exception:
        _data = {}
    uid = ((_data or {}).get("user_id") or "").strip()
    # Preferimos o email/nome vindos do admin: a lista de usuários do admin lê
    # via RPC SECURITY DEFINER (enxerga tudo), então sempre tem o dado correto.
    # Isso funciona mesmo sem SUPABASE_SERVICE_ROLE_KEY (o endpoint é admin-only).
    email = ((_data or {}).get("email") or "").strip()
    name = (_data or {}).get("name") or ""
    # Fallback: se o admin não mandou email, tenta ler do profiles direto
    # (precisa service_role setada — senão o RLS bloqueia e volta vazio).
    if not email and uid:
        try:
            import urllib.request as _ura
            _qa = f"{SUPABASE_URL}/rest/v1/profiles?user_id=eq.{uid}&select=full_name,email"
            _ra = _ura.Request(_qa, method="GET")
            _ra.add_header("apikey", SUPABASE_KEY)
            _ra.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            _rows = _json.loads(_ura.urlopen(_ra, timeout=8).read().decode("utf-8"))
            if _rows:
                name = name or (_rows[0].get("full_name") or "")
                email = _rows[0].get("email") or ""
        except Exception as _e:
            raise HTTPException(500, f"Erro lendo perfil: {_e}")
    if not email:
        raise HTTPException(400, "Sem email pra enviar (mande 'email' no corpo ou configure a service_role)")
    if not name.strip() and uid:
        name = _name_from_auth(uid)
    sent = _send_welcome_email(email, name)
    return {"status": "ok", "sent": sent, "email": email, "name": name}


@app.post("/api/admin/send-nudge")
async def admin_send_nudge(request: Request):
    """Lembrete (admin-only) com login de 1 clique. kind='cadastro' pra quem
    ficou incompleto; kind='onboarding' pra quem tem conta mas 0 projetos.
    Email vem da lista do admin (não depende de ler profiles)."""
    _require_admin(request)
    try:
        _data = await request.json()
    except Exception:
        _data = {}
    email = ((_data or {}).get("email") or "").strip()
    name = (_data or {}).get("name") or ""
    kind = ((_data or {}).get("kind") or "cadastro").strip()
    uid = ((_data or {}).get("user_id") or "").strip()
    if not email:
        raise HTTPException(400, "email requerido")
    # Sempre tenta nome: se o profile não tinha, puxa do metadata do auth.
    if not name.strip() and uid:
        name = _name_from_auth(uid)
    if kind not in ("cadastro", "onboarding", "feedback"):
        raise HTTPException(400, "kind inválido (use 'cadastro', 'onboarding' ou 'feedback')")
    redirect = "https://ai.arq.br/feedback.html" if kind == "feedback" else "https://ai.arq.br/login.html"
    link = _generate_magic_link(email, redirect)
    if not link:
        raise HTTPException(502, "Não consegui gerar o link de login (verifique a service_role)")
    sent = _send_nudge_email(email, name, kind, link)
    return {"status": "ok", "sent": sent, "email": email, "kind": kind}


# ═══════════════════ EMAILS AUTOMÁTICOS (tick horário via pg_cron) ═══════════════════
# Automatiza o que o Pedro fazia na mão (decisão 07/07): reenviar convite de
# cadastro incompleto, lembrar quem nunca subiu prancha, chamar de volta quem
# sumiu há 30 dias. Anti-spam DURO: cada (email, kind) sai UMA vez na vida,
# registrado em email_auto_log (unique). Janela de recência evita ressuscitar
# cadastros antigos. EMAILS_AUTO=0 no ambiente desliga tudo. ?dry=1 = ensaio.

def _email_auto_ja_enviado(email: str, kind: str, ref: str = "") -> bool:
    """True se este lembrete já saiu (ou se a checagem falhar — na dúvida, NÃO envia)."""
    import urllib.request as _u, urllib.parse as _up, json as _j
    try:
        q = (f"{SUPABASE_URL}/rest/v1/email_auto_log?select=id"
             f"&email=eq.{_up.quote(email)}&kind=eq.{_up.quote(kind)}&ref=eq.{_up.quote(ref)}&limit=1")
        req = _u.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _j.loads(_u.urlopen(req, timeout=10).read().decode("utf-8"))
        return bool(rows)
    except Exception as e:
        print(f"[emails-auto] dedup check falhou ({kind}/{email}): {e} — NÃO enviando por segurança")
        return True


def _ja_recebeu_kind(email: str, kind: str) -> bool:
    """True se a pessoa já recebeu um e-mail deste tipo alguma vez (consulta o
    email_sent_log, que registra TODO envio — transacional ou automático).

    Usado pelo resgate do boas-vindas (02/08): o welcome só saía quando a
    pessoa abria o dashboard na 1ª hora de conta; quem criava conta e voltava
    no dia seguinte NUNCA recebia. Como aquele envio grava kind='boas_vindas'
    aqui, esta checagem cobre os dois caminhos e evita e-mail duplicado.
    Na dúvida (erro), diz True — não envia."""
    import urllib.request as _u, urllib.parse as _up, json as _j
    try:
        q = (f"{SUPABASE_URL}/rest/v1/email_sent_log?select=id"
             f"&email=ilike.{_up.quote(email)}&kind=eq.{_up.quote(kind)}&limit=1")
        req = _u.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        return bool(_j.loads(_u.urlopen(req, timeout=10).read().decode("utf-8")))
    except Exception as e:
        print(f"[emails-auto] checagem de kind falhou ({kind}/{email}): {e} — NÃO enviando")
        return True


def _email_auto_recente(email: str, dias: int = 7) -> bool:
    """True se a pessoa recebeu QUALQUER e-mail automático nos últimos `dias`.

    Cooldown global entre tipos: o dedup vida-inteira é por tipo, então nada
    impedia dois tipos DIFERENTES caírem na mesma caixa em dias seguidos (ex.:
    nudge de cadastro num dia, convite pro 2º projeto no outro). Regra do Pedro:
    não encher a caixa de ninguém — no máximo 1 automático por pessoa por semana.
    Lembrete que perdeu a vez espera o próximo tick; se a janela dele passou,
    paciência — silêncio é melhor que spam. Na dúvida (erro na checagem), NÃO envia.
    """
    import urllib.request as _u, urllib.parse as _up, json as _j
    from datetime import datetime as _dt, timedelta as _td, timezone as _tz
    try:
        corte = (_dt.now(_tz.utc) - _td(days=dias)).isoformat()
        q = (f"{SUPABASE_URL}/rest/v1/email_auto_log?select=id"
             f"&email=eq.{_up.quote(email)}&sent_at=gte.{_up.quote(corte)}&limit=1")
        req = _u.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        return bool(_j.loads(_u.urlopen(req, timeout=10).read().decode("utf-8")))
    except Exception as e:
        print(f"[emails-auto] cooldown check falhou ({email}): {e} — NÃO enviando por segurança")
        return True


def _email_auto_registrar(email: str, kind: str, ref: str = "") -> None:
    try:
        _supabase_insert("email_auto_log", {"email": email, "kind": kind, "ref": ref})
    except Exception as e:
        print(f"[emails-auto] registrar falhou ({kind}/{email}): {e}")


def _auth_admin_list_users(max_pages: int = 5) -> list[dict]:
    """Lista usuários via Auth Admin API (service_role). Paginado, best-effort."""
    import urllib.request as _u, json as _j
    users: list[dict] = []
    for page in range(1, max_pages + 1):
        try:
            req = _u.Request(f"{SUPABASE_URL}/auth/v1/admin/users?page={page}&per_page=200", method="GET")
            req.add_header("apikey", SUPABASE_SERVICE_ROLE_KEY)
            req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            data = _j.loads(_u.urlopen(req, timeout=15).read().decode("utf-8"))
            batch = data.get("users", data if isinstance(data, list) else [])
            if not batch:
                break
            users.extend(batch)
            if len(batch) < 200:
                break
        except Exception as e:
            print(f"[emails-auto] list users p{page} falhou: {e}")
            break
    return users


def _build_retorno30_email(name: str):
    """Monta (subject, html) do email de retorno após 30 dias sem projeto —
    isca: cronograma grátis. Separado do envio pra reuso no preview."""
    import html as _hn
    greet = _greeting_line(_hn.escape(name or ""))
    body = (_email_img("retorno-foto.jpg", "Interior de projeto de arquitetura", margem="2px 0 14px")
            + f"{greet}<br><br>Faz um tempinho que você não aparece por aqui — e desde a sua "
            f"última visita o AI.arq melhorou bastante: a leitura das pranchas ficou mais "
            f"precisa e agora todo projeto vira também <b>cronograma de obra</b> e <b>memorial "
            f"descritivo em rascunho</b> (editável na tela, sai em Word ou PDF) — de graça."
            f"<br><br>Se tiver um projeto na mesa, manda a prancha (PDF, DWG ou DXF) "
            f"que em minutos você recebe a planilha de quantitativos. Nessa fase de beta está "
            f"<b>grátis e ilimitado</b>.")
    subject = "Seu próximo quantitativo sai em minutos — e o cronograma é grátis"
    html = _email_wrap("Sentimos sua falta por aqui", body,
                       "Subir um projeto", "https://ai.arq.br/dashboard.html",
                       reason=("Você está recebendo este e-mail porque tem uma conta no AI.arq. "
                               "Se não quiser mais lembretes, é só responder avisando."))
    return subject, html


def _send_email_retorno30(email: str, name: str) -> bool:
    """Retorno após 30 dias sem projeto — isca: cronograma grátis."""
    name = _resolve_client_name(email, hint=name)
    subject, html = _build_retorno30_email(name)
    return _send_email_smtp(email, subject, html, log_kind="retorno_30d")


def _build_proximo_projeto_email(name: str, project_name: str):
    """Monta (subject, html) do convite pro 2º projeto, dias 3-10 após o 1º.

    Fecha o vão morto do funil: entre o "planilha pronta" (dia 0) e o
    retorno_30d (dia 30+) o produto ficava MUDO — e o pico de intenção pro 2º
    projeto é na primeira semana, logo depois da primeira entrega de valor.
    Convida o 2º projeto e, de carona, pede a planilha revisada (calibração).
    SEM mencionar dinheiro/cashback (beta grátis — mesma regra do calibracao)."""
    import html as _hp
    pn = _hp.escape(project_name or "seu primeiro projeto")
    greet = _greeting_line(_hp.escape(name or ""))
    body = (_email_img("proximo-foto.jpg", "Escritório de arquitetura moderno", margem="2px 0 14px")
            + f"{greet}<br><br>"
            f"O quantitativo do <b>{pn}</b> te ajudou? Se tiver outro projeto na mesa — "
            f"mesmo que seja um estudo ou uma reforma pequena — manda a prancha (PDF, DWG "
            f"ou DXF): o próximo sai em minutos, e nessa fase de beta está <b>grátis e "
            f"ilimitado</b>.<br><br>"
            f"E se você chegou a <b>revisar</b> a planilha do primeiro, subir a sua versão "
            f"corrigida na página do projeto afina o motor — os seus próximos quantitativos "
            f"saem medindo melhor exatamente o que você corrigiu.")
    subject = "Tem outro projeto na mesa? O próximo sai em minutos"
    html = _email_wrap("Bora fazer o próximo?", body,
                       "Subir outro projeto", "https://ai.arq.br/dashboard.html",
                       reason=("Você está recebendo este e-mail porque processou um projeto "
                               "no AI.arq. Se não quiser mais lembretes, é só responder avisando."))
    return subject, html


def _send_email_proximo_projeto(email: str, name: str, project_name: str) -> bool:
    name = _resolve_client_name(email, hint=name)
    subject, html = _build_proximo_projeto_email(name, project_name)
    return _send_email_smtp(email, subject, html, log_kind="proximo_projeto")


def _build_cronograma_checkin_email(name: str, project_name: str, semana: int, job_id: str):
    """Monta (subject, html) do check-in semanal de obra em andamento.

    O cronograma é a âncora de retorno do produto ([[feedback_cronograma_gratis]]):
    uma obra dura meses, mas a tela era one-shot — nada puxava a pessoa de volta
    pra atualizar o % executado. Este e-mail dispara pra obra EM ANDAMENTO
    (entre início e fim previstos), no máximo 1x/semana por projeto (ref traz a
    semana) — e o cooldown global de 7 dias garante que ninguém recebe junto com
    outro automático."""
    import html as _hk
    pn = _hk.escape(project_name or "sua obra")
    greet = _greeting_line(_hk.escape(name or ""))
    body = (f"{greet}<br><br>"
            f"Pelo cronograma, a obra <b>{pn}</b> está na <b>semana {semana}</b>. "
            f"Como estão as fases no canteiro?"
            + _email_img("checkin-gantt.png", "Cronograma com percentual executado por fase")
            + "<br>"
            f"Atualizar o % executado leva um minuto e mantém a <b>Curva S</b> e o "
            f"avanço real em dia — bom pra você enxergar atraso cedo, e pronto pra "
            f"mostrar ao cliente (dá pra exportar em PDF direto da tela).")
    subject = f"Semana {semana} da obra {project_name or 'do seu projeto'} — como está o avanço?"
    html = _email_wrap(f"Semana {semana} — bora atualizar?", body,
                       "Atualizar o avanço", f"https://ai.arq.br/cronograma.html?job={job_id}",
                       reason=("Você está recebendo este e-mail porque tem um cronograma de obra "
                               "em andamento no AI.arq. Se não quiser mais estes lembretes, é só "
                               "responder avisando."))
    return subject, html


def _send_email_cronograma_checkin(email: str, name: str, project_name: str,
                                   semana: int, job_id: str) -> bool:
    name = _resolve_client_name(email, hint=name)
    subject, html = _build_cronograma_checkin_email(name, project_name, semana, job_id)
    return _send_email_smtp(email, subject, html, log_kind="cronograma_checkin")


def _email_eh_interno(email: str) -> bool:
    """True pra contas do Pedro/teste — inclusive aliases (+smoke etc.).
    O dry-run de 07/07 pegou zarelalopes+smoke@ escapando do filtro exato."""
    e = (email or "").lower()
    if e == ADMIN_EMAIL:
        return True
    try:
        local, dom = e.split("@", 1)
        a_local, a_dom = ADMIN_EMAIL.split("@", 1)
        return dom == a_dom and local.split("+")[0] == a_local.split("+")[0]
    except ValueError:
        return False


TICK_SECRET = os.getenv("TICK_SECRET", "")


def _require_tick_secret(request):
    """Gate dos ticks de cron (achado da auditoria de 01/08/2026).

    Os 3 ticks (emails, newsletter, instagram) eram POST públicos sem segredo —
    qualquer um na internet podia disparar ciclo de envio/publicação, e o
    force_slot do IG permitia publicar post agendado ANTES da data.

    Retrocompatível: sem TICK_SECRET no ambiente, o gate fica aberto (nada
    quebra antes de setar a env var no Render). Com a env setada, exige o
    header X-Tick-Secret — que os pg_cron já enviam desde 01/08."""
    if TICK_SECRET and request.headers.get("X-Tick-Secret", "") != TICK_SECRET:
        raise HTTPException(401, "Tick não autorizado")


@app.post("/api/emails/auto/tick")
async def emails_auto_tick(request: Request, dry: int = 0):
    """Varredura horária (pg_cron): decide e envia os lembretes automáticos.
    dry=1 → só lista o que ENVIARIA, sem mandar nada (ensaio)."""
    _require_tick_secret(request)
    if os.environ.get("EMAILS_AUTO", "1") == "0":
        return {"status": "off"}
    from datetime import datetime as _dt, timezone as _tz
    import urllib.request as _u, json as _j

    now = _dt.now(_tz.utc)

    def _parse(ts):
        try:
            return _dt.fromisoformat(str(ts).replace("Z", "+00:00"))
        except Exception:
            return None

    users = _auth_admin_list_users()

    # 🪤 28/07/2026 — BUG QUE SEGUROU ESTE LEMBRETE DESDE SEMPRE:
    # "cadastro incompleto" NÃO é "e-mail não confirmado". A condição antiga era
    # `not confirmado`, mas no fluxo atual (Supabase + Google) TODA conta nasce
    # com email_confirmed_at preenchido — as 42 contas do banco estão assim.
    # Resultado: a regra nunca foi verdadeira e o template "Falta pouco pra
    # terminar seu cadastro" marcava 0 enviados no painel, com gente parada há
    # 21 dias sem receber nada.
    # Incompleto de verdade = criou login e NÃO criou o perfil em `profiles`.
    ids_com_perfil: set[str] = set()
    emails_com_perfil: set[str] = set()
    try:
        qp = f"{SUPABASE_URL}/rest/v1/profiles?select=user_id,email&limit=5000"
        rp = _u.Request(qp, method="GET")
        rp.add_header("apikey", SUPABASE_KEY)
        rp.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        for row in _j.loads(_u.urlopen(rp, timeout=15).read().decode("utf-8")):
            if row.get("user_id"):
                ids_com_perfil.add(str(row["user_id"]))
            if row.get("email"):
                emails_com_perfil.add(str(row["email"]).lower())
    except Exception as e:
        # Sem essa lista, todo mundo pareceria "sem perfil" e levaria cutucada
        # indevida. Na dúvida, NÃO envia.
        print(f"[emails-auto] profiles falhou: {e}")
        return {"status": "erro", "detail": "não consegui ler perfis — abortando por segurança"}

    # Quem JÁ recebeu o boas-vindas (por qualquer caminho: 1º acesso ao
    # dashboard ou resgate deste tick). Uma query só, mesmo padrão dos perfis.
    # 🪤 Se a leitura falhar, a lista fica "cheia" (None) e o resgate é PULADO
    # neste tick — melhor não mandar do que mandar boas-vindas repetido.
    emails_com_welcome: set[str] = set()
    _welcome_ok = True
    try:
        qw = (f"{SUPABASE_URL}/rest/v1/email_sent_log?select=email"
              f"&kind=eq.boas_vindas&limit=5000")
        rw = _u.Request(qw, method="GET")
        rw.add_header("apikey", SUPABASE_KEY)
        rw.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        for row in _j.loads(_u.urlopen(rw, timeout=15).read().decode("utf-8")):
            if row.get("email"):
                emails_com_welcome.add(str(row["email"]).lower())
    except Exception as e:
        print(f"[emails-auto] histórico de boas-vindas falhou: {e} — pulando resgate neste tick")
        _welcome_ok = False

    # projetos por email (pra saber quem nunca subiu / quem sumiu)
    proj_by_email: dict[str, list] = {}
    try:
        q = (f"{SUPABASE_URL}/rest/v1/projects?select=user_email,created_at,status,project_name,user_name,job_id"
             f"&order=created_at.desc&limit=2000")
        req = _u.Request(q, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        for p in _j.loads(_u.urlopen(req, timeout=15).read().decode("utf-8")):
            e = (p.get("user_email") or "").lower()
            if e:
                proj_by_email.setdefault(e, []).append(p)
    except Exception as e:
        print(f"[emails-auto] projetos falhou: {e}")
        return {"status": "erro", "detail": "não consegui ler projetos — abortando por segurança"}

    acoes: list[dict] = []
    H = 3600.0
    for u in users:
        email = (u.get("email") or "").lower()
        if not email or _email_eh_interno(email):
            continue
        created = _parse(u.get("created_at"))
        if not created:
            continue
        idade_h = (now - created).total_seconds() / H
        recente = idade_h <= 14 * 24  # janela: não ressuscitar cadastro antigo
        # Janela maior SÓ pro cadastro incompleto (decisão do Pedro, 28/07/2026):
        # como o gatilho estava quebrado desde sempre, essa gente nunca recebeu
        # nada — ficaram parados meses sem uma única mensagem. A janela de 14
        # dias é pra não ressuscitar quem já foi contatado; aqui ninguém foi.
        # O anti-spam continua valendo: cada pessoa recebe este lembrete UMA vez
        # na vida (dedup por email+kind) e no máximo 1 automático por semana.
        recente_cadastro = idade_h <= 60 * 24
        confirmado = bool(u.get("email_confirmed_at"))
        nome = ((u.get("user_metadata") or {}).get("full_name")
                or (u.get("user_metadata") or {}).get("name") or "")

        tem_perfil = (str(u.get("id") or "") in ids_com_perfil) or (email in emails_com_perfil)

        # 0) RESGATE DO BOAS-VINDAS (02/08/2026). O welcome só saía quando a
        #    pessoa abria o dashboard na PRIMEIRA HORA de conta — quem criava
        #    conta à noite e voltava no dia seguinte nunca recebia o cartão de
        #    visita. Aqui ele alcança todo mundo: 3h de espera (dá tempo do
        #    caminho normal acontecer) e conta de até 14 dias. Vem ANTES dos
        #    lembretes de propósito: primeiro se apresenta, depois cutuca.
        #    🪤 A checagem "já recebeu?" tem que estar NESTA condição (e não só
        #    no filtro lá embaixo): como o encadeamento é elif, um candidato
        #    descartado depois bloquearia o lembrete da mesma pessoa pra sempre.
        if _welcome_ok and 3 <= idade_h and recente and email not in emails_com_welcome:
            acoes.append({"kind": "boas_vindas", "email": email, "nome": nome})
        # 1) Parou ANTES de completar o cadastro (sem perfil). Vale tanto pra
        #    quem não confirmou o e-mail quanto pra quem confirmou, logou e
        #    abandonou a segunda etapa — que é o caso real que estava passando
        #    batido. 24h de espera pra não cutucar quem ainda está preenchendo.
        elif not tem_perfil and 24 <= idade_h and recente_cadastro:
            acoes.append({"kind": "nudge_cadastro", "email": email, "nome": nome})
        # 2) Completou o cadastro, mas nunca subiu prancha.
        elif tem_perfil and confirmado and 48 <= idade_h and recente and email not in proj_by_email:
            acoes.append({"kind": "nudge_onboarding", "email": email, "nome": nome})

    # retorno 30d: tem projeto concluído, último movimento entre 30 e 60 dias atrás
    for email, plist in proj_by_email.items():
        if _email_eh_interno(email):
            continue
        ultimo = max((_parse(p.get("created_at")) for p in plist if _parse(p.get("created_at"))), default=None)
        tem_done = any(p.get("status") == "done" for p in plist)
        if ultimo and tem_done:
            dias = (now - ultimo).total_seconds() / (24 * H)
            if 30 <= dias <= 60:
                acoes.append({"kind": "retorno_30d", "email": email, "nome": ""})
            # próximo projeto (dias 3-10): fecha o vão morto entre o "planilha
            # pronta" (dia 0) e o retorno_30d (dia 30) — o pico de intenção pro
            # 2º projeto é na 1ª semana. SÓ pra quem tem exatamente 1 projeto:
            # quem já subiu o 2º se ativou sozinho e não precisa de cutucada.
            # Dedup vida-inteira (sem ref) → cada pessoa recebe este no máximo
            # UMA vez, no primeiro projeto da conta.
            elif 3 <= dias <= 10 and len(plist) == 1 and plist[0].get("status") == "done":
                acoes.append({"kind": "proximo_projeto", "email": email,
                              "nome": plist[0].get("user_name") or "",
                              "projeto": plist[0].get("project_name") or ""})
            # calibração (dias 12-30, LIGADO 19/07): pede a PLANILHA REVISADA do
            # projeto — cada correção real afina o motor (loop de aprendizado).
            # Janela deliberadamente DEPOIS do proximo_projeto (3-10, que já
            # pede a revisada de carona) e ANTES do retorno_30d. Só se o job
            # ainda NÃO tem planilha_upload (checado no envio, 1 query); ref =
            # job_id → máx 1x por projeto, + cooldown global de 7d por pessoa.
            elif 12 <= dias <= 30:
                _dones = [p for p in plist if p.get("status") == "done" and p.get("job_id")]
                if _dones:
                    _dj = max(_dones, key=lambda p: p.get("created_at") or "")
                    acoes.append({"kind": "calibracao", "email": email,
                                  "nome": _dj.get("user_name") or "",
                                  "projeto": _dj.get("project_name") or "",
                                  "job_id": _dj["job_id"],
                                  "ref": _dj["job_id"]})

    # check-in de cronograma: obra EM ANDAMENTO (hoje entre início e fim previsto),
    # que ninguém mexeu nos últimos 7 dias (quem atualizou está engajado — não
    # precisa de cutucada). O ref carrega a SEMANA da obra, então o dedup
    # vida-inteira vira "1x por semana por projeto" — e o cooldown global de 7
    # dias abaixo garante no máximo 1 automático por pessoa por semana.
    try:
        proj_by_job = {}
        for plist in proj_by_email.values():
            for p in plist:
                if p.get("job_id"):
                    proj_by_job[p["job_id"]] = p
        qc = (f"{SUPABASE_URL}/rest/v1/cronogramas"
              f"?select=job_id,data_inicio,duracao_meses,updated_at&limit=500")
        rc = _u.Request(qc, method="GET")
        rc.add_header("apikey", SUPABASE_KEY)
        rc.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        for c in _j.loads(_u.urlopen(rc, timeout=15).read().decode("utf-8")):
            p = proj_by_job.get(c.get("job_id"))
            if not p:
                continue
            email = (p.get("user_email") or "").lower()
            if not email or _email_eh_interno(email):
                continue
            ini = _parse(str(c.get("data_inicio")) + "T00:00:00+00:00")
            upd = _parse(c.get("updated_at"))
            if not ini:
                continue
            dias_de_obra = (now - ini).total_seconds() / (24 * H)
            fim_dias = (c.get("duracao_meses") or 0) * 30
            mexeu_ha_pouco = upd and (now - upd).total_seconds() / (24 * H) < 7
            # começa a cutucar a partir da semana 2 (semana 1 = acabou de planejar)
            if 7 <= dias_de_obra <= fim_dias and not mexeu_ha_pouco:
                semana = int(dias_de_obra // 7) + 1
                acoes.append({"kind": "cronograma_checkin", "email": email,
                              "nome": p.get("user_name") or "",
                              "projeto": p.get("project_name") or "",
                              "job_id": c["job_id"], "semana": semana,
                              "ref": f"{c['job_id']}:w{semana}"})
    except Exception as e:
        print(f"[emails-auto] check-in cronograma falhou (pulando): {e}")

    # Três camadas anti-spam, nesta ordem (regra do Pedro: caixa de ninguém vira lixo):
    # 1) máx. 1 ação por PESSOA por tick (mesmo que ela caia em 2 regras diferentes);
    # 2) dedup vida-inteira por tipo (o ref, quando existe, limita por semana/projeto);
    # 3) cooldown global de 7 dias — nenhum automático pra quem recebeu QUALQUER
    #    automático na última semana (checado por último por custar 1 query/pessoa).
    # Teto de 5 por tick continua: goteja, nunca rajada.
    _vistos: set[str] = set()
    acoes = [a for a in acoes if not (a["email"] in _vistos or _vistos.add(a["email"]))]
    acoes = [a for a in acoes if not _email_auto_ja_enviado(a["email"], a["kind"], ref=a.get("ref", ""))]
    acoes = [a for a in acoes if not _email_auto_recente(a["email"], dias=7)][:5]

    # ── Alerta INTERNO pro Pedro: chegou gente nova ────────────────────────
    # 28/07/2026: chega por e-mail em até 1h, dizendo se a pessoa completou o
    # cadastro ou parou no meio.
    # ⚠️ CORREÇÃO 29/07: o comentário aqui dizia que não existia alerta de cadastro
    # nenhum. Existia — o "Novo cliente cadastrado", disparado quando a pessoa abre
    # o painel. Resultado: 2 e-mails pela mesma pessoa (caso Walter). Os dois ficam,
    # porque se completam — o imediato só dispara pra quem ABRE o painel, este pega
    # quem criou conta e sumiu. Agora aquele registra o mesmo kind/ref aqui embaixo,
    # então quem chegar primeiro cala o outro.
    # Dedup por e-mail → 1 alerta por pessoa, na vida. NÃO entra no cooldown dos
    # automáticos do usuário (aquilo protege a caixa do cliente; isto é interno).
    novos = []
    for _u2 in users:
        _e2 = (_u2.get("email") or "").lower()
        if not _e2 or _email_eh_interno(_e2):
            continue
        _c2 = _parse(_u2.get("created_at"))
        if not _c2 or (now - _c2).total_seconds() / H > 36:
            continue
        if _email_auto_ja_enviado(NOTIFY_EMAIL, "alerta_novo_cadastro", ref=_e2):
            continue
        novos.append({
            "email": _e2,
            "nome": ((_u2.get("user_metadata") or {}).get("full_name")
                     or (_u2.get("user_metadata") or {}).get("name") or ""),
            "tem_perfil": (str(_u2.get("id") or "") in ids_com_perfil) or (_e2 in emails_com_perfil),
        })

    if not dry and novos:
        import html as _ha
        for _n in novos:
            _sit = ("✅ completou o cadastro" if _n["tem_perfil"]
                    else "⚠️ <b>parou antes de completar o cadastro</b>")
            _corpo = (
                f"<b>{_ha.escape(_n['nome']) or '(ainda sem nome)'}</b> criou uma conta.<br>"
                f"E-mail: {_ha.escape(_n['email'])}<br>"
                f"Situação: {_sit}<br><br>"
                f'<a href="https://ai.arq.br/admin.html#usuarios">Abrir no painel</a>'
            )
            if _notify_admin(f"Cadastro novo — {_n['email']}", _corpo):
                _email_auto_registrar(NOTIFY_EMAIL, "alerta_novo_cadastro", ref=_n["email"])
            else:
                # Não registra: assim tenta de novo no próximo tick em vez de
                # perder o aviso porque o SMTP piscou.
                print(f"[emails-auto] alerta de cadastro novo NÃO entregue: {_n['email']}")

    # ── Lembrete da newsletter no ÚLTIMO DIA ÚTIL do mês ───────────────────
    # Decisão do Pedro (28/07/2026, opção B): o sistema NÃO dispara newsletter
    # sozinho — ele lembra, e o Pedro revisa a prévia e dispara na mão.
    # Motivo de existir: a edição de julho quase passou em branco. A última que
    # saiu foi 30/06 e ninguém percebeu que a de julho não tinha sido agendada.
    # Dedup por mês (ref='YYYY-MM') → no máximo 1 lembrete por mês.
    lembrete_news = 0
    try:
        from datetime import timedelta as _td_nl
        import calendar as _cal_nl
        import html as _ha_nl
        _agora_br = now - _td_nl(hours=3)  # Brasília = UTC-3 (sem horário de verão)
        _hoje_br = _agora_br.date()
        _ult = _hoje_br.replace(day=_cal_nl.monthrange(_hoje_br.year, _hoje_br.month)[1])
        while _ult.weekday() > 4:  # 5=sábado, 6=domingo
            _ult -= _td_nl(days=1)
        _ref_mes = _hoje_br.strftime("%Y-%m")
        # Janela de 1h pra não depender de o tick cair num minuto exato.
        if _hoje_br == _ult and 9 <= _agora_br.hour < 10:
            if not _email_auto_ja_enviado(NOTIFY_EMAIL, "lembrete_newsletter", ref=_ref_mes):
                # Se a edição do mês JÁ foi enviada, não enche o saco.
                _ja_saiu = False
                try:
                    _qn = (f"{SUPABASE_URL}/rest/v1/newsletter_scheduled"
                           f"?select=sent_at&status=eq.sent&sent_at=gte.{_hoje_br.replace(day=1)}T00:00:00Z")
                    _rn = _u.Request(_qn, method="GET")
                    _rn.add_header("apikey", SUPABASE_KEY)
                    _rn.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                    _ja_saiu = bool(_j.loads(_u.urlopen(_rn, timeout=15).read().decode("utf-8")))
                except Exception as _en:
                    print(f"[newsletter] checagem do mês falhou: {_en}")
                if not _ja_saiu:
                    lembrete_news = 1
                    if not dry:
                        _corpo_nl = (
                            f"Hoje é o <b>último dia útil de {_hoje_br.strftime('%m/%Y')}</b> e a "
                            f"newsletter do mês ainda não saiu.<br><br>"
                            f"Assunto atual: <b>{_ha_nl.escape(_NEWSLETTER_SUBJECT)}</b><br><br>"
                            f"Abra a prévia, confira e dispare — nada sai sem você clicar."
                            f'<br><br><a href="https://ai.arq.br/admin.html#newsletter">Abrir a aba Newsletter</a>'
                        )
                        if _notify_admin("Newsletter do mês — revisar e disparar hoje", _corpo_nl):
                            _email_auto_registrar(NOTIFY_EMAIL, "lembrete_newsletter", ref=_ref_mes)
    except Exception as _enl:
        print(f"[newsletter] lembrete do último dia útil falhou: {_enl}")

    if dry:
        # Endpoint é aberto (mesma mecânica dos outros ticks de cron). O modo dry
        # NÃO pode devolver email/nome — vazava PII de usuários beta pra qualquer
        # chamador anônimo. Devolve só a contagem por tipo. (Revisão 16/07.)
        from collections import Counter as _Counter
        return {
            "status": "dry",
            "total": len(acoes),
            "por_tipo": dict(_Counter(a["kind"] for a in acoes)),
            "alertas_novo_cadastro": len(novos),
            "lembrete_newsletter": lembrete_news,
        }

    enviados = []
    for a in acoes:
        ok = False
        try:
            if a["kind"] == "boas_vindas":
                ok = _send_welcome_email(a["email"], a["nome"])
            elif a["kind"] in ("nudge_cadastro", "nudge_onboarding"):
                link = _generate_magic_link(a["email"])
                if link:
                    ok = _send_nudge_email(a["email"], a["nome"],
                                           "cadastro" if a["kind"] == "nudge_cadastro" else "onboarding",
                                           link)
            elif a["kind"] == "proximo_projeto":
                ok = _send_email_proximo_projeto(a["email"], a["nome"], a.get("projeto") or "")
            elif a["kind"] == "calibracao":
                # se o cliente JÁ subiu a revisada deste job, não pede (silêncio)
                if _job_tem_planilha_revisada(a.get("job_id") or ""):
                    ok = False
                else:
                    ok = _send_email_calibracao(a["email"], a["nome"], a.get("projeto") or "",
                                                a.get("job_id") or "")
            elif a["kind"] == "cronograma_checkin":
                ok = _send_email_cronograma_checkin(a["email"], a["nome"],
                                                    a.get("projeto") or "",
                                                    a.get("semana") or 0,
                                                    a.get("job_id") or "")
            else:
                ok = _send_email_retorno30(a["email"], a["nome"])
        except Exception as e:
            print(f"[emails-auto] envio {a['kind']} -> {a['email']} falhou: {e}")
        if ok:
            _email_auto_registrar(a["email"], a["kind"], ref=a.get("ref", ""))
            enviados.append(a)
    if enviados:
        print(f"[emails-auto] tick enviou {len(enviados)}: {[(a['kind'], a['email']) for a in enviados]}")
    return {"status": "ok", "enviados": len(enviados), "detalhe": enviados}


# ── NEWSLETTER MENSAL ──
# Edição atual (revisada, com fonte). Pra trocar de mês: edita SUBJECT + HTML aqui
# e deploya; o admin dispara pelos botões do painel. {{UNSUB}} é trocado pelo link
# de descadastro por destinatário no envio.
_NEWSLETTER_SUBJECT = "IA que só olha a planta erra quase metade"
_NEWSLETTER_HTML = """<div style="background:#eaeef3;padding:24px 12px;font-family:Arial,Helvetica,sans-serif;color:#334155;">
<div style="max-width:560px;margin:0 auto;background:#ffffff;border-radius:14px;overflow:hidden;border:1px solid #e2e8f0;">
<div style="height:5px;background:#4F46E5;"></div>
<img src="https://ai.arq.br/assets/email/news-banner.png" width="560" alt="News do AI.arq — uma vez por m&ecirc;s: uma ideia &uacute;til e o que mudou por aqui" style="width:100%;max-width:560px;height:auto;display:block;border:0;">
<div style="padding:24px 30px;">
<p style="margin:0 0 14px;font-size:15px;line-height:1.6;">{{SAUDACAO}}</p>
<div style="padding:14px 16px;background:#F8FAFC;border-left:3px solid #4F46E5;border-radius:6px;font-size:14px;line-height:1.6;margin:0 0 18px;">
<b style="color:#0F172A;">&#129504; IA na arquitetura &mdash; o que ela ainda n&atilde;o faz bem</b><br>
O <b>AECV-Bench</b>, estudo de 2026, testou os melhores modelos de IA do mundo lendo pranchas. Em tarefas de s&iacute;mbolo &mdash; como <b>contar portas e janelas</b> &mdash; o acerto ficou <b>entre 40% e 55%</b>. Quase metade errada, s&oacute; &ldquo;olhando&rdquo; o desenho.<br><br>
&Eacute; por isso que o AI.arq <b>mede a geometria do arquivo</b> em vez de interpretar a imagem: no DWG e no DXF, cada parede e cada bloco t&ecirc;m coordenadas. O que n&atilde;o d&aacute; pra medir sai como <b>estimativa</b> &mdash; nunca em branco.
<div style="color:#94a3b8;font-size:12px;margin-top:6px;">Fonte: AECV-Bench (2026) &mdash; <a href="https://arxiv.org/abs/2601.04819" style="color:#8b93f6;">arxiv.org/abs/2601.04819</a></div>
</div>
<b style="color:#0F172A;font-size:15px;">&#10024; O que melhorou em julho</b>
<ul style="margin:8px 0 14px;padding-left:20px;font-size:14px;line-height:1.6;">
<li style="margin-bottom:8px;"><b>Projeto grande n&atilde;o trava no meio.</b> O motor salva o progresso prancha por prancha e retoma de onde parou. E o envio mostra o progresso real &mdash; arquivo pesado n&atilde;o parece mais travado.</li>
<li style="margin-bottom:8px;"><b>Planta sem cota n&atilde;o volta vazia.</b> Sem cotas no desenho, voc&ecirc; informa a <b>&aacute;rea total</b> no envio. Ela vira base pros c&aacute;lculos, sempre rotulada como <i>informada por voc&ecirc;</i> &mdash; n&atilde;o medida.</li>
<li><b>Menos n&uacute;mero estranho na planilha.</b> Um bebedouro chegou a aparecer com a &aacute;rea do pavimento inteiro: o motor leu a observa&ccedil;&atilde;o &ldquo;H=1,00m do piso&rdquo; e, pela palavra <i>piso</i>, tratou o item como o piso todo. Noutro projeto, uma TV veio em m&sup2; em vez de contada por pe&ccedil;a. Os dois corrigidos &mdash; se vir um n&uacute;mero estranho, responde esse e-mail. &Eacute; assim que a gente acha.</li>
</ul>
<p style="margin:0 0 18px;padding:10px 14px;background:#FFFBEB;border-radius:8px;font-size:14px;line-height:1.6;color:#78350F;">&#128161; <b>Dica:</b> exportando do Revit, do AutoCAD ou do ArchiCAD, salve em <b>DXF</b> &mdash; &eacute; o formato que mede melhor. PDF a gente l&ecirc;, mas vira estimativa.</p>
<div style="background:#F1F5F9;border-radius:10px;padding:16px;text-align:center;margin:0 0 18px;">
<p style="margin:0 0 12px;font-size:15px;color:#0F172A;"><b>Tem projeto novo?</b> Durante o beta &eacute; <b>gr&aacute;tis</b>, quantos projetos voc&ecirc; quiser.</p>
<a href="https://ai.arq.br/dashboard.html" style="display:inline-block;background:#4F46E5;color:#fff;text-decoration:none;padding:12px 24px;border-radius:10px;font-size:15px;font-weight:600;">Abrir o AI.arq</a>
</div>
<div style="border-top:1px solid #eef2f7;padding-top:16px;font-size:14px;color:#475569;">Um abra&ccedil;o,<br><b style="color:#0F172A;">Pedro</b> <span style="color:#94a3b8;">&mdash; AI.arq</span>
<div style="margin-top:12px;"><a href="{{WPP}}" style="display:inline-block;border:1px solid #25D366;color:#128C4A;text-decoration:none;padding:8px 14px;border-radius:8px;font-size:13px;font-weight:600;">&#128172; Falar no WhatsApp</a></div></div>
<div style="margin-top:14px;font-size:11px;color:#aab4c0;line-height:1.6;">Voc&ecirc; recebe porque tem conta no AI.arq. <a href="{{UNSUB}}" style="color:#8b93f6;">Sair da lista</a> &middot; <a href="https://ai.arq.br/privacidade.html" style="color:#8b93f6;">Privacidade</a></div>
</div></div></div>"""

import hmac as _hmac_nl, hashlib as _hashlib_nl


def _newsletter_token(email: str) -> str:
    key = (SUPABASE_SERVICE_ROLE_KEY or "fallback").encode()
    return _hmac_nl.new(key, (email or "").strip().lower().encode(), _hashlib_nl.sha256).hexdigest()[:24]


def _newsletter_recipients() -> list:
    """(email, nome) distintos de projects, menos optout e contas de teste/admin."""
    import urllib.request as _ur
    by_email = {}
    try:
        r = _ur.Request(f"{SUPABASE_URL}/rest/v1/projects?select=user_email,user_name&order=created_at.desc", method="GET")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        for row in _json.loads(_ur.urlopen(r, timeout=20).read().decode("utf-8")):
            e = (row.get("user_email") or "").strip().lower()
            n = (row.get("user_name") or "").strip()
            if not (e and "@" in e):
                continue
            if e not in by_email or (not by_email[e] and n):
                by_email[e] = n  # nome mais recente; preenche se o recente vier vazio
        r2 = _ur.Request(f"{SUPABASE_URL}/rest/v1/newsletter_optout?select=email", method="GET")
        r2.add_header("apikey", SUPABASE_KEY)
        r2.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        opt = {(o.get("email") or "").strip().lower() for o in _json.loads(_ur.urlopen(r2, timeout=20).read().decode("utf-8"))}
    except Exception as _e:
        print(f"[newsletter] recipients erro: {_e}")
        return []
    # 28/07: o repositório é PÚBLICO. E-mails de pessoas reais saíram daqui —
    # ficam só as impressões digitais (SHA-256). Mesma lista de antes:
    # a conta de smoke test e um descadastro manual.
    import hashlib as _hl_block
    _block_hashes = {
        "8efe29cfa06aab88dfd0bb7cc63e7d0aab297c3670d4c6626b15092b0119c35c",
        "f4100dd59b6bd35a3ee43afa1bc53540ec191a04cce22287a87664189c074c57",
    }

    def _blocked(e):
        return _hl_block.sha256(e.strip().lower().encode()).hexdigest() in _block_hashes

    return [
        (e, by_email[e])
        for e in sorted(by_email)
        if e not in opt and e != ADMIN_EMAIL and not _blocked(e) and "+smoke" not in e
    ]


def _newsletter_blast(subject, html_template, recipients):
    """Manda html_template (com {{SAUDACAO}} e {{UNSUB}}) pra cada (email, nome).
    Retorna (sent, fail). Best-effort por destinatário. Usado pelo envio manual e
    pelo tick agendado."""
    import urllib.parse as _up
    sent = 0
    fail = 0
    for email, name in recipients:
        try:
            first = (name or "").strip().split(" ")[0] if (name or "").strip() else ""
            greet = f"Olá, {first}!" if first else "Olá, tudo bem?"
            unsub = (f"https://ai-arq.onrender.com/api/newsletter/unsub"
                     f"?e={_up.quote(email)}&t={_newsletter_token(email)}")
            html = (html_template.replace("{{SAUDACAO}}", greet)
                    .replace("{{UNSUB}}", unsub)
                    .replace("{{WPP}}", _WHATSAPP_LINK))
            ok = _send_email_smtp(email, subject, html)
            sent += 1 if ok else 0
            fail += 0 if ok else 1
        except Exception as _e:
            print(f"[newsletter] envio {email}: {_e}")
            fail += 1
    return sent, fail


@app.post("/api/admin/newsletter/send")
async def admin_newsletter_send(request: Request):
    """Dispara a newsletter mensal (admin-only). body {test_only: bool}.
    test_only=true manda só pro ADMIN_EMAIL (pré-visualização). Senão manda pra
    lista (projects distinct, menos optout/teste). Cada email leva link de
    descadastro 1-clique personalizado."""
    _require_admin(request)
    try:
        data = await request.json()
    except Exception:
        data = {}
    test_only = bool((data or {}).get("test_only"))
    recipients = [(ADMIN_EMAIL, "Pedro")] if test_only else _newsletter_recipients()
    sent, fail = _newsletter_blast(_NEWSLETTER_SUBJECT, _NEWSLETTER_HTML, recipients)
    print(f"[newsletter] test_only={test_only} recipients={len(recipients)} sent={sent} fail={fail}")
    if not test_only and sent:
        # registra o envio manual no histórico (mesma tabela dos agendamentos)
        try:
            from datetime import datetime as _dt, timezone as _tz
            _now = _dt.now(_tz.utc).isoformat()
            _supabase_insert("newsletter_scheduled", {
                "subject": _NEWSLETTER_SUBJECT, "html_template": _NEWSLETTER_HTML,
                "scheduled_for": _now, "status": "sent",
                "recipients": len(recipients), "sent": sent, "sent_at": _now,
            })
        except Exception as _e:
            print(f"[newsletter] log histórico falhou: {_e}")
    return {"status": "ok", "test_only": test_only,
            "recipients": len(recipients), "sent": sent, "fail": fail}


@app.get("/api/admin/newsletter/preview")
async def admin_newsletter_preview(request: Request):
    """Edição atual renderizada (subject + html) pra prévia no admin."""
    _require_admin(request)
    html = (_NEWSLETTER_HTML.replace("{{SAUDACAO}}", "Olá, Pedro!")
            .replace("{{UNSUB}}", "#").replace("{{WPP}}", _WHATSAPP_LINK))
    return {"subject": _NEWSLETTER_SUBJECT, "html": html}


@app.get("/api/newsletter/unsub")
async def newsletter_unsub(e: str = "", t: str = ""):
    """Descadastro 1-clique (público, sem login). Valida token = HMAC(email) e
    grava em newsletter_optout. Retorna página HTML simples."""
    email = (e or "").strip().lower()
    shell = ('<!doctype html><html lang="pt-br"><head><meta charset="utf-8">'
             '<meta name="viewport" content="width=device-width,initial-scale=1"><title>AI.arq</title></head>'
             '<body style="font-family:Arial,Helvetica,sans-serif;background:#eaeef3;padding:60px 20px;'
             'text-align:center;color:#334155;"><div style="max-width:420px;margin:0 auto;background:#fff;'
             'border-radius:14px;padding:36px 28px;border:1px solid #e2e8f0;">{msg}</div></body></html>')
    if not email or not t or t != _newsletter_token(email):
        return HTMLResponse(shell.format(msg='<h2 style="color:#0F172A;">Link inv&aacute;lido</h2>'
            '<p>N&atilde;o consegui validar esse link. Se quiser sair da lista, responda o e-mail '
            'que a gente tira na hora.</p>'), status_code=400)
    try:
        _supabase_insert("newsletter_optout", {"email": email, "source": "link"})
    except Exception:
        pass
    return HTMLResponse(shell.format(msg='<div style="height:4px;background:#4F46E5;border-radius:4px;'
        'margin-bottom:18px;"></div><h2 style="color:#0F172A;">Pronto, voc&ecirc; saiu da lista</h2>'
        '<p>N&atilde;o vamos mais te mandar a newsletter. Suas planilhas e avisos de projeto continuam normais.</p>'
        '<p style="font-size:13px;color:#94a3b8;">Mudou de ideia? &Eacute; s&oacute; responder um e-mail nosso. &mdash; AI.arq</p>'))


@app.post("/api/admin/newsletter/schedule")
async def admin_newsletter_schedule(request: Request):
    """Agenda a edição ATUAL da newsletter pra uma data/hora (admin). Snapshot do
    conteúdo embutido vai pro banco; o tick dispara na hora."""
    _require_admin(request)
    try:
        data = await request.json()
    except Exception:
        data = {}
    when = ((data or {}).get("scheduled_for") or "").strip()
    if not when:
        raise HTTPException(400, "scheduled_for requerido (ISO 8601 com timezone)")
    ok = _supabase_insert("newsletter_scheduled", {
        "subject": _NEWSLETTER_SUBJECT,
        "html_template": _NEWSLETTER_HTML,
        "scheduled_for": when,
        "status": "pending",
    })
    if not ok:
        raise HTTPException(502, "Não consegui agendar (verifique a service_role)")
    return {"status": "ok", "scheduled_for": when}


@app.get("/api/admin/blog/agenda")
async def admin_blog_agenda(request: Request):
    """Agenda do blog (slug + título + data) pro Calendário de Conteúdo.

    🪤 POR QUE ESTE ENDPOINT EXISTE (31/07/2026): o calendário lia
    `/blog/posts.json` do site — mas o workflow de deploy APAGA esse arquivo de
    propósito (limpeza de 28/07, pra não expor post não publicado). Resultado:
    404, o fetch caía no catch, e o painel anunciava "a fila secou" com 10 posts
    agendados até outubro. A decisão de apagar do site está certa; o painel é
    que não podia depender dela. Aqui o backend lê o arquivo do próprio
    repositório, que ele tem em disco.

    Devolve só o necessário pro calendário — nunca o corpo dos posts.
    """
    _require_admin(request)
    # 🪤 O container do Render é construído SÓ com a pasta backend/ (`COPY . .`
    # com o contexto em backend/), então `../blog/posts.json` NÃO existe lá — a
    # 1ª versão deste endpoint respondia 503 em produção e funcionava só local.
    # Ordem: arquivo local (desenvolvimento) → repositório (produção).
    _p = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                      "blog", "posts.json")
    _raw = None
    _erros = []
    try:
        with open(_p, "r", encoding="utf-8") as _f:
            _raw = _json.load(_f)
    except Exception as e:
        _erros.append(f"local: {type(e).__name__}")
    if _raw is None:
        try:
            import urllib.request as _urb
            _u = ("https://raw.githubusercontent.com/pedrozellmer/ai-arq/"
                  "main/blog/posts.json")
            _req = _urb.Request(_u, headers={"User-Agent": "aiarq-admin"})
            _raw = _json.loads(_urb.urlopen(_req, timeout=12).read().decode("utf-8"))
        except Exception as e:
            _erros.append(f"repo: {type(e).__name__}")
    if _raw is None:
        # Falha explícita: o painel PRECISA distinguir "não consegui ler" de
        # "não há nada agendado" — foi essa confusão que gerou o alarme falso.
        raise HTTPException(503, "Não consegui ler a agenda do blog (" + "; ".join(_erros) + ")")
    _posts = _raw.get("posts", _raw) if isinstance(_raw, dict) else _raw
    _out = []
    for _p2 in (_posts or []):
        if not isinstance(_p2, dict):
            continue
        _out.append({
            "slug": _p2.get("slug", ""),
            "title": _p2.get("title", "") or _p2.get("slug", ""),
            "publish_date": _p2.get("publish_date", ""),
        })
    _out.sort(key=lambda x: x.get("publish_date") or "")
    return {"posts": _out, "total": len(_out)}


@app.get("/api/admin/newsletter/scheduled")
async def admin_newsletter_scheduled(request: Request):
    """Lista os agendamentos (admin)."""
    _require_admin(request)
    import urllib.request as _ur
    try:
        q = (f"{SUPABASE_URL}/rest/v1/newsletter_scheduled"
             f"?select=id,subject,scheduled_for,status,recipients,sent,sent_at"
             f"&order=scheduled_for.desc&limit=20")
        r = _ur.Request(q, method="GET")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _json.loads(_ur.urlopen(r, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        print(f"[newsletter] scheduled list erro: {_e}")
        rows = []
    return {"scheduled": rows}


@app.get("/api/admin/motor-health")
async def admin_motor_health(request: Request):
    """Painel de SAÚDE DO MOTOR (admin): % medido por fonte (CAD vs PDF), distribuição
    de % medido por projeto, itens medidos x total e top motivos de falha. Agrega via
    RPC admin_motor_health (SECURITY DEFINER, só service_role) pra o Pedro parar de
    otimizar às cegas (board 15/07)."""
    _require_admin(request)
    import urllib.request as _ur
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/admin_motor_health"
        req = _ur.Request(url, data=b"{}", method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        data = _json.loads(_ur.urlopen(req, timeout=20).read().decode("utf-8"))
    except Exception as _e:
        print(f"[motor-health] erro: {_e}")
        raise HTTPException(502, "Não consegui carregar a saúde do motor")

    # #3 PDF SOMBRA: o medidor vetorial de PDF roda em modo sombra (loga em
    # error_log stage 'pdfvec:shadow') e NÃO vira item medido — o item de PDF é
    # sempre estimado por design. Este readout mostra o que o vetor TERIA medido:
    # quantos PDFs pegaram escala, quanto de cômodo/parede foi achado. É o portão
    # pra um dia graduar PDF de sombra→medido com dado, não chute (board 15/07).
    _shadow = {"n_pdfs": 0, "com_escala": 0, "pulados": 0, "rooms_m2_total": 0.0,
               "walls_m_total": 0.0, "amostras": []}
    try:
        _sq = (f"{SUPABASE_URL}/rest/v1/error_log?stage=like.pdfvec*"
               f"&select=message,job_id,created_at&order=created_at.desc&limit=120")
        _sr = _ur.Request(_sq, method="GET")
        _sr.add_header("apikey", SUPABASE_KEY)
        _sr.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _slogs = _json.loads(_ur.urlopen(_sr, timeout=20).read().decode("utf-8"))
        for _row in _slogs:
            try:
                _pg = (_json.loads(_row.get("message") or "{}").get("pages") or [{}])[0]
            except Exception:
                continue
            _shadow["n_pdfs"] += 1
            if _pg.get("skip"):
                _shadow["pulados"] += 1
                continue
            if _pg.get("scale"):
                _shadow["com_escala"] += 1
            _rm = float(_pg.get("rooms_m2") or 0)
            _wm = float(_pg.get("walls_m") or 0)
            _shadow["rooms_m2_total"] += _rm
            _shadow["walls_m_total"] += _wm
            if (_rm > 0 or _wm > 0) and len(_shadow["amostras"]) < 8:
                _shadow["amostras"].append({
                    "job_id": _row.get("job_id"),
                    "escala": _pg.get("scale"),
                    "n_rooms": _pg.get("n_rooms") or 0,
                    "rooms_m2": round(_rm, 1),
                    "walls_m": round(_wm, 1),
                })
        _shadow["rooms_m2_total"] = round(_shadow["rooms_m2_total"], 1)
        _shadow["walls_m_total"] = round(_shadow["walls_m_total"], 1)
    except Exception as _se:
        print(f"[motor-health] pdf-shadow erro: {_se}")
    data["pdf_shadow"] = _shadow
    return data


# ── PAINEL FINANCEIRO (custos) ─────────────────────────────────────────
_COST_UUID_RE = __import__("re").compile(r"^[0-9a-fA-F-]{32,40}$")


@app.get("/api/admin/costs")
async def admin_costs_list(request: Request):
    """Lista as linhas de custo (admin). Tabela financial_costs = só service_role."""
    _require_admin(request)
    import urllib.request as _ur
    try:
        q = (f"{SUPABASE_URL}/rest/v1/financial_costs"
             f"?select=id,servico,categoria,valor,moeda,periodo,confirmado,obs,ordem"
             f"&order=ordem.asc,servico.asc")
        r = _ur.Request(q, method="GET")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _json.loads(_ur.urlopen(r, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        print(f"[costs] list erro: {_e}")
        rows = []
    # projetos concluídos nos últimos 30d → base pro "custo por projeto" no painel
    _proj30 = 0
    try:
        _since = (datetime.utcnow() - timedelta(days=30)).isoformat()
        _pq = (f"{SUPABASE_URL}/rest/v1/projects?status=eq.done&created_at=gte.{_since}&select=job_id")
        _pr = _ur.Request(_pq, method="GET")
        _pr.add_header("apikey", SUPABASE_KEY)
        _pr.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _proj30 = len(_json.loads(_ur.urlopen(_pr, timeout=15).read().decode("utf-8")))
    except Exception as _pe:
        print(f"[costs] contagem projetos 30d erro: {_pe}")
    return {"costs": rows, "projetos_30d": _proj30}


def _cost_to_float(v):
    """float robusto pro valor de custo. Aceita número, '1.234,56' (BR) ou
    '1234.56'. (Antes usava sf(), que só existe DENTRO de process_job → dava
    NameError e 500 ao salvar custo. Bug pego na revisão 16/07.)"""
    if v is None:
        return 0.0
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip().replace("R$", "").replace(" ", "")
    if not s:
        return 0.0
    if "," in s:                       # BR: vírgula é decimal; ponto é milhar
        s = s.replace(".", "").replace(",", ".")
    try:
        return float(s)
    except Exception:
        return 0.0


@app.post("/api/admin/costs")
async def admin_costs_upsert(request: Request):
    """Adiciona (sem id) ou edita (com id) uma linha de custo (admin)."""
    _require_admin(request)
    try:
        d = await request.json()
    except Exception:
        d = {}
    _cid = str(d.get("id") or "").strip()
    _fields = {
        "servico": str(d.get("servico") or "").strip()[:120],
        "categoria": str(d.get("categoria") or "outro").strip()[:40],
        "valor": _cost_to_float(d.get("valor", 0)),
        "moeda": str(d.get("moeda") or "BRL").strip()[:8],
        "periodo": str(d.get("periodo") or "mensal").strip()[:12],
        "confirmado": bool(d.get("confirmado", False)),
        "obs": str(d.get("obs") or "")[:500],
    }
    if not _fields["servico"]:
        raise HTTPException(400, "servico requerido")
    if _cid:
        if not _COST_UUID_RE.match(_cid):
            raise HTTPException(400, "id inválido")
        _fields["updated_at"] = datetime.utcnow().isoformat()
        ok = _supabase_update("financial_costs", "id", _cid, _fields)
    else:
        ok = _supabase_insert("financial_costs", _fields)
    if not ok:
        raise HTTPException(502, "Não consegui salvar o custo")
    return {"status": "ok"}


@app.post("/api/admin/costs/delete")
async def admin_costs_delete(request: Request):
    """Exclui uma linha de custo (admin)."""
    _require_admin(request)
    try:
        d = await request.json()
    except Exception:
        d = {}
    _cid = str(d.get("id") or "").strip()
    if not _COST_UUID_RE.match(_cid):
        raise HTTPException(400, "id inválido")
    import urllib.request as _ur
    try:
        q = f"{SUPABASE_URL}/rest/v1/financial_costs?id=eq.{_cid}"
        r = _ur.Request(q, method="DELETE")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        r.add_header("Prefer", "return=minimal")
        _ur.urlopen(r, timeout=15)
    except Exception as _e:
        print(f"[costs] delete erro: {_e}")
        raise HTTPException(502, "Não consegui excluir")
    return {"status": "ok"}


@app.get("/api/admin/ops")
async def admin_ops_panel(request: Request):
    """Painel de OPERAÇÃO (admin): erros do motor (error_log), falhas recentes com
    contato do dono (recuperar 1-a-1) e projetos só-PDF que mediram pouco (puxar pro
    CAD). Tudo via RPC admin_ops (SECURITY DEFINER, só service_role)."""
    _require_admin(request)
    import urllib.request as _ur
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/admin_ops"
        req = _ur.Request(url, data=b"{}", method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        data = _json.loads(_ur.urlopen(req, timeout=20).read().decode("utf-8"))
    except Exception as _e:
        print(f"[ops] erro: {_e}")
        raise HTTPException(502, "Não consegui carregar a operação")
    return data


@app.post("/api/admin/newsletter/cancel")
async def admin_newsletter_cancel(request: Request):
    """Cancela um agendamento pendente (admin)."""
    _require_admin(request)
    try:
        data = await request.json()
    except Exception:
        data = {}
    sid = ((data or {}).get("id") or "").strip()
    if not sid:
        raise HTTPException(400, "id requerido")
    import urllib.request as _ur
    try:
        q = f"{SUPABASE_URL}/rest/v1/newsletter_scheduled?id=eq.{sid}&status=eq.pending"
        r = _ur.Request(q, data=_json.dumps({"status": "canceled"}).encode("utf-8"), method="PATCH")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        r.add_header("Content-Type", "application/json")
        _ur.urlopen(r, timeout=15)
    except Exception as _e:
        raise HTTPException(502, f"Não consegui cancelar: {str(_e)[:120]}")
    return {"status": "ok"}


# ══════════════════════════════════════════════════
#  Instagram — painel admin (fila de posts)
#  Lê/gerencia instagram_scheduled_posts. O cron (/scheduler/tick) só publica
#  status='pending' com publish_at<=now; 'pending_approval' é rascunho (ignorado).
# ══════════════════════════════════════════════════

_IG_ADMIN_STATUSES = {"pending_approval", "pending", "canceled"}
_IG_UUID_RE = __import__("re").compile(r"^[0-9a-fA-F-]{32,40}$")


def _ig_has_media(post: dict) -> bool:
    """True se o post tem mídia suficiente pra publicar, conforme o tipo.
    Guarda a regra dura: não agendar (pending) sem imagem/vídeo — senão o cron falha."""
    mt = (post.get("media_type") or "feed").lower()
    if mt == "carousel":
        imgs = post.get("image_urls")
        if isinstance(imgs, str):
            try:
                imgs = _json.loads(imgs)
            except Exception:
                imgs = []
        return isinstance(imgs, list) and len([u for u in imgs if str(u).strip()]) >= 2
    if mt == "reel":
        return bool((post.get("video_url") or "").strip())
    return bool((post.get("image_url") or "").strip())  # feed / story


@app.get("/api/admin/instagram/posts")
async def admin_instagram_posts(request: Request, limit: int = 150):
    """Lista a fila do Instagram (admin): rascunhos, agendados e publicados."""
    _require_admin(request)
    from datetime import datetime as _dt, timezone as _tz
    import urllib.request as _ur
    limit = max(1, min(limit, 300))
    try:
        q = (f"{SUPABASE_URL}/rest/v1/instagram_scheduled_posts"
             f"?select=id,slot_key,media_type,caption,notes,publish_at,status,"
             f"image_url,image_urls,video_url,thumbnail_url,media_id,error_message,"
             f"attempts,published_at,created_at"
             f"&order=publish_at.desc&limit={limit}")
        r = _ur.Request(q, method="GET")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _json.loads(_ur.urlopen(r, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        print(f"[ig-admin] list erro: {_e}")
        raise HTTPException(502, "Não consegui listar os posts do Instagram")
    return {"posts": rows, "now": _dt.now(_tz.utc).isoformat()}


@app.post("/api/admin/instagram/posts/status")
async def admin_instagram_post_status(request: Request):
    """Muda o status de um post (admin): aprovar (->pending), pausar (->pending_approval), cancelar."""
    _require_admin(request)
    import urllib.request as _ur
    try:
        data = await request.json()
    except Exception:
        data = {}
    pid = str((data or {}).get("id") or "").strip()
    new_status = str((data or {}).get("status") or "").strip()
    if not pid or not _IG_UUID_RE.match(pid):
        raise HTTPException(400, "id inválido")
    if new_status not in _IG_ADMIN_STATUSES:
        raise HTTPException(400, f"status inválido (use: {', '.join(sorted(_IG_ADMIN_STATUSES))})")

    # busca o post pra validar mídia e estado atual
    try:
        q = (f"{SUPABASE_URL}/rest/v1/instagram_scheduled_posts"
             f"?id=eq.{pid}&select=id,media_type,image_url,image_urls,video_url,status")
        r = _ur.Request(q, method="GET")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        found = _json.loads(_ur.urlopen(r, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        raise HTTPException(502, f"Não consegui ler o post: {str(_e)[:120]}")
    if not found:
        raise HTTPException(404, "post não encontrado")
    post = found[0]

    if (post.get("status") or "") in ("published", "publishing"):
        raise HTTPException(409, "Post já publicado ou publicando — não dá pra mudar.")
    # GUARDA DURA: não entra na fila (pending) sem mídia — senão o cron publica errado/falha
    if new_status == "pending" and not _ig_has_media(post):
        raise HTTPException(400, "Esse post ainda não tem imagem/vídeo. Adicione a mídia antes de agendar.")

    try:
        u = f"{SUPABASE_URL}/rest/v1/instagram_scheduled_posts?id=eq.{pid}"
        body = _json.dumps({"status": new_status, "error_message": None}).encode("utf-8")
        r = _ur.Request(u, data=body, method="PATCH")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        r.add_header("Content-Type", "application/json")
        r.add_header("Prefer", "return=minimal")
        _ur.urlopen(r, timeout=15)
    except Exception as _e:
        raise HTTPException(502, f"Não consegui atualizar: {str(_e)[:120]}")
    return {"status": "ok", "new_status": new_status}


@app.post("/api/admin/instagram/posts/update")
async def admin_instagram_post_update(request: Request):
    """Edita campos de um post (admin): legenda, data, mídia, notas."""
    _require_admin(request)
    import re as _re_ig
    import urllib.request as _ur
    try:
        data = await request.json()
    except Exception:
        data = {}
    pid = str((data or {}).get("id") or "").strip()
    if not pid or not _IG_UUID_RE.match(pid):
        raise HTTPException(400, "id inválido")
    patch = {}
    for k in ("caption", "publish_at", "image_url", "video_url", "thumbnail_url", "notes"):
        if k in (data or {}):
            patch[k] = data[k]
    if "image_urls" in (data or {}):
        iu = data["image_urls"]
        if isinstance(iu, str):
            iu = [x.strip() for x in _re_ig.split(r"[\n,]+", iu) if x.strip()]
        patch["image_urls"] = iu
    if not patch:
        raise HTTPException(400, "nada pra atualizar")

    # lê o post atual: protege publicados e reavalia a regra dura de mídia
    try:
        q = (f"{SUPABASE_URL}/rest/v1/instagram_scheduled_posts"
             f"?id=eq.{pid}&select=status,media_type,image_url,image_urls,video_url")
        rq = _ur.Request(q, method="GET")
        rq.add_header("apikey", SUPABASE_KEY)
        rq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        cur = _json.loads(_ur.urlopen(rq, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        raise HTTPException(502, f"Não consegui ler o post: {str(_e)[:120]}")
    if not cur:
        raise HTTPException(404, "post não encontrado")
    cur = cur[0]
    if (cur.get("status") or "") in ("published", "publishing"):
        raise HTTPException(409, "Post já publicado ou publicando — não dá pra editar.")

    # GUARDA DURA: se o post está AGENDADO (pending) e a edição o deixa sem mídia,
    # rebaixa pra rascunho — senão o cron publica sem imagem e falha.
    demoted = False
    if (cur.get("status") or "") == "pending":
        merged = dict(cur)
        merged.update(patch)
        if not _ig_has_media(merged):
            patch["status"] = "pending_approval"
            demoted = True

    try:
        u = f"{SUPABASE_URL}/rest/v1/instagram_scheduled_posts?id=eq.{pid}"
        r = _ur.Request(u, data=_json.dumps(patch).encode("utf-8"), method="PATCH")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        r.add_header("Content-Type", "application/json")
        r.add_header("Prefer", "return=minimal")
        _ur.urlopen(r, timeout=15)
    except Exception as _e:
        raise HTTPException(502, f"Não consegui salvar: {str(_e)[:120]}")
    return {"status": "ok", "demoted": demoted}


@app.post("/api/newsletter/tick")
async def newsletter_tick(request: Request):
    """Chamado pelo pg_cron. Dispara 1 newsletter agendada vencida por vez, com
    claim atômico anti-duplicação. Gate por X-Tick-Secret desde 01/08."""
    _require_tick_secret(request)
    import urllib.request as _ur, urllib.parse as _up
    from datetime import datetime as _dt, timezone as _tz
    now_iso = _dt.now(_tz.utc).isoformat()

    def _patch(qs, payload, want_repr=False):
        req = _ur.Request(f"{SUPABASE_URL}/rest/v1/newsletter_scheduled?{qs}",
                          data=_json.dumps(payload).encode("utf-8"), method="PATCH")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        if want_repr:
            req.add_header("Prefer", "return=representation")
        resp = _ur.urlopen(req, timeout=25).read().decode("utf-8")
        return _json.loads(resp) if want_repr else None

    try:
        q = (f"{SUPABASE_URL}/rest/v1/newsletter_scheduled?status=eq.pending"
             f"&scheduled_for=lte.{_up.quote(now_iso)}&select=id&order=scheduled_for.asc&limit=1")
        r = _ur.Request(q, method="GET")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        due = _json.loads(_ur.urlopen(r, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        return {"ok": False, "error": str(_e)[:200]}
    if not due:
        return {"ok": True, "message": "nada agendado vencido"}
    sid = due[0]["id"]
    try:
        claimed = _patch(f"id=eq.{sid}&status=eq.pending", {"status": "sending"}, want_repr=True)
    except Exception as _e:
        return {"ok": False, "error": f"claim: {str(_e)[:150]}"}
    if not claimed:
        return {"ok": True, "message": "já reivindicado por outro tick"}
    row = claimed[0]
    recipients = _newsletter_recipients()
    sent, fail = _newsletter_blast(row.get("subject") or _NEWSLETTER_SUBJECT,
                                   row.get("html_template") or _NEWSLETTER_HTML, recipients)
    try:
        _patch(f"id=eq.{sid}", {"status": "sent", "recipients": len(recipients),
                                "sent": sent, "sent_at": _dt.now(_tz.utc).isoformat()})
    except Exception as _e:
        print(f"[newsletter-tick] update final falhou: {_e}")
    print(f"[newsletter-tick] {sid} enviado: {sent}/{len(recipients)} (fail {fail})")
    return {"ok": True, "sent": sent, "recipients": len(recipients)}


@app.get("/api/debug/service-role")
async def debug_service_role(request: Request):
    _require_admin(request)
    """Diagnóstico (sem PII): confirma se a service_role consegue ler dados
    protegidos por RLS (profiles). Se 'service_role_reads_profiles' for false,
    os emails automáticos pra usuário REAL (welcome/planilha-pronta/erro) não
    funcionam — falta setar SUPABASE_SERVICE_ROLE_KEY no Render. Retorna só
    booleans."""
    import urllib.request as _urd

    def _can_read(key):
        try:
            _r = _urd.Request(f"{SUPABASE_URL}/rest/v1/profiles?select=user_id&limit=1", method="GET")
            _r.add_header("apikey", SUPABASE_KEY)
            _r.add_header("Authorization", f"Bearer {key}")
            _rows = _json.loads(_urd.urlopen(_r, timeout=8).read().decode("utf-8"))
            return len(_rows) > 0
        except Exception:
            return False

    service_ok = _can_read(SUPABASE_SERVICE_ROLE_KEY)
    return {
        "service_role_set": SUPABASE_SERVICE_ROLE_KEY != SUPABASE_KEY,
        "service_role_reads_profiles": service_ok,
        "anon_reads_profiles": _can_read(SUPABASE_KEY),
        "emails_to_real_users_ok": service_ok,
    }


@app.get("/api/debug/email-preview")
async def email_preview(request: Request):
    """Manda 1 amostra de cada email transacional pro NOTIFY_EMAIL (pessoal do
    dono) pra pré-visualizar o layout. Admin-only. Recipiente FIXO."""
    _require_admin(request)
    to = NOTIFY_EMAIL
    fake_link = "https://ai.arq.br/login.html"
    out = {}
    out["1_boas_vindas"] = _send_welcome_email(to, "Pedro")
    _s1, _h1 = _build_falha_email("Pedro", "Residencial Vila Nova", True)
    out["2_erro_reprocessar"] = _send_email_smtp(to, _s1, _h1)
    _s2, _h2 = _build_falha_email("Pedro", "Residencial Vila Nova", False)
    out["3_erro_trocar_arquivo"] = _send_email_smtp(to, _s2, _h2)
    out["4_nudge_cadastro"] = _send_nudge_email(to, "Pedro", "cadastro", fake_link)
    out["5_nudge_onboarding"] = _send_nudge_email(to, "Pedro", "onboarding", fake_link)
    out["6_feedback"] = _send_nudge_email(to, "Pedro", "feedback", "")
    return {"to": to, "sent": out, "obs": "planilha-pronta voce ja viu (projeto Template Novo)"}


# ═══════════════════ CENTRAL DE EMAILS (painel admin) ═══════════════════
# Catálogo dos emails transacionais + preview (sem enviar). Cada tipo tem um
# helper _build_* que devolve (subject, html) — a Central de Emails renderiza o
# preview e mostra o volume por tipo (email_sent_log.kind).

# Catálogo: fonte única da verdade dos tipos que a Central de Emails mostra.
# grupo "auto" = disparado sozinho pelo sistema; "manual" = você dispara em
# Usuários. Tipos que são auto E manual ficam em "auto" com o gatilho anotando
# "também manual" (evita card duplicado).
_EMAIL_CATALOG = [
    {"key": "boas_vindas", "nome": "Boas-vindas", "grupo": "auto",
     "gatilho": "auto: 1º acesso ao painel · também manual (reenvio em Usuários)"},
    {"key": "planilha_pronta", "nome": "Planilha pronta", "grupo": "auto",
     "gatilho": "auto: quando o quantitativo termina de processar"},
    {"key": "erro_reprocessar", "nome": "Erro — reprocessar", "grupo": "auto",
     "gatilho": "auto: falha passageira (reprocessar resolve)"},
    {"key": "erro_trocar", "nome": "Erro — trocar arquivo", "grupo": "auto",
     "gatilho": "auto: arquivo não-quantificável (precisa de outro arquivo)"},
    {"key": "proximo_projeto", "nome": "Convite pro 2º projeto (dias 3-10)", "grupo": "auto",
     "gatilho": "auto: 1º projeto done há 3-10 dias, sem 2º projeto (tick horário, 1x por pessoa)"},
    {"key": "cronograma_checkin", "nome": "Check-in de obra (cronograma)", "grupo": "auto",
     "gatilho": "auto: obra em andamento, sem atualização há 7+ dias (máx 1x/semana por projeto)"},
    {"key": "retorno_30d", "nome": "Retorno após 30 dias", "grupo": "auto",
     "gatilho": "auto: 30 dias sem subir projeto (tick horário)"},
    {"key": "nudge_onboarding", "nome": "Lembrar de subir a 1ª prancha", "grupo": "auto",
     "gatilho": "auto: tem conta, 0 projetos · também manual (Usuários)"},
    {"key": "nudge_cadastro", "nome": "Lembrar de terminar o cadastro", "grupo": "auto",
     "gatilho": "auto: cadastro incompleto · também manual (Usuários)"},
    {"key": "calibracao", "nome": "Pedir planilha revisada (calibração)", "grupo": "auto",
     "gatilho": "auto: projeto concluído há 12-30 dias sem planilha revisada (1x por projeto)"},
    {"key": "feedback", "nome": "Como foi? (feedback)", "grupo": "manual",
     "gatilho": "manual: botão em Usuários (usuário com projeto)"},
]


def _render_email_by_type(key: str):
    """Devolve (subject, html) de um tipo de email com dados de EXEMPLO. NÃO
    envia nada. Usado só pelo preview do painel. Levanta KeyError se o tipo não
    existe."""
    nome = "Pedro"
    projeto = "Residencial Vila Nova"
    fake_link = "https://ai.arq.br/login.html"
    fake_job = "exemplo-1234"
    if key == "boas_vindas":
        return _build_welcome_email(nome)
    if key == "planilha_pronta":
        # Exemplo dos "próximos passos" (bloco real, dados fictícios) pro preview
        # ficar representativo — sem tocar em all_items reais.
        _extra = _next_steps_html(fake_job, 30, 42, False)
        return _build_planilha_pronta_email(nome, projeto, fake_job, 42, _extra)
    if key == "erro_reprocessar":
        return _build_falha_email(nome, projeto, True)
    if key == "erro_trocar":
        return _build_falha_email(nome, projeto, False)
    if key == "nudge_cadastro":
        return _build_nudge_email(nome, "cadastro", fake_link)
    if key == "nudge_onboarding":
        return _build_nudge_email(nome, "onboarding", fake_link)
    if key == "feedback":
        return _build_nudge_email(nome, "feedback", fake_link)
    if key == "retorno_30d":
        return _build_retorno30_email(nome)
    if key == "proximo_projeto":
        return _build_proximo_projeto_email(nome, projeto)
    if key == "cronograma_checkin":
        return _build_cronograma_checkin_email(nome, projeto, 6, fake_job)
    if key == "calibracao":
        return _build_calibracao_email(nome, projeto)
    raise KeyError(key)


@app.get("/api/admin/email-render")
async def admin_email_render(request: Request, type: str = ""):
    """Renderiza (subject, html) de UM tipo de email com dados de EXEMPLO.
    Admin-only. NÃO envia nada — é só pra pré-visualizar no painel."""
    _require_admin(request)
    key = (type or "").strip()
    try:
        subject, html = _render_email_by_type(key)
    except KeyError:
        raise HTTPException(400, f"tipo de email desconhecido: {key}")
    return {"type": key, "subject": subject, "html": html}


@app.get("/api/admin/email-catalog")
async def admin_email_catalog(request: Request):
    """Lista os tipos de email da Central de Emails com {key, nome, grupo,
    gatilho, volume}. `volume` = quantos já saíram (email_sent_log por kind).
    Admin-only."""
    _require_admin(request)
    # Agrega volume por kind num único fetch (baixo volume no beta).
    volumes = {}
    try:
        _st, _rows = _supa_rest_service("GET", "/rest/v1/email_sent_log?select=kind")
        if _rows:
            for _r in _rows:
                _k = (_r or {}).get("kind") or ""
                if _k:
                    volumes[_k] = volumes.get(_k, 0) + 1
    except Exception as _e:
        print(f"[email-catalog] volume falhou (nao critico): {_e}")
    items = []
    for c in _EMAIL_CATALOG:
        items.append({**c, "volume": int(volumes.get(c["key"], 0))})
    return {"items": items}


@app.get("/api/health")
async def health():
    """Health check com métricas do sistema."""
    try:
        import psutil
    except ImportError:
        psutil = None
    api_key = os.getenv("ANTHROPIC_API_KEY", "")
    stripe_key = os.getenv("STRIPE_SECRET_KEY", "")

    # Métricas de sistema
    if psutil:
        mem = psutil.virtual_memory()
        disk = psutil.disk_usage('/')
    else:
        mem = None
        disk = None

    # Contar projetos hoje
    today_count = 0
    try:
        import urllib.request, json
        url = f"{SUPABASE_URL}/rest/v1/projects?select=id&created_at=gte.{datetime.utcnow().strftime('%Y-%m-%d')}T00:00:00"
        req = urllib.request.Request(url)
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Prefer', 'count=exact')
        resp = urllib.request.urlopen(req, timeout=3)
        count_header = resp.headers.get('content-range', '')
        if '/' in count_header:
            today_count = int(count_header.split('/')[1])
        else:
            today_count = len(json.loads(resp.read()))
    except:
        pass

    # Contar totais via Supabase count=exact.
    # FIX 2026-05-14: antes usava locals()[var_name] = ... que é noop em
    # CPython (locals() não tem write-back). Resultava em total_projects=0
    # e total_users=0 sempre, quebrando dashboard de métrica admin.
    total_projects = 0
    total_users = 0
    try:
        import urllib.request, json
        def _count_table(table):
            url = f"{SUPABASE_URL}/rest/v1/{table}?select=id"
            req = urllib.request.Request(url)
            req.add_header('apikey', SUPABASE_KEY)
            req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
            req.add_header('Prefer', 'count=exact')
            resp = urllib.request.urlopen(req, timeout=3)
            count_header = resp.headers.get('content-range', '')
            if '/' in count_header:
                return int(count_header.split('/')[1])
            return len(json.loads(resp.read()))
        try: total_projects = _count_table('projects')
        except Exception: pass
        try: total_users = _count_table('profiles')
        except Exception: pass
    except Exception:
        pass

    return {
        "status": "healthy",
        "api_key_configured": bool(api_key and api_key.startswith("sk-")),
        "stripe_configured": bool(stripe_key),
        "timestamp": datetime.utcnow().isoformat(),
        "system": {
            "ram_used_pct": round(mem.percent, 1) if mem else 0,
            "ram_used_mb": round(mem.used / 1024 / 1024) if mem else 0,
            "ram_total_mb": round(mem.total / 1024 / 1024) if mem else 0,
            "disk_used_pct": round(disk.percent, 1) if disk else 0,
            "cpu_pct": psutil.cpu_percent(interval=0.5) if psutil else 0,
        },
        "stats": {
            "projects_today": today_count,
            "total_projects": total_projects,
            "total_users": total_users,
        },
        "features": {
            "pdf": True,
            "dxf": True,
            "dwg": shutil.which("ODAFileConverter") is not None,
            "calibrator": HAS_CALIBRATOR if 'HAS_CALIBRATOR' in dir() else False,
            # Fallback pros DWG que o ODA recusa. Duas coisas separadas: o
            # binário existir e a chave estar ligada no Render. Ficam à vista
            # porque "liguei lá" e "o backend está usando" já se desencontraram
            # antes — e sem isso só dá pra conferir com token de admin.
            "libredwg_instalado": shutil.which("dwg2dxf") is not None,
            "libredwg_ligado": os.getenv("LIBREDWG_FALLBACK", "0").strip().lower()
                               in ("1", "true", "on", "sim"),
        }
    }


# ── TCPO BIM: busca técnica de composições ──
@app.get("/api/tcpo/search")
async def tcpo_search(q: str, limit: int = 5):
    """Busca composições TCPO BIM por similaridade de descrição.

    Exemplo: /api/tcpo/search?q=luminaria+fluorescente&limit=5

    Retorna top N composições ordenadas por similaridade, com o 1º match
    enriquecido com a lista de insumos (material/mão de obra/equipamento).
    Usado pela UI pra permitir que o revisor busque a referência TCPO ideal
    pra cada item.
    """
    try:
        from tcpo_matcher import match_item, get_insumos
    except ImportError:
        return {"results": [], "error": "tcpo_matcher indisponível"}

    if not q or len(q.strip()) < 2:
        return {"results": []}

    matches = match_item(q.strip(), limit=min(int(limit or 5), 15))
    # Enriquece o 1º com insumos pra UI poder mostrar direto
    if matches:
        matches[0]['insumos'] = get_insumos(matches[0]['id'])
    return {"query": q, "count": len(matches), "results": matches}


@app.get("/api/tcpo/details/{composicao_id}")
async def tcpo_details(composicao_id: str):
    """Retorna detalhes completos de uma composição TCPO — todos os insumos,
    metadados (conteúdo, critério, normas). Usado quando o usuário expande uma
    referência pra ver a composição inteira.
    """
    try:
        from tcpo_matcher import get_insumos, _supabase_rpc
    except ImportError:
        return {"error": "tcpo_matcher indisponível"}

    rows = _supabase_rpc("get_tcpo_details", {"p_id": composicao_id}) or []
    if not rows:
        return {"error": "composição não encontrada"}
    return rows[0] if isinstance(rows, list) else rows


# ── HEURÍSTICAS DE MERCADO: alertas por categoria ──
@app.get("/api/heuristics/check")
async def heuristics_check(description: str, unit: str = "",
                            typology: str = "office"):
    """Retorna alertas de heurística pra um item (descrição + unidade).

    Exemplo: /api/heuristics/check?description=demolicao+de+drywall&typology=office

    Retorna: {
        'category': 'demolicao',
        'alertas': ['💡 variação 103%...'],
        'metrics': {'dispersion': {...}, 'mat_mo_share': {...}, 'coverage': {...}}
    }
    """
    try:
        from market_heuristics import (
            categorize_item, check_item_anomaly,
            get_dispersion_for_category, get_mat_mo_share_for_category,
            get_coverage_pattern_for_category,
        )
    except ImportError:
        return {"error": "market_heuristics indisponível"}

    if not description or len(description.strip()) < 3:
        return {"error": "descrição muito curta"}

    category = categorize_item(description)
    item_dict = {"description": description.strip(), "unit": unit.strip()}
    alertas = check_item_anomaly(item_dict, typology=typology)

    return {
        "category": category,
        "typology": typology,
        "alertas": alertas,
        "metrics": {
            "dispersion": get_dispersion_for_category(category, typology),
            "mat_mo_share": get_mat_mo_share_for_category(category, typology),
            "coverage": get_coverage_pattern_for_category(category, typology),
        },
    }


@app.get("/api/heuristics/summary")
async def heuristics_summary():
    """Retorna resumo do que tem na base de heurísticas (pra debug/transparência)."""
    try:
        from market_heuristics import get_summary
        return get_summary()
    except ImportError:
        return {"error": "market_heuristics indisponível"}


# ═══════════════════════════════════════════════════════════════
#  CHAT PÚBLICO (widget de vendas/dúvidas no site) + CAPTURA DE LEADS
# ═══════════════════════════════════════════════════════════════

# Rate limit simples em memória: max 20 mensagens / 10min por IP
_PUBLIC_CHAT_HITS: dict = {}  # ip → [timestamps]


# ── Rate limit genérico por IP + guarda anti-SSRF (fix segurança 2026-07-22) ──
_RATE_BUCKETS: dict = {}  # (bucket, ip) → [timestamps]


def _client_ip(request) -> str:
    _xff = [p.strip() for p in (request.headers.get("x-forwarded-for") or "").split(",") if p.strip()]
    return (_xff[-1] if _xff else "") or (request.client.host if request.client else "unknown")


def _rate_limit_ok(bucket: str, request, limit: int, window_s: int) -> bool:
    """True se dentro do limite; False se estourou. Best-effort em memória (por IP)."""
    ip = _client_ip(request)
    now_ts = datetime.utcnow().timestamp()
    key = (bucket, ip)
    hist = [t for t in _RATE_BUCKETS.get(key, []) if now_ts - t < window_s]
    if len(hist) >= limit:
        _RATE_BUCKETS[key] = hist
        return False
    hist.append(now_ts)
    _RATE_BUCKETS[key] = hist
    # Anti memory-leak: limpa buckets velhos de vez em quando.
    if len(_RATE_BUCKETS) > 500 and (int(now_ts) % 100) == 0:
        for k in [k for k, v in _RATE_BUCKETS.items() if not v or (now_ts - max(v)) > window_s]:
            _RATE_BUCKETS.pop(k, None)
    return True


def _url_is_safe_public(url: str) -> bool:
    """Guarda anti-SSRF: aceita só http/https cujo host resolve pra IP PÚBLICO.
    Barra file://, localhost, 169.254.169.254 (metadata da nuvem) e redes
    privadas — impede que uma logo_url arbitrária faça o servidor ler recurso
    interno."""
    import ipaddress
    import socket
    from urllib.parse import urlparse
    try:
        p = urlparse(url or "")
    except Exception:
        return False
    if p.scheme not in ("http", "https") or not p.hostname:
        return False
    try:
        port = p.port or (443 if p.scheme == "https" else 80)
        infos = socket.getaddrinfo(p.hostname, port, proto=socket.IPPROTO_TCP)
    except Exception:
        return False
    for info in infos:
        try:
            ip = ipaddress.ip_address(info[4][0])
        except ValueError:
            return False
        if (ip.is_private or ip.is_loopback or ip.is_link_local
                or ip.is_reserved or ip.is_multicast or ip.is_unspecified):
            return False
    return True


@app.post("/api/public/chat/lead")
async def save_chat_lead(request: Request):
    """Salva (ou atualiza) um lead do chat widget.

    Chamado quando o visitante preenche nome + email ANTES da primeira
    mensagem. Usa email como chave de dedupe — se já existir, atualiza.
    """
    try:
        body = await request.json()
    except Exception:
        body = {}

    name = str(body.get("name", "")).strip()[:120]
    email = str(body.get("email", "")).strip().lower()[:200]
    phone = str(body.get("phone", "")).strip()[:40]
    source_page = str(body.get("source_page", "")).strip()[:80]
    first_question = str(body.get("first_question", "")).strip()[:500]
    user_agent = request.headers.get("user-agent", "")[:300]

    if not name or not email or "@" not in email:
        return {"error": "Nome e email válidos são obrigatórios"}

    try:
        import urllib.parse as _up
        # Upsert por email (se já existe, atualiza last_message_at e n_messages)
        # Primeiro tenta achar
        find_url = (f"{SUPABASE_URL}/rest/v1/chat_leads"
                    f"?email=eq.{_up.quote(email)}&select=id,n_messages")
        req = urllib.request.Request(find_url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = urllib.request.urlopen(req, timeout=8)
        existing = json.loads(resp.read().decode("utf-8"))

        now_iso = datetime.utcnow().isoformat()

        if existing:
            # Atualiza existente
            lead_id = existing[0]["id"]
            n_msgs = (existing[0].get("n_messages") or 0) + 1
            upd_url = f"{SUPABASE_URL}/rest/v1/chat_leads?id=eq.{lead_id}"
            upd_payload = {
                "name": name,
                "phone": phone or None,
                "n_messages": n_msgs,
                "last_message_at": now_iso,
            }
            if first_question:
                upd_payload["first_question"] = first_question[:500]
            body_bytes = json.dumps(upd_payload).encode("utf-8")
            upd_req = urllib.request.Request(upd_url, data=body_bytes, method="PATCH")
            upd_req.add_header("apikey", SUPABASE_KEY)
            upd_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            upd_req.add_header("Content-Type", "application/json")
            upd_req.add_header("Prefer", "return=minimal")
            urllib.request.urlopen(upd_req, timeout=8)
            return {"status": "ok", "lead_id": lead_id, "new": False}

        # Insere novo
        ins_payload = {
            "name": name,
            "email": email,
            "phone": phone or None,
            "source_page": source_page,
            "user_agent": user_agent,
            "first_question": first_question[:500],
            "n_messages": 1,
        }
        ins_url = f"{SUPABASE_URL}/rest/v1/chat_leads"
        body_bytes = json.dumps(ins_payload).encode("utf-8")
        ins_req = urllib.request.Request(ins_url, data=body_bytes, method="POST")
        ins_req.add_header("apikey", SUPABASE_KEY)
        ins_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        ins_req.add_header("Content-Type", "application/json")
        ins_req.add_header("Prefer", "return=representation")
        ins_resp = urllib.request.urlopen(ins_req, timeout=8)
        inserted = json.loads(ins_resp.read().decode("utf-8"))
        return {
            "status": "ok",
            "lead_id": inserted[0]["id"] if inserted else None,
            "new": True,
        }
    except Exception as e:
        print(f"[chat_lead] erro: {e}")
        return {"error": str(e)}


@app.get("/api/admin/chat/leads")
async def admin_list_chat_leads(request: Request, limit: int = 200):
    """Lista leads do chat pra painel admin. Requer auth admin."""
    _require_admin(request)
    try:
        url = (f"{SUPABASE_URL}/rest/v1/chat_leads"
               f"?select=*&order=last_message_at.desc&limit={int(limit)}")
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = urllib.request.urlopen(req, timeout=15)
        return {"leads": json.loads(resp.read().decode("utf-8"))}
    except Exception as e:
        return {"error": str(e)}


# ══════════════════════════════════════════════════
#  Formulário de contato (público)
# ══════════════════════════════════════════════════

@app.post("/api/contact")
async def submit_contact(request: Request):
    """Recebe envio do formulário de contato (JSON ou multipart/form-data).

    Aceita arquivo opcional (campo 'file', máx 10MB) que vai pro Supabase
    Storage no bucket 'contact-attachments'. URL pública é salva junto.
    """
    # Anti-spam/flood: público sem login — limita por IP (fix segurança 2026-07-22).
    if not _rate_limit_ok("contact", request, limit=8, window_s=600):
        return {"ok": False, "error": "Muitas mensagens em pouco tempo. Espere alguns minutos e tente de novo."}
    content_type = (request.headers.get("content-type") or "").lower()
    upload_file = None
    upload_filename = None

    # Parse: JSON ou multipart
    if "multipart/form-data" in content_type:
        try:
            form = await request.form()
        except Exception as e:
            return {"ok": False, "error": f"Form inválido: {e}"}
        body = dict(form)
        # Extrai arquivo se vier
        file_obj = form.get("file")
        if file_obj is not None and hasattr(file_obj, "read"):
            try:
                upload_filename = (getattr(file_obj, "filename", "") or "anexo.bin")[:200]
                upload_file = await file_obj.read()
                if len(upload_file) > 10 * 1024 * 1024:
                    return {"ok": False, "error": "Arquivo grande demais (máx 10MB)"}
            except Exception as e:
                print(f"[contact] erro lendo arquivo: {e}")
                upload_file = None
    else:
        try:
            body = await request.json()
        except Exception:
            return {"ok": False, "error": "JSON inválido"}

    name    = str(body.get("name", "")).strip()[:200]
    email   = str(body.get("email", "")).strip().lower()[:200]
    phone   = str(body.get("phone", "")).strip()[:50]
    msg_type = str(body.get("type", "duvida")).strip().lower()
    subject = str(body.get("subject", "")).strip()[:300]
    message = str(body.get("message", "")).strip()[:5000]

    if not name or not email or not message:
        return {"ok": False, "error": "Nome, email e mensagem são obrigatórios"}
    if "@" not in email or "." not in email:
        return {"ok": False, "error": "Email inválido"}
    if msg_type not in ("reclamacao", "sugestao", "duvida", "parceria", "elogio", "outro"):
        msg_type = "outro"

    # Metadados de origem
    source_page = request.headers.get("referer", "")[:500]
    user_agent  = request.headers.get("user-agent", "")[:500]

    # Upload do anexo (se houver) pro Supabase Storage
    attachment_url = None
    attachment_size_kb = None
    if upload_file and upload_filename:
        ext = ""
        if "." in upload_filename:
            ext = "." + upload_filename.rsplit(".", 1)[-1][:10]
        # 🔒 28/07/2026: só aceita os tipos que a gente realmente espera num
        # ticket de suporte. Antes, qualquer extensão passava e caía em
        # application/octet-stream — inclusive .html e .svg, que executam script
        # se abertos direto do Storage.
        # 🪤 Esta checagem fica FORA do try de baixo de propósito: lá os erros são
        # engolidos ("segue sem anexo"), e recusa de arquivo o usuário precisa ver.
        _ext_ok = {
            ".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp",
            ".xlsx", ".xls", ".csv", ".doc", ".docx", ".txt",
            ".dxf", ".dwg", ".zip",
        }
        if ext.lower() not in _ext_ok:
            raise HTTPException(
                400,
                "Tipo de arquivo não aceito. Envie PDF, imagem, planilha, "
                "documento de texto, DXF, DWG ou ZIP."
            )
        try:
            import uuid as _uuid
            object_key = f"contact/{datetime.utcnow().strftime('%Y%m')}/{_uuid.uuid4()}{ext}"

            up_url = f"{SUPABASE_URL}/storage/v1/object/contact-attachments/{object_key}"
            req = urllib.request.Request(up_url, data=upload_file, method="POST")
            req.add_header("apikey", SUPABASE_KEY)
            req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            # Detecta content type básico
            ct_map = {
                ".pdf": "application/pdf", ".png": "image/png", ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg", ".gif": "image/gif", ".webp": "image/webp",
                ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ".xls": "application/vnd.ms-excel", ".csv": "text/csv",
                ".doc": "application/msword",
                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".txt": "text/plain",
            }
            req.add_header("Content-Type", ct_map.get(ext.lower(), "application/octet-stream"))
            req.add_header("x-upsert", "true")
            urllib.request.urlopen(req, timeout=30)

            # 🔒 28/07/2026: o bucket era PÚBLICO e a URL abaixo era
            # /object/public/... — qualquer um com o endereço abria o anexo que
            # um cliente mandou no formulário de contato. O bucket virou privado
            # e agora geramos um link ASSINADO, com validade. Se a assinatura
            # falhar, preferimos ficar SEM link a devolver um link público.
            attachment_size_kb = round(len(upload_file) / 1024)
            try:
                import json as _json_sign  # 🪤 não existe `json` global neste módulo
                sign_url = f"{SUPABASE_URL}/storage/v1/object/sign/contact-attachments/{object_key}"
                sign_req = urllib.request.Request(
                    sign_url,
                    data=_json_sign.dumps({"expiresIn": 60 * 60 * 24 * 365}).encode("utf-8"),
                    method="POST",
                )
                sign_req.add_header("apikey", SUPABASE_KEY)
                sign_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                sign_req.add_header("Content-Type", "application/json")
                signed = _json_sign.loads(urllib.request.urlopen(sign_req, timeout=20).read().decode("utf-8"))
                signed_path = signed.get("signedURL") or signed.get("signedUrl") or ""
                if signed_path:
                    attachment_url = f"{SUPABASE_URL}/storage/v1{signed_path}"
                    print(f"[contact] anexo enviado (link assinado), {attachment_size_kb} KB")
                else:
                    print("[contact] anexo enviado, mas assinatura veio vazia — sem link")
            except Exception as _se:
                print(f"[contact] anexo enviado, mas falhou ao assinar link: {_se}")
        except Exception as e:
            print(f"[contact] erro upload Storage: {e}")
            # Não falha a request por causa do anexo — segue sem ele

    payload = {
        "name": name,
        "email": email,
        "phone": phone or None,
        "message_type": msg_type,
        "subject": subject or None,
        "message": message,
        "source_page": source_page or None,
        "user_agent": user_agent or None,
        "status": "new",
        "attachment_url": attachment_url,
        "attachment_filename": upload_filename if attachment_url else None,
        "attachment_size_kb": attachment_size_kb,
    }

    ok = _supabase_insert("contact_messages", payload)
    if not ok:
        return {"ok": False, "error": "Falha ao salvar mensagem"}

    print(f"[contact] nova mensagem: {email} ({msg_type}) — {subject[:50]}")
    return {"ok": True, "message": "Mensagem recebida! Vamos responder em breve."}


@app.get("/api/admin/contact-messages")
async def admin_list_contact_messages(request: Request, limit: int = 200, status: str = ""):
    """Lista mensagens de contato pro admin. Requer auth admin."""
    _require_admin(request)
    try:
        url = f"{SUPABASE_URL}/rest/v1/contact_messages?select=*&order=created_at.desc&limit={int(limit)}"
        if status:
            url += f"&status=eq.{status}"
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = urllib.request.urlopen(req, timeout=15)
        return {"messages": json.loads(resp.read().decode("utf-8"))}
    except Exception as e:
        return {"error": str(e)}

PUBLIC_CHAT_SYSTEM_PROMPT = """Você é o assistente virtual do AI.arq — uma plataforma brasileira que acelera o levantamento de quantitativos em projetos de arquitetura.

**Sua missão:** responder perguntas de visitantes do site (ai.arq.br) com honestidade e brevidade. Você ajuda a tirar dúvidas sobre o produto, preços e LGPD. Foco em conversão mas SEM vender gordura — o AI.arq é honesto sobre o que faz e o que não faz.

## O QUE O AI.ARQ FAZ

1. **Quantitativo a partir de CAD** — usuário envia PDF/DWG/DXF, IA lê e em ~5 min gera planilha Excel com 18 disciplinas (demolição, alvenaria, elétrica, hidráulica, pintura, pisos, forro, etc). Cada item cita a prancha de origem e tem código de cor:
   - BRANCO = medido direto do CAD (confiável)
   - LARANJA = estimado pela IA (revisar antes de usar)
   - CINZA = metadado do projeto
   - ROXO = custo indireto / gestão (checklist — não sai no CAD mas toda obra tem)

2. **Memória técnica SINAPI + TCPO BIM** — aba extra na planilha com referência oficial pra cada item (composição + insumos + coeficientes). Base carregada: 10.284 composições SINAPI + 54.529 insumos, e 1.333 composições TCPO BIM + 6.733 insumos.

3. **Comparativo de fornecedores** (opcional) — depois que o usuário recebe propostas dos fornecedores, pode fazer upload das planilhas deles e o AI.arq gera XLSX + PPT executivo com ranking, discrepâncias, itens esquecidos.

4. **PPT com a marca do escritório** — logo e cor do escritório aparecem na apresentação. Envio direto pro cliente final via WhatsApp.

## O QUE O AI.ARQ **NÃO** FAZ

- **NÃO precifica** e **NÃO fecha orçamento** — a planilha sai SEM preços preenchidos. Precificação é trabalho do orçamentista do usuário (humano, não IA).
- **NÃO substitui profissional habilitado** — é ferramenta de apoio ao levantamento inicial. Revisão por engenheiro/orçamentista é obrigatória.
- **NÃO adivinha o que não está no CAD** — itens estimados são sinalizados em laranja pra revisão.

## PREÇOS

- **1º projeto: GRÁTIS**
- **Pequeno** (1-5 pranchas): R$ 97
- **Médio** (6-10 pranchas): R$ 157 (mais comum — -19% por prancha vs Pequeno)
- **Grande** (11-20 pranchas): R$ 247 (-36% por prancha)
- **Acima de 20**: R$ 247 + R$ 10/prancha extra
- Pagamento via Stripe (cartão/PIX). **Sem mensalidade.** Paga só quando usa.
- Todas as outras features (memória técnica, comparativo, PPT) estão INCLUÍDAS no preço do projeto. Sem taxas extras. Cada feature é OPCIONAL — use só se fizer sentido. Se perguntarem de cashback: no beta está tudo grátis, não há o que abater; créditos ficam guardados pra quando a cobrança ativar.

## CASHBACK

Programa granular — cada ação que ajuda a calibrar a IA gera crédito automático abatível no próximo projeto:
- Revisão inline (validar item no navegador): crédito por ação
- Upload de planilha revisada offline: +R$ 20
- Upload de cotação de fornecedor: +R$ 5
- Feedback NPS: crédito ao responder

## LGPD E PRIVACIDADE

- Usuário é controlador LGPD dos dados do cliente final (nome, telefone, endereço). AI.arq é operador.
- Valores em R$ de cotações **NUNCA** viram referência pra outros projetos. Só métricas adimensionais anônimas (coeficientes de variação, share MAT/MO) alimentam a base de mercado.
- Dados isolados por projeto. Arquivos processados com criptografia. Transferências internacionais (Anthropic/Supabase/Render/Stripe) com bases legais do art. 33 da LGPD.
- Detalhes em ai.arq.br/privacidade.html e ai.arq.br/termos.html

## COMO USAR

1. Crie conta (grátis) em ai.arq.br → login
2. No dashboard: clique "Novo projeto" → faça upload dos PDFs/DWGs
3. Aguarde ~5 min o processamento
4. Baixe a planilha ou revise no navegador (suas correções afinam o motor)
5. Mande pros fornecedores, receba as propostas
6. (Opcional) Faça upload das propostas no projeto → gere comparativo + PPT

## TOM DE VOZ

- Responda em **português brasileiro, informal-profissional**. Sem jargão técnico desnecessário. Sem "caro cliente".
- Seja **objetivo e curto** (2-4 frases geralmente é o suficiente).
- Se a dúvida for específica e técnica (ex: "como o piso é calculado?"), explique resumido e aponte pro FAQ completo: ai.arq.br/faq.html
- Se o usuário quiser testar: direcione pra ai.arq.br → "Comece Grátis"
- Se não souber a resposta com certeza, DIGA QUE NÃO SABE e sugira contato: contato@ai.arq.br

## NÃO INVENTE

Se perguntarem sobre coisas que você não tem certeza (integrações específicas, recursos futuros, planos corporativos), responda honestamente que não tem essa informação e sugira email de contato. NUNCA invente funcionalidades.
"""


@app.post("/api/public/chat")
async def public_chat(request: Request):
    """Chat público do widget no site — sem auth, com rate limit por IP.

    Usado pra visitantes tirarem dúvidas sobre o produto antes de cadastrar.
    """
    # Rate limit: 20 msgs / 10min por IP
    # FIX 2026-05-14: parênteses pra precedência correta de `or ... if ... else`
    # (antes parseava como `(A or B) if C else D` e podia explodir quando request.client é None).
    # Segurança (16/07): usa o ÚLTIMO hop do X-Forwarded-For. No Render, o proxy
    # confiável ANEXA o IP real do cliente no fim; os valores da esquerda podem ser
    # FORJADOS pelo cliente. Antes [0] deixava burlar o rate-limit trocando o header
    # (abuso de custo da IA na chat pública).
    _xff = [p.strip() for p in request.headers.get("x-forwarded-for", "").split(",") if p.strip()]
    fwd = _xff[-1] if _xff else ""
    client_ip = fwd or (request.client.host if request.client else "unknown")
    now_ts = datetime.utcnow().timestamp()
    history = _PUBLIC_CHAT_HITS.get(client_ip, [])
    history = [t for t in history if now_ts - t < 600]  # últimos 10min
    if len(history) >= 20:
        return {
            "error": "rate_limit",
            "message": "Muitas mensagens em pouco tempo. Espere uns minutos e tente de novo.",
        }
    history.append(now_ts)
    _PUBLIC_CHAT_HITS[client_ip] = history

    # FIX 2026-05-14: anti memory-leak. Periodicamente limpa IPs cuja última
    # hit foi há mais de 10min. Sem isso, o dict cresce pra sempre (cada IP
    # novo = entrada nova nunca limpa) e esgota RAM do Render free tier.
    # Trigger barato: 1× a cada ~50 requests.
    if len(_PUBLIC_CHAT_HITS) > 50 and (int(now_ts) % 50) == 0:
        stale = [ip for ip, hits in _PUBLIC_CHAT_HITS.items()
                 if not hits or (now_ts - max(hits)) > 600]
        for ip in stale:
            _PUBLIC_CHAT_HITS.pop(ip, None)

    try:
        body = await request.json()
    except Exception:
        body = {}

    messages = body.get("messages", [])
    if not isinstance(messages, list) or not messages:
        return {"error": "Precisa mandar { messages: [{role, content}, ...] }"}

    # Sanitização: limita tamanho e quantidade de turns
    messages = messages[-10:]  # só os últimos 10 turns (proteção contra prompt inflation)
    clean_msgs = []
    for m in messages:
        role = m.get("role")
        content = str(m.get("content", ""))[:2000]  # max 2k chars por msg
        if role in ("user", "assistant") and content.strip():
            clean_msgs.append({"role": role, "content": content})

    if not clean_msgs:
        return {"error": "Nenhuma mensagem válida"}

    # Chama Claude Haiku (mais barato e rápido pra chat de vendas)
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {"error": "Chat indisponível no momento."}

    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key)
        resp = client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=600,
            system=PUBLIC_CHAT_SYSTEM_PROMPT,
            messages=clean_msgs,
        )
        reply = resp.content[0].text if resp.content else "Desculpe, não consegui responder agora."
        return {
            "status": "ok",
            "reply": reply,
            "tokens_used": (resp.usage.input_tokens if resp.usage else 0) + (resp.usage.output_tokens if resp.usage else 0),
        }
    except Exception as e:
        print(f"[public_chat] erro: {e}")
        return {
            "error": "Erro temporário",
            "message": "Desculpe, algo deu errado. Tente de novo em instantes ou mande e-mail pra contato@ai.arq.br",
        }


# ═══════════════════════════════════════════════════════════════
#  CHAT DO PROJETO — "Pergunte sobre o seu quantitativo"
#  Chat ancorado nos itens DESTE projeto. Guarda-corpo DURO contra
#  preço (regra nº5: AI.arq entrega quantitativo, NÃO precifica).
# ═══════════════════════════════════════════════════════════════
_PROJECT_CHAT_HITS: dict = {}

PROJECT_CHAT_SYSTEM = """Você é o assistente do AI.arq. Ajuda o cliente a ENTENDER o quantitativo dele — a planilha que a nossa IA gerou lendo o projeto CAD. Responda em português do Brasil, de forma clara, curta e cordial.

REGRAS DURAS (nunca violar):
1. Fale SÓ sobre ESTE projeto e os itens listados abaixo. Nunca invente um item que não está na lista.
2. NUNCA dê preço, custo, valor em R$, ou BDI. O AI.arq entrega QUANTIDADE, não preço. Se perguntarem de preço/custo/orçamento, responda gentil: "O AI.arq gera o quantitativo (as quantidades) — a precificação é com você e seu orçamentista. Mas posso te ajudar a entender as quantidades e as referências SINAPI."
3. Você NÃO substitui o profissional nem dá parecer técnico definitivo — você ajuda a LER e entender a planilha.
4. Explique bem a diferença: um item MEDIDO foi extraído direto da geometria do CAD (confiável); uma ESTIMATIVA é quando o desenho não deixou claro e o cliente precisa revisar. Para ter MAIS itens medidos, oriente enviar o projeto em DWG ou DXF (PDF a IA lê, mas vira estimativa).
5. Se não souber, ou o dado não estiver na planilha, seja honesto e diga que não consta.
6. FORMATO: texto corrido e curto, com listas de traços quando ajudar. NÃO use títulos markdown (#, ##) nem tabelas — sua resposta aparece num balão de chat simples.

DADOS DO PROJETO:
{context}
"""


@app.post("/api/projeto/{job_id}/chat")
async def project_chat(job_id: str, request: Request):
    """Chat ancorado no quantitativo DESTE projeto. Auth: dono. Contexto: os
    itens da planilha. Modelo Haiku (barato). Não precifica (regra nº5)."""
    _require_project_owner(request, job_id)

    import urllib.request as _u, json as _j
    now_ts = datetime.utcnow().timestamp()
    hits = [t for t in _PROJECT_CHAT_HITS.get(job_id, []) if now_ts - t < 600]
    if len(hits) >= 30:  # rate limit leve: 30 msgs / 10min por projeto
        return {"error": "rate_limit", "reply": "Você mandou bastante pergunta em pouco tempo — dá um respiro de uns minutinhos. 🙂"}
    hits.append(now_ts)
    _PROJECT_CHAT_HITS[job_id] = hits
    if len(_PROJECT_CHAT_HITS) > 200 and int(now_ts) % 50 == 0:
        for _k in [k for k, v in _PROJECT_CHAT_HITS.items() if not v or now_ts - max(v) > 600]:
            _PROJECT_CHAT_HITS.pop(_k, None)

    try:
        body = await request.json()
    except Exception:
        body = {}
    messages = body.get("messages", [])
    if not isinstance(messages, list) or not messages:
        return {"error": "Precisa mandar messages"}
    clean = []
    for m in messages[-8:]:
        role = m.get("role")
        content = str(m.get("content", ""))[:1500]
        if role in ("user", "assistant") and content.strip():
            clean.append({"role": role, "content": content})
    if not clean:
        return {"error": "Nenhuma mensagem válida"}

    # Contexto: metadados + itens da planilha (service role; ownership já validado)
    ctx_lines = []
    try:
        q = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
             f"&select=project_name,total_area,typology,warnings")
        rq = _u.Request(q, method="GET")
        rq.add_header("apikey", SUPABASE_KEY)
        rq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        prow = _j.loads(_u.urlopen(rq, timeout=10).read().decode("utf-8"))
        if prow:
            if prow[0].get("total_area"):
                ctx_lines.append(f"Área total do projeto: {prow[0]['total_area']} m²")
            if prow[0].get("typology"):
                ctx_lines.append(f"Tipologia: {prow[0]['typology']}")
    except Exception as e:
        print(f"[project_chat] meta erro: {e}")

    items = []
    try:
        q = (f"{SUPABASE_URL}/rest/v1/project_items?job_id=eq.{job_id}"
             f"&select=description,unit,quantity,confidence,discipline&limit=400")
        rq = _u.Request(q, method="GET")
        rq.add_header("apikey", SUPABASE_KEY)
        rq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        items = _j.loads(_u.urlopen(rq, timeout=15).read().decode("utf-8"))
    except Exception as e:
        print(f"[project_chat] itens erro: {e}")

    by_disc: dict = {}
    for it in items:
        by_disc.setdefault(it.get("discipline") or "Outros", []).append(it)
    ctx_lines.append(f"Total de itens na planilha: {len(items)}")
    for _d, _lst in by_disc.items():
        ctx_lines.append(f"\n[{_d}] ({len(_lst)} itens)")
        for it in _lst[:40]:
            _marca = "MEDIDO" if (it.get("confidence") == "confirmado") else "estimativa"
            ctx_lines.append(f"  - {str(it.get('description',''))[:80]}: {it.get('quantity',0)} {it.get('unit','')} ({_marca})")
    context = "\n".join(ctx_lines)[:9000]

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {"error": "Chat indisponível no momento."}
    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key, timeout=60.0)
        resp = client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=500,
            system=PROJECT_CHAT_SYSTEM.format(context=context),
            messages=clean,
        )
        reply = resp.content[0].text if resp.content else "Não consegui responder agora."
        return {"status": "ok", "reply": reply}
    except Exception as e:
        print(f"[project_chat] erro: {e}")
        return {"error": "Erro temporário", "reply": "Ops, algo deu errado. Tenta de novo em instantes."}


# ═══════════════════════════════════════════════════════════════
#  COTAÇÕES DE FORNECEDORES (supplier quotes)
# ═══════════════════════════════════════════════════════════════

@app.post("/api/projects/{job_id}/quotes/upload")
async def upload_supplier_quote(
    request: Request,
    job_id: str,
    supplier_name: str = Form(...),
    file: UploadFile = File(...),
    parser_mode: str = Form("auto"),
    user_id: str = Form(""),
):
    """Recebe uma planilha de orçamento de fornecedor, parseia e salva no banco.

    Quem chama tem que ser o dono do projeto (ou admin) — senão atacante
    sobe cotação no projeto alheio + ganha R$5 de cashback no nome dele.
    """
    # Valida que JWT == dono do projeto (ou admin)
    _require_project_owner(request, job_id)
    # Se passou user_id no Form, ele tem que casar com o JWT (cashback fica
    # registrado no nome certo, não no que o atacante quiser)
    if user_id:
        jwt_user = _get_user_from_request(request)
        if jwt_user and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
            if jwt_user.get("id") != user_id:
                raise HTTPException(403, "user_id não corresponde ao token")

    try:
        from supplier_quote_parser import parse_supplier_quote
    except ImportError:
        return {"error": "supplier_quote_parser indisponível"}

    # Só .xlsx e com teto de tamanho (o parse carrega o arquivo inteiro em
    # memória — Render tem 2GB pra tudo)
    if not (file.filename or "").lower().endswith(".xlsx"):
        return {"status": "error",
                "error": "Envie a cotação em .xlsx (Excel). Se estiver em .xls "
                         "antigo ou PDF, abra no Excel e salve como .xlsx."}
    _quote_bytes = await file.read()
    if len(_quote_bytes) > 15 * 1024 * 1024:
        return {"status": "error",
                "error": "Arquivo acima de 15MB — exporte só a planilha de "
                         "orçamento, sem imagens."}

    # Salva temporariamente (anti path-traversal)
    work_dir = os.path.join(WORK_DIR, job_id, "quotes")
    os.makedirs(work_dir, exist_ok=True)
    safe_name = _safe_local_filename(file.filename)
    temp_path = os.path.join(work_dir, safe_name)
    with open(temp_path, "wb") as f:
        f.write(_quote_bytes)

    # Parseia (arquivo corrompido/renomeado lançava BadZipFile → HTTP 500 cru)
    try:
        parsed = parse_supplier_quote(temp_path, supplier_name, mode=parser_mode)
    except Exception as e:
        print(f"[quotes/upload] parse falhou ({file.filename}): {e}")
        return {"status": "error",
                "error": "Não consegui abrir esse arquivo como planilha Excel. "
                         "Confira se é um .xlsx válido e tente de novo."}
    if "error" in parsed:
        return {"status": "error", "error": parsed["error"]}

    # Grava no Supabase
    payload = {
        "job_id": job_id,
        "supplier_name": supplier_name,
        "original_filename": file.filename,
        "storage_path": temp_path,
        "parser_mode": parsed.get("parser_mode_used", "strict"),
        "n_items_quoted": parsed["n_items_quoted"],
        "total_bruto": parsed["total_bruto"],
        "total_material": parsed["total_material"],
        "total_mao_obra": parsed["total_mao_obra"],
        "items": parsed["items"],
        "status": "parsed",
        "uploaded_by": user_id or "",
    }

    try:
        st, inserted = _supa_rest_as_user(
            request, "POST",
            "/project_supplier_quotes",
            body=payload, prefer="return=representation", timeout=20,
        )
        if st >= 400 or inserted is None:
            return {"status": "error", "error": f"Erro ao gravar (status {st})"}
        quote_id = inserted[0]["id"] if inserted else None
    except Exception as e:
        return {"status": "error", "error": f"Erro ao gravar: {e}"}

    # Cria evento de cashback: R$ 10 por cotação, MAX 3 cotações por projeto.
    # Política nova (2026-05-13): subiu de R$ 5 → R$ 10 mas com cap de 3 — incentiva
    # diversificação de cotações sem virar farm.
    cashback_status = None
    if user_id:
        # Conta quantos uploads já tiveram cashback nesse projeto
        try:
            _, cnt_rows = _supa_rest_as_user(
                request, "GET",
                f"/project_cashback_events?job_id=eq.{job_id}"
                f"&event_type=eq.supplier_quote_upload&select=id",
                timeout=8,
            )
            existing_count = len(cnt_rows or [])
        except Exception:
            existing_count = 0

        if existing_count >= 3:
            cashback_status = "capped"
        else:
            try:
                cb_payload = {
                    "job_id": job_id,
                    "user_id": user_id,
                    "event_type": "supplier_quote_upload",
                    "credit_cents": 1000,  # R$ 10 por cotação (max 3 = R$ 30)
                    "description": f"Upload cotação {supplier_name}",
                    "ref_id": quote_id,
                    "approved_at": datetime.utcnow().isoformat(),
                    "approved_by": "auto",
                }
                # service_role: a escrita de cashback é do sistema (não do user).
                # Com JWT do user a RLS negava (não há policy de INSERT authenticated)
                # — o crédito vinha falhando. (auditoria 2026-07-13)
                _cb_ok = _supabase_insert("project_cashback_events", cb_payload)
                cashback_status = "credited" if _cb_ok else "error"
            except Exception:
                cashback_status = "error"  # best-effort, não bloqueia upload

    return {
        "status": "ok",
        "quote_id": quote_id,
        "supplier_name": supplier_name,
        "n_items_quoted": parsed["n_items_quoted"],
        "total_bruto": parsed["total_bruto"],
        "parser_mode_used": parsed.get("parser_mode_used"),
        "cashback": cashback_status,  # 'credited' | 'capped' | 'error' | None
    }


@app.get("/api/projects/{job_id}/quotes")
async def list_supplier_quotes(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Lista cotações de fornecedores de um projeto."""
    try:
        _, rows = _supa_rest_as_user(
            request, "GET",
            f"/project_supplier_quotes?job_id=eq.{job_id}"
            f"&select=id,supplier_name,original_filename,parser_mode,"
            f"n_items_quoted,total_bruto,total_material,total_mao_obra,"
            f"status,uploaded_at"
            f"&order=uploaded_at.asc",
            timeout=8,
        )
        return {"quotes": rows or []}
    except Exception as e:
        return {"error": str(e)}


@app.delete("/api/projects/{job_id}/quotes/{quote_id}")
async def delete_supplier_quote(job_id: str, quote_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Remove uma cotação."""
    try:
        st, _ = _supa_rest_as_user(
            request, "DELETE",
            f"/project_supplier_quotes?id=eq.{quote_id}&job_id=eq.{job_id}",
            timeout=8,
        )
        if st >= 400:
            return {"error": f"Falha ao remover (status {st})"}
        return {"status": "ok"}
    except Exception as e:
        return {"error": str(e)}


@app.get("/api/projects/{job_id}/quotes/compare")
async def compare_supplier_quotes(job_id: str, request: Request, include_reference: int = 1):
    _require_project_owner(request, job_id)
    """Compara todas as cotações de um projeto, gera XLSX+PPT, retorna URLs."""
    try:
        from supplier_quote_compare import (compare_quotes,
                                             generate_comparative_xlsx,
                                             generate_comparative_pptx)
    except ImportError:
        return {"error": "supplier_quote_compare indisponível"}

    # Busca quotes
    try:
        _, quotes_raw = _supa_rest_as_user(
            request, "GET",
            f"/project_supplier_quotes?job_id=eq.{job_id}"
            f"&select=*&order=uploaded_at.asc",
            timeout=10,
        )
        quotes_raw = quotes_raw or []
    except Exception as e:
        return {"error": f"erro buscando quotes: {e}"}

    if len(quotes_raw) < 2:
        return {"error": "precisa de pelo menos 2 cotações pra comparar"}

    # Converte pra formato esperado pelo compare
    quotes = []
    for q in quotes_raw:
        quotes.append({
            "supplier_name": q["supplier_name"],
            "n_items_quoted": q.get("n_items_quoted", 0),
            "total_bruto": float(q.get("total_bruto") or 0),
            "total_material": float(q.get("total_material") or 0),
            "total_mao_obra": float(q.get("total_mao_obra") or 0),
            "items": q.get("items") or [],
        })

    # Busca reference items do quantitativo original
    reference_items = None
    if include_reference:
        try:
            _, reference_items = _supa_rest_as_user(
                request, "GET",
                f"/project_items?job_id=eq.{job_id}"
                f"&select=description,unit,quantity&limit=500",
                timeout=10,
            )
        except Exception:
            pass

    analysis = compare_quotes(quotes, reference_items=reference_items)

    # Busca project_clients + project data + logo do escritório pro cabeçalho
    project_name = ""
    architect_name = ""
    client_name = ""
    logo_url = ""
    project_user_id = ""
    try:
        _, pj_data = _supa_rest_as_user(
            request, "GET",
            f"/projects?job_id=eq.{job_id}&select=project_name,user_name,user_id",
            timeout=8,
        )
        pj_data = pj_data or []
        if pj_data:
            project_name = pj_data[0].get("project_name", "")
            architect_name = pj_data[0].get("user_name", "")
            project_user_id = pj_data[0].get("user_id", "")

        _, cl_data = _supa_rest_as_user(
            request, "GET",
            f"/project_clients?job_id=eq.{job_id}&select=client_name,client_company",
            timeout=8,
        )
        cl_data = cl_data or []
        if cl_data:
            client_name = cl_data[0].get("client_name", "")
            if cl_data[0].get("client_company"):
                client_name = f"{client_name} ({cl_data[0]['client_company']})" \
                    if client_name else cl_data[0]["client_company"]

        # Logo e cor de marca do escritório (perfil do dono do projeto)
        if project_user_id:
            _, pr_data = _supa_rest_as_user(
                request, "GET",
                f"/profiles?user_id=eq.{project_user_id}"
                f"&select=logo_url,company,company_brand_color",
                timeout=8,
            )
            pr_data = pr_data or []
            if pr_data:
                logo_url = pr_data[0].get("logo_url", "") or ""
                brand_color = pr_data[0].get("company_brand_color", "") or ""
                if not architect_name and pr_data[0].get("company"):
                    architect_name = pr_data[0]["company"]
            else:
                brand_color = ""
        else:
            brand_color = ""
    except Exception:
        brand_color = ""

    # Download logo temporariamente pra embutir no PPT
    logo_path = None
    # Guarda anti-SSRF (fix 2026-07-22): só baixa logo de host público.
    if logo_url and _url_is_safe_public(logo_url):
        try:
            work_dir_lg = os.path.join(WORK_DIR, job_id)
            os.makedirs(work_dir_lg, exist_ok=True)
            logo_path = os.path.join(work_dir_lg, "logo_escritorio.png")
            lg_req = urllib.request.Request(logo_url, method="GET")
            with urllib.request.urlopen(lg_req, timeout=10) as resp:
                with open(logo_path, "wb") as f:
                    f.write(resp.read(20 * 1024 * 1024))  # cap 20MB
        except Exception:
            logo_path = None

    # Gera XLSX
    work_dir = os.path.join(WORK_DIR, job_id)
    os.makedirs(work_dir, exist_ok=True)
    xlsx_path = os.path.join(work_dir, f"comparativo_{job_id}.xlsx")
    generate_comparative_xlsx(
        analysis, xlsx_path,
        project_name=project_name,
        architect_name=architect_name,
        client_name=client_name,
    )

    # Gera PPT
    pptx_path = os.path.join(work_dir, f"comparativo_{job_id}.pptx")
    try:
        generate_comparative_pptx(
            analysis, pptx_path,
            project_name=project_name,
            architect_name=architect_name,
            client_name=client_name,
            logo_path=logo_path,
            brand_color_hex=brand_color,
        )
        pptx_ok = True
    except Exception as e:
        print(f"[quotes/compare] erro PPT: {e}")
        pptx_ok = False

    # Persiste no Storage (19/07): o disco do Render é EFÊMERO — o comparativo
    # sumia a cada deploy/restart e o download virava 404 "gere de novo".
    # Best-effort: sem Storage, o fluxo antigo (disco) continua valendo.
    try:
        _supabase_storage_upload(xlsx_path, f"quotes/{job_id}/comparativo.xlsx")
        if pptx_ok:
            _supabase_storage_upload(pptx_path, f"quotes/{job_id}/comparativo.pptx")
    except Exception as _se:
        print(f"[quotes/compare] upload Storage falhou (segue com disco): {_se}")

    # ═══ ALIMENTA MOTOR ANONIMAMENTE ═══
    # Extrai heurísticas do comparativo (sem dados do projeto) e insere
    # em market_heuristics. Loop de aprendizado do motor.
    try:
        _extract_and_save_heuristics(analysis, "office")
    except Exception as e:
        print(f"[quotes/compare] erro extração heurísticas: {e}")

    return {
        "status": "ok",
        "xlsx_url": f"/api/projects/{job_id}/quotes/download/xlsx",
        "pptx_url": f"/api/projects/{job_id}/quotes/download/pptx" if pptx_ok else None,
        "summary": {
            "n_suppliers": len(analysis["suppliers"]),
            "paired_count": analysis.get("paired_count", 0),
            "ranking_confiavel": analysis.get("ranking_confiavel", False),
            "ranking": [{"supplier": s, "paired_total": t}
                        for s, t in analysis["ranking"]],
            "n_discrepancies_above_100pct":
                len([d for d in analysis["biggest_discrepancies"]
                     if d["pct_diff"] > 100]),
            "reference_check_summary":
                analysis.get("reference_check", {}).get("summary")
                if analysis.get("reference_check") else None,
        },
    }


def _quotes_download_path(job_id: str, ext: str) -> Optional[str]:
    """Caminho local do comparativo; se o disco efêmero perdeu (deploy/restart),
    re-baixa do Storage (quotes/{job}/comparativo.{ext}, persistido no compare).
    None = não existe em lugar nenhum → caller devolve 404."""
    path = os.path.join(WORK_DIR, job_id, f"comparativo_{job_id}.{ext}")
    if os.path.exists(path):
        return path
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        if _supabase_storage_download(f"quotes/{job_id}/comparativo.{ext}", path) \
                and os.path.exists(path) and os.path.getsize(path) > 0:
            print(f"[quotes/download] {job_id}.{ext}: resgatado do Storage")
            return path
    except Exception as e:
        print(f"[quotes/download] Storage {job_id}.{ext}: {e}")
    try:
        if os.path.exists(path):
            os.remove(path)  # não deixar arquivo vazio/corrompido pra trás
    except Exception:
        pass
    return None


@app.get("/api/projects/{job_id}/quotes/download/xlsx")
async def download_quotes_xlsx(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Baixa o comparativo XLSX gerado (disco → fallback Storage)."""
    path = _quotes_download_path(job_id, "xlsx")
    if not path:
        # 404, não 200 com JSON: o downloadProtected salva o corpo da resposta
        # como arquivo — 200 com {"error":...} virava .xlsx corrompido no Excel.
        raise HTTPException(404, "Comparativo não está mais disponível — clique em "
                                 "'Gerar comparativo' novamente.")
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=f"comparativo_fornecedores_{job_id}.xlsx",
    )


@app.get("/api/projects/{job_id}/quotes/download/pptx")
async def download_quotes_pptx(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Baixa o comparativo PPT gerado (disco → fallback Storage)."""
    path = _quotes_download_path(job_id, "pptx")
    if not path:
        # Mesmo motivo do XLSX acima: 200 com JSON virava .pptx corrompido.
        raise HTTPException(404, "Apresentação não está mais disponível — clique em "
                                 "'Gerar comparativo' novamente.")
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"comparativo_fornecedores_{job_id}.pptx",
    )


def _extract_and_save_heuristics(analysis: dict, typology: str = "office"):
    """Extrai heurísticas adimensionais do comparativo e salva em market_heuristics.

    REGRA DURA: zero valores absolutos, zero nomes de projeto/fornecedor.
    Fornecedores são anonimizados como fornecedor_1/2/3 (ordem aleatória).

    Inclui:
      - dispersion: CV por categoria dos top itens divergentes
      - coverage_pattern: contagem por categoria (anonimizada)
      - mat_mo_share: ratio médio MAT/MO por categoria
    """
    from collections import defaultdict
    from statistics import mean, stdev
    import unicodedata
    import re as _re
    import hashlib
    import random

    def _norm(s):
        if not s: return ""
        s = str(s).strip().lower()
        s = unicodedata.normalize("NFKD", s)
        s = "".join(c for c in s if not unicodedata.combining(c))
        s = _re.sub(r"[^\w\s]", " ", s)
        return _re.sub(r"\s+", " ", s).strip()

    categorias_kw = {
        "demolicao": ["demolicao", "remocao", "retirada", "demolir"],
        "drywall": ["drywall", "gesso", "septo", "sept"],
        "eletrica": ["eletr", "tomada", "interruptor", "luminaria"],
        "hidraulica": ["hidr", "tubulacao", "bacia", "torneira"],
        "piso": ["piso", "carpete", "laminado", "vinil"],
        "pintura": ["pintura", "tinta", "selador", "verniz"],
        "forro": ["forro"],
        "esquadria": ["porta", "janela", "esquadria"],
        "ar_condicionado": ["ar condicionado", "hvac", "condicionado"],
        "mobiliario": ["mobiliar", "armario", "marcenar"],
        "preliminares": ["mobilizacao", "canteiro", "tapume", "protecao"],
    }

    def _cat(desc):
        n = _norm(desc)
        for cat, keys in categorias_kw.items():
            for k in keys:
                if k in n:
                    return cat
        return "outros"

    # Source hash anônimo pra não repetir dados
    ranking_key = "_".join(sorted(analysis["paired_totals"].keys()))
    source = "auto_" + hashlib.md5(ranking_key.encode()).hexdigest()[:10]

    rows = []
    n_sup = len(analysis["suppliers"])

    # 1. Dispersão por categoria
    cat_disp = defaultdict(list)
    for d in analysis["biggest_discrepancies"][:30]:
        cat = _cat(d["desc"])
        cat_disp[cat].append(d["pct_diff"] / 100)  # normaliza pra 0-1

    for cat, cvs in cat_disp.items():
        if len(cvs) < 1:
            continue
        avg_cv = mean(cvs)
        if avg_cv < 0.2:
            continue
        rows.append({
            "heuristic_type": "dispersion",
            "typology": typology,
            "category": cat,
            "unit": "",
            "keywords": "",
            "metric_name": "cv_total",
            "metric_value": round(avg_cv, 3),
            "stddev": round(stdev(cvs), 3) if len(cvs) > 1 else 0,
            "n_observations": len(cvs),
            "source_anonimo": source,
        })

    # 2. Cobertura por categoria (n itens por fornecedor anônimo)
    cov_by_cat = defaultdict(lambda: defaultdict(int))
    suppliers_shuffled = list(analysis["suppliers"])
    random.shuffle(suppliers_shuffled)  # anonimiza posição
    sup_alias = {s: i + 1 for i, s in enumerate(suppliers_shuffled)}

    for m in analysis["merged_items"]:
        cat = _cat(m["desc"])
        for sup in analysis["suppliers"]:
            if m.get(sup) and m[sup].get("total"):
                cov_by_cat[cat][sup_alias[sup]] += 1

    for cat, supcov in cov_by_cat.items():
        cobertura = sum(1 for i in range(1, n_sup + 1) if supcov.get(i, 0) > 0)
        rows.append({
            "heuristic_type": "coverage_pattern",
            "typology": typology,
            "category": cat,
            "unit": "",
            "keywords": "",
            "metric_name": "cobertura_completa",
            "metric_value": cobertura,
            "stddev": 0,
            "n_observations": n_sup,
            "source_anonimo": source,
        })

    # 3. Share MAT/MO por categoria
    cat_shares = defaultdict(list)
    for m in analysis["merged_items"]:
        cat = _cat(m["desc"])
        for sup in analysis["suppliers"]:
            it = m.get(sup)
            if not it or not it.get("qtd"):
                continue
            mat = (it.get("unit_mat") or 0) * it["qtd"]
            mo = (it.get("unit_mo") or 0) * it["qtd"]
            soma = mat + mo
            if soma <= 0:
                continue
            cat_shares[cat].append(mat / soma)

    for cat, shares in cat_shares.items():
        if len(shares) < 3:
            continue
        share_mat_avg = mean(shares)
        stddev_mat = stdev(shares) if len(shares) > 1 else 0
        rows.append({
            "heuristic_type": "mat_mo_share",
            "typology": typology,
            "category": cat,
            "unit": "",
            "keywords": "",
            "metric_name": "share_mat",
            "metric_value": round(share_mat_avg, 3),
            "stddev": round(stddev_mat, 3),
            "n_observations": len(shares),
            "source_anonimo": source,
        })
        rows.append({
            "heuristic_type": "mat_mo_share",
            "typology": typology,
            "category": cat,
            "unit": "",
            "keywords": "",
            "metric_name": "share_mo",
            "metric_value": round(1 - share_mat_avg, 3),
            "stddev": 0,
            "n_observations": len(shares),
            "source_anonimo": source,
        })

    if not rows:
        return

    # Audita antes de inserir
    for r in rows:
        v = r.get("metric_value")
        if v and abs(v) > 100 and r["metric_name"] not in (
                "cobertura_completa",) and not r["metric_name"].startswith("n_itens"):
            print(f"[heuristics] ALERTA valor suspeito: {r}")
            return  # não insere se tem valor absoluto vazado

    # Insere
    try:
        url = f"{SUPABASE_URL}/rest/v1/market_heuristics"
        body = json.dumps(rows).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        req.add_header("Prefer", "return=minimal")
        urllib.request.urlopen(req, timeout=15)
        print(f"[heuristics] +{len(rows)} métricas anônimas inseridas")
    except Exception as e:
        print(f"[heuristics] erro insert: {e}")


# ═══════════════════════════════════════════════════════════════
#  UPLOAD DE PLANILHA REVISADA (cashback)
# ═══════════════════════════════════════════════════════════════

@app.post("/api/projects/{job_id}/revised-sheet/upload")
async def upload_revised_sheet(
    request: Request,
    job_id: str,
    file: UploadFile = File(...),
    user_id: str = Form(""),
):
    """Recebe planilha revisada offline. Salva na pasta do job e dá cashback.

    Dono do projeto (ou admin) apenas — senão atacante ganha R$20 de
    cashback no nome de outro user.
    """
    _require_project_owner(request, job_id)
    if user_id:
        jwt_user = _get_user_from_request(request)
        if jwt_user and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
            if jwt_user.get("id") != user_id:
                raise HTTPException(403, "user_id não corresponde ao token")

    # Salva o arquivo (anti path-traversal)
    work_dir = os.path.join(WORK_DIR, job_id, "revised")
    os.makedirs(work_dir, exist_ok=True)
    safe_name = _safe_local_filename(file.filename)
    dst = os.path.join(work_dir, safe_name)
    with open(dst, "wb") as f:
        f.write(await file.read())

    # Evento de cashback: R$ 30 (2026-05-13: subiu de R$ 20 → R$ 30 pra
    # incentivar mais uploads de planilha revisada — feedback alimenta calibração)
    credit_cents = 3000
    if user_id:
        try:
            cb_payload = {
                "job_id": job_id,
                "user_id": user_id,
                "event_type": "planilha_upload",
                "credit_cents": credit_cents,
                "description": f"Upload planilha revisada ({file.filename})",
                "approved_at": datetime.utcnow().isoformat(),
                "approved_by": "auto",
            }
            # service_role (escrita do sistema; RLS negava com JWT do user).
            _supabase_insert("project_cashback_events", cb_payload)
        except Exception as e:
            print(f"[revised-sheet] erro cashback: {e}")

    # Feedback de revisão (best-effort, em thread): compara a planilha revisada
    # com os itens originais (project_items) e grava em revision_feedback —
    # é o que alimenta o painel "Onde a IA mais erra" no admin. Nunca pode
    # atrasar nem derrubar a resposta do upload.
    try:
        import threading
        import revision_feedback as _rf
        threading.Thread(
            target=_rf.processar_revisao, args=(job_id, dst), daemon=True
        ).start()
    except Exception as _rfe:
        print(f"[revision-feedback] não iniciou thread job={job_id}: {_rfe}")

    return {
        "status": "ok",
        "filename": file.filename,
        "credit_cents": credit_cents,
    }


@app.get("/api/admin/revision-feedback")
async def admin_revision_feedback(request: Request):
    """Onde a IA mais erra — agregado das planilhas revisadas pelos clientes.

    Lê a tabela revision_feedback (service_role; RLS sem policy) e devolve o
    resumo pronto pro painel: por disciplina + recorte medido×estimado."""
    _require_admin(request)
    import urllib.request as _url_rf
    import json as _json_rf
    import revision_feedback as _rf
    try:
        url = (f"{SUPABASE_URL}/rest/v1/revision_feedback"
               f"?select=*&order=created_at.desc&limit=200")
        req = _url_rf.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = _url_rf.urlopen(req, timeout=10)
        rows = _json_rf.loads(resp.read().decode("utf-8"))
    except Exception as _e:
        print(f"[revision-feedback] admin erro: {_e}")
        raise HTTPException(502, "Não consegui carregar o feedback de revisões")
    resumo = _rf.resumo_para_admin(rows if isinstance(rows, list) else [])

    # 🪤 01/08/2026: o painel só lia revision_feedback — que depende do cliente
    # REENVIAR um XLSX revisado, gesto que ninguém fez em 3 meses (0 linhas).
    # Enquanto isso a revisão INLINE gravava em item_reviews (24 sinais em 5
    # projetos) e ninguém olhava. Agora as duas fontes aparecem juntas:
    # approve = "o número está certo" (valida a medição); edit = correção real.
    try:
        url2 = (f"{SUPABASE_URL}/rest/v1/item_reviews"
                f"?select=job_id,item_id,action,edits,comment,reviewed_at"
                f"&order=reviewed_at.desc&limit=500")
        req2 = _url_rf.Request(url2, method="GET")
        req2.add_header("apikey", SUPABASE_KEY)
        req2.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows2 = _json_rf.loads(_url_rf.urlopen(req2, timeout=10).read().decode("utf-8"))
        if isinstance(rows2, list):
            aprov = [r for r in rows2 if r.get("action") == "approve"]
            edits = [r for r in rows2 if r.get("action") == "edit"]
            resumo["revisao_inline"] = {
                "aprovacoes": len(aprov),
                "edicoes": len(edits),
                "projetos": len({r.get("job_id") for r in rows2 if r.get("job_id")}),
                "ultimos_edits": [
                    {"job_id": r.get("job_id"),
                     "quando": (r.get("reviewed_at") or "")[:10],
                     "edits": r.get("edits"),
                     "comment": r.get("comment")}
                    for r in edits[:10]
                ],
            }
    except Exception as _e2:
        print(f"[revision-feedback] item_reviews erro: {_e2}")
        resumo["revisao_inline"] = {"erro": "não consegui ler item_reviews"}
    return resumo


# ═══════════════════════════════════════════════════════════════
#  CASHBACK EVENTS POR PROJETO
# ═══════════════════════════════════════════════════════════════

@app.get("/api/user/{user_id}/cashback-all")
async def get_user_cashback_all(request: Request, user_id: str):
    """Retorna cashback agregado de TODOS os projetos do usuário.

    Usa pra mostrar o resumo na tela de cashback (total geral + por projeto).
    Só dono (ou admin) — saldo de cashback é dado financeiro.
    """
    jwt_user = _get_user_from_request(request)
    if not jwt_user:
        raise HTTPException(401, "Autenticação requerida")
    if jwt_user.get("id") != user_id and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
        raise HTTPException(403, "Só é possível consultar seu próprio cashback")
    try:
        # Busca todos eventos do user
        url = (f"{SUPABASE_URL}/rest/v1/project_cashback_events"
               f"?user_id=eq.{user_id}"
               f"&select=*&order=created_at.desc&limit=1000")
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        resp = urllib.request.urlopen(req, timeout=8)
        events = json.loads(resp.read().decode("utf-8"))

        # Busca nome dos projetos envolvidos
        job_ids = list({e["job_id"] for e in events})
        projects_map = {}
        if job_ids:
            ids_filter = ",".join(f'"{j}"' for j in job_ids)
            pj_url = (f"{SUPABASE_URL}/rest/v1/projects"
                      f"?job_id=in.({','.join(job_ids)})&select=job_id,project_name")
            pj_req = urllib.request.Request(pj_url, method="GET")
            pj_req.add_header("apikey", SUPABASE_KEY)
            pj_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            try:
                pj_resp = urllib.request.urlopen(pj_req, timeout=8)
                for p in json.loads(pj_resp.read().decode("utf-8")):
                    projects_map[p["job_id"]] = p.get("project_name", "Projeto sem nome")
            except Exception:
                pass

        # Agrupa por projeto
        from collections import defaultdict
        by_proj = defaultdict(lambda: {
            "job_id": "", "project_name": "", "total_cents": 0,
            "n_events": 0, "event_types": set(),
        })
        for e in events:
            p = by_proj[e["job_id"]]
            p["job_id"] = e["job_id"]
            p["project_name"] = projects_map.get(e["job_id"], "Projeto sem nome")
            p["total_cents"] += e.get("credit_cents", 0)
            p["n_events"] += 1
            p["event_types"].add(e.get("event_type", ""))

        by_project_list = []
        for p in by_proj.values():
            by_project_list.append({
                "job_id": p["job_id"],
                "project_name": p["project_name"],
                "total_cents": p["total_cents"],
                "n_events": p["n_events"],
                "event_types": sorted(p["event_types"]),
            })
        by_project_list.sort(key=lambda x: -x["total_cents"])

        return {
            "total_cents": sum(e.get("credit_cents", 0) for e in events),
            "n_events": len(events),
            "by_project": by_project_list,
        }
    except Exception as e:
        return {"error": str(e)}


@app.get("/api/projects-confidence")
async def projects_confidence_summary(request: Request):
    """Selo de confiança em 'Meus Projetos': medido/total por projeto do usuário,
    numa query só (RPC user_project_confidence). Evita N fetches pesados de itens."""
    user = _get_user_from_request(request)
    if not user:
        raise HTTPException(401, "Autenticação requerida")
    uid = str(user.get("id") or "")
    if not uid:
        raise HTTPException(401, "Sessão inválida")
    import urllib.request as _ur
    try:
        u = f"{SUPABASE_URL}/rest/v1/rpc/user_project_confidence"
        body = _json.dumps({"p_user_id": uid}).encode("utf-8")
        r = _ur.Request(u, data=body, method="POST")
        r.add_header("apikey", SUPABASE_KEY)
        r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        r.add_header("Content-Type", "application/json")
        rows = _json.loads(_ur.urlopen(r, timeout=15).read().decode("utf-8"))
    except Exception as _e:
        print(f"[proj-conf] erro: {_e}")
        return {"summary": {}}
    summary = {}
    for row in (rows or []):
        jid = row.get("job_id")
        if jid:
            summary[jid] = {"total": int(row.get("total") or 0), "medido": int(row.get("medido") or 0)}
    return {"summary": summary}


@app.get("/api/projects/{job_id}/cashback")
async def get_project_cashback(job_id: str, request: Request):
    """Retorna eventos de cashback + total acumulado desse projeto.

    Ownership check adicionado em 2026-05-27: antes o endpoint ficava
    aberto — qualquer um com job_id válido lia o saldo. Não é PII forte,
    mas vazava histórico financeiro do projeto. Bug bônus achado durante
    auditoria RLS/REST (commit b5ddd07).
    """
    _require_project_owner(request, job_id)
    try:
        status, events = _supa_rest_as_user(
            request, "GET",
            f"project_cashback_events?job_id=eq.{job_id}"
            f"&select=*&order=created_at.desc"
        )
        if status >= 400 or events is None:
            return {"error": f"HTTP {status}", "events": [], "total_cents": 0,
                    "total_reais": 0.0, "count": 0}
        total_cents = sum(e.get("credit_cents", 0) for e in events)
        return {
            "events": events,
            "total_cents": total_cents,
            "total_reais": total_cents / 100.0,
            "count": len(events),
        }
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════
#  DADOS DO CLIENTE FINAL POR PROJETO
# ═══════════════════════════════════════════════════════════════

@app.get("/api/projects/{job_id}/client")
async def get_project_client(job_id: str, request: Request):
    """Retorna dados do cliente final de um projeto.

    LGPD: project_clients contém PII (nome, email, telefone, endereço).
    Só o dono do projeto (ou admin) pode ler.
    """
    _require_project_owner(request, job_id)
    try:
        status, data = _supa_rest_as_user(
            request, "GET",
            f"/project_clients?job_id=eq.{job_id}&select=*",
            timeout=8,
        )
        if status >= 400 or data is None:
            return {} if status == 200 else {"error": "Falha interna ao buscar cliente"}
        return data[0] if data else {}
    except Exception as e:
        print(f"[get_project_client] erro: {e}")
        return {"error": "Falha interna ao buscar cliente"}


class ProjectMetaPayload(BaseModel):
    project_name: Optional[str] = None
    typology: Optional[str] = None
    address: Optional[str] = None
    phase: Optional[str] = None


@app.post("/api/projects/{job_id}/meta")
async def update_project_meta(job_id: str, payload: ProjectMetaPayload, request: Request):
    """Edita metadados do projeto: nome amigável, tipologia, endereço, fase.

    Só dono do projeto (ou admin) pode editar. Validação de tipologia, e
    aceita campos opcionais — só atualiza o que foi informado.
    """
    _require_project_owner(request, job_id)

    updates: dict = {}
    if payload.project_name is not None:
        name = payload.project_name.strip()
        if not name:
            raise HTTPException(400, "Nome do projeto não pode ser vazio")
        if len(name) > 120:
            raise HTTPException(400, "Nome do projeto excede 120 caracteres")
        updates["project_name"] = name
    if payload.typology is not None:
        if payload.typology not in _VALID_TYPOLOGIES:
            raise HTTPException(400, f"Tipologia inválida. Aceitas: {sorted(_VALID_TYPOLOGIES)}")
        updates["typology"] = payload.typology
    if payload.address is not None:
        updates["address"] = (payload.address or "").strip()[:240]
    if payload.phase is not None:
        updates["phase"] = (payload.phase or "").strip()[:60]

    if not updates:
        raise HTTPException(400, "Nenhum campo a atualizar")

    ok = _supabase_update("projects", "job_id", job_id, updates)
    if not ok:
        raise HTTPException(500, "Erro ao salvar no banco")
    return {"status": "ok", "job_id": job_id, "updated": list(updates.keys())}


@app.post("/api/projects/{job_id}/client")
async def upsert_project_client(
    job_id: str,
    request: Request,
    client_name: str = Form(""),
    client_company: str = Form(""),
    client_email: str = Form(""),
    client_phone: str = Form(""),
    address_site: str = Form(""),
    internal_notes: str = Form(""),
):
    """Cria ou atualiza dados do cliente final de um projeto.

    LGPD: só o dono do projeto (ou admin) pode escrever — antes era aberto,
    qualquer atacante podia sobrescrever PII alheia conhecendo o job_id.
    """
    _require_project_owner(request, job_id)
    payload = {
        "job_id": job_id,
        "client_name": client_name,
        "client_company": client_company,
        "client_email": client_email,
        "client_phone": client_phone,
        "address_site": address_site,
        "internal_notes": internal_notes,
        "updated_at": datetime.utcnow().isoformat(),
    }
    try:
        # Checa se já existe
        status, existing = _supa_rest_as_user(
            request, "GET",
            f"/project_clients?job_id=eq.{job_id}&select=id",
            timeout=8,
        )
        existing = existing or []

        if existing:
            # UPDATE
            up_status, _ = _supa_rest_as_user(
                request, "PATCH",
                f"/project_clients?id=eq.{existing[0]['id']}",
                body=payload, prefer="return=minimal", timeout=8,
            )
            if up_status >= 400:
                raise HTTPException(500, "Falha ao atualizar os dados do cliente")
        else:
            # INSERT
            ins_status, _ = _supa_rest_as_user(
                request, "POST",
                "/project_clients",
                body=payload, prefer="return=minimal", timeout=8,
            )
            if ins_status >= 400:
                raise HTTPException(500, "Falha ao gravar os dados do cliente")
        return {"status": "ok"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Erro ao salvar cliente: {e}")


# ── STRIPE CHECKOUT ──
@app.post("/api/estimate-price")
async def estimate_price(request: Request, files: list[UploadFile] = File(...)):
    """Conta pranchas REAIS (páginas dentro de PDFs + layouts dentro de DWG/DXF)
    e devolve preço calculado. Usado antes do checkout pra mostrar preview
    transparente pro cliente.
    """
    if not files:
        raise HTTPException(400, "Nenhum arquivo enviado")
    # Anti-DoS: endpoint PÚBLICO — limita rajadas por IP (fix segurança 2026-07-22).
    if not _rate_limit_ok("estimate", request, limit=15, window_s=600):
        raise HTTPException(429, "Muitas estimativas em pouco tempo. Espere alguns minutos e tente de novo.")
    # Anti-OOM: este endpoint é PÚBLICO (sem login) — sem teto vira porta anônima
    # de estouro de memória. Teto de request + streaming pro disco (mesmo padrão
    # do /api/process). Cap por-arquivo abaixo evita um único CAD gigante.
    _clen = request.headers.get("content-length") or request.headers.get("Content-Length")
    if _clen and _clen.isdigit() and int(_clen) > 450 * 1024 * 1024:
        raise HTTPException(413, "Arquivos muito grandes (máx. ~450 MB no total).")
    tmp_dir = os.path.join(WORK_DIR, "_estimate_tmp", str(uuid.uuid4())[:8])
    os.makedirs(tmp_dir, exist_ok=True)
    saved_paths = []
    try:
        for f in files:
            if not f.filename:
                continue
            safe_name = _safe_local_filename(f.filename)
            p = os.path.join(tmp_dir, safe_name)
            # Grava em pedaços (não bufferiza o arquivo inteiro na RAM).
            n_written, _ = await _stream_upload_to_disk(f, p)
            if n_written > 150 * 1024 * 1024:
                try: os.remove(p)
                except OSError: pass
                raise HTTPException(413, f"Arquivo '{f.filename}' grande demais (máx. ~150 MB por prancha).")
            saved_paths.append(p)
        from pricing import estimate_for_files, precheck_warnings
        result = estimate_for_files(saved_paths)
        # QW5 (20/07): precheck barato ANTES de pagar, reaproveitando os arquivos
        # que já estão em disco aqui (antes do finally apagar). Best-effort —
        # se o precheck falhar, o preço sai igual (warnings=[]).
        try:
            warnings = precheck_warnings(saved_paths)
        except Exception as _pe:
            print(f"[estimate] precheck falhou (não crítico): {_pe}")
            warnings = []
        return {"status": "ok", **result, "warnings": warnings}
    finally:
        # Limpa
        for p in saved_paths:
            try: os.remove(p)
            except Exception: pass
        try: os.rmdir(tmp_dir)
        except Exception: pass


# ═══════════════════════════════════════════════════════════════
#  Créditos de usuário (cashback, cupom, referral)
# ═══════════════════════════════════════════════════════════════

def _get_available_credits(user_id: str) -> list[dict]:
    """Retorna créditos disponíveis (used_at IS NULL e não expirados),
    ordenados por expires_at asc (expira primeiro consome primeiro)."""
    if not user_id or user_id == "anonymous":
        return []
    try:
        import urllib.request as _ur, urllib.parse as _up
        query = (f"select=id,amount_cents,source,source_ref,description,expires_at,created_at"
                 f"&user_id=eq.{_up.quote(user_id)}&used_at=is.null"
                 f"&order=expires_at.asc.nullslast,created_at.asc")
        url = f"{SUPABASE_URL}/rest/v1/user_credits?{query}"
        req = _ur.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Accept", "application/json")
        resp = _ur.urlopen(req, timeout=8)
        rows = _json.loads(resp.read().decode("utf-8"))
        now = datetime.utcnow().isoformat()
        # Filtra expirados
        return [r for r in rows if not r.get("expires_at") or r["expires_at"] > now]
    except Exception as e:
        print(f"credits query error: {e}")
        return []


def _total_available_credit(user_id: str) -> int:
    """Soma dos créditos disponíveis em centavos."""
    return sum(int(c.get("amount_cents") or 0) for c in _get_available_credits(user_id))


def _grant_credit(user_id: str, amount_cents: int, source: str,
                  source_ref: str = "", description: str = "",
                  expires_days: Optional[int] = None) -> bool:
    """Cria um crédito novo pro user."""
    if not user_id or user_id == "anonymous" or amount_cents <= 0:
        return False
    record = {
        "user_id": user_id,
        "amount_cents": int(amount_cents),
        "source": source,
        "source_ref": (source_ref or "")[:100],
        "description": (description or "")[:200],
    }
    if expires_days:
        record["expires_at"] = (datetime.utcnow() + __import__("datetime").timedelta(
            days=expires_days)).isoformat() + "Z"
    return _supabase_insert("user_credits", record)


def _consume_credits(user_id: str, amount_cents: int, job_id: str) -> int:
    """Consome créditos até somar amount_cents. Marca used_at nos usados.
    Retorna o TOTAL efetivamente consumido (pode ser menor se saldo insuficiente)."""
    available = _get_available_credits(user_id)
    if not available:
        return 0
    consumed = 0
    now_iso = datetime.utcnow().isoformat() + "Z"
    import urllib.request as _ur
    for cr in available:
        if consumed >= amount_cents:
            break
        # Marca used_at
        try:
            url = f"{SUPABASE_URL}/rest/v1/user_credits?id=eq.{cr['id']}&used_at=is.null"
            body = _json.dumps({"used_at": now_iso, "used_on_job_id": job_id}).encode("utf-8")
            req = _ur.Request(url, data=body, method="PATCH")
            req.add_header("apikey", SUPABASE_KEY)
            req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            req.add_header("Content-Type", "application/json")
            req.add_header("Prefer", "return=minimal")
            _ur.urlopen(req, timeout=8)
            consumed += int(cr.get("amount_cents") or 0)
        except Exception as e:
            print(f"consume credit error: {e}")
    return consumed


@app.get("/api/credits/balance")
async def credits_balance(request: Request, user_id: str):
    """Retorna saldo de crédito disponível pro usuário.

    Auth: exige JWT do próprio usuário (ou admin). Sem isso, qualquer um que
    descubra um user_id veria o saldo financeiro alheio (IDOR — auditoria
    2026-07-13). O frontend já chama com authFetch (manda o Bearer)."""
    jwt_user = _get_user_from_request(request)
    if not jwt_user:
        raise HTTPException(401, "Autenticação requerida")
    if jwt_user.get("id") != user_id and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
        raise HTTPException(403, "Só é possível consultar seu próprio saldo")
    credits = _get_available_credits(user_id)
    total = sum(int(c.get("amount_cents") or 0) for c in credits)
    return {
        "status": "ok",
        "user_id": user_id,
        "total_cents": total,
        "total_brl": round(total / 100, 2),
        "credits": credits,
    }


@app.post("/api/checkout")
async def create_checkout(request: Request, num_pranchas: int = 1, num_files: int = 0,
                          user_id: str = ""):
    """Cria sessão de pagamento no Stripe baseado em PRANCHAS REAIS.

    Aplica automaticamente créditos disponíveis do usuário (cashback,
    cupom, etc.) reduzindo o valor cobrado. Se o saldo cobrir 100%,
    pula Stripe e retorna direto `is_free=true`.

    Aceita `num_pranchas` (preferido — vem de /api/estimate-price). Mantém
    `num_files` como fallback compatibilidade.

    Auth: se `user_id` é informado e não é 'anonymous', exige JWT do dono.
    Senão atacante consome créditos da vítima via `_total_available_credit`.
    """
    # Validação JWT vs user_id (mesmo padrão do /api/process)
    if user_id and user_id != "anonymous":
        jwt_user = _get_user_from_request(request)
        if not jwt_user:
            raise HTTPException(401, "Autenticação requerida quando user_id é informado")
        if jwt_user.get("id") != user_id and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
            raise HTTPException(403, "user_id não corresponde ao token de autenticação")

    import stripe
    stripe.api_key = os.getenv("STRIPE_SECRET_KEY")
    if not stripe.api_key:
        raise HTTPException(500, "Stripe não configurado")

    n = num_pranchas if num_pranchas > 0 else (num_files or 1)
    from pricing import calculate_price, get_tier_for
    price_cents = calculate_price(n)
    tier = get_tier_for(n)
    plan_name = f"Plano {tier['name']} — {n} pranchas"

    # Aplica créditos disponíveis (se houver user_id)
    credit_cents = min(price_cents, _total_available_credit(user_id)) if user_id else 0
    final_cents = max(0, price_cents - credit_cents)

    # Se crédito cobriu 100%, pula Stripe
    if final_cents == 0 and credit_cents > 0:
        return {
            "is_free": True,
            "checkout_url": None,
            "session_id": None,
            "num_pranchas": n,
            "price_cents": price_cents,
            "credit_applied_cents": credit_cents,
            "credit_applied_brl": round(credit_cents / 100, 2),
            "final_cents": 0,
            "final_brl": 0.0,
            "message": "Coberto pelo saldo de créditos — processe direto sem pagamento",
        }

    # Descrição do produto inclui o desconto aplicado se houver
    description = plan_name
    if credit_cents > 0:
        description += f" (R$ {credit_cents/100:.2f} de desconto aplicado)"

    line_items = [{
        "price_data": {
            "currency": "brl",
            "product_data": {
                "name": f"AI.arq — Planilha de quantitativos",
                "description": description,
            },
            "unit_amount": final_cents,
        },
        "quantity": 1,
    }]
    base_params = {
        "line_items": line_items,
        "mode": "payment",
        "success_url": "https://ai.arq.br/dashboard.html?payment=success&session_id={CHECKOUT_SESSION_ID}",
        "cancel_url": "https://ai.arq.br/dashboard.html?payment=cancelled",
        "metadata": {
            "user_id": user_id or "anonymous",
            "num_pranchas": str(n),
            "price_cents": str(price_cents),
            "credit_applied_cents": str(credit_cents),
        },
    }

    # Tenta com PIX+card primeiro. Se o Stripe não tiver PIX ativado
    # (típico em conta nova ou ainda não configurada), cai pra card-only
    # automaticamente — evita quebrar o fluxo do cliente.
    session = None
    methods_used = []
    try:
        session_params = {**base_params, "payment_method_types": ["card", "pix"]}
        session = stripe.checkout.Session.create(**session_params)
        methods_used = ["card", "pix"]
    except stripe.error.InvalidRequestError as e:
        msg = str(e).lower()
        if "pix" in msg or "payment_method_types" in msg:
            print(f"[checkout] PIX indisponível ({e}); caindo pra card-only")
            try:
                session_params = {**base_params, "payment_method_types": ["card"]}
                session = stripe.checkout.Session.create(**session_params)
                methods_used = ["card"]
            except Exception as e2:
                raise HTTPException(500, f"Erro Stripe (card-only fallback): {str(e2)}")
        else:
            raise HTTPException(500, f"Erro ao criar checkout: {str(e)}")
    except Exception as e:
        raise HTTPException(500, f"Erro ao criar checkout: {str(e)}")

    if not session:
        raise HTTPException(500, "Erro ao criar checkout: session não foi criada")
    return {
        "is_free": False,
        "checkout_url": session.url,
        "session_id": session.id,
        "num_pranchas": n,
        "price_cents": price_cents,
        "credit_applied_cents": credit_cents,
        "credit_applied_brl": round(credit_cents / 100, 2),
        "final_cents": final_cents,
        "final_brl": round(final_cents / 100, 2),
        "payment_methods": methods_used,
    }


@app.get("/api/checkout/verify/{session_id}")
async def verify_payment(session_id: str, request: Request):
    """Verifica se o pagamento foi concluído.

    Exige login e confere que a sessão Stripe pertence ao próprio usuário
    (metadata.user_id gravado no create-checkout). Fix 2026-07-22: antes era
    aberto — qualquer um com um session_id via status/valor do pagamento."""
    jwt_user = _get_user_from_request(request)
    if not jwt_user:
        raise HTTPException(401, "Autenticação requerida")
    import stripe
    stripe.api_key = os.getenv("STRIPE_SECRET_KEY")
    try:
        session = stripe.checkout.Session.retrieve(session_id)
    except Exception as e:
        raise HTTPException(404, f"Sessão não encontrada: {str(e)}")
    try:
        sess_user = (session.metadata or {}).get("user_id")
    except Exception:
        sess_user = None
    if (jwt_user.get("email", "").lower() != ADMIN_EMAIL
            and sess_user and sess_user != jwt_user.get("id")):
        raise HTTPException(403, "Sessão de pagamento não pertence a este usuário")
    return {
        "paid": session.payment_status == "paid",
        "status": session.payment_status,
        "amount": session.amount_total,
    }


# ── CALIBRATION ENDPOINTS ──

@app.post("/api/cashback/upload")
async def cashback_upload(
    request: Request,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    job_id: str = "",
):
    """Cashback: recebe a planilha revisada do cliente e alimenta os
    benchmarks de densidade da mesma tipologia.

    A tipologia e a área de referência são resolvidas consultando o
    projeto no Supabase. A lógica é a mesma da ingestão do admin —
    o sistema aprende **proporções típicas** (qty/m²) e apenas
    alerta sobre anomalias em projetos futuros. Nunca copia valores
    absolutos entre projetos.
    """
    if not file.filename or not file.filename.lower().endswith((".xlsx", ".xls")):
        raise HTTPException(400, "Arquivo deve ser .xlsx")
    if not job_id:
        raise HTTPException(400, "job_id é obrigatório")
    _require_project_owner(request, job_id)
    if not HAS_DENSITY_CAL:
        raise HTTPException(500, "Módulo density_calibration não carregado")

    project = _get_project_from_supabase(job_id)
    if not project:
        raise HTTPException(404, f"Projeto não encontrado pro job_id={job_id}")

    typology = project.get("typology") or "office"
    ref_area = project.get("layout_area") or project.get("total_area") or 0
    try:
        ref_area = float(ref_area or 0)
    except Exception:
        ref_area = 0
    if ref_area <= 0:
        raise HTTPException(
            400,
            "Projeto não tem área de referência (layout_area ou total_area) — "
            "ainda não foi possível computar a área pelo DWG. "
            "Aguarde o processamento terminar completamente antes de enviar a revisão.",
        )

    label = (project.get("project_name") or f"job_{job_id}")[:100]

    tmp_dir = os.path.join(WORK_DIR, "_density_tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    safe_name = _safe_local_filename(file.filename)
    revised_path = os.path.join(tmp_dir, f"cashback_{job_id}_{safe_name}")
    try:
        content = await file.read()
        with open(revised_path, "wb") as f:
            f.write(content)

        from density_calibration import ingest_budget as _ingest
        summary = _ingest(
            revised_path, area_m2=ref_area, typology=typology, project_label=label,
        )

        # Concede crédito de R$ 20 de cashback pro usuário — idempotente
        # por job_id (só libera uma vez por projeto revisado). Expira em
        # 180 dias pra evitar credit acumulando pra sempre sem uso.
        CASHBACK_CENTS = 2000  # R$ 20
        CASHBACK_EXPIRE_DAYS = 180
        cashback_granted = False
        user_id = project.get("user_id") or ""
        if user_id and user_id != "anonymous" and summary.get("items_parsed", 0) > 0:
            # Evita duplicar crédito pra mesma revisão
            already = _get_available_credits(user_id)
            duplicate = any(c.get("source") == "cashback" and c.get("source_ref") == job_id
                            for c in already)
            if not duplicate:
                cashback_granted = _grant_credit(
                    user_id=user_id,
                    amount_cents=CASHBACK_CENTS,
                    source="cashback",
                    source_ref=job_id,
                    description=f"Cashback revisão do projeto {label}",
                    expires_days=CASHBACK_EXPIRE_DAYS,
                )

        return {
            "status": "ok",
            "source": "cashback",
            "job_id": job_id,
            "typology": typology,
            "area_m2": ref_area,
            "project_label": label,
            "items_parsed": summary.get("items_parsed", 0),
            "benchmarks_updated": summary.get("benchmarks_updated", 0),
            "new_item_types": summary.get("new_item_types", 0),
            # Legacy fields for UI compat (dashboard shows these)
            "items_compared": summary.get("items_parsed", 0),
            "items_saved": summary.get("benchmarks_updated", 0),
            "avg_deviation_pct": 0,
            "cashback_granted": cashback_granted,
            "cashback_amount_brl": CASHBACK_CENTS / 100 if cashback_granted else 0,
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Erro na ingestão da revisão: {str(e)}")
    finally:
        if os.path.exists(revised_path):
            try:
                os.remove(revised_path)
            except Exception:
                pass


# Endpoints /api/calibration/manual e /api/calibration/factors foram removidos
# porque alimentavam o modelo de "fator absoluto" (correction_factor = real/ai),
# que contraria a regra de isolamento de projetos. A mesma ingestão agora é
# feita via /api/calibration/ingest (admin) e /api/cashback/upload (cliente),
# que aprendem RATIOS de densidade (qty/m²) por tipologia e só geram alertas.


# ═══════════════════════════════════════════════
#  Calibração por DENSIDADE (ratios qty/área)
#  — regra: aprende padrões proporcionais pra ALERTAR anomalias,
#  nunca copia valores absolutos entre projetos.
# ═══════════════════════════════════════════════

try:
    from density_calibration import (
        ingest_budget as density_ingest_budget,
        get_benchmarks as density_get_benchmarks,
    )
    HAS_DENSITY_CAL = True
except ImportError:
    HAS_DENSITY_CAL = False
    print("density_calibration.py não disponível — calibração por densidade desativada")


@app.post("/api/calibration/ingest")
async def calibration_ingest_density(
    request: Request,
    xlsx: UploadFile = File(...),
    area_m2: float = 0,
    typology: str = "office",
    project_label: str = "",
):
    """Ingere um orçamento-fonte histórico pra enriquecer os benchmarks
    de densidade. `area_m2` é a área de referência do projeto-fonte (layout
    ou laje) usada pra computar qty/área por item.

    Benchmarks são agregados por (typology, item_type, unit) — projetos
    novos só recebem ALERTAS, nunca valores copiados.
    """
    _require_admin(request)  # escreve no motor de TODOS os projetos — só admin
    if not HAS_DENSITY_CAL:
        raise HTTPException(500, "Módulo density_calibration não carregado")
    if area_m2 <= 0:
        raise HTTPException(400, "area_m2 deve ser > 0 (área de referência do projeto-fonte)")
    if not typology or len(typology) < 3:
        raise HTTPException(400, "typology obrigatória (ex.: 'office', 'residential')")

    tmp_dir = os.path.join(WORK_DIR, "_density_tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    safe_name = _safe_local_filename(xlsx.filename)
    xlsx_path = os.path.join(tmp_dir, f"ingest_{safe_name}")
    try:
        content = await xlsx.read()
        with open(xlsx_path, "wb") as f:
            f.write(content)

        label = project_label or xlsx.filename
        summary = density_ingest_budget(
            xlsx_path, area_m2=area_m2, typology=typology,
            project_label=label,
        )
        return {"status": "ok", "area_m2": area_m2, "typology": typology,
                "project_label": label, **summary}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Erro na ingestão: {str(e)}")
    finally:
        if os.path.exists(xlsx_path):
            try:
                os.remove(xlsx_path)
            except Exception:
                pass


def _get_project_from_supabase(job_id: str) -> dict:
    """Busca um projeto pelo job_id — usado pelo cashback pra resolver
    typology + área do projeto revisado."""
    try:
        import urllib.request, json as _json
        url = f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}&select=*"
        req = urllib.request.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Accept", "application/json")
        resp = urllib.request.urlopen(req, timeout=8)
        rows = _json.loads(resp.read().decode("utf-8"))
        return rows[0] if rows else {}
    except Exception as e:
        print(f"Supabase select project error: {e}")
        return {}


# REMOVIDO 2026-05-14: /api/calibration/ingest-from-review era duplicata
# zumbi de /api/cashback/upload (mesma lógica de calibração por densidade).
# O frontend nunca chamou /calibration/ingest-from-review — só usa o
# /cashback/upload. Eliminado pra reduzir surface area e clarear o código.


# ═══════════════════════════════════════════════════════════════
#  Agente "tira-dúvidas" — Q&A sobre uma planilha gerada
# ═══════════════════════════════════════════════════════════════

@app.post("/api/agent/ask")
async def agent_ask(request: Request, job_id: str, question: str = ""):
    """Cliente faz uma pergunta sobre o quantitativo de UM job. O agente
    investiga (lê planilha, busca itens, lê DXFs, checa calibração) e
    responde em linguagem natural com referências aos itens.

    Body opcional (JSON): {"history": [{"role": "user|assistant", "content": "..."}]}
    pra manter contexto de conversação contínua.
    """
    if not job_id:
        raise HTTPException(400, "job_id obrigatório")
    if not question or len(question.strip()) < 2:
        raise HTTPException(400, "pergunta vazia")
    _require_project_owner(request, job_id)

    # History opcional via JSON body
    history = None
    try:
        body_bytes = await request.body()
        if body_bytes:
            import json as _j
            body = _j.loads(body_bytes.decode("utf-8"))
            if isinstance(body, dict) and isinstance(body.get("history"), list):
                history = body["history"]
    except Exception:
        history = None

    try:
        from agent import ask
        result = ask(job_id=job_id, question=question.strip(), history=history)
        return {"status": "ok", **result}
    except Exception as e:
        raise HTTPException(500, f"Erro do agente: {type(e).__name__}: {e}")


@app.get("/api/agent/conversations")
async def agent_conversations(request: Request, job_id: Optional[str] = None, limit: int = 50):
    # Se passou job_id, valida ownership; senão, exige admin (lista global de conversas)
    if job_id:
        _require_project_owner(request, job_id)
    else:
        _require_admin(request)
    """Lista conversas do agente — usado pelo admin pra acompanhar uso."""
    try:
        import urllib.request as _ur
        query = f"select=*&order=created_at.desc&limit={int(limit)}"
        if job_id:
            query += f"&job_id=eq.{_ur.quote(job_id)}"
        url = f"{SUPABASE_URL}/rest/v1/agent_conversations?{query}"
        req = _ur.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Accept", "application/json")
        resp = _ur.urlopen(req, timeout=10)
        rows = _json.loads(resp.read().decode("utf-8"))
        # Enriquece com QUEM perguntou e QUAL projeto: a conversa grava
        # user_id='anonymous', entao a identidade vem do projeto ao qual o chat
        # esta ancorado (job_id -> projects). Uma query so pra todos os job_ids.
        _jids = sorted({(r.get("job_id") or "") for r in rows if r.get("job_id")})
        _pmap = {}
        if _jids:
            try:
                _inlist = ",".join(_jids)
                _purl = (f"{SUPABASE_URL}/rest/v1/projects?job_id=in.({_ur.quote(_inlist, safe=',')})"
                         f"&select=job_id,user_email,user_name,project_name")
                _preq = _ur.Request(_purl, method="GET")
                _preq.add_header("apikey", SUPABASE_KEY)
                _preq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                _preq.add_header("Accept", "application/json")
                for _p in _json.loads(_ur.urlopen(_preq, timeout=10).read().decode("utf-8")):
                    _pmap[_p.get("job_id")] = _p
            except Exception:
                _pmap = {}
        for _r in rows:
            _p = _pmap.get(_r.get("job_id")) or {}
            _r["user_email"] = _p.get("user_email") or ""
            _r["user_name"] = _p.get("user_name") or ""
            _r["project_name"] = _p.get("project_name") or ""
        return {"status": "ok", "count": len(rows), "conversations": rows}
    except Exception as e:
        raise HTTPException(500, f"Erro ao listar conversas: {str(e)}")


@app.get("/api/agent/stats")
async def agent_stats(request: Request):
    """Estatísticas agregadas do uso do agente — pra dashboard admin."""
    _require_admin(request)
    try:
        # FIX 2026-05-14: antes definia URL pra RPC `agent_stats_summary` e
        # sobrescrevia logo na linha seguinte com fallback. RPC nunca era
        # chamada. Mantém apenas o select agregado (consistente com a RPC
        # se um dia existir).
        import urllib.request as _ur
        url = (f"{SUPABASE_URL}/rest/v1/agent_conversations"
               "?select=id,iterations,duration_ms,error,job_id,created_at")
        req = _ur.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Accept", "application/json")
        resp = _ur.urlopen(req, timeout=10)
        rows = _json.loads(resp.read().decode("utf-8"))
        total = len(rows)
        with_error = sum(1 for r in rows if r.get("error"))
        unique_jobs = len({r.get("job_id") for r in rows if r.get("job_id")})
        durations = [r.get("duration_ms") or 0 for r in rows]
        avg_dur = int(sum(durations)/len(durations)) if durations else 0
        avg_iter = (sum(r.get("iterations") or 0 for r in rows) / total) if total else 0
        return {
            "status": "ok",
            "total_conversations": total,
            "unique_jobs": unique_jobs,
            "errors": with_error,
            "avg_duration_ms": avg_dur,
            "avg_iterations": round(avg_iter, 2),
        }
    except Exception as e:
        raise HTTPException(500, f"Erro stats: {str(e)}")


@app.post("/api/calibration/reclassify-raws")
async def calibration_reclassify_raws(
    request: Request,
    typology: Optional[str] = None,
    limit: Optional[int] = None,
    only_unclassified: bool = True,
):
    """Classifica linhas raw existentes via LLM e recompila benchmarks.

    Requer auth admin — chamadas em massa consomem créditos Anthropic.

    Útil pra "ativar" raws antigos ingeridos antes do classificador
    existir. Idempotente: por padrão só toca raws sem familia_id.

    Cada raw vira ~3s (chamada Claude Haiku). Lote de 264 leva ~10min
    e custa ~$0.50. Use `limit` pra testar incrementalmente.
    """
    _require_admin(request)
    if not HAS_DENSITY_CAL:
        raise HTTPException(500, "Módulo density_calibration não carregado")
    try:
        from density_calibration import reclassify_raws
        result = reclassify_raws(
            typology=typology, limit=limit, only_unclassified=only_unclassified,
        )
        return {"status": "ok", **result}
    except Exception as e:
        raise HTTPException(500, f"Erro na reclassificação: {str(e)}")


@app.get("/api/calibration/benchmarks")
async def calibration_benchmarks(request: Request, typology: Optional[str] = None):
    """Lista os benchmarks de densidade agregados (mean ± stddev por
    tipologia × item_type × unit). Usado pelo admin pra auditar os
    padrões aprendidos."""
    _require_admin(request)
    if not HAS_DENSITY_CAL:
        raise HTTPException(500, "Módulo density_calibration não carregado")
    try:
        raw = density_get_benchmarks(typology=typology)
        rows = []
        for (item_type, unit), data in raw.items():
            rows.append({
                "typology": data.get("typology"),
                "item_type": item_type,
                "unit": unit,
                "mean": data.get("mean"),
                "stddev": data.get("stddev"),
                "min_value": data.get("min_value"),
                "max_value": data.get("max_value"),
                "n_projects": data.get("n_projects"),
            })
        rows.sort(key=lambda r: (-(r["n_projects"] or 0), r["item_type"]))
        return {"status": "ok", "count": len(rows), "benchmarks": rows}
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar benchmarks: {str(e)}")


# ═══════════════════════════════════════════════════════════════
#  REVISÃO INLINE DE ITENS (feedback loop)
# ═══════════════════════════════════════════════════════════════

@app.get("/api/items/{job_id}")
async def get_project_items(job_id: str, request: Request):
    """Retorna lista de itens individuais de um job pra revisão inline.
    Usa RPC `list_project_items` pra bypassar RLS."""
    _require_project_owner(request, job_id)
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
        body = json.dumps({"p_job_id": job_id}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        items = json.loads(resp.read().decode('utf-8'))
        # Meta do projeto (datas, status) — pra páginas que NÃO acham o projeto no
        # by-user usarem como fallback (caso avaliação/eval: user_id != dono logado,
        # então a linha não vem na lista e o tempo/datas ficavam "--"). Best-effort.
        _meta = {}
        try:
            _murl = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
                     f"&select=project_name,status,typology,files_count,items_count,"
                     f"total_area,user_total_area,created_at,completed_at,phase&limit=1")
            _mreq = urllib.request.Request(_murl, method="GET")
            _mreq.add_header("apikey", SUPABASE_KEY)
            _mreq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            _mrows = json.loads(urllib.request.urlopen(_mreq, timeout=10).read().decode("utf-8"))
            if _mrows:
                _meta = _mrows[0]
        except Exception as _me:
            print(f"[items] meta do projeto falhou (não crítico): {_me}")
        return {"status": "ok", "job_id": job_id, "items": items,
                "count": len(items), "project": _meta}
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar itens: {str(e)}")


@app.get("/api/memorial/{job_id}")
async def memorial_docx(job_id: str, request: Request):
    """Memorial descritivo (RASCUNHO) em .docx.

    v1 determinística (01/08/2026) — texto = template + itens; zero IA, zero
    invenção. v1.1: se o cliente EDITOU na tela (memorial.html), o .docx sai
    da versão salva em project_memorial. Download exige downloadProtected no
    frontend (armadilha nº9: <a href> não manda Authorization)."""
    _require_project_owner(request, job_id)
    import tempfile
    try:
        from memorial import estrutura_para_docx
        salvo = _memorial_carregar_salvo(job_id)
        if salvo:
            estrutura = salvo
        else:
            projeto, items = _memorial_dados_frescos(job_id)
            from memorial import montar_estrutura
            estrutura = montar_estrutura(projeto, items)
        tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        tmp.close()
        resumo = estrutura_para_docx(tmp.name, estrutura)
        print(f"[memorial] job={job_id} salvo={bool(salvo)} → {resumo}")
        fname = f"memorial_descritivo_rascunho_{_slug_filename(estrutura.get('obra') or job_id)}.docx"
        return FileResponse(
            tmp.name,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=fname)
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[memorial] erro: {e}")
        print(traceback.format_exc())
        _log_error("memorial:gerar", str(e), job_id=job_id)
        raise HTTPException(500, f"Erro ao gerar memorial: {e}")


def _assinatura_quantitativo(items) -> str:
    """Impressão digital do quantitativo. Cronograma e memorial nascem DAQUI —
    se o cliente editar/confirmar/excluir um item depois, a assinatura muda e o
    site avisa que o entregável salvo ficou velho, em vez de entregar número
    desatualizado calado. Regra do Pedro (02/08/2026): tudo interligado.

    Entram só os campos que mudam o resultado: descrição e observação viram
    texto do memorial; quantidade/unidade viram esforço e duração de fase;
    disciplina vira fase; confiança vira o selo medido/estimado."""
    import hashlib
    partes = []
    for it in (items or []):
        try:
            _q = round(float(it.get("quantity") or 0), 4)
        except (TypeError, ValueError):
            _q = 0.0
        partes.append("|".join([
            str(it.get("id") or ""),
            (it.get("description") or "").strip(),
            f"{_q:.4f}",
            (it.get("unit") or "").strip(),
            (it.get("discipline") or "").strip(),
            (it.get("confidence") or "").strip(),
            (it.get("observations") or "").strip(),
        ]))
    partes.sort()  # ordem do banco não pode mudar a assinatura
    bruto = f"n={len(partes)}\n" + "\n".join(partes)
    return hashlib.sha256(bruto.encode("utf-8")).hexdigest()[:32]


_ASSINATURA_CACHE: dict = {}   # job_id -> (timestamp, assinatura)
_ASSINATURA_TTL = 60           # seg. O autosave do memorial salva a cada 1,5s;
                               # sem cache seria uma leitura de itens por tecla.


def _assinatura_invalidar(job_id: str) -> None:
    """Chamado quando um item muda — a próxima leitura recalcula na hora."""
    _ASSINATURA_CACHE.pop(job_id, None)


def _assinatura_atual(job_id: str) -> str:
    """Assinatura do quantitativo que está no banco AGORA. String vazia se não
    der pra ler — nesse caso o chamador não deve acusar desatualização (não
    inventar alarme em cima de falha de rede)."""
    import json
    import time as _time
    import urllib.request as _url_req
    _hit = _ASSINATURA_CACHE.get(job_id)
    if _hit and (_time.time() - _hit[0]) < _ASSINATURA_TTL:
        return _hit[1]
    try:
        _u = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
        _b = json.dumps({"p_job_id": job_id}).encode("utf-8")
        _r = _url_req.Request(_u, data=_b, method="POST")
        _r.add_header("apikey", SUPABASE_KEY)
        _r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _r.add_header("Content-Type", "application/json")
        items = json.loads(_url_req.urlopen(_r, timeout=15).read().decode("utf-8")) or []
        _sig = _assinatura_quantitativo(items)
        _ASSINATURA_CACHE[job_id] = (_time.time(), _sig)
        return _sig
    except Exception as e:
        print(f"[coerencia] não consegui assinar o quantitativo de {job_id}: {e}")
        return ""


def _carimbar_planilha(job_id: str) -> None:
    """Registra de qual quantitativo o .xlsx atual nasceu. Chamar SEMPRE depois
    de gerar/subir a planilha e depois dos itens já estarem no banco — o .xlsx é
    servido de arquivo salvo, então sem este carimbo o cliente que corrige itens
    e não clica em 'Finalizar revisão' baixa os números velhos sem saber.
    Best-effort: falhar aqui não pode derrubar um job que deu certo."""
    try:
        _assinatura_invalidar(job_id)   # itens acabaram de mudar
        _supabase_update("projects", "job_id", job_id, {
            "planilha_assinatura": _assinatura_atual(job_id),
            "planilha_gerada_em": datetime.utcnow().isoformat() + "Z",
        })
    except Exception as e:
        print(f"[coerencia] carimbo da planilha falhou ({job_id}), não crítico: {e}")


def _memorial_dados_frescos(job_id: str):
    """(projeto, items) pros geradores do memorial — mesma RPC do /api/items."""
    import json
    import urllib.request as _url_req
    _u = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
    _b = json.dumps({"p_job_id": job_id}).encode("utf-8")
    _r = _url_req.Request(_u, data=_b, method="POST")
    _r.add_header("apikey", SUPABASE_KEY)
    _r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
    _r.add_header("Content-Type", "application/json")
    items = json.loads(_url_req.urlopen(_r, timeout=15).read().decode("utf-8")) or []
    if not items:
        raise HTTPException(404, "Projeto sem itens — gere o quantitativo primeiro")
    projeto = {}
    try:
        _mu = (f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}"
               f"&select=project_name,typology,total_area,user_total_area&limit=1")
        _mr = _url_req.Request(_mu, method="GET")
        _mr.add_header("apikey", SUPABASE_KEY)
        _mr.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _rows = json.loads(_url_req.urlopen(_mr, timeout=10).read().decode("utf-8"))
        if _rows:
            projeto = _rows[0]
    except Exception as _me:
        print(f"[memorial] meta do projeto falhou (não crítico): {_me}")
    return projeto, items


def _memorial_carregar_salvo(job_id: str):
    """Estrutura editada salva em project_memorial, ou None. Tupla do REST
    helper SEMPRE desempacotada (lição do bug de 01/08 no dedupe de reviews)."""
    import urllib.parse
    _st, _rows = _supa_rest_service(
        "GET", f"project_memorial?job_id=eq.{urllib.parse.quote(job_id)}&select=conteudo&limit=1")
    if _st == 200 and isinstance(_rows, list) and _rows:
        return _rows[0].get("conteudo") or None
    return None


@app.get("/api/memorial/{job_id}/pdf")
async def memorial_pdf(job_id: str, request: Request):
    """Memorial em PDF (WeasyPrint, mesmo motor do cronograma). Prefere a
    versão editada/salva, igual ao .docx."""
    _require_project_owner(request, job_id)
    import tempfile
    try:
        salvo = _memorial_carregar_salvo(job_id)
        if salvo:
            estrutura = salvo
        else:
            projeto, items = _memorial_dados_frescos(job_id)
            from memorial import montar_estrutura
            estrutura = montar_estrutura(projeto, items)
        from memorial import estrutura_para_pdf_bytes
        pdf = estrutura_para_pdf_bytes(estrutura)
        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        tmp.write(pdf)
        tmp.close()
        fname = f"memorial_descritivo_rascunho_{_slug_filename(estrutura.get('obra') or job_id)}.pdf"
        return FileResponse(tmp.name, media_type="application/pdf", filename=fname)
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[memorial pdf] erro: {e}")
        print(traceback.format_exc())
        _log_error("memorial:pdf", str(e), job_id=job_id)
        raise HTTPException(500, f"Erro ao gerar PDF: {e}")


@app.get("/api/memorial/{job_id}/estrutura")
async def memorial_estrutura(job_id: str, request: Request):
    """Estrutura editável do memorial pra tela memorial.html.
    Devolve a versão SALVA se existir; senão monta fresca dos itens.

    `?fresco=1` ignora a salva e remonta com os números de agora — é o botão
    'atualizar' de quando o cliente corrige o quantitativo depois. Não grava
    nada: o texto novo só entra no lugar quando ele salvar."""
    _require_project_owner(request, job_id)
    fresco = str(request.query_params.get("fresco") or "") in ("1", "true", "sim")
    try:
        salvo = None if fresco else _memorial_carregar_salvo(job_id)
        if salvo:
            return {"status": "ok", "salvo": True, "estrutura": salvo}
        projeto, items = _memorial_dados_frescos(job_id)
        from memorial import montar_estrutura
        return {"status": "ok", "salvo": False, "estrutura": montar_estrutura(projeto, items)}
    except HTTPException:
        raise
    except Exception as e:
        _log_error("memorial:estrutura", str(e), job_id=job_id)
        raise HTTPException(500, f"Erro ao montar memorial: {e}")


@app.post("/api/memorial/{job_id}/redigir")
async def memorial_redigir_ia(job_id: str, request: Request):
    """Reescreve com IA os parágrafos de abertura das seções do memorial.

    Coleira (02/08): a IA vê só as DESCRIÇÕES dos itens (nunca as quantidades)
    e devolve um parágrafo por disciplina; o validador rejeita qualquer texto
    com número, norma ou juízo de valor. Rejeitado = fica o texto
    determinístico. Item, quantidade, rótulo medido/estimado e [A PREENCHER]
    não passam pela IA em momento nenhum."""
    _require_project_owner(request, job_id)
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(503, "Redação por IA indisponível agora")
    try:
        estrutura = _memorial_carregar_salvo(job_id)
        if not estrutura:
            projeto, items = _memorial_dados_frescos(job_id)
            from memorial import montar_estrutura
            estrutura = montar_estrutura(projeto, items)

        import anthropic
        _client = anthropic.Anthropic(api_key=api_key, timeout=60.0)

        def _chamar(system, user):
            r = _client.messages.create(
                model="claude-haiku-4-5", max_tokens=300,
                system=system, messages=[{"role": "user", "content": user}])
            return r.content[0].text if r.content else ""

        from memorial import redigir_intros_ia, validar_estrutura_editada
        resumo = redigir_intros_ia(estrutura, _chamar)
        if resumo["puladas"]:
            print(f"[memorial ia] job={job_id} puladas: {resumo['puladas']}")
        limpa = validar_estrutura_editada(estrutura)
        _st, _resp = _supa_rest_service(
            "POST", "project_memorial?on_conflict=job_id",
            {"job_id": job_id, "conteudo": limpa, "updated_at": "now()",
             "itens_assinatura": _assinatura_atual(job_id)},
            prefer="resolution=merge-duplicates")
        if _st not in (200, 201, 204):
            _log_error("memorial:redigir-salvar", f"REST {_st}", job_id=job_id)
            raise HTTPException(500, "Texto gerado, mas falhou ao salvar")
        return {"status": "ok", "reescritas": resumo["reescritas"],
                "puladas": len(resumo["puladas"]), "estrutura": limpa}
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[memorial ia] erro: {e}")
        print(traceback.format_exc())
        _log_error("memorial:redigir", str(e), job_id=job_id)
        raise HTTPException(500, f"Erro ao redigir: {e}")


@app.post("/api/memorial/{job_id}/estrutura")
async def memorial_salvar(job_id: str, request: Request):
    """Salva a estrutura EDITADA na tela. Valida/sanitiza tudo que vem do
    navegador (memorial.validar_estrutura_editada) — só campos conhecidos
    sobrevivem, textos com teto de tamanho."""
    _require_project_owner(request, job_id)
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(400, "JSON inválido")
    from memorial import validar_estrutura_editada
    try:
        limpa = validar_estrutura_editada((body or {}).get("estrutura"))
    except ValueError as ve:
        raise HTTPException(400, f"Estrutura inválida: {ve}")
    _st, _resp = _supa_rest_service(
        "POST", "project_memorial?on_conflict=job_id",
        {"job_id": job_id, "conteudo": limpa, "updated_at": "now()",
         "itens_assinatura": _assinatura_atual(job_id)},
        prefer="resolution=merge-duplicates")
    if _st not in (200, 201, 204):
        _log_error("memorial:salvar", f"REST {_st}: {str(_resp)[:300]}", job_id=job_id)
        raise HTTPException(500, "Erro ao salvar o memorial")
    return {"status": "ok"}


# ═══════════════════════════════════════════════════════════════
#  COERÊNCIA ENTRE OS ENTREGÁVEIS  (regra do Pedro, 02/08/2026)
#  Quantitativo, cronograma e memorial saem do MESMO projeto. Mexeu
#  no quantitativo, os outros dois envelheceram — e o cliente tem que
#  ficar sabendo, não descobrir na obra. Aqui a gente detecta e avisa;
#  refazer é sempre escolha dele (memorial e cronograma têm edição
#  manual dentro, ninguém sobrescreve o trabalho do cliente sozinho).
# ═══════════════════════════════════════════════════════════════

def _revisoes_depois_de(job_id: str, quando: str) -> dict:
    """Quantos itens foram editados/excluídos depois de `quando` (ISO). É o
    'o que mudou' que a assinatura sozinha não conta."""
    import urllib.parse
    out = {"editados": 0, "excluidos": 0}
    if not quando:
        return out
    _st, _rows = _supa_rest_service(
        "GET",
        f"item_reviews?job_id=eq.{urllib.parse.quote(job_id)}"
        f"&reviewed_at=gt.{urllib.parse.quote(quando)}"
        f"&action=in.(edit,reject)&select=action,item_id")
    if _st != 200 or not isinstance(_rows, list):
        return out
    # Mesmo item mexido 3× continua sendo 1 item mudado.
    editados, excluidos = set(), set()
    for r in _rows:
        (excluidos if r.get("action") == "reject" else editados).add(r.get("item_id"))
    out["editados"] = len(editados - excluidos)
    out["excluidos"] = len(excluidos)
    return out


def _frase_mudanca(m: dict) -> str:
    p = []
    if m.get("editados"):
        n = m["editados"]
        p.append(f"{n} item corrigido" if n == 1 else f"{n} itens corrigidos")
    if m.get("excluidos"):
        n = m["excluidos"]
        p.append(f"{n} item excluído" if n == 1 else f"{n} itens excluídos")
    return " e ".join(p)


def _coerencia_do_projeto(job_id: str) -> dict:
    """Estado de sincronia dos três entregáveis do projeto."""
    import urllib.parse
    atual = _assinatura_atual(job_id)

    def _avaliar(nome: str, row: Optional[dict]) -> dict:
        if not row:
            return {"existe": False, "desatualizado": False}
        quando = row.get("updated_at") or row.get("created_at") or ""
        salva = row.get("itens_assinatura")
        mudancas = _revisoes_depois_de(job_id, quando)
        # Duas provas independentes. A assinatura pega qualquer diferença de
        # números; a contagem de revisões cobre os entregáveis antigos, salvos
        # antes de existir assinatura (8 cronogramas em 02/08/2026).
        por_assinatura = bool(atual and salva and salva != atual)
        por_revisao = bool(mudancas["editados"] or mudancas["excluidos"])
        return {
            "existe": True,
            "desatualizado": por_assinatura or por_revisao,
            "em": quando,
            "mudancas": mudancas,
            "frase": _frase_mudanca(mudancas),
        }

    _st, _mem = _supa_rest_service(
        "GET", f"project_memorial?job_id=eq.{urllib.parse.quote(job_id)}"
               f"&select=updated_at,itens_assinatura&limit=1")
    memorial = _avaliar("memorial", (_mem or [None])[0] if _st == 200 and _mem else None)
    cronograma = _avaliar("cronograma", _supabase_get_cronograma(job_id))

    # A planilha .xlsx é servida de arquivo salvo — envelhece igual aos outros.
    _st, _prj = _supa_rest_service(
        "GET", f"projects?job_id=eq.{urllib.parse.quote(job_id)}"
               f"&select=planilha_assinatura,planilha_gerada_em&limit=1")
    _row = (_prj or [None])[0] if _st == 200 and _prj else None
    planilha = _avaliar("planilha", {
        "updated_at": _row.get("planilha_gerada_em"),
        "itens_assinatura": _row.get("planilha_assinatura"),
    } if _row and _row.get("planilha_gerada_em") else None)

    desatualizados = [n for n, v in (("planilha", planilha),
                                     ("cronograma", cronograma),
                                     ("memorial", memorial)) if v.get("desatualizado")]
    return {
        "assinatura": atual,
        "planilha": planilha,
        "cronograma": cronograma,
        "memorial": memorial,
        "desatualizados": desatualizados,
        "tudo_em_dia": not desatualizados,
    }


@app.get("/api/projetos/desatualizados")
async def projetos_desatualizados(request: Request):
    """Lista, de uma vez, os projetos do usuário logado com entregável velho —
    pro chip de aviso em Meus Projetos. Uma consulta pra todos: recalcular a
    assinatura projeto por projeto significaria reler os itens de cada um."""
    user = _get_user_from_request(request)
    if not user or not user.get("id"):
        raise HTTPException(401, "Autenticação requerida")
    try:
        _st, _rows = _supa_rest_service(
            "POST", "rpc/projetos_desatualizados", {"p_user_id": user["id"]})
        if _st != 200 or not isinstance(_rows, list):
            # Devolver lista vazia calado esconderia permissão errada no RPC:
            # o chip some e parece que está tudo em dia. Registra pra aparecer
            # no error_log em vez de virar "funcionou".
            _log_error("coerencia:lista", f"RPC {_st}: {str(_rows)[:300]}")
            return {"projetos": []}
        return {"projetos": _rows}
    except Exception as e:
        _log_error("coerencia:lista", str(e))
        return {"projetos": []}


@app.get("/api/projeto/{job_id}/coerencia")
async def projeto_coerencia(job_id: str, request: Request):
    """Diz quais entregáveis salvos ficaram velhos depois que o cliente mexeu
    no quantitativo. Só leitura — quem refaz é o cliente, com um clique."""
    _require_project_owner(request, job_id)
    try:
        return _coerencia_do_projeto(job_id)
    except Exception as e:
        # Nunca derruba a tela por causa do aviso: sem resposta confiável, o
        # front simplesmente não mostra banner nenhum. Mas registra — "nenhum
        # aviso" por falha é indistinguível de "está tudo em dia".
        _log_error("coerencia:projeto", str(e), job_id=job_id)
        _vazio = {"existe": False, "desatualizado": False}
        return {"assinatura": "", "planilha": _vazio, "cronograma": _vazio,
                "memorial": _vazio, "desatualizados": [], "tudo_em_dia": True,
                "erro": True}


# ═══════════════════════════════════════════════════════════════
#  CRONOGRAMA FÍSICO-FINANCEIRO (Fase 2 do roadmap)
# ═══════════════════════════════════════════════════════════════

class CronogramaPayload(BaseModel):
    data_inicio: str       # ISO YYYY-MM-DD
    duracao_meses: int     # 1..36


class CronogramaSavePayload(BaseModel):
    data_inicio: str
    duracao_meses: int
    k_sigmoid: Optional[int] = 10
    fases_custom: Optional[list] = []   # array de {label, inicio, fim, cor?, ambiente?, ordem?, manual?}
    notas: Optional[str] = ""


@app.get("/api/cronograma/{job_id}/sugestao")
async def cronograma_sugestao(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Sugere duração de obra baseada em tipologia + área + n disciplinas
    detectadas. Chamado pelo frontend ANTES do "Gerar cronograma" pra
    preencher o slider com valor inteligente.
    """
    import urllib.request, json as _json

    # Busca o projeto pra pegar typology + area + files_count
    try:
        _, rows = _supa_rest_as_user(
            request, "GET",
            f"/projects?job_id=eq.{job_id}"
            "&select=typology,layout_area,total_area,files_count",
            timeout=10,
        )
        rows = rows or []
        if not rows:
            raise HTTPException(404, "Projeto não encontrado")
        proj = rows[0]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar projeto: {e}")

    # Busca disciplinas ativas via items
    try:
        url2 = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
        body = _json.dumps({"p_job_id": job_id}).encode('utf-8')
        req2 = urllib.request.Request(url2, data=body, method='POST')
        req2.add_header('apikey', SUPABASE_KEY)
        req2.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req2.add_header('Content-Type', 'application/json')
        resp2 = urllib.request.urlopen(req2, timeout=15)
        items = _json.loads(resp2.read().decode('utf-8'))
    except Exception:
        items = []

    disciplinas = set()
    for it in items:
        d = (it.get('discipline') or '').strip().upper()
        if d and 'PREMISSA' not in d:
            disciplinas.add(d)
    n_disc = len(disciplinas)

    # Aplica heurística
    from cronograma import sugerir_duracao
    typology = proj.get('typology')
    area = proj.get('layout_area') or proj.get('total_area') or 0
    files_count = proj.get('files_count') or 0

    sugestao = sugerir_duracao(
        typology=typology,
        area_m2=area,
        files_count=files_count,
        n_disciplinas=n_disc,
    )
    return {"status": "ok", "job_id": job_id, **sugestao}


@app.post("/api/cronograma/{job_id}/generate")
async def generate_cronograma(job_id: str, payload: CronogramaPayload, request: Request):
    _require_project_owner(request, job_id)
    """Gera cronograma FÍSICO (prazo e avanço, sem valores) a partir do quantitativo.

    Devolve JSON com fases + Gantt + curva S, pronto pra renderizar no
    frontend. Não persiste (recalcula a cada chamada — UX live).

    NÃO precifica. Distribui esforço no tempo seguindo sequenciamento
    construtivo padrão BR (16 etapas — convenção de mercado; a NBR 16636 foi
    citada aqui no passado por engano: ela trata de fases de PROJETO, não de obra).
    """
    # 1. Busca os items do projeto via mesma RPC do get_project_items
    import urllib.request, json as _json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
        body = _json.dumps({"p_job_id": job_id}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        items = _json.loads(resp.read().decode('utf-8'))
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar itens do projeto: {e}")

    if not items:
        raise HTTPException(404, "Projeto sem itens — gere a planilha primeiro.")

    # 2. Valida inputs
    try:
        from datetime import date as _date
        _date.fromisoformat(payload.data_inicio)
    except Exception:
        raise HTTPException(400, "data_inicio deve estar no formato YYYY-MM-DD")

    if payload.duracao_meses < 1 or payload.duracao_meses > 60:
        raise HTTPException(400, "duracao_meses deve estar entre 1 e 60")

    # 3. Gera cronograma
    try:
        from cronograma import gerar_cronograma
        resultado = gerar_cronograma(
            items=items,
            data_inicio=payload.data_inicio,
            duracao_meses=payload.duracao_meses,
        )
        return {"status": "ok", "job_id": job_id, **resultado}
    except Exception as e:
        import traceback
        print(f"[cronograma] erro: {e}")
        print(traceback.format_exc())
        raise HTTPException(500, f"Erro ao gerar cronograma: {e}")


def _get_branding_context(job_id: str, request=None) -> dict:
    """Padrão de co-branding pra TODAS as exportações (PDF/PPTX/XLSX).

    Centraliza a busca de:
    - project_name (sempre amigável, nunca o hash do job_id)
    - architect_name (do projeto OU company do profile)
    - client_name (project_clients.client_name + client_company)
    - logo_url + logo_local_path (baixa pra arquivo temp pra embedar)
    - brand_color (hex; default = indigo AI.arq se vazio)
    - company (nome do escritório)

    `request` é opcional — se passado, repassa o JWT do user pras queries
    REST (respeita RLS). Sem request, usa anon key (caso de uso interno
    quando RLS já foi validada upstream).

    Retorna dict pronto pra passar pra funções de export.
    """
    import urllib.request as _urlreq, json as _json
    import tempfile

    ctx = {
        'job_id': job_id,
        'project_name': '',
        'architect_name': '',
        'client_name': '',
        'company': '',
        'logo_url': '',
        'logo_local_path': None,
        'brand_color': '#4F46E5',  # default indigo AI.arq
    }

    project_user_id = ''

    def _fetch(path):
        """Helper local: usa o JWT do request se houver, senão service_role
        (bypassa RLS) — caso de uso interno quando branding é montado fora de
        um endpoint autenticado."""
        if request is not None:
            _, data = _supa_rest_as_user(request, "GET", path, timeout=8)
            return data or []
        # fallback service_role (bypassa RLS)
        try:
            url = f"{SUPABASE_URL}/rest/v1{path if path.startswith('/') else '/' + path}"
            req_ = _urlreq.Request(url, method="GET")
            req_.add_header("apikey", SUPABASE_KEY)
            req_.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            resp_ = _urlreq.urlopen(req_, timeout=8)
            return _json.loads(resp_.read().decode("utf-8"))
        except Exception:
            return []

    # 1. projects: nome amigável + arquiteto
    try:
        pj_data = _fetch(f"/projects?job_id=eq.{job_id}"
                          "&select=project_name,user_name,user_id")
        if pj_data:
            ctx['project_name'] = (pj_data[0].get('project_name') or '').strip()
            ctx['architect_name'] = (pj_data[0].get('user_name') or '').strip()
            project_user_id = pj_data[0].get('user_id') or ''
    except Exception as e:
        print(f"[branding] erro projects: {e}")

    # 2. project_clients: cliente final
    try:
        cl_data = _fetch(f"/project_clients?job_id=eq.{job_id}"
                          "&select=client_name,client_company")
        if cl_data:
            cn = (cl_data[0].get('client_name') or '').strip()
            cc = (cl_data[0].get('client_company') or '').strip()
            if cn and cc:
                ctx['client_name'] = f"{cn} ({cc})"
            else:
                ctx['client_name'] = cn or cc
    except Exception as e:
        print(f"[branding] erro project_clients: {e}")

    # 3. profiles: logo + cor + company
    if project_user_id:
        try:
            pr_data = _fetch(f"/profiles?user_id=eq.{project_user_id}"
                              "&select=logo_url,company,company_brand_color")
            if pr_data:
                ctx['logo_url'] = (pr_data[0].get('logo_url') or '').strip()
                ctx['company'] = (pr_data[0].get('company') or '').strip()
                bc = (pr_data[0].get('company_brand_color') or '').strip()
                if bc and bc.startswith('#') and len(bc) in (4, 7):
                    ctx['brand_color'] = bc
                if not ctx['architect_name'] and ctx['company']:
                    ctx['architect_name'] = ctx['company']
        except Exception as e:
            print(f"[branding] erro profiles: {e}")

    # 4. Fallback project_name
    if not ctx['project_name']:
        ctx['project_name'] = 'Projeto sem nome'

    # 5. Baixa logo pra arquivo temp (se houver).
    # Guarda anti-SSRF (fix 2026-07-22): logo_url vem do perfil do usuário e
    # poderia apontar pra arquivo local/rede interna. Só baixa de host público.
    if ctx['logo_url'] and _url_is_safe_public(ctx['logo_url']):
        try:
            req = urllib.request.Request(ctx['logo_url'], method="GET",
                                          headers={'User-Agent': 'AI.arq/1.0'})
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = resp.read(20 * 1024 * 1024)  # cap 20MB
            ext = '.png'
            if '.jpg' in ctx['logo_url'].lower() or '.jpeg' in ctx['logo_url'].lower():
                ext = '.jpg'
            tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
            tmp.write(data)
            tmp.close()
            ctx['logo_local_path'] = tmp.name
        except Exception as e:
            print(f"[branding] erro download logo: {e}")
            ctx['logo_local_path'] = None

    return ctx


def _supabase_get_cronograma(job_id: str, request=None) -> Optional[dict]:
    """Busca cronograma salvo do job. Retorna None se não existir.

    `request` opcional repassa JWT pra respeitar RLS quando chamado por endpoint."""
    import urllib.request, json as _json
    try:
        if request is not None:
            _, rows = _supa_rest_as_user(
                request, "GET",
                f"/cronogramas?job_id=eq.{job_id}&select=*&limit=1",
                timeout=10,
            )
            rows = rows or []
        else:
            url = (f"{SUPABASE_URL}/rest/v1/cronogramas"
                   f"?job_id=eq.{job_id}&select=*&limit=1")
            req = urllib.request.Request(url, method='GET')
            req.add_header('apikey', SUPABASE_KEY)
            req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
            req.add_header('Accept', 'application/json')
            resp = urllib.request.urlopen(req, timeout=10)
            rows = _json.loads(resp.read().decode('utf-8'))
        return rows[0] if rows else None
    except Exception as e:
        print(f"[cronograma get] erro: {e}")
        return None


def _supabase_upsert_cronograma(job_id: str, data: dict, request=None) -> bool:
    """Insere ou atualiza cronograma. Usa Prefer: resolution=merge-duplicates.

    `request` opcional repassa JWT pra respeitar RLS."""
    import urllib.request, json as _json
    # Carimba de qual quantitativo este cronograma nasceu. Se o cliente mexer
    # num item depois, /api/projeto/{job}/coerencia acusa a diferença.
    payload = {"job_id": job_id, "itens_assinatura": _assinatura_atual(job_id), **data}
    try:
        if request is not None:
            st, _ = _supa_rest_as_user(
                request, "POST",
                "/cronogramas",
                body=payload,
                prefer="resolution=merge-duplicates,return=minimal",
                timeout=10,
            )
            return st < 400
        url = f"{SUPABASE_URL}/rest/v1/cronogramas"
        body = _json.dumps(payload, default=str).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        req.add_header('Prefer', 'resolution=merge-duplicates,return=minimal')
        urllib.request.urlopen(req, timeout=10)
        return True
    except Exception as e:
        print(f"[cronograma upsert] erro: {e}")
        return False


@app.get("/api/cronograma/{job_id}")
async def get_cronograma(job_id: str, request: Request):
    """Retorna cronograma salvo do projeto (config + fases custom se houver).

    Se não houver salvo, devolve sugestão de duração + flag indicando que
    cliente ainda não gerou cronograma. Frontend usa essa info pra decidir
    se mostra inputs (gerar primeiro) ou já renderiza o salvo.
    """
    _require_project_owner(request, job_id)
    saved = _supabase_get_cronograma(job_id, request=request)
    if not saved:
        return {"status": "empty", "job_id": job_id, "saved": None}
    return {"status": "ok", "job_id": job_id, "saved": saved}


@app.post("/api/cronograma/{job_id}/save")
async def save_cronograma(job_id: str, payload: CronogramaSavePayload, request: Request):
    """Salva (upsert) cronograma com config + fases editadas."""
    _require_project_owner(request, job_id)
    # Valida data
    try:
        from datetime import date as _date
        _date.fromisoformat(payload.data_inicio)
    except Exception:
        raise HTTPException(400, "data_inicio deve estar no formato YYYY-MM-DD")
    if payload.duracao_meses < 1 or payload.duracao_meses > 60:
        raise HTTPException(400, "duracao_meses entre 1 e 60")

    data = {
        "data_inicio": payload.data_inicio,
        "duracao_meses": payload.duracao_meses,
        "k_sigmoid": payload.k_sigmoid or 10,
        "fases_custom": payload.fases_custom or [],
        "notas": payload.notas or "",
        "updated_at": "now()",
    }
    ok = _supabase_upsert_cronograma(job_id, data, request=request)
    if not ok:
        raise HTTPException(500, "Erro ao salvar no banco")
    return {"status": "ok", "job_id": job_id}


@app.post("/api/cronograma/{job_id}/preview")
async def preview_cronograma(job_id: str, payload: CronogramaSavePayload, request: Request):
    """Recalcula o cronograma a partir das fases editadas SEM salvar.

    Existe porque editar fases na tela mudava o Gantt mas deixava "Curva S",
    "Distribuição por mês" e "Caminho crítico" com os números de ANTES da edição:
    eles são calculados no servidor (curva_s / meses / matriz_pct) e o navegador
    só re-desenhava os dados velhos. Refazer essa conta em JS seria duplicar a
    matemática do cronograma.py e abrir espaço pra os dois lados divergirem —
    então a tela pergunta pra mesma função que gera o PDF/PPTX.
    """
    _require_project_owner(request, job_id)
    if not payload.fases_custom:
        raise HTTPException(400, "fases_custom vazio")
    try:
        from datetime import date as _date
        _date.fromisoformat(payload.data_inicio)
    except Exception:
        raise HTTPException(400, "data_inicio deve estar no formato YYYY-MM-DD")
    try:
        from cronograma import gerar_cronograma_de_fases_custom
        cron = gerar_cronograma_de_fases_custom(
            fases_custom=payload.fases_custom,
            data_inicio=payload.data_inicio,
            duracao_meses=payload.duracao_meses,
        )
    except Exception as e:
        raise HTTPException(500, f"Erro ao recalcular cronograma: {e}")
    return {"status": "ok", "job_id": job_id, **cron}


@app.get("/api/cronograma/{job_id}/full")
async def get_cronograma_full(job_id: str, request: Request):
    """Retorna o cronograma já renderizado (Gantt+CurvaS+Matriz+PPC) pra
    abrir a página sem precisar clicar em 'Gerar'.

    Se o cliente já salvou (com ou sem fases_custom), reusa a config salva.
    Se não há salvo, retorna 404 — frontend mostra o form pra gerar pela
    primeira vez.
    """
    _require_project_owner(request, job_id)
    saved = _supabase_get_cronograma(job_id)
    if not saved:
        raise HTTPException(404, "Cronograma ainda não gerado para este projeto")
    try:
        cron, _branding = _build_cronograma_for_export(job_id, request=request)
        return cron
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[cronograma full] erro: {e}")
        print(traceback.format_exc())
        raise HTTPException(500, f"Erro ao montar cronograma: {e}")


def _build_cronograma_for_export(job_id: str, request=None) -> tuple:
    """Monta o JSON do cronograma usando fases_custom se houver, senão gera
    automaticamente. Retorna (cronograma_dict, branding_context).

    `request` é opcional — quando passado, repassa o JWT do user pras
    leituras REST (RLS). Sem request, usa anon key (compat legado).
    """
    import urllib.request, json as _json
    saved = _supabase_get_cronograma(job_id, request=request)

    # Branding co-branded (nome projeto + cliente + logo + cor + arquiteto)
    branding = _get_branding_context(job_id, request=request)

    from cronograma import gerar_cronograma, gerar_cronograma_de_fases_custom

    if saved and saved.get('fases_custom'):
        # Usa fases editadas pelo cliente
        cron = gerar_cronograma_de_fases_custom(
            fases_custom=saved['fases_custom'],
            data_inicio=str(saved['data_inicio']),
            duracao_meses=saved['duracao_meses'],
        )
    else:
        # Gera automaticamente a partir dos items
        try:
            url = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
            body = _json.dumps({"p_job_id": job_id}).encode('utf-8')
            req = urllib.request.Request(url, data=body, method='POST')
            req.add_header('apikey', SUPABASE_KEY)
            req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
            req.add_header('Content-Type', 'application/json')
            resp = urllib.request.urlopen(req, timeout=15)
            items = _json.loads(resp.read().decode('utf-8'))
        except Exception as e:
            raise HTTPException(500, f"Erro ao buscar items: {e}")

        if not items:
            raise HTTPException(404, "Projeto sem itens")

        # Se há config salva mas sem fases_custom, usa config; senão default
        if saved:
            data_inicio = str(saved['data_inicio'])
            duracao = saved['duracao_meses']
        else:
            from datetime import date as _date, timedelta as _td
            data_inicio = (_date.today() + _td(days=30)).isoformat()
            duracao = 6
        cron = gerar_cronograma(items, data_inicio, duracao)

    return cron, branding


def _slug_filename(name: str) -> str:
    """Converte nome do projeto pra filename amigável (sem acentos/especiais)."""
    import re, unicodedata
    s = unicodedata.normalize('NFKD', name)
    s = ''.join(c for c in s if not unicodedata.combining(c))
    s = re.sub(r'[^a-zA-Z0-9_-]+', '_', s).strip('_')
    return s[:60] or 'cronograma'


# Templates de cronograma disponiveis (seletor no cronograma.html). Default escuro.
_CRONO_TEMPLATES = {"escuro", "blueprint", "claro", "editorial", "bold"}


async def _cronograma_preview_png_impl(job_id: str, request: Request,
                                       template: str = "", accent: str = ""):
    """Prévia do modelo como IMAGEM (PNG da 1ª página).

    Por que imagem e não o PDF embutido (02/08): o Chrome pode estar
    configurado pra BAIXAR PDF em vez de exibir (aconteceu no PC do Pedro) e
    celular quase nunca renderiza PDF em iframe — nos dois casos a prévia
    virava um cartão "Abrir". PNG aparece em qualquer lugar.
    Reusa o mesmo render do download, então a prévia continua sendo o
    arquivo real; só muda o formato de exibição."""
    _require_project_owner(request, job_id)
    import tempfile
    tmpl = (template or "").strip().lower()
    if tmpl not in _CRONO_TEMPLATES:
        tmpl = "escuro"
    try:
        cron, branding = _build_cronograma_for_export(job_id, request=request)
        from cronograma_render import render_pdf_bytes, render_png_paginas
        pdf = render_pdf_bytes(cron, branding, tmpl, (accent or "").strip() or None)
        if not pdf:
            raise RuntimeError("PDF vazio")
        pngs = render_png_paginas(pdf, scale=1.6)
        if not pngs:
            raise RuntimeError("rasterização falhou")
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp.write(pngs[0])
        tmp.close()
        return FileResponse(tmp.name, media_type="image/png",
                            headers={"Cache-Control": "no-store"})
    except Exception as e:
        import traceback
        print(f"[crono preview png] {tmpl}: {e}")
        print(traceback.format_exc())
        raise HTTPException(500, f"Erro ao gerar prévia: {e}")


@app.get("/api/cronograma/{job_id}/preview.png")
async def cronograma_preview_png(job_id: str, request: Request,
                                 template: str = "", accent: str = ""):
    return await _cronograma_preview_png_impl(job_id, request, template, accent)


@app.get("/api/cronograma/{job_id}/export/pdf")
async def export_cronograma_pdf(job_id: str, request: Request,
                                template: str = "", accent: str = ""):
    """Exporta cronograma como PDF co-branded. Usa os novos templates (WeasyPrint,
    5 direcoes, cor da marca); se falhar, cai no gerador antigo (reportlab)."""
    _require_project_owner(request, job_id)
    import tempfile
    from fastapi.responses import FileResponse
    cron, branding = _build_cronograma_for_export(job_id, request=request)
    tmpl = (template or "").strip().lower()
    if tmpl not in _CRONO_TEMPLATES:
        tmpl = "escuro"
    acc = (accent or "").strip() or None
    fname = f"cronograma_{_slug_filename(branding['project_name'])}.pdf"
    # 1) Novo render HTML->PDF (fiel aos 5 templates do handoff)
    try:
        from cronograma_render import render_pdf_bytes
        pdf = render_pdf_bytes(cron, branding, tmpl, acc)
        if pdf:
            tmp = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
            tmp.write(pdf)
            tmp.close()
            return FileResponse(tmp.name, media_type='application/pdf', filename=fname)
        raise RuntimeError("render_pdf_bytes vazio")
    except Exception as e:
        import traceback
        print(f"[export pdf] novo render falhou ({tmpl}), fallback reportlab: {e}")
        print(traceback.format_exc())
    # 2) Fallback: gerador antigo (reportlab)
    try:
        from cronograma_export import exportar_pdf
        tmp = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
        tmp.close()
        exportar_pdf(cron, tmp.name, branding=branding)
        return FileResponse(tmp.name, media_type='application/pdf', filename=fname)
    except Exception as e:
        import traceback
        print(f"[export pdf] erro: {e}")
        print(traceback.format_exc())
        raise HTTPException(500, f"Erro ao gerar PDF: {e}")


@app.get("/api/cronograma/{job_id}/export/pptx")
async def export_cronograma_pptx(job_id: str, request: Request,
                                 template: str = "", accent: str = ""):
    """Exporta cronograma como PPTX (5 slides). Novo: renderiza o PDF dos templates
    e insere 1 imagem full-bleed por slide (A4 paisagem). Fallback: gerador antigo."""
    _require_project_owner(request, job_id)
    import tempfile
    from fastapi.responses import FileResponse
    cron, branding = _build_cronograma_for_export(job_id, request=request)
    tmpl = (template or "").strip().lower()
    if tmpl not in _CRONO_TEMPLATES:
        tmpl = "escuro"
    acc = (accent or "").strip() or None
    fname = f"cronograma_{_slug_filename(branding['project_name'])}.pptx"
    _pptx_mime = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    # 1) Novo: PDF -> PNGs -> slides full-bleed
    try:
        from cronograma_render import render_pdf_bytes, render_png_paginas
        pdf = render_pdf_bytes(cron, branding, tmpl, acc)
        pngs = render_png_paginas(pdf) if pdf else []
        if pngs:
            from pptx import Presentation
            from pptx.util import Inches
            import io as _io
            prs = Presentation()
            prs.slide_width = Inches(11.69)   # A4 paisagem — mesma proporcao das paginas
            prs.slide_height = Inches(8.27)
            blank = prs.slide_layouts[6]
            for png in pngs:
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(_io.BytesIO(png), 0, 0,
                                         width=prs.slide_width, height=prs.slide_height)
            tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            tmp.close()
            prs.save(tmp.name)
            return FileResponse(tmp.name, media_type=_pptx_mime, filename=fname)
        raise RuntimeError("render_png_paginas vazio")
    except Exception as e:
        import traceback
        print(f"[export pptx] novo render falhou ({tmpl}), fallback: {e}")
        print(traceback.format_exc())
    # 2) Fallback: gerador antigo
    try:
        from cronograma_export import exportar_pptx
        tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        tmp.close()
        exportar_pptx(cron, tmp.name, branding=branding)
        return FileResponse(tmp.name, media_type=_pptx_mime, filename=fname)
    except Exception as e:
        import traceback
        print(f"[export pptx] erro: {e}")
        print(traceback.format_exc())
        raise HTTPException(500, f"Erro ao gerar PPTX: {e}")


# (endpoint temporário /api/debug/cronograma-sample removido após validar o
#  WeasyPrint no Render em 2026-07-14 — os 5 templates renderizaram OK)


# (endpoint temporário /api/debug/ig-insights removido após a análise de posts
#  de 2026-07-14 — dados puxados da Meta Graph API. Se virar painel fixo, criar
#  como rota admin autenticada + voltar a popular instagram_post_insights.)


class ReviewPayload(BaseModel):
    action: str                              # 'approve' | 'reject' | 'edit'
    edits: Optional[dict] = None             # {description?, unit?, quantity?, discipline?, observations?}
    comment: Optional[str] = ""
    reviewed_by: Optional[str] = ""


@app.post("/api/items/{job_id}/review/{item_id}")
async def submit_item_review(job_id: str, item_id: str, payload: ReviewPayload, request: Request):
    _require_project_owner(request, job_id)
    """Registra uma revisão pra um item específico.
    Se action='edit', também aplica os edits à row em project_items.
    Insere sempre uma linha em item_reviews pra histórico/aprendizado."""
    import urllib.request, urllib.error, json

    action = (payload.action or "").strip().lower()
    if action not in ("approve", "reject", "edit"):
        raise HTTPException(400, "action inválida (use approve/reject/edit)")

    # 1) Log da revisão — com DEDUPE de clique repetido (01/08/2026: o mesmo
    # item aparecia aprovado 6× na tabela; duplicata infla o painel e não é
    # sinal novo). Mesmo action pro mesmo item = atualiza, não insere.
    review_row = {
        "job_id": job_id,
        "item_id": item_id,
        "action": action,
        "edits": payload.edits or None,
        "comment": payload.comment or "",
        "reviewed_by": payload.reviewed_by or "",
    }
    _ja_tem = False
    try:
        import urllib.parse as _upq
        _q = (f"item_reviews?job_id=eq.{_upq.quote(job_id)}"
              f"&item_id=eq.{_upq.quote(item_id)}&action=eq.{action}&select=id&limit=1")
        # 🪤 (status, dados): bool da TUPLA era sempre True — o dedupe passou a
        # engolir TODA revisao nova. Regressao de 01/08, viva por ~2 horas.
        _st_rv, _rows_rv = _supa_rest_service("GET", _q)
        _ja_tem = bool(_rows_rv) and isinstance(_rows_rv, list) and len(_rows_rv) > 0
    except Exception:
        pass   # na dúvida, insere (comportamento antigo)
    if not _ja_tem:
        _supabase_insert("item_reviews", review_row)

    # 2) Aplicar edits à row original (se action=edit)
    if action == "edit" and payload.edits:
        safe_edits = {}
        allowed = {"description", "unit", "quantity", "discipline", "observations"}
        for k, v in payload.edits.items():
            if k in allowed and v is not None:
                safe_edits[k] = v
        if safe_edits:
            try:
                # Segurança: amarra o item ao job_id JÁ validado como do dono
                # (_require_project_owner acima). Sem o &job_id, um dono passaria
                # o próprio job_id + o item_id de OUTRO projeto e editaria item
                # alheio (IDOR, escrita cross-tenant com service_role).
                url = f"{SUPABASE_URL}/rest/v1/project_items?id=eq.{item_id}&job_id=eq.{job_id}"
                body = json.dumps(safe_edits).encode('utf-8')
                req = urllib.request.Request(url, data=body, method='PATCH')
                req.add_header('apikey', SUPABASE_KEY)
                req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
                req.add_header('Content-Type', 'application/json')
                req.add_header('Prefer', 'return=minimal')
                urllib.request.urlopen(req, timeout=15)
            except Exception as e:
                _supa_log(f"REVIEW edit item={item_id} ERR {e}")

    # 3) Se rejeitado, deleta row do item (mantém review pra histórico)
    if action == "reject":
        try:
            # Segurança (IDOR): só apaga item que pertence ao job_id do dono.
            url = f"{SUPABASE_URL}/rest/v1/project_items?id=eq.{item_id}&job_id=eq.{job_id}"
            req = urllib.request.Request(url, method='DELETE')
            req.add_header('apikey', SUPABASE_KEY)
            req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
            req.add_header('Prefer', 'return=minimal')
            urllib.request.urlopen(req, timeout=15)
        except Exception as e:
            _supa_log(f"REVIEW reject item={item_id} ERR {e}")

    # 4) O quantitativo mudou → cronograma e memorial salvos podem ter ficado
    # velhos. Derruba o cache pra próxima leitura de coerência ser a real.
    if action in ("edit", "reject"):
        _assinatura_invalidar(job_id)

    return {"status": "ok", "action": action}


@app.post("/api/items/{job_id}/finalize")
async def finalize_review(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Regenera o .xlsx com as revisões aplicadas e sobe pro Storage.
    Busca items atuais do Supabase (já com edits/rejects aplicados) e
    passa pro generate_spreadsheet. Retorna URL de download."""
    import urllib.request, urllib.error, json
    from models import BudgetItem, Confidence, ProjectData
    from spreadsheet import generate_spreadsheet

    # 1) Buscar project metadata
    try:
        _, projects = _supa_rest_as_user(
            request, "GET",
            f"/projects?job_id=eq.{job_id}&select=*",
            timeout=15,
        )
        projects = projects or []
        if not projects:
            raise HTTPException(404, "Projeto não encontrado")
        proj = projects[0]
    except HTTPException:
        raise
    except urllib.error.HTTPError as e:
        raise HTTPException(500, f"Erro Supabase: {e}")

    # 2) Buscar items atuais (já revisados) — RPC SECURITY DEFINER, anon ok
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
        body = json.dumps({"p_job_id": job_id}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        rows = json.loads(resp.read().decode('utf-8'))
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar itens: {e}")

    # 3) Reconstituir ProjectData + BudgetItems
    pd = ProjectData(
        name=proj.get("project_name", "") or "Projeto",
        total_area=proj.get("total_area") or 0,
        layout_area=proj.get("layout_area") or 0,
    )
    items = []
    for r in rows:
        try:
            items.append(BudgetItem(
                item_num=r.get("item_num", "") or "",
                description=r.get("description", "") or "",
                unit=r.get("unit", "vb") or "vb",
                quantity=float(r.get("quantity") or 0),
                observations=r.get("observations", "") or "",
                ref_sheet=r.get("ref_sheet", "") or "",
                confidence=Confidence(r.get("confidence", "estimado") or "estimado"),
                discipline=r.get("discipline", "Complementares") or "Complementares",
            ))
        except Exception:
            continue

    # 4) Gerar xlsx revisado — também enriquece com matches TCPO + heurísticas
    try:
        from tcpo_matcher import match_item, get_insumos
        for it in items:
            try:
                ms = match_item(it.description, limit=3)
                if ms:
                    ms[0]['insumos'] = get_insumos(ms[0]['id'])
                    it.tcpo_matches = ms
            except Exception:
                pass
    except ImportError:
        pass

    typology = proj.get("typology") or "office"

    # Heurísticas de mercado (alertas agregados anônimos)
    try:
        from market_heuristics import check_item_anomaly
        for it in items:
            try:
                alertas = check_item_anomaly(it, typology=typology)
                if alertas:
                    sep = " | " if it.observations else ""
                    it.observations = (it.observations or "") + sep + " ".join(alertas)
            except Exception:
                pass
    except ImportError:
        pass

    work_dir = os.path.join(WORK_DIR, job_id)
    os.makedirs(work_dir, exist_ok=True)
    output_path = os.path.join(work_dir, f"orcamento_{job_id}_revisado.xlsx")
    generate_spreadsheet(pd, items, output_path, typology=typology)

    # 5) Subir pra Storage sobrescrevendo o antigo
    _storage_ok = _supabase_storage_upload(output_path, f"{job_id}.xlsx")
    # A planilha agora bate com os itens revisados — tira o aviso de velha.
    _carimbar_planilha(job_id)

    return {
        "status": "ok",
        "job_id": job_id,
        "items_count": len(items),
        "download_url": f"/api/download/{job_id}",
        "storage_uploaded": _storage_ok,
    }


class InformAreaPayload(BaseModel):
    area: float


@app.post("/api/project/{job_id}/inform-area")
async def inform_project_area(job_id: str, payload: InformAreaPayload, request: Request):
    """Cliente informa a metragem DEPOIS do processamento, pra completar itens de
    área que ficaram em branco porque a planta não tinha cota de área (caso layout,
    ex. Catarina). NÃO reprocessa do zero (sem custo de IA, sem re-erro): reaproveita
    os itens já salvos, preenche as superfícies horizontais (piso/forro/laje) com a
    área informada — SEMPRE estimado (laranja) e rotulado 'informado por você, não
    medido' (regra dura nº1, via _apply_area_honesty) — e regenera a planilha in-place
    (mesmo projeto, sem criar cópia). Itens que não escalam com piso (pintura de
    parede, rodapé) NÃO são preenchidos."""
    _require_project_owner(request, job_id)
    import urllib.request, urllib.error, json
    from models import BudgetItem, Confidence, ProjectData
    from spreadsheet import generate_spreadsheet

    # 1) Validar área informada
    try:
        area = round(float(payload.area or 0), 2)
    except (TypeError, ValueError):
        area = 0
    if area <= 0 or area > 1_000_000:
        raise HTTPException(400, "Informe uma área válida em m² (maior que 0).")

    # 2) Projeto + itens atuais (mesmo padrão do finalize_review)
    try:
        _, projects = _supa_rest_as_user(
            request, "GET", f"/projects?job_id=eq.{job_id}&select=*", timeout=15)
        projects = projects or []
        if not projects:
            raise HTTPException(404, "Projeto não encontrado")
        proj = projects[0]
    except HTTPException:
        raise
    except urllib.error.HTTPError as e:
        raise HTTPException(500, f"Erro Supabase: {e}")

    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_project_items"
        body = json.dumps({"p_job_id": job_id}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        rows = json.loads(resp.read().decode('utf-8'))
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar itens: {e}")

    # 3) Reconstituir ProjectData (área = informada) + BudgetItems na MESMA ordem
    #    (list_project_items vem ordenado por sort_order → re-persistir preserva ordem)
    pd = ProjectData(
        name=proj.get("project_name", "") or "Projeto",
        total_area=area,
        layout_area=proj.get("layout_area") or 0,
    )
    try:
        pd.total_area_source = "informado"
    except Exception:
        pass

    items = []
    for r in rows:
        try:
            items.append(BudgetItem(
                item_num=r.get("item_num", "") or "",
                description=r.get("description", "") or "",
                unit=r.get("unit", "vb") or "vb",
                quantity=float(r.get("quantity") or 0),
                observations=r.get("observations", "") or "",
                ref_sheet=r.get("ref_sheet", "") or "",
                confidence=Confidence(r.get("confidence", "estimado") or "estimado"),
                discipline=r.get("discipline", "Complementares") or "Complementares",
            ))
        except Exception:
            continue

    # 4) Preenche piso/forro/laje com a área informada (estimado, rotulado) — mesma
    #    regra do motor. Itens em branco não-horizontais seguem em branco.
    filled, _ = _apply_area_honesty(items, area, "informado")

    # 5) Enriquecimento (TCPO + heurísticas) igual ao finalize + gerar xlsx in-place
    typology = proj.get("typology") or "office"
    try:
        from tcpo_matcher import match_item, get_insumos
        for it in items:
            try:
                ms = match_item(it.description, limit=3)
                if ms:
                    ms[0]['insumos'] = get_insumos(ms[0]['id'])
                    it.tcpo_matches = ms
            except Exception:
                pass
    except ImportError:
        pass
    try:
        from market_heuristics import check_item_anomaly
        for it in items:
            try:
                alertas = check_item_anomaly(it, typology=typology)
                if alertas:
                    sep = " | " if it.observations else ""
                    it.observations = (it.observations or "") + sep + " ".join(alertas)
            except Exception:
                pass
    except ImportError:
        pass

    work_dir = os.path.join(WORK_DIR, job_id)
    os.makedirs(work_dir, exist_ok=True)
    output_path = os.path.join(work_dir, f"orcamento_{job_id}.xlsx")
    generate_spreadsheet(pd, items, output_path, typology=typology)
    _storage_ok = _supabase_storage_upload(output_path, f"{job_id}.xlsx")

    # 6) Persistir itens atualizados (revisão/planilha na tela refletem o preenchimento)
    _persist_items_to_supabase(job_id, items)
    _carimbar_planilha(job_id)

    # 7) Atualizar projeto: área informada + warning honesto (troca avisos antigos
    #    de "informe a área" pra não duplicar)
    _warn = (f"Área total de {area:.0f} m² INFORMADA POR VOCÊ (a planta não trazia cota "
             f"pra medir). Preenchemos os itens de piso/forro/laje com essa base — "
             f"confira antes de orçar. Pra medir de verdade, envie o DXF.")
    _existing = proj.get("warnings") or []
    if not isinstance(_existing, list):
        _existing = []
    _existing = [w for w in _existing
                 if "não trazia" not in str(w) and "informe a área" not in str(w).lower()
                 and "informe a metragem" not in str(w).lower()]
    _supabase_update("projects", "job_id", job_id, {
        "total_area": area,
        "user_total_area": area,
        "warnings": _existing + [_warn],
    })

    return {
        "status": "ok",
        "job_id": job_id,
        "area": area,
        "filled_count": filled,
        "items_count": len(items),
        "download_url": f"/api/download/{job_id}",
    }


# ═══════════════════════════════════════════════════════════════
#  STATUS EDITÁVEL + NOTAS POR ITEM
# ═══════════════════════════════════════════════════════════════

class StatusPayload(BaseModel):
    user_status: str


@app.post("/api/project/{job_id}/status")
async def update_project_user_status(job_id: str, payload: StatusPayload, request: Request):
    """Atualiza o status editável do projeto (em_analise, enviado_cliente,
    aprovado, fechado, arquivado)."""
    _require_project_owner(request, job_id)
    import urllib.request, urllib.error, json
    valid = {"em_analise", "enviado_cliente", "aprovado", "fechado", "arquivado"}
    if payload.user_status not in valid:
        raise HTTPException(400, f"status inválido. Aceito: {valid}")
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/update_user_status"
        body = json.dumps({
            "p_job_id": job_id,
            "p_user_status": payload.user_status,
        }).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        return {"status": "ok", "user_status": payload.user_status}
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")


class NotePayload(BaseModel):
    note: str
    author: Optional[str] = ""


@app.post("/api/items/{job_id}/note/{item_id}")
async def save_item_note(job_id: str, item_id: str, payload: NotePayload, request: Request):
    """Salva (upsert) nota de um item. Vazio deleta a nota."""
    _require_project_owner(request, job_id)
    import urllib.request, urllib.error, json
    text = (payload.note or "").strip()
    try:
        if not text:
            # Delete
            _supa_rest_as_user(
                request, "DELETE",
                f"/item_notes?item_id=eq.{item_id}&job_id=eq.{job_id}",
                timeout=10,
            )
            return {"status": "ok", "deleted": True}

        # Upsert: try update first, insert on 0 rows
        patch_body = {
            "note": text,
            "author": payload.author or "",
            "updated_at": datetime.utcnow().isoformat(),
        }
        _, rows = _supa_rest_as_user(
            request, "PATCH",
            f"/item_notes?item_id=eq.{item_id}&job_id=eq.{job_id}",
            body=patch_body, prefer="return=representation", timeout=10,
        )
        rows = rows or []
        if not rows:
            # Insert novo
            _supabase_insert("item_notes", {
                "job_id": job_id,
                "item_id": item_id,
                "note": text,
                "author": payload.author or "",
            })
        return {"status": "ok", "saved": True}
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")


@app.get("/api/items/{job_id}/notes")
async def list_job_notes(job_id: str, request: Request):
    _require_project_owner(request, job_id)
    """Lista todas as notas de itens de um job — restaura estado na revisão."""
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_item_notes"
        body = json.dumps({"p_job_id": job_id}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        notes = json.loads(resp.read().decode('utf-8'))
        # Retorna como mapa {item_id: {note, author, updated_at}}
        state = {n["item_id"]: {
            "note": n.get("note", ""),
            "author": n.get("author", ""),
            "updated_at": n.get("updated_at"),
        } for n in notes}
        return {"status": "ok", "notes": state}
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")


@app.get("/api/projects/by-user/{user_id}")
async def list_my_projects(user_id: str, request: Request):
    """Lista projetos de um usuário (pra tela 'Meus projetos').
    Usa RPC list_user_projects (SECURITY DEFINER).

    Auth: sempre exige JWT Bearer do próprio usuário (ou admin).
    28/07/2026: 'anonymous' deixou de ser livre — listava os 53 projetos órfãos
    do começo do beta, com nome de projeto de cliente real, sem login."""
    user = _get_user_from_request(request)
    if not user:
        raise HTTPException(401, "Autenticação requerida")
    is_admin = user.get("email", "").lower() == ADMIN_EMAIL
    if user_id == "anonymous" and not is_admin:
        raise HTTPException(403, "Lista indisponível")
    if user_id != "anonymous" and user.get("id") != user_id and not is_admin:
        raise HTTPException(403, "Só é possível listar seus próprios projetos")
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_user_projects"
        body = json.dumps({"p_user_id": user_id}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=15)
        projects = json.loads(resp.read().decode('utf-8'))
        return {"status": "ok", "user_id": user_id, "projects": projects, "count": len(projects)}
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar projetos: {str(e)}")


# ═══════════════════════════════════════════════════════════════
#  NPS (Net Promoter Score)
# ═══════════════════════════════════════════════════════════════

class NPSPayload(BaseModel):
    user_id: str
    score: int                               # 0-10
    comment: Optional[str] = ""
    context: Optional[str] = "manual"        # 'after_download' | 'after_review' | 'manual' | 'first_project'
    user_email: Optional[str] = ""
    user_name: Optional[str] = ""
    job_id: Optional[str] = ""


class NPSDetailedPayload(BaseModel):
    recommend: int                           # 0-10 (NPS clássico)
    stage_ratings: dict = {}                 # {"upload":1-5, "tempo":1-5, "precisao":1-5, "planilha":1-5}
    comment: Optional[str] = ""
    job_id: Optional[str] = ""


@app.post("/api/nps")
async def submit_nps(payload: NPSPayload):
    """Registra uma resposta NPS. Score 0-10 obrigatório, comentário opcional."""
    if payload.score < 0 or payload.score > 10:
        raise HTTPException(400, "score deve estar entre 0 e 10")
    if not payload.user_id:
        raise HTTPException(400, "user_id obrigatório")

    row = {
        "user_id": payload.user_id,
        "user_email": payload.user_email or "",
        "user_name": payload.user_name or "",
        "score": int(payload.score),
        "comment": (payload.comment or "")[:2000],
        "context": (payload.context or "manual")[:50],
        "job_id": payload.job_id or "",
    }
    ok = _supabase_insert("nps_responses", row)
    category = "promoter" if payload.score >= 9 else "passive" if payload.score >= 7 else "detractor"
    return {"status": "ok" if ok else "error", "category": category}


_NPS_STAGES = ("upload", "tempo", "precisao", "planilha")


@app.post("/api/nps/detailed")
async def submit_nps_detailed(payload: NPSDetailedPayload, request: Request):
    """Feedback detalhado: nota 0-10 (recomendaria) + nota 1-5 por etapa +
    comentário. Autenticado (JWT) — só grava no nome do próprio usuário."""
    user = _get_user_from_request(request)
    if not user or not user.get("id"):
        raise HTTPException(401, "Autenticação requerida")
    if payload.recommend < 0 or payload.recommend > 10:
        raise HTTPException(400, "recommend deve estar entre 0 e 10")
    stages = {}
    for k in _NPS_STAGES:
        v = (payload.stage_ratings or {}).get(k)
        if isinstance(v, (int, float)) and 1 <= int(v) <= 5:
            stages[k] = int(v)
    cat = "promoter" if payload.recommend >= 9 else "passive" if payload.recommend >= 7 else "detractor"
    # OBS: NÃO incluir 'category' — é coluna GERADA (o banco calcula do score).
    # Inserir valor nela faz o insert inteiro falhar.
    row = {
        "user_id": user["id"],
        "user_email": user.get("email", ""),
        "user_name": _name_from_auth(user["id"]),
        "score": int(payload.recommend),
        "comment": (payload.comment or "")[:2000],
        "context": "feedback_detailed",
        "job_id": payload.job_id or "",
        "stage_ratings": stages,
    }
    ok = _supabase_insert("nps_responses", row)
    return {"status": "ok" if ok else "error", "category": cat}


@app.get("/api/nps/check/{user_id}")
async def should_show_nps(user_id: str, request: Request):
    """Verifica se o usuário já respondeu NPS recentemente (últimos 60 dias).
    Frontend usa isso pra não mostrar o widget repetidamente.

    Auth: exige JWT do próprio usuário (ou admin) — evita checar/inferir NPS de
    terceiros (auditoria 2026-07-13). O frontend já chama com authFetch."""
    jwt_user = _get_user_from_request(request)
    if not jwt_user:
        raise HTTPException(401, "Autenticação requerida")
    if jwt_user.get("id") != user_id and jwt_user.get("email", "").lower() != ADMIN_EMAIL:
        raise HTTPException(403, "Só é possível checar o próprio NPS")
    import urllib.request, urllib.error, json
    try:
        url = (f"{SUPABASE_URL}/rest/v1/nps_responses"
               f"?user_id=eq.{user_id}&select=created_at"
               f"&order=created_at.desc&limit=1")
        req = urllib.request.Request(url, method='GET')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        resp = urllib.request.urlopen(req, timeout=10)
        rows = json.loads(resp.read().decode('utf-8'))
        if not rows:
            return {"should_show": True, "last_answered": None}
        last = rows[0].get("created_at")
        # 60 dias de cooldown
        from datetime import timedelta
        try:
            dt = datetime.fromisoformat(last.replace("Z", "+00:00"))
            days_since = (datetime.utcnow().replace(tzinfo=dt.tzinfo) - dt).days
            return {"should_show": days_since >= 60, "last_answered": last, "days_since": days_since}
        except Exception:
            return {"should_show": False, "last_answered": last}
    except Exception as e:
        return {"should_show": False, "error": str(e)}


@app.get("/api/admin/nps/summary")
async def admin_nps_summary(request: Request, days: int = 30):
    """Dashboard admin: resumo agregado de NPS."""
    _require_admin(request)
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/get_nps_summary"
        body = json.dumps({"p_days": days}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=10)
        rows = json.loads(resp.read().decode('utf-8'))
        return rows[0] if rows else {}
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")


@app.get("/api/admin/nps/responses")
async def admin_nps_responses(request: Request, limit: int = 50):
    """Dashboard admin: respostas recentes com comentários (insights qualitativos)."""
    _require_admin(request)
    import urllib.request, urllib.error, json
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_nps_responses"
        body = json.dumps({"p_limit": limit}).encode('utf-8')
        req = urllib.request.Request(url, data=body, method='POST')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        req.add_header('Content-Type', 'application/json')
        resp = urllib.request.urlopen(req, timeout=10)
        return {"responses": json.loads(resp.read().decode('utf-8'))}
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")


@app.get("/api/admin/nps/stages")
async def admin_nps_stages(request: Request):
    """Médias por etapa do feedback detalhado (upload, tempo, precisão, planilha).
    Lê as respostas e agrega em Python (volume pequeno)."""
    _require_admin(request)
    import urllib.request as _urs
    try:
        url = (f"{SUPABASE_URL}/rest/v1/nps_responses"
               f"?context=eq.feedback_detailed&select=stage_ratings,score,comment,user_name,created_at"
               f"&order=created_at.desc")
        req = _urs.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _json.loads(_urs.urlopen(req, timeout=15).read().decode("utf-8"))
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")
    stages = ("upload", "tempo", "precisao", "planilha")
    sums = {k: 0.0 for k in stages}
    counts = {k: 0 for k in stages}
    rec_sum, rec_n = 0.0, 0
    for r in rows:
        sr = r.get("stage_ratings") or {}
        for k in stages:
            v = sr.get(k)
            if isinstance(v, (int, float)):
                sums[k] += v
                counts[k] += 1
        sc = r.get("score")
        if isinstance(sc, (int, float)):
            rec_sum += sc
            rec_n += 1
    avgs = {k: (round(sums[k] / counts[k], 2) if counts[k] else None) for k in stages}
    return {
        "total": len(rows),
        "stages": avgs,
        "counts": counts,
        "recommend_avg": round(rec_sum / rec_n, 2) if rec_n else None,
    }


# ═══════════════════════════════════════════════════════════════
#  USAGE TRACKING (Painel de Atividade) — quem usa o quê, onde para
# ═══════════════════════════════════════════════════════════════

class TrackPayload(BaseModel):
    event: str
    user_id: Optional[str] = ""
    user_email: Optional[str] = ""
    job_id: Optional[str] = ""
    path: Optional[str] = ""
    meta: dict = {}


# Lista branca de eventos aceitos pelo /api/track. Endpoint é ABERTO (sem auth),
# então SÓ nomes conhecidos entram — mata XSS armazenado (nome de evento cru
# renderizado no admin) e evita poluição com eventos arbitrários. Ao criar
# evento novo no front, adicione aqui também.
_TRACK_ALLOWED = {
    "view_landing", "view_cadastro", "signup_done", "view_dashboard",
    "start_project", "open_project", "download_xlsx",
    "use_cronograma", "use_comparativo", "review_item", "review_finish",
}


@app.post("/api/track")
async def track_event(payload: TrackPayload):
    """Registra um evento de uso (best-effort, nunca falha pro cliente).
    Chamado pelo trackEvent() do front. Aberto (sem auth), MAS só aceita nomes
    de evento da allowlist e só as chaves de meta conhecidas (cid/type) —
    segurança: nada de HTML/JS arbitrário chega ao painel admin."""
    ev = (payload.event or "").strip()
    if ev not in _TRACK_ALLOWED:
        return {"status": "ignored"}
    # meta capado: só cid (id anônimo), type e src (origem first-touch), curtos.
    _meta = {}
    if isinstance(payload.meta, dict):
        _cid = str(payload.meta.get("cid") or "")[:40]
        _type = str(payload.meta.get("type") or "")[:40]
        _src = str(payload.meta.get("src") or "")[:40]
        if _cid:
            _meta["cid"] = _cid
        if _type:
            _meta["type"] = _type
        if _src:
            _meta["src"] = _src
    row = {
        "event": ev,  # já validado contra a allowlist
        "user_id": (payload.user_id or "")[:80],
        "user_email": (payload.user_email or "")[:200],
        "job_id": (payload.job_id or "")[:80],
        "path": (payload.path or "")[:200],
        "meta": _meta,
    }
    try:
        _supabase_insert("usage_events", row)
    except Exception:
        pass  # telemetria nunca quebra nada
    return {"status": "ok"}


@app.get("/api/admin/activity")
async def admin_activity(request: Request, days: int = 30, limit: int = 200):
    """Painel de Atividade: eventos recentes + agregados (por evento, por usuário,
    ativos 7/30d). Agrega em Python — volume pequeno."""
    _require_admin(request)
    import urllib.request as _urs
    from datetime import datetime, timedelta, timezone
    days = max(1, min(int(days or 30), 365))
    limit = max(1, min(int(limit or 200), 1000))
    since = (datetime.now(timezone.utc) - timedelta(days=days)).isoformat()
    since_url = since.replace("+00:00", "Z")
    cutoff_7d = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
    try:
        url = (f"{SUPABASE_URL}/rest/v1/usage_events"
               f"?created_at=gte.{since_url}"
               f"&select=event,user_email,user_id,job_id,path,meta,created_at"
               f"&order=created_at.desc&limit=5000")
        req = _urs.Request(url, method="GET")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        rows = _json.loads(_urs.urlopen(req, timeout=20).read().decode("utf-8"))
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")

    # ── Mistura o USO REAL do produto (tabela projects) ───────────────────
    # usage_events é opt-in (só quem aceitou o cookie de analytics), então
    # some gente que de fato USA o produto (ex.: quem ignorou o banner). A
    # tabela projects tem TODOS os uploads/processamentos — é a fonte completa
    # e é dado operacional do serviço (não cookie de navegação), então entra
    # sempre. Assim o painel para de mostrar quase só o admin.
    try:
        purl = (f"{SUPABASE_URL}/rest/v1/projects"
                f"?created_at=gte.{since_url}"
                f"&select=user_email,user_id,job_id,status,created_at,completed_at"
                f"&order=created_at.desc&limit=5000")
        preq = _urs.Request(purl, method="GET")
        preq.add_header("apikey", SUPABASE_KEY)
        preq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _projects = _json.loads(_urs.urlopen(preq, timeout=20).read().decode("utf-8"))
    except Exception:
        _projects = []
    # "Subiu projeto" passa a vir de projects (completo) — descarta o
    # start_project opt-in pra não contar o mesmo upload duas vezes.
    rows = [r for r in rows if (r.get("event") or "") != "start_project"]
    for _p in _projects:
        _base = {
            "user_email": (_p.get("user_email") or "").strip(),
            "user_id": _p.get("user_id") or "",
            "job_id": _p.get("job_id") or "",
            "path": "", "meta": {},
        }
        _ca = _p.get("created_at") or ""
        _st = (_p.get("status") or "").strip()
        _cp = _p.get("completed_at") or ""
        rows.append({**_base, "event": "start_project", "created_at": _ca})
        if _st == "done" and _cp:
            rows.append({**_base, "event": "project_done", "created_at": _cp})
        elif _st == "error":
            rows.append({**_base, "event": "project_error", "created_at": _cp or _ca})

    # ── Signups do AUTH (3ª fonte) ────────────────────────────────────────
    # usage_events e projects não cobrem quem CRIOU CONTA mas não fez mais nada
    # (cadastro incompleto: conta Google/email criada, perfil nunca completado) —
    # ficava INVISÍVEL na Atividade. Puxa os usuários do auth e injeta um evento
    # 'signup_created' por conta criada na janela. Best-effort — nunca derruba o
    # painel. (feedback Pedro 21/07: um incompleto de hoje não aparecia.)
    _since_dt = datetime.now(timezone.utc) - timedelta(days=days)
    try:
        aurl = f"{SUPABASE_URL}/auth/v1/admin/users?per_page=200&page=1"
        areq = _urs.Request(aurl, method="GET")
        areq.add_header("apikey", SUPABASE_KEY)
        areq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        _adata = _json.loads(_urs.urlopen(areq, timeout=20).read().decode("utf-8"))
        _ausers = _adata.get("users", []) if isinstance(_adata, dict) else (_adata or [])
    except Exception as _ae:
        print(f"[activity] auth users falhou (não crítico): {_ae}")
        _ausers = []
    for _au in _ausers:
        _acreated = _au.get("created_at") or ""
        try:
            _adt = datetime.fromisoformat(_acreated.replace("Z", "+00:00"))
            if _adt.tzinfo is None:
                _adt = _adt.replace(tzinfo=timezone.utc)
        except Exception:
            continue
        if _adt < _since_dt:
            continue
        rows.append({
            "event": "signup_created",
            "user_email": (_au.get("email") or "").strip(),
            "user_id": _au.get("id") or "",
            "job_id": "", "path": "", "meta": {},
            "created_at": _acreated,
        })

    rows.sort(key=lambda r: r.get("created_at") or "", reverse=True)

    # Tira as contas internas/do Pedro do painel: ele entra todo dia e dominaria
    # (poluia a Atividade). _email_eh_interno pega o email dele + aliases (+smoke).
    # Anonimos (sem email) continuam contando.
    rows = [r for r in rows if not _email_eh_interno(r.get("user_email") or "")]

    by_event, by_user = {}, {}
    seen_7d, seen_30d = set(), set()
    # Funil do topo por VISITANTE único (cid anônimo no meta): landing → cadastro
    # → completou. Responde "quanta gente visita e não converte".
    _FUNNEL = ("view_landing", "view_cadastro", "signup_done")
    funnel_cids = {k: set() for k in _FUNNEL}
    for r in rows:
        ev = r.get("event") or "?"
        by_event[ev] = by_event.get(ev, 0) + 1
        if ev in funnel_cids:
            _m = r.get("meta")
            _cid = _m.get("cid") if isinstance(_m, dict) else ""
            if _cid:
                funnel_cids[ev].add(_cid)
        email = (r.get("user_email") or "").strip() or (r.get("user_id") or "anonymous")
        u = by_user.setdefault(email, {"email": email, "events": 0,
                                       "last_seen": None, "first_seen": None, "kinds": {}})
        u["events"] += 1
        u["kinds"][ev] = u["kinds"].get(ev, 0) + 1
        ts = r.get("created_at") or ""
        if u["last_seen"] is None or ts > u["last_seen"]:
            u["last_seen"] = ts
        if u["first_seen"] is None or ts < u["first_seen"]:
            u["first_seen"] = ts
        if email and email != "anonymous":
            seen_30d.add(email)
            if ts >= cutoff_7d:
                seen_7d.add(email)
    users = sorted(by_user.values(), key=lambda x: x.get("last_seen") or "", reverse=True)
    return {
        "total_events": len(rows),
        "window_days": days,
        "active_7d": len(seen_7d),
        "active_30d": len(seen_30d),
        "by_event": by_event,
        "funnel": {k: len(funnel_cids[k]) for k in _FUNNEL},
        "users": users,
        "recent": rows[:limit],
    }


# ═══════════════════════════════════════════════════════════════
#  SERVIR PRANCHA (pra revisão inline abrir em nova aba)
# ═══════════════════════════════════════════════════════════════

def _find_prancha_file(job_id: str, ref: str) -> Optional[str]:
    """Dado um ref_sheet (que pode vir como nome exato, descrição da IA, ou
    concatenação 'filename (hint)'), encontra o filename real do PDF:

    1. Se `ref` termina em .pdf e existe no hot cache → usa direto.
    2. Se `ref` tem formato "filename.pdf (hint)" → extrai o filename antes do "(".
    3. Fuzzy match: lista PDFs do job no /tmp OU no Storage, escolhe o que
       tem maior substring overlap com `ref`.

    Retorna o filename encontrado ou None.
    """
    import urllib.request
    ref_clean = (ref or "").strip()
    if not ref_clean:
        return None

    # Remover sufixo "(hint da IA)" se houver
    if "(" in ref_clean and ref_clean.endswith(")"):
        ref_clean = ref_clean.split("(")[0].strip()

    # Caso direto: já é um filename .pdf válido
    if ref_clean.lower().endswith(".pdf"):
        candidate = os.path.basename(ref_clean)
        # Verificar se existe
        if os.path.exists(os.path.join(WORK_DIR, job_id, candidate)):
            return candidate
        # Senão, tenta no Storage (via list)
        try:
            url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
            import json as _j
            body = _j.dumps({"prefix": f"{job_id}/", "limit": 200}).encode("utf-8")
            req = urllib.request.Request(url, data=body, method="POST")
            req.add_header("apikey", SUPABASE_KEY)
            req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            req.add_header("Content-Type", "application/json")
            resp = urllib.request.urlopen(req, timeout=10)
            files = _j.loads(resp.read().decode("utf-8"))
            names = [f.get("name", "") for f in files]
            for n in names:
                # Storage keys são URL-encoded
                from urllib.parse import unquote
                if unquote(n) == candidate:
                    return candidate
        except Exception:
            pass

    # Fuzzy match: lista PDFs locais + Storage, acha melhor match
    candidates = set()
    local_dir = os.path.join(WORK_DIR, job_id)
    if os.path.isdir(local_dir):
        for fn in os.listdir(local_dir):
            if fn.lower().endswith(".pdf"):
                candidates.add(fn)
    try:
        from urllib.parse import unquote
        url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        import json as _j2
        body = _j2.dumps({"prefix": f"{job_id}/", "limit": 200}).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=10)
        for f in _j2.loads(resp.read().decode("utf-8")):
            n = unquote(f.get("name", ""))
            if n.lower().endswith(".pdf"):
                candidates.add(n)
    except Exception:
        pass

    if not candidates:
        return None

    # Tokeniza o ref e cada candidato; retorna o candidato com mais tokens
    # em comum (case-insensitive, ignorando caracteres não-alfanum).
    import re as _re
    def _tokens(s):
        return set(t for t in _re.split(r"[^a-z0-9]+", s.lower()) if len(t) > 1)
    ref_tok = _tokens(ref_clean)
    if not ref_tok:
        return None
    best, best_score = None, 0
    for c in candidates:
        score = len(ref_tok & _tokens(c))
        if score > best_score:
            best_score = score
            best = c
    return best if best_score >= 1 else None


@app.get("/api/sheet/{job_id}")
async def get_sheet_pdf(job_id: str, request: Request, ref: str = ""):
    """Serve a prancha (PDF, PNG, etc) inline pro viewer. Pra DWG/DXF
    tenta servir o PNG renderizado (render server-side) antes.

    Segurança: exige que quem chama seja dono do projeto (ou admin).
    Antes (até 2026-06-02) o endpoint era público — qualquer um com job_id
    válido baixava a prancha. Fix IDOR aplicado adicionando
    _require_project_owner."""
    _require_project_owner(request, job_id)
    from fastapi.responses import Response

    if not ref:
        raise HTTPException(400, "parâmetro 'ref' obrigatório")

    filename = _find_prancha_file(job_id, ref)
    if not filename:
        # Não ecoar o 'ref' cru (evita refletir input do usuário no viewer —
        # defesa em profundidade contra XSS; fix 2026-07-22). Mensagem estática.
        raise HTTPException(404, "Prancha não encontrada")

    ext = os.path.splitext(filename.lower())[1]
    preview_filename = filename
    mime = _PRANCHA_MIME.get(ext, "application/octet-stream")

    # DWG/DXF: preferir PNG renderizado se existe
    if ext in ('.dwg', '.dxf'):
        _png_name = os.path.splitext(filename)[0] + '.png'
        _png_local = os.path.join(WORK_DIR, job_id, _png_name)
        if os.path.exists(_png_local):
            preview_filename = _png_name
            mime = "image/png"
        else:
            _png_data = _supabase_storage_download_prancha(job_id, _png_name)
            if _png_data:
                return Response(
                    content=_png_data, media_type="image/png",
                    headers={
                        "Content-Disposition": "inline",
                        "X-Filename": _png_name,
                        "Cache-Control": "private, max-age=3600",
                    }
                )

    # Hot cache local
    local_path = os.path.join(WORK_DIR, job_id, preview_filename)
    if os.path.exists(local_path):
        try:
            with open(local_path, "rb") as f:
                data = f.read()
            return Response(
                content=data, media_type=mime,
                headers={
                    "Content-Disposition": "inline",
                    "X-Filename": preview_filename,
                    "Cache-Control": "private, max-age=3600",
                }
            )
        except Exception:
            pass

    # Storage
    data = _supabase_storage_download_prancha(job_id, preview_filename)
    if data:
        return Response(
            content=data, media_type=mime,
            headers={
                "Content-Disposition": "inline",
                "X-Filename": preview_filename,
                "Cache-Control": "private, max-age=3600",
            }
        )

    raise HTTPException(404, f"Prancha '{preview_filename}' não encontrada no Storage")


# ═══════════════════════════════════════════════════════════════
#  REPROCESSAR PROJETO (motor atualizado)
# ═══════════════════════════════════════════════════════════════

REPROCESS_FREE_LIMIT = 1  # 1 reprocessamento grátis por projeto


@app.post("/api/project/{job_id}/reprocess")
async def reprocess_project(job_id: str, request: Request):
    """Baixa os arquivos originais do Storage e cria novo job com os mesmos
    arquivos + tipologia. Usa a última versão do motor (prompts + regras).

    Política: 1 reprocessamento grátis por projeto. Tentativas adicionais
    retornam 402 (Payment Required) com mensagem orientando o user."""
    _require_project_owner(request, job_id)
    # Freio anti-abuso de custo (auditoria 27/07): reprocessar dispara o motor de
    # IA de novo. Teto por-projeto — humano reprocessa 1-2×; loop é barrado.
    if not _rate_limit_ok(f"reprocess:{job_id}", request, limit=6, window_s=600):
        raise HTTPException(429, "Muitos reprocessamentos em pouco tempo. Espere alguns minutos.")
    import urllib.request, urllib.error, json, shutil

    # 1) Buscar projeto original + checar contador.
    # Bug fix 2026-05-25: antes a query usava anon key como Bearer e RLS
    # bloqueava (a tabela projects exige auth.uid() = user_id). Resultado:
    # rows=[] → falso 404 "Projeto original não encontrado", mesmo o usuário
    # sendo dono. Agora usa o JWT do request (mesmo padrão dos endpoints
    # outros que sofreram do mesmo problema na onda Daniela 2026-05-18).
    auth_header = request.headers.get("Authorization", "") or request.headers.get("authorization", "")
    user_token = auth_header[7:].strip() if auth_header.lower().startswith("bearer ") else None
    try:
        url = f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}&select=*"
        req = urllib.request.Request(url, method='GET')
        req.add_header('apikey', SUPABASE_KEY)
        # JWT do user quando disponível (passa RLS por auth.uid()); service_role
        # como fallback pra projetos legados anônimos (bypassa RLS). _require_project_owner
        # já validou ownership upstream — fallback aqui só serve pra ler row anônima.
        req.add_header('Authorization', f'Bearer {user_token or SUPABASE_SERVICE_ROLE_KEY}')
        resp = urllib.request.urlopen(req, timeout=15)
        projects = json.loads(resp.read().decode('utf-8'))
        if not projects:
            # Diagnóstico extra: se _require_project_owner passou mas a query
            # devolveu vazio, é quase certamente bug de RLS — não "não existe".
            print(f"[reprocess] /projects?job_id={job_id} vazio com token={'sim' if user_token else 'nao'}")
            raise HTTPException(404, "Projeto original não encontrado no banco. Se o erro persistir, recarregue a página (Ctrl+Shift+R) e tente de novo.")
        orig = projects[0]
    except urllib.error.HTTPError as e:
        try:
            body = e.read().decode('utf-8', errors='replace')[:300]
        except Exception:
            body = ''
        print(f"[reprocess] /projects HTTP {e.code}: {body}")
        raise HTTPException(500, f"Erro ao buscar projeto: HTTP {e.code}")

    # Política: 1 reprocessamento grátis por projeto
    current_count = int(orig.get("reprocess_count") or 0)
    if current_count >= REPROCESS_FREE_LIMIT:
        raise HTTPException(
            402,  # Payment Required
            f"Este projeto já foi reprocessado {current_count}× — o limite "
            f"gratuito é {REPROCESS_FREE_LIMIT} por projeto. Crie um novo "
            f"projeto pra processar com o motor atualizado."
        )

    # 2) Listar arquivos originais no Storage
    from urllib.parse import unquote
    try:
        list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        body = json.dumps({"prefix": f"{job_id}/", "limit": 100}).encode("utf-8")
        req = urllib.request.Request(list_url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=20)
        storage_files = json.loads(resp.read().decode("utf-8"))
    except Exception as e:
        raise HTTPException(500, f"Erro ao listar Storage: {e}")

    valid_ext = ('.pdf', '.dwg', '.dxf')
    original_filenames = []
    for f in storage_files:
        n = unquote(f.get("name", ""))
        if n.lower().endswith(valid_ext):
            original_filenames.append(n)
    if not original_filenames:
        raise HTTPException(400,
            "Arquivos originais não disponíveis no Storage. "
            "Projetos anteriores a 21/04/2026 não tiveram upload automático — "
            "suba novamente pra reprocessar.")

    # 3) Criar novo job_id + work dir + baixar arquivos
    new_job_id = str(uuid.uuid4())[:8]
    new_work_dir = os.path.join(WORK_DIR, new_job_id)
    os.makedirs(new_work_dir, exist_ok=True)

    new_file_paths = []
    file_types = {'pdf': 0, 'dwg': 0, 'dxf': 0}
    for fname in original_filenames:
        local_path = os.path.join(new_work_dir, fname)
        data = _supabase_storage_download_prancha(job_id, fname)
        if not data:
            continue
        with open(local_path, "wb") as f:
            f.write(data)
        new_file_paths.append(local_path)
        ext = fname.lower().rsplit('.', 1)[-1]
        file_types[ext] = file_types.get(ext, 0) + 1

    if not new_file_paths:
        shutil.rmtree(new_work_dir, ignore_errors=True)
        raise HTTPException(500, "Falha ao baixar arquivos do Storage")

    # 4) Criar novo projeto + status
    types_summary = ", ".join(f"{v} {k.upper()}" for k, v in file_types.items() if v > 0)
    jobs[new_job_id] = ProcessingStatus(
        job_id=new_job_id, status="queued", progress=0,
        current_step=f"Reprocessamento: {len(new_file_paths)} arquivos recuperados ({types_summary})",
        total_steps=3,
    )

    typology = orig.get("typology") or "office"
    ptype = orig.get("project_type") or "arquitetura"
    try:
        _uta_orig = float(orig.get("user_total_area") or 0)
    except (TypeError, ValueError):
        _uta_orig = 0
    _supabase_insert("projects", {
        "job_id": new_job_id,
        "user_id": orig.get("user_id") or "anonymous",
        "user_email": orig.get("user_email") or "",
        "user_name": orig.get("user_name") or "",
        "project_name": f"{orig.get('project_name','Projeto')} (reprocessado)",
        "typology": typology,
        "project_type": ptype,  # propaga tipo: projeto ESTRUTURAL não vira arquitetura
        "files_count": len(new_file_paths),
        "file_types": file_types,
        "status": "queued",
        "parent_job_id": job_id,  # rastreabilidade: novo projeto é filho do original
        "user_total_area": _uta_orig if _uta_orig > 0 else None,  # propaga área informada
    })

    # Incrementar contador do ORIGINAL via RPC atômica
    try:
        inc_url = f"{SUPABASE_URL}/rest/v1/rpc/increment_reprocess_count"
        inc_body = json.dumps({"p_job_id": job_id}).encode('utf-8')
        inc_req = urllib.request.Request(inc_url, data=inc_body, method='POST')
        inc_req.add_header('apikey', SUPABASE_KEY)
        inc_req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        inc_req.add_header('Content-Type', 'application/json')
        urllib.request.urlopen(inc_req, timeout=10)
    except Exception as _inc_e:
        print(f"[reprocess] Erro ao incrementar contador: {_inc_e}")

    # 5) Disparar em thread (throttled pelo semáforo — 1 job por vez)
    import threading
    t = threading.Thread(
        target=_process_job_throttled,
        args=(new_job_id, new_file_paths, new_work_dir),
        kwargs={"typology": typology, "project_type": ptype,
                "user_total_area": _uta_orig},
        daemon=True,
    )
    t.start()

    return {
        "status": "ok",
        "original_job_id": job_id,
        "new_job_id": new_job_id,
        "files_count": len(new_file_paths),
        "typology": typology,
    }


@app.post("/api/admin/eval-reprocess/{job_id}")
async def admin_eval_reprocess(job_id: str, request: Request):
    """ADMIN — "modo avaliação": re-roda os arquivos de um projeto num job
    ISOLADO, sem tocar no projeto do cliente. Serve pra comparar leituras sem
    risco (ex.: DWG vs PDF, ou motor novo vs antigo).

    Isolamento (por quê é seguro):
    - is_eval=true  → fora das varreduras (auto-retry/alerta) e marcado no admin;
    - user_email='' → TODO email de cliente (sucesso e falha) no-opa (o
      _send_email_smtp/_email_falha_cliente checam email não-vazio);
    - parent_job_id → guarda extra contra email de falha;
    - user_id='eval'→ não aparece no dashboard de nenhum cliente;
    - NÃO incrementa o reprocess_count do original (não gasta o grátis do cliente).
    Mesmo pipeline do processamento real → avaliação fiel. Só ADMIN_EMAIL acessa."""
    _require_admin(request)
    import urllib.request, urllib.error, json, shutil
    from urllib.parse import unquote

    # 1) Projeto original (service role — admin já validado no _require_admin)
    try:
        url = f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}&select=*"
        req = urllib.request.Request(url, method='GET')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        projects = json.loads(urllib.request.urlopen(req, timeout=15).read().decode('utf-8'))
        if not projects:
            raise HTTPException(404, "Projeto original não encontrado.")
        orig = projects[0]
    except urllib.error.HTTPError as e:
        raise HTTPException(500, f"Erro ao buscar projeto: HTTP {e.code}")

    if orig.get("is_eval"):
        raise HTTPException(400, "Este já é um projeto de avaliação — avalie o original.")

    # 2) Arquivos originais no Storage do job ORIGINAL
    try:
        list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        body = json.dumps({"prefix": f"{job_id}/", "limit": 100}).encode("utf-8")
        req = urllib.request.Request(list_url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        storage_files = json.loads(urllib.request.urlopen(req, timeout=20).read().decode("utf-8"))
    except Exception as e:
        raise HTTPException(500, f"Erro ao listar Storage: {e}")

    valid_ext = ('.pdf', '.dwg', '.dxf')
    original_filenames = [unquote(f.get("name", "")) for f in storage_files
                          if unquote(f.get("name", "")).lower().endswith(valid_ext)]
    if not original_filenames:
        raise HTTPException(400,
            "Arquivos originais não estão no Storage (projeto antigo, pré-upload "
            "automático de 21/04/2026). Sem os arquivos não dá pra avaliar.")

    # 3) Novo job de avaliação + download dos arquivos
    eval_job_id = "ev" + str(uuid.uuid4())[:6]
    eval_work_dir = os.path.join(WORK_DIR, eval_job_id)
    os.makedirs(eval_work_dir, exist_ok=True)
    eval_file_paths = []
    file_types = {'pdf': 0, 'dwg': 0, 'dxf': 0}
    for fname in original_filenames:
        data = _supabase_storage_download_prancha(job_id, fname)
        if not data:
            continue
        lp = os.path.join(eval_work_dir, fname)
        with open(lp, "wb") as f:
            f.write(data)
        eval_file_paths.append(lp)
        ext = fname.lower().rsplit('.', 1)[-1]
        file_types[ext] = file_types.get(ext, 0) + 1
    if not eval_file_paths:
        shutil.rmtree(eval_work_dir, ignore_errors=True)
        raise HTTPException(500, "Falha ao baixar arquivos do Storage")

    # 4) Row de avaliação ISOLADA
    typology = orig.get("typology") or "office"
    ptype = orig.get("project_type") or "arquitetura"
    types_summary = ", ".join(f"{v} {k.upper()}" for k, v in file_types.items() if v > 0)
    jobs[eval_job_id] = ProcessingStatus(
        job_id=eval_job_id, status="queued", progress=0,
        current_step=f"Avaliação (teste): {len(eval_file_paths)} arquivo(s) ({types_summary})",
        total_steps=3,
    )
    _supabase_insert("projects", {
        "job_id": eval_job_id,
        "user_id": "eval",
        "user_email": "",  # vazio → nenhum email de cliente dispara
        "user_name": "",
        "project_name": f"[TESTE] {orig.get('project_name', 'Projeto')} — avaliação",
        "typology": typology,
        "project_type": ptype,
        "files_count": len(eval_file_paths),
        "file_types": file_types,
        "status": "queued",
        "parent_job_id": job_id,  # rastreabilidade + guarda extra contra email
        "user_total_area": orig.get("user_total_area"),  # espelha a área informada do original
        "is_eval": True,
    })

    # 5) Dispara — MESMO pipeline (avaliação fiel). NÃO incrementa reprocess do original.
    import threading
    threading.Thread(
        target=_process_job_throttled,
        args=(eval_job_id, eval_file_paths, eval_work_dir),
        # Passa a área INFORMADA pelo cliente no original — sem isso a avaliação
        # não é fiel (cai no ramo que zera todo m²). 0 se não houve área informada.
        kwargs={"typology": typology, "project_type": ptype,
                "user_total_area": float(orig.get("user_total_area") or 0)},
        daemon=True,
    ).start()

    print(f"[eval] avaliação {eval_job_id} disparada a partir de {job_id} "
          f"({len(eval_file_paths)} arquivo(s), {types_summary})")
    return {
        "status": "ok",
        "original_job_id": job_id,
        "eval_job_id": eval_job_id,
        "files_count": len(eval_file_paths),
        "view_url": f"/projeto.html?job_id={eval_job_id}",
    }


def _ai_suggest_combine(files_meta: list, base_type: str = "") -> dict:
    """Pede pra Claude Haiku olhar os arquivos (nome/tipo/tamanho/disciplina) e
    sugerir QUAIS combinar (mesma obra) e quais deixar de fora (outra obra,
    grande demais >150MB, duplicado). Best-effort — se a IA falhar, cai num
    fallback determinístico (recomenda todos até 150MB, exclui os gigantes)."""
    def _fallback():
        rec, exc = [], []
        for f in files_meta:
            if f.get("size_mb", 0) and f["size_mb"] > 150:
                exc.append({"filename": f["filename"],
                            "reason": f"grande demais ({f['size_mb']:.0f} MB, limite 150) — não abre"})
            else:
                rec.append(f["filename"])
        return {"recommended": rec, "excluded": exc, "ai": False,
                "summary": "Sugestão automática (IA indisponível): todos os arquivos até 150 MB."}
    api_key = os.getenv("ANTHROPIC_API_KEY", "")
    if not api_key or not files_meta:
        return _fallback()
    try:
        import anthropic, json as _j
        _lines = [
            f"{i+1}. {f['filename']} | tipo={f.get('ext','?').upper()} | {f.get('size_mb',0):.1f} MB "
            f"| projeto=\"{f.get('project_name','')}\" | disciplina={f.get('project_type','?')}"
            for i, f in enumerate(files_meta)
        ]
        prompt = (
            "Você ajuda um admin a COMBINAR arquivos CAD que pertencem à MESMA obra num "
            "quantitativo só. Abaixo, os arquivos que um cliente subiu (em vários envios):\n\n"
            + "\n".join(_lines) +
            "\n\nDecida QUAIS combinar (são a mesma obra/edifício) e quais DEIXAR DE FORA. "
            "Deixe de fora: arquivo claramente de OUTRA obra (pelo nome), duplicado óbvio, ou "
            "grande demais (>150 MB, não abre). DWG, DXF e PDF da MESMA obra DEVEM ser combinados "
            "(o motor mede pelo CAD e completa pelo PDF), mesmo em disciplinas diferentes "
            "(arquitetura + estrutura da MESMA obra = ok juntar).\n\n"
            "Responda SÓ com JSON, sem texto fora: {\"recommended\":[\"nome1\",...],"
            "\"excluded\":[{\"filename\":\"nome\",\"reason\":\"motivo curto\"}],"
            "\"summary\":\"1-2 frases pro admin (pt-br): o que combinar e por quê\"}"
        )
        client = anthropic.Anthropic(api_key=api_key)
        resp = client.messages.create(model="claude-haiku-4-5", max_tokens=900,
                                      messages=[{"role": "user", "content": prompt}])
        txt = resp.content[0].text if resp.content else ""
        _s, _e = txt.find("{"), txt.rfind("}")
        data = _j.loads(txt[_s:_e + 1]) if (_s >= 0 and _e > _s) else {}
        rec = [str(x) for x in (data.get("recommended") or [])]
        exc = [{"filename": str(x.get("filename", "")), "reason": str(x.get("reason", ""))}
               for x in (data.get("excluded") or []) if isinstance(x, dict)]
        if not rec and not exc:
            return _fallback()
        return {"recommended": rec, "excluded": exc, "ai": True,
                "summary": (str(data.get("summary") or "").strip() or "A IA analisou os arquivos.")}
    except Exception as _e:
        print(f"[combine-preview] IA falhou, fallback: {_e}")
        return _fallback()


@app.get("/api/admin/combine-preview/{job_id}")
async def admin_combine_preview(job_id: str, request: Request):
    """PRÉVIA do Combinar: lista os arquivos do MESMO cliente (nome/tipo/tamanho/
    disciplina) e pede pra IA sugerir quais juntar (mesma obra). O admin confere
    e ajusta antes de rodar. Não processa nada — só lê. Só ADMIN_EMAIL."""
    _require_admin(request)
    import urllib.request, urllib.error, json
    from urllib.parse import unquote, quote as _q
    try:
        url = f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}&select=*"
        req = urllib.request.Request(url, method='GET')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        base = json.loads(urllib.request.urlopen(req, timeout=15).read().decode('utf-8'))
        if not base:
            raise HTTPException(404, "Projeto não encontrado.")
        base = base[0]
    except urllib.error.HTTPError as e:
        raise HTTPException(500, f"Erro ao buscar projeto: HTTP {e.code}")
    _email = (base.get("user_email") or "").strip()
    if not _email:
        raise HTTPException(400, "Projeto sem cliente (email vazio) — não dá pra combinar.")
    try:
        pq = (f"{SUPABASE_URL}/rest/v1/projects?user_email=eq.{_q(_email)}"
              f"&is_eval=not.is.true&select=job_id,project_name,project_type,created_at&limit=40")
        preq = urllib.request.Request(pq, method='GET')
        preq.add_header('apikey', SUPABASE_KEY)
        preq.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        projs = json.loads(urllib.request.urlopen(preq, timeout=15).read().decode('utf-8'))
    except Exception as e:
        raise HTTPException(500, f"Erro ao listar projetos: {e}")
    valid_ext = ('.pdf', '.dwg', '.dxf')
    files_meta, seen_names = [], set()
    for p in projs or []:
        pj = p.get("job_id")
        try:
            list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
            body = json.dumps({"prefix": f"{pj}/", "limit": 100}).encode("utf-8")
            r = urllib.request.Request(list_url, data=body, method="POST")
            r.add_header("apikey", SUPABASE_KEY)
            r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            r.add_header("Content-Type", "application/json")
            objs = json.loads(urllib.request.urlopen(r, timeout=20).read().decode("utf-8"))
        except Exception:
            continue
        for o in objs or []:
            n = unquote(o.get("name", ""))
            if not n.lower().endswith(valid_ext):
                continue
            bn = os.path.basename(n)
            if bn in seen_names:
                continue
            seen_names.add(bn)
            try:
                _sz = int(((o.get("metadata") or {}).get("size")) or 0)
            except Exception:
                _sz = 0
            files_meta.append({
                "filename": bn, "ext": bn.lower().rsplit('.', 1)[-1],
                "size_mb": round(_sz / 1048576, 1), "job_id": pj,
                "project_name": p.get("project_name", ""),
                "project_type": p.get("project_type", ""),
            })
    if not files_meta:
        raise HTTPException(400, "Nenhum arquivo no Storage desse cliente.")
    suggestion = _ai_suggest_combine(files_meta, base.get("project_type", ""))
    return {"status": "ok", "client_email": _email, "base_job_id": job_id,
            "base_project_type": (base.get("project_type") or ""),
            "files": files_meta, "suggestion": suggestion}


@app.post("/api/admin/eval-combine/{job_id}")
async def admin_eval_combine(job_id: str, request: Request):
    """ADMIN — "combinar e avaliar": junta os arquivos de TODOS os projetos do
    mesmo cliente com o MESMO nome (caso comum: o cliente subiu a mesma obra em
    DWG, DXF e PDF separados) num teste ISOLADO. O motor mede o esqueleto pelo
    DXF/DWG e completa a cobertura pelo PDF — o melhor dos dois. Mesma blindagem
    do eval-reprocess (is_eval, user_email='', não avisa o cliente, não gasta
    reprocesso). Só ADMIN_EMAIL acessa."""
    _require_admin(request)
    import urllib.request, urllib.error, json, shutil
    from urllib.parse import unquote, quote as _q

    # 1) Projeto base + irmãos (mesmo email + mesmo nome, não-eval)
    try:
        url = f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}&select=*"
        req = urllib.request.Request(url, method='GET')
        req.add_header('apikey', SUPABASE_KEY)
        req.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
        base = json.loads(urllib.request.urlopen(req, timeout=15).read().decode('utf-8'))
        if not base:
            raise HTTPException(404, "Projeto não encontrado.")
        base = base[0]
    except urllib.error.HTTPError as e:
        raise HTTPException(500, f"Erro ao buscar projeto: HTTP {e.code}")

    _email = (base.get("user_email") or "").strip()
    _pname = (base.get("project_name") or "").strip()

    # Lista EXPLÍCITA de arquivos (vinda da PRÉVIA — o admin conferiu e escolheu).
    # Formato: {"files":[{"job_id":"...","filename":"..."}]}. Se vier, MANDA.
    explicit = []
    try:
        _body = await request.json()
        if isinstance(_body, dict):
            explicit = [x for x in (_body.get("files") or [])
                        if isinstance(x, dict) and x.get("job_id") and x.get("filename")]
    except Exception:
        explicit = []

    valid_ext = ('.pdf', '.dwg', '.dxf')
    eval_job_id = "ev" + str(uuid.uuid4())[:6]
    eval_work_dir = os.path.join(WORK_DIR, eval_job_id)
    os.makedirs(eval_work_dir, exist_ok=True)
    seen = {}  # basename -> source_job_id

    if explicit:
        # TRAVA DE ISOLAMENTO (regra dura): todo job referenciado TEM que ser do
        # MESMO cliente que o base. Blinda contra combinar clientes diferentes
        # mesmo que o front tenha bug.
        src_ids = list(dict.fromkeys(str(x["job_id"]) for x in explicit))
        try:
            vq = (f"{SUPABASE_URL}/rest/v1/projects?job_id=in.({','.join(src_ids)})"
                  f"&select=job_id,user_email")
            vreq = urllib.request.Request(vq, method='GET')
            vreq.add_header('apikey', SUPABASE_KEY)
            vreq.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
            owners = json.loads(urllib.request.urlopen(vreq, timeout=15).read().decode('utf-8'))
        except Exception as _ve:
            shutil.rmtree(eval_work_dir, ignore_errors=True)
            raise HTTPException(500, f"Erro ao validar donos dos arquivos: {_ve}")
        _emails = {(o.get("user_email") or "").strip().lower() for o in owners}
        if (not _email) or _emails != {_email.lower()} or len(owners) != len(src_ids):
            shutil.rmtree(eval_work_dir, ignore_errors=True)
            raise HTTPException(400, "Não dá pra combinar: os arquivos têm que ser TODOS do "
                                     "mesmo cliente (isolamento). Revise a seleção.")
        for x in explicit:
            bn = os.path.basename(str(x["filename"]))
            if bn.lower().endswith(valid_ext):
                seen.setdefault(bn, str(x["job_id"]))
    else:
        # Fallback (sem prévia): auto por email + nome (mesma obra, mesmo dia)
        sibling_ids = [job_id]
        if _email and _pname:
            try:
                sq = (f"{SUPABASE_URL}/rest/v1/projects?user_email=eq.{_q(_email)}"
                      f"&project_name=eq.{_q(_pname)}&is_eval=not.is.true&select=job_id&limit=20")
                sreq = urllib.request.Request(sq, method='GET')
                sreq.add_header('apikey', SUPABASE_KEY)
                sreq.add_header('Authorization', f'Bearer {SUPABASE_SERVICE_ROLE_KEY}')
                rows = json.loads(urllib.request.urlopen(sreq, timeout=15).read().decode('utf-8'))
                sibling_ids = list(dict.fromkeys([job_id] + [r.get("job_id") for r in rows if r.get("job_id")]))
            except Exception as _se:
                print(f"[combine] busca de irmãos falhou (usa só o base): {_se}")
        for src in sibling_ids:
            try:
                list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
                body = json.dumps({"prefix": f"{src}/", "limit": 100}).encode("utf-8")
                r = urllib.request.Request(list_url, data=body, method="POST")
                r.add_header("apikey", SUPABASE_KEY)
                r.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
                r.add_header("Content-Type", "application/json")
                files = json.loads(urllib.request.urlopen(r, timeout=20).read().decode("utf-8"))
            except Exception:
                continue
            for f in files:
                n = unquote(f.get("name", ""))
                if n.lower().endswith(valid_ext):
                    seen.setdefault(os.path.basename(n), src)

    eval_file_paths = []
    file_types = {'pdf': 0, 'dwg': 0, 'dxf': 0}
    for bn, src in seen.items():
        data = _supabase_storage_download_prancha(src, bn)
        if not data:
            continue
        lp = os.path.join(eval_work_dir, bn)
        with open(lp, "wb") as f:
            f.write(data)
        eval_file_paths.append(lp)
        ext = bn.lower().rsplit('.', 1)[-1]
        file_types[ext] = file_types.get(ext, 0) + 1
    if not eval_file_paths:
        shutil.rmtree(eval_work_dir, ignore_errors=True)
        raise HTTPException(400, "Nenhum arquivo encontrado no Storage desses projetos.")

    # 3) Row de avaliação ISOLADA (mesma blindagem do eval-reprocess)
    typology = base.get("typology") or "office"
    ptype = base.get("project_type") or "arquitetura"
    types_summary = ", ".join(f"{v} {k.upper()}" for k, v in file_types.items() if v > 0)
    jobs[eval_job_id] = ProcessingStatus(
        job_id=eval_job_id, status="queued", progress=0,
        current_step=f"Avaliação combinada: {len(eval_file_paths)} arquivo(s) ({types_summary})",
        total_steps=3,
    )
    _supabase_insert("projects", {
        "job_id": eval_job_id, "user_id": "eval", "user_email": "", "user_name": "",
        "project_name": f"[TESTE] {_pname or 'Projeto'} — combinado ({types_summary})",
        "typology": typology, "project_type": ptype,
        "files_count": len(eval_file_paths), "file_types": file_types,
        "status": "queued", "parent_job_id": job_id, "is_eval": True,
    })
    import threading
    threading.Thread(target=_process_job_throttled,
                     args=(eval_job_id, eval_file_paths, eval_work_dir),
                     kwargs={"typology": typology, "project_type": ptype}, daemon=True).start()
    print(f"[combine] avaliação combinada {eval_job_id}: {len(sibling_ids)} projeto(s) → "
          f"{len(eval_file_paths)} arquivo(s) ({types_summary})")
    return {
        "status": "ok",
        "eval_job_id": eval_job_id,
        "combined_from": sibling_ids,
        "files_count": len(eval_file_paths),
        "file_types": file_types,
        "view_url": f"/projeto.html?job_id={eval_job_id}",
    }


@app.post("/api/project/{job_id}/add-file")
async def add_file_and_reprocess(job_id: str, request: Request, files: list[UploadFile] = File(...)):
    """Anexa UM OU MAIS arquivos (DWG/DXF/PDF) a um projeto EXISTENTE e reprocessa NO
    MESMO job_id — pro caso "veio PDF, deu tudo estimado (ou poucas pranchas); agora
    manda o CAD ou mais pranchas pra medir/completar". Diferente de /reprocess (que cria
    um job FILHO), aqui é IN-PLACE: sobe os arquivos pro Storage do job, baixa TODOS (os
    novos + os originais) e re-dispara process_job no MESMO job (throttled pelo semáforo).
    Grátis — é melhoria de insumo. A limpeza dos itens antigos só acontece no SUCESSO do
    reprocesso (_persist), então um reprocesso que falhe preserva a planilha anterior."""
    _require_project_owner(request, job_id)
    # Freio anti-abuso de custo (auditoria 27/07): anexar+reprocessar dispara o
    # motor de IA. Teto por-projeto — loop é barrado, uso normal passa folgado.
    if not _rate_limit_ok(f"addfile:{job_id}", request, limit=6, window_s=600):
        raise HTTPException(429, "Muitos envios de arquivo em pouco tempo. Espere alguns minutos.")
    import urllib.request, urllib.error, json
    from urllib.parse import unquote

    if not files:
        raise HTTPException(400, "Envie ao menos um arquivo: DWG, DXF ou PDF.")

    # Valida extensões de TODOS antes de ler nada
    for _f in files:
        _rn = _f.filename or ""
        _ext = _rn.lower().rsplit(".", 1)[-1] if "." in _rn else ""
        if _ext not in ("dwg", "dxf", "pdf"):
            raise HTTPException(400, f"'{_rn or 'arquivo'}': envie só DWG, DXF ou PDF.")

    # Teto de request (anti-OOM). O upload agora é gravado em pedaços no disco
    # (_stream_upload_to_disk, mesmo do /api/process) — não bufferiza mais o
    # arquivo inteiro na RAM. Unificado em 450 MB com o /api/process (4 GB).
    _clen = request.headers.get("content-length") or request.headers.get("Content-Length")
    if _clen and _clen.isdigit() and int(_clen) > 450 * 1024 * 1024:
        raise HTTPException(413, "Arquivos muito grandes (máx. ~450 MB no total). Envie só as pranchas necessárias.")

    # Projeto original: tipologia/tipo + guarda contra reprocesso concorrente
    typology, ptype = "office", "arquitetura"
    try:
        purl = f"{SUPABASE_URL}/rest/v1/projects?job_id=eq.{job_id}&select=typology,project_type,status"
        preq = urllib.request.Request(purl, method="GET")
        preq.add_header("apikey", SUPABASE_KEY)
        preq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        prows = json.loads(urllib.request.urlopen(preq, timeout=15).read().decode("utf-8"))
        if prows:
            typology = prows[0].get("typology") or "office"
            ptype = prows[0].get("project_type") or "arquitetura"
            if (prows[0].get("status") or "") in ("queued", "processing"):
                raise HTTPException(409, "Esse projeto ainda está processando. Espera terminar e tenta de novo.")
    except HTTPException:
        raise
    except Exception as _pe:
        print(f"[add-file] leitura projeto {job_id}: {_pe}")

    work_dir = os.path.join(WORK_DIR, job_id)
    os.makedirs(work_dir, exist_ok=True)
    saved = 0
    for _f in files:
        safe_local = _safe_local_filename(_f.filename or f"arquivo_{saved}")
        new_local = os.path.join(work_dir, safe_local)
        # Grava em pedaços (não bufferiza o arquivo inteiro na RAM — mesmo padrão
        # anti-OOM do /api/process; este arquivo alimenta o ezdxf logo em seguida).
        n_written, _ = await _stream_upload_to_disk(_f, new_local)
        if not n_written:
            try: os.remove(new_local)
            except OSError: pass
            continue
        if n_written > 150 * 1024 * 1024:
            try: os.remove(new_local)
            except OSError: pass
            raise HTTPException(413, f"'{_f.filename}' passa de 150 MB. Exporte só a prancha necessária.")
        if not _supabase_storage_upload_prancha(new_local, job_id, safe_local):
            raise HTTPException(500, "Não consegui guardar um dos arquivos. Tenta de novo.")
        saved += 1
    if not saved:
        raise HTTPException(400, "Arquivo(s) vazio(s).")

    # Lista + baixa TODOS os arquivos do job (o novo + os originais)
    try:
        list_url = f"{SUPABASE_URL}/storage/v1/object/list/{PRANCHAS_BUCKET}"
        lbody = json.dumps({"prefix": f"{job_id}/", "limit": 200}).encode("utf-8")
        lreq = urllib.request.Request(list_url, data=lbody, method="POST")
        lreq.add_header("apikey", SUPABASE_KEY)
        lreq.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        lreq.add_header("Content-Type", "application/json")
        objs = json.loads(urllib.request.urlopen(lreq, timeout=20).read().decode("utf-8"))
    except Exception as e:
        raise HTTPException(500, f"Erro ao listar Storage: {e}")

    names = [unquote(o.get("name", "")) for o in objs if o.get("name")]
    names = [n for n in names if n.lower().endswith((".pdf", ".dwg", ".dxf"))]
    file_paths = []
    for n in names:
        bn = os.path.basename(n)
        lp = os.path.join(work_dir, _safe_local_filename(bn))
        d = _supabase_storage_download_prancha(job_id, bn)
        if not d:
            continue
        with open(lp, "wb") as f:
            f.write(d)
        file_paths.append(lp)
    if not file_paths:
        raise HTTPException(500, "Falha ao preparar os arquivos pra reprocessar.")

    # Se há CAD (dwg/dxf), processa SÓ o CAD e descarta os PDFs. O PDF era o
    # stand-in que deu "estimado"; misturar os dois DUPLICA a quantidade (PDF
    # estima 100, CAD mede 95 → 2 linhas). Com CAD presente, o CAD manda.
    _cads = [p for p in file_paths if p.lower().endswith((".dwg", ".dxf"))]
    if _cads:
        # Se existe o DXF re-exportado de um DWG que falhava (MESMO nome), prefere
        # o DXF e descarta o DWG. Senão o DWG velho no Storage re-falha a cada
        # reprocesso e gera aviso "planilha INCOMPLETA" enganoso, mesmo o DXF tendo
        # medido tudo (caso forro MEP do Pedro, 15/07 — DWG AEC não abre, DXF sim).
        _dxf_stems = {os.path.splitext(os.path.basename(p))[0].lower()
                      for p in _cads if p.lower().endswith(".dxf")}
        _cads = [p for p in _cads
                 if not (p.lower().endswith(".dwg")
                         and os.path.splitext(os.path.basename(p))[0].lower() in _dxf_stems)]
        file_paths = _cads

    # 🚨 TRAVA ANTI-PERDA (caso Walter, 30/07/2026). O complemento refaz o projeto
    # a partir do que está no STORAGE. Se o CAD original não subiu (o upload tem
    # timeout de 60s e falha em SILÊNCIO), sobram só os PDFs — e o passo final
    # APAGA os itens antigos antes de gravar os novos. Aconteceu de verdade: um
    # projeto com 40 itens MEDIDOS virou 57 estimados, sem erro nenhum na tela.
    # Regra dura nº1: nunca trocar medição por estimativa pelas nossas costas.
    if not _cads:
        _medidos_antes = _job_medidos_count(job_id)
        if _medidos_antes > 0:
            raise HTTPException(409, (
                f"Esse projeto tem {_medidos_antes} itens medidos do CAD, e o arquivo CAD "
                f"original não está mais guardado aqui. Refazer só com PDF apagaria essa "
                f"medição e deixaria tudo como estimativa. Anexe o DXF (ou DWG) junto que "
                f"a gente refaz medindo — seus arquivos novos já ficaram salvos."
            ))

    # NÃO apaga os itens antigos aqui: a limpeza acontece no _persist (só no
    # SUCESSO do reprocesso). Se o CAD falhar (não-convertível/0 itens), a
    # planilha estimada anterior é preservada. Aqui só marca 'queued'.
    #
    # 🪤 E ATUALIZA A COMPOSIÇÃO DO PROJETO (31/07/2026). Até aqui, anexar
    # arquivo não mexia em files_count/file_types: o projeto ficava congelado no
    # que o cliente mandou PRIMEIRO. O painel mostrava "1 PDF" num projeto que
    # já tinha CAD e 30 itens MEDIDOS (caso Fernando, 31/07) — e, pior, as
    # estatísticas de "PDF mede X, CAD mede Y" contavam esses sucessos do
    # complemento na coluna do PDF, inflando o PDF e escondendo o valor do
    # anexo. É o mesmo vício de "contar o que entrou, não o que o projeto virou".
    # Conta pelo que está no STORAGE (a composição real), não pelo que este run
    # processa — o run descarta os PDFs quando há CAD, mas eles seguem no projeto.
    _comp = {"dwg": 0, "dxf": 0, "pdf": 0}
    for _n in names:
        _e = os.path.splitext(_n.lower())[1].lstrip(".")
        if _e in _comp:
            _comp[_e] += 1
    _upd = {"status": "queued", "error_message": None}
    if sum(_comp.values()) > 0:
        _upd["file_types"] = _comp
        _upd["files_count"] = sum(_comp.values())
    _supabase_update("projects", "job_id", job_id, _upd)

    jobs[job_id] = ProcessingStatus(
        job_id=job_id, status="queued", progress=0,
        current_step=f"Refazendo com o novo arquivo ({len(file_paths)} no total)",
        total_steps=3,
    )
    import threading
    threading.Thread(
        target=_process_job_throttled,
        args=(job_id, file_paths, work_dir),
        kwargs={"typology": typology, "project_type": ptype, "is_complement": True},
        daemon=True,
    ).start()

    return {"status": "ok", "job_id": job_id, "files_count": len(file_paths)}


@app.post("/api/items/{job_id}/review-finalize")
async def finalize_review(job_id: str, request: Request):
    """Marca a revisão inline como concluída e credita cashback (R$0,10 por
    ação, teto R$ 20,00 por projeto).

    Chamado quando o usuário clica "Concluir revisão" na tela de revisão.
    """
    _require_project_owner(request, job_id)  # fix IDOR 2026-07-22 (igual ao review-state)
    try:
        body = await request.json()
    except Exception:
        body = {}
    user_id = body.get("user_id", "")

    # Cashback inline removido em 2026-05-13. Política nova:
    # - Revisar itens dentro do site: agradecimento, sem cashback (treinava IA
    #   mas era cognição confusa: R$ 0,10/item raramente atingia o teto).
    # - Cashback fica concentrado em ações de MAIOR sinal pra IA:
    #     * Upload de planilha revisada offline → R$ 30
    #     * Upload de cotação de fornecedor → R$ 10 (max 3 = R$ 30)
    # Total possível por projeto: R$ 60.
    try:
        # Ainda conta reviews pra retornar n_actions ao frontend (pra UI mostrar
        # progresso do usuário), mas não cria evento de cashback.
        _, reviews = _supa_rest_as_user(
            request, "GET",
            f"/item_reviews?job_id=eq.{job_id}&select=item_id&limit=5000",
            timeout=10,
        )
        reviews = reviews or []
        n_actions = len({r["item_id"] for r in reviews if r.get("item_id")})
    except Exception as e:
        n_actions = 0

    return {
        "status": "ok",
        "n_actions": n_actions,
        "credit_cents": 0,
        "credit_brl": 0.0,
        "message": "Obrigado pela revisão — suas correções treinam a IA pro próximo projeto.",
    }


@app.get("/api/items/{job_id}/review-state")
async def get_review_state(job_id: str, request: Request):
    """Retorna as revisões já feitas nesse job pra restaurar o estado no
    browser quando o user volta pra terminar depois.

    Retorna mapa {item_id: {action, edits, comment, reviewed_at}} com a
    última review de cada item."""
    _require_project_owner(request, job_id)
    import urllib.request, urllib.error, json
    try:
        _, reviews = _supa_rest_as_user(
            request, "GET",
            f"/item_reviews?job_id=eq.{job_id}"
            f"&select=item_id,action,edits,comment,reviewed_at"
            f"&order=reviewed_at.desc&limit=2000",
            timeout=15,
        )
        reviews = reviews or []
        # Manter só a MAIS RECENTE por item (supabase retorna ordenado desc)
        state = {}
        for r in reviews:
            iid = r.get("item_id")
            if iid and iid not in state:
                state[iid] = {
                    "action": r.get("action"),
                    "comment": r.get("comment") or "",
                    "reviewed_at": r.get("reviewed_at"),
                }
        return {"status": "ok", "job_id": job_id, "state": state, "count": len(state)}
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar state: {e}")


# ═══════════════════════════════════════════════════════════════
#  INSIGHTS DE REVISÕES (pra admin entender o que ajustar no motor)
# ═══════════════════════════════════════════════════════════════

@app.get("/api/admin/review-insights")
async def admin_review_insights(request: Request, days: int = 30, limit: int = 30):
    """Agrega revisões recentes pra identificar padrões de erro do motor:
    - Quais tipos de item são mais rejeitados (alucinação)?
    - Quais edits são mais comuns (qty errada, unidade errada)?
    - Quais comentários o usuário deixou (insights qualitativos)?

    Serve pra eu (humano ou Claude) olhar e decidir o que ajustar no
    SYSTEM_PROMPT ou nas regras de consolidação."""
    _require_admin(request)
    import urllib.request, urllib.error, json

    # 1) Busca reviews recentes (com JOIN manual pra pegar description do item)
    try:
        # service_role: endpoint admin (agrega reviews de TODOS os projetos).
        # Com JWT do admin a RLS owner-scoped voltaria só os do próprio admin.
        _rev_url = (f"{SUPABASE_URL}/rest/v1/item_reviews"
                    f"?select=id,job_id,item_id,action,edits,comment,reviewed_at"
                    f"&reviewed_at=gte.{(datetime.utcnow() - timedelta(days=days)).isoformat()}"
                    f"&order=reviewed_at.desc&limit=1000")
        _rev_req = urllib.request.Request(_rev_url, method="GET")
        _rev_req.add_header("apikey", SUPABASE_KEY)
        _rev_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        reviews = json.loads(urllib.request.urlopen(_rev_req, timeout=20).read().decode("utf-8")) or []
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar reviews: {e}")

    # 2) Enrich com descrição original do item (se ainda existe).
    # FIX 2026-05-14: N+1 eliminado — antes fazia 1 GET por review (até 1000
    # requests sequenciais ao Supabase, timeout em pico). Agora 1 GET batch
    # com `item_id=in.(id1,id2,...)`. Reduz pra 1 chamada (até 100 ids).
    from collections import Counter
    rejected_patterns = Counter()
    edit_patterns = Counter()
    comments = []
    by_typology = Counter()

    # Pré-busca descrições em batch
    item_descs = {}
    item_ids_uniq = [r.get("item_id") for r in reviews if r.get("item_id")]
    item_ids_uniq = list({i for i in item_ids_uniq if i})
    for i in range(0, len(item_ids_uniq), 100):
        chunk = item_ids_uniq[i:i+100]
        try:
            ids_csv = ",".join(str(x) for x in chunk)
            # service_role: endpoint admin (itens de projetos de qualquer user).
            _pi_url = (f"{SUPABASE_URL}/rest/v1/project_items?id=in.({ids_csv})"
                       f"&select=id,description,discipline")
            _pi_req = urllib.request.Request(_pi_url, method="GET")
            _pi_req.add_header("apikey", SUPABASE_KEY)
            _pi_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            rows_ = json.loads(urllib.request.urlopen(_pi_req, timeout=10).read().decode("utf-8"))
            for row in (rows_ or []):
                item_descs[row.get("id")] = row.get("description", "")
        except Exception:
            pass

    for r in reviews:
        action = r.get("action", "")
        comment = (r.get("comment") or "").strip()
        if comment:
            comments.append({
                "comment": comment[:300],
                "action": action,
                "job_id": r.get("job_id"),
                "at": r.get("reviewed_at"),
            })

        item_id = r.get("item_id")
        item_desc = item_descs.get(item_id) if item_id else None

        if action == "reject" and item_desc:
            # Normaliza pra primeiras 3 palavras significativas
            key = " ".join(item_desc.lower().split()[:3])
            rejected_patterns[key] += 1

        if action == "edit":
            edits = r.get("edits") or {}
            if isinstance(edits, dict):
                for field in edits.keys():
                    edit_patterns[field] += 1

    return {
        "period_days": days,
        "total_reviews": len(reviews),
        "rejected_top": [
            {"pattern": k, "count": v}
            for k, v in rejected_patterns.most_common(limit)
        ],
        "edit_fields_top": [
            {"field": k, "count": v}
            for k, v in edit_patterns.most_common(limit)
        ],
        "comments_recent": comments[:limit],
    }


# ═══════════════════════════════════════════════════════════════
#  CLEANUP AUTOMÁTICO — retenção de 90 dias (LGPD)
#  Chamado diariamente via GitHub Actions cron. Deleta arquivos do
#  Storage e marca projetos como archived=true.
# ═══════════════════════════════════════════════════════════════

CLEANUP_RETENTION_DAYS = 90
CLEANUP_SECRET = os.getenv("CLEANUP_SECRET", "")


def _supabase_storage_delete(bucket: str, object_path: str) -> bool:
    """Deleta um objeto do Storage. Retorna True se OK (ou se já não existia)."""
    import urllib.request, urllib.error
    try:
        url = f"{SUPABASE_URL}/storage/v1/object/{bucket}/{object_path}"
        req = urllib.request.Request(url, method="DELETE")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        urllib.request.urlopen(req, timeout=15)
        return True
    except urllib.error.HTTPError as e:
        # 404 = já não existe (tudo ok)
        if e.code == 404:
            return True
        return False
    except Exception:
        return False


def _supabase_storage_list(bucket: str, prefix: str) -> list:
    """Lista objetos no bucket com o prefix dado. Retorna lista de nomes."""
    import urllib.request, json as _j
    try:
        url = f"{SUPABASE_URL}/storage/v1/object/list/{bucket}"
        body = _j.dumps({"prefix": prefix, "limit": 500}).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=15)
        files = _j.loads(resp.read().decode("utf-8"))
        return [f.get("name", "") for f in files if f.get("name")]
    except Exception as e:
        print(f"[cleanup] list {bucket}/{prefix}: {e}")
        return []


@app.get("/api/admin/cleanup-secret-check")
async def cleanup_secret_check(request: Request):
    """Debug: confirma se CLEANUP_SECRET está configurado no Render (sem
    expor o valor). Requer auth admin — anteriormente vazava first_4/last_4
    do segredo, reduzindo entropia pra brute-force do header X-Cleanup-Secret."""
    _require_admin(request)
    if not CLEANUP_SECRET:
        return {
            "configured": False,
            "message": "CLEANUP_SECRET NÃO está setado no Render. "
                       "Adicione em Environment e dê 'Save, rebuild, and deploy'."
        }
    return {
        "configured": True,
        "length": len(CLEANUP_SECRET),
        "message": "Secret configurada no Render."
    }


@app.post("/api/admin/cleanup-old-projects")
async def cleanup_old_projects(request: Request):
    """Deleta arquivos originais + planilha de projetos com 90+ dias.
    Marca projeto como archived=true pra preservar histórico no DB.

    Autenticação: header 'X-Cleanup-Secret' deve bater com env var
    CLEANUP_SECRET. Se CLEANUP_SECRET não estiver setado no Render, o
    endpoint responde 503.

    Chamado por GitHub Action cron diariamente."""
    # Auth
    provided_secret = request.headers.get("X-Cleanup-Secret", "")
    if not CLEANUP_SECRET:
        raise HTTPException(503, "CLEANUP_SECRET não configurado no ambiente")
    if provided_secret != CLEANUP_SECRET:
        raise HTTPException(401, "Secret inválido")

    # Query days pode customizar pra testes (default 90)
    import urllib.request, json as _j
    try:
        days = int(request.query_params.get("days", CLEANUP_RETENTION_DAYS))
    except Exception:
        days = CLEANUP_RETENTION_DAYS

    # 1) Lista projetos elegíveis via RPC
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_expired_projects"
        body = _j.dumps({"p_days": days}).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=15)
        expired = _j.loads(resp.read().decode("utf-8"))
    except Exception as e:
        raise HTTPException(500, f"Erro ao buscar projetos expirados: {e}")

    stats = {
        "days_threshold": days,
        "total_expired": len(expired),
        "archived": 0,
        "files_deleted": 0,
        "errors": [],
        "jobs": [],
    }

    # 2) Pra cada projeto: deletar arquivos do Storage + marcar archived
    for proj in expired:
        job_id = proj.get("job_id")
        if not job_id:
            continue

        # 2a) Arquivos originais (bucket aiarq-pranchas/{job_id}/)
        pranchas_list = _supabase_storage_list(PRANCHAS_BUCKET, f"{job_id}/")
        files_ok = 0
        for obj in pranchas_list:
            # list retorna "nome.pdf" (sem prefix). Delete precisa do path completo.
            path = f"{job_id}/{obj}"
            if _supabase_storage_delete(PRANCHAS_BUCKET, path):
                files_ok += 1

        # 2b) Planilha (bucket aiarq-planilhas/{job_id}.xlsx)
        if _supabase_storage_delete(PLANILHAS_BUCKET, f"{job_id}.xlsx"):
            files_ok += 1

        # 2c) Marcar projeto como archived
        try:
            mark_url = f"{SUPABASE_URL}/rest/v1/rpc/mark_project_archived"
            mark_body = _j.dumps({"p_job_id": job_id}).encode("utf-8")
            mark_req = urllib.request.Request(mark_url, data=mark_body, method="POST")
            mark_req.add_header("apikey", SUPABASE_KEY)
            mark_req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
            mark_req.add_header("Content-Type", "application/json")
            urllib.request.urlopen(mark_req, timeout=10)
            stats["archived"] += 1
            stats["files_deleted"] += files_ok
            stats["jobs"].append({"job_id": job_id, "files_deleted": files_ok})
        except Exception as e:
            stats["errors"].append({"job_id": job_id, "error": str(e)[:200]})

    _supa_log(f"CLEANUP archived={stats['archived']} files={stats['files_deleted']} "
              f"errors={len(stats['errors'])}")
    print(f"[cleanup] {stats['archived']} projetos arquivados, "
          f"{stats['files_deleted']} arquivos deletados, "
          f"{len(stats['errors'])} erros")

    # Registra execução no cleanup_log pra monitoramento no admin
    _supabase_insert("cleanup_log", {
        "days_threshold": days,
        "total_expired": stats["total_expired"],
        "archived": stats["archived"],
        "files_deleted": stats["files_deleted"],
        "errors_count": len(stats["errors"]),
        "details": {"jobs": stats["jobs"][:50]},  # limita tamanho
    })

    return stats


@app.get("/api/admin/cleanup-log")
async def admin_cleanup_log(request: Request, limit: int = 30):
    """Retorna últimas execuções do cleanup pra dashboard admin."""
    _require_admin(request)
    import urllib.request, json as _j
    try:
        url = f"{SUPABASE_URL}/rest/v1/rpc/list_cleanup_runs"
        body = _j.dumps({"p_limit": limit}).encode("utf-8")
        req = urllib.request.Request(url, data=body, method="POST")
        req.add_header("apikey", SUPABASE_KEY)
        req.add_header("Authorization", f"Bearer {SUPABASE_SERVICE_ROLE_KEY}")
        req.add_header("Content-Type", "application/json")
        resp = urllib.request.urlopen(req, timeout=10)
        return {"runs": _j.loads(resp.read().decode("utf-8"))}
    except Exception as e:
        raise HTTPException(500, f"Erro: {e}")


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
# Trigger autodeploy Mon Apr 13 10:58:57     2026
