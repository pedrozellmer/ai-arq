# -*- coding: utf-8 -*-
"""Exportação profissional do cronograma — PDF + PPTX.

Aplica princípios de diagramação profissional:
- Hierarquia visual (F1 hero / F2 apoio / F3 diagramas / F4 texto)
- Grid 12-col com margens 2cm
- Tipografia hierárquica (H1 36 / H2 18 / body 10 / caption 8)
- Paleta AI.arq consistente (indigo + cyan + dark + cream)
- Espaço em branco generoso
- Cards / boxes pra agrupar info
- Carimbo de prancha (nº + data + autoria)
"""
import os
import tempfile
from datetime import datetime as _dt, date as _date
from typing import Dict, Optional


# Paleta sóbria — alinhada com o frontend (2026-05-25).
# Antes era indigo/cyan/amber gritante (parecia material de apresentação
# escolar). Agora é escala slate puro com 1 cor de acento sutil baseada
# na cor de marca do escritório (que entra como hairline ou label, não
# como faixa fullbleed). Mais profissional, mais respeitoso ao conteúdo
# técnico do cronograma.
COLOR_DARK = '#0F172A'        # slate-900 — títulos
COLOR_DARK_2 = '#1E293B'      # slate-800 — body forte
COLOR_GRAY_TX = '#475569'     # slate-600 — body suave
COLOR_GRAY_MID = '#64748B'    # slate-500 — secundário
COLOR_GRAY_LIGHT = '#94A3B8'  # slate-400 — caption
COLOR_BORDER = '#E2E8F0'      # slate-200 — divisores
COLOR_BG_SOFT = '#F8FAFC'     # slate-50 — fundo sutil pra cards
COLOR_CREAM = '#FAF7F0'       # cream — capa
COLOR_WHITE = '#FFFFFF'
# Acentos mantidos só pra retrocompat (alguns lugares ainda referenciam)
COLOR_INDIGO = '#475569'      # remapeado pra slate-600 (era #4F46E5)
COLOR_INDIGO_DARK = '#334155' # slate-700 (era #3730A3)
COLOR_CYAN = '#64748B'        # slate-500 (era #06B6D4 vibrante)
COLOR_AMBER = '#92400E'       # amber-800 (era #F59E0B saturado)


# ─── PNGs auxiliares (Gantt + Curva S) ────────────────────────────

def gerar_gantt_png(cronograma: Dict, output_path: str,
                     titulo: str = '') -> str:
    """Gantt profissional. Sem título embutido (vai no layout do PDF/PPT)."""
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates

    fases = cronograma.get('fases', [])
    if not fases:
        fig, ax = plt.subplots(figsize=(12, 4))
        ax.text(0.5, 0.5, 'Cronograma sem fases', ha='center', va='center',
                fontsize=14, color=COLOR_GRAY_LIGHT)
        ax.axis('off')
        plt.savefig(output_path, dpi=160, bbox_inches='tight',
                    facecolor='white', edgecolor='none')
        plt.close()
        return output_path

    h = max(4.5, len(fases) * 0.45)
    fig, ax = plt.subplots(figsize=(14, h))
    fig.patch.set_facecolor('white')

    for i, f in enumerate(fases):
        ini = _dt.fromisoformat(f['inicio'])
        fim = _dt.fromisoformat(f['fim'])
        cor = f.get('cor', COLOR_INDIGO)
        # Barra com borda suave
        ax.barh(i, (fim - ini).days, left=mdates.date2num(ini),
                color=cor, alpha=0.95, edgecolor='white',
                linewidth=1.5, height=0.65)
        # Label dias dentro
        meio = ini + (fim - ini) / 2
        dur_txt = f"{f['dur_dias']}d"
        # Só mostra se barra suficientemente larga
        if (fim - ini).days >= 14:
            ax.text(mdates.date2num(meio), i, dur_txt,
                    ha='center', va='center', color='white',
                    fontsize=10, fontweight='600',
                    family='sans-serif')

    ax.set_yticks(list(range(len(fases))))
    ax.set_yticklabels([f['label'] for f in fases], fontsize=11,
                       color=COLOR_DARK, family='sans-serif')
    ax.invert_yaxis()
    ax.xaxis_date()
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b/%y'))
    ax.tick_params(axis='x', colors=COLOR_GRAY_TX, labelsize=10)
    ax.tick_params(axis='y', colors=COLOR_DARK, labelsize=10)

    # Remove spines
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color(COLOR_BORDER)
    ax.spines['bottom'].set_color(COLOR_BORDER)

    # Grid suave só horizontal
    ax.grid(axis='x', alpha=0.25, linestyle='-', linewidth=0.6,
            color=COLOR_GRAY_LIGHT)
    ax.set_axisbelow(True)

    plt.tight_layout()
    plt.savefig(output_path, dpi=160, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    return output_path


def gerar_curva_s_png(cronograma: Dict, output_path: str) -> str:
    """Curva S sigmoidal profissional."""
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates

    curva = cronograma.get('curva_s', [])
    if not curva:
        fig, ax = plt.subplots(figsize=(12, 6))
        ax.text(0.5, 0.5, 'Sem curva S', ha='center', va='center',
                fontsize=14, color=COLOR_GRAY_LIGHT)
        ax.axis('off')
        plt.savefig(output_path, dpi=160, bbox_inches='tight',
                    facecolor='white', edgecolor='none')
        plt.close()
        return output_path

    fig, ax = plt.subplots(figsize=(13, 6))
    fig.patch.set_facecolor('white')

    datas = [_dt.fromisoformat(p['data_fim_mes']) for p in curva]
    pcts = [p['pct_acumulado'] for p in curva]

    # Linha + área
    ax.fill_between(datas, pcts, alpha=0.12, color=COLOR_INDIGO)
    ax.plot(datas, pcts, color=COLOR_INDIGO, linewidth=3,
            solid_joinstyle='round', solid_capstyle='round')

    # Marcos com anotação
    for pct_alvo in (25, 50, 75, 100):
        for i, p in enumerate(pcts):
            if p >= pct_alvo:
                ax.axhline(pct_alvo, color=COLOR_BORDER, alpha=0.5,
                           linestyle='-', linewidth=0.8, zorder=0)
                ax.scatter([datas[i]], [pct_alvo], color=COLOR_CYAN,
                           s=90, zorder=5, edgecolor='white', linewidth=2)
                ax.annotate(f'{pct_alvo}%',
                            xy=(datas[i], pct_alvo),
                            xytext=(8, -16), textcoords='offset points',
                            fontsize=10, color=COLOR_INDIGO_DARK,
                            fontweight='bold', family='sans-serif')
                break

    ax.set_ylabel('% Avanço acumulado', fontsize=11, color=COLOR_GRAY_TX,
                  family='sans-serif')
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b/%y'))
    ax.tick_params(axis='x', colors=COLOR_GRAY_TX, labelsize=10)
    ax.tick_params(axis='y', colors=COLOR_GRAY_TX, labelsize=10)
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color(COLOR_BORDER)
    ax.spines['bottom'].set_color(COLOR_BORDER)
    ax.set_ylim(0, 105)
    ax.set_axisbelow(True)
    ax.grid(alpha=0.15, linestyle='-', linewidth=0.6,
            color=COLOR_GRAY_LIGHT)

    plt.tight_layout()
    plt.savefig(output_path, dpi=160, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    return output_path


def _format_br(iso_date: Optional[str]) -> str:
    if not iso_date or not isinstance(iso_date, str):
        return ''
    parts = iso_date.split('-')
    if len(parts) == 3:
        return f'{parts[2]}/{parts[1]}/{parts[0]}'
    return iso_date


# ═════════════════════════════════════════════════════════════════
#  PDF — Layout profissional A4 paisagem
# ═════════════════════════════════════════════════════════════════

def _hex_to_rgb01(hex_color: str) -> tuple:
    """#4F46E5 → (0.31, 0.27, 0.90) pra reportlab/colors."""
    h = (hex_color or '#4F46E5').lstrip('#')
    if len(h) == 3:
        h = ''.join(c*2 for c in h)
    return tuple(int(h[i:i+2], 16)/255.0 for i in (0, 2, 4))


def _hex_to_rgb_int(hex_color: str) -> tuple:
    """#4F46E5 → (79, 70, 229) pra pptx."""
    h = (hex_color or '#4F46E5').lstrip('#')
    if len(h) == 3:
        h = ''.join(c*2 for c in h)
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def exportar_pdf(cronograma: Dict, output_path: str,
                  branding: Optional[Dict] = None,
                  # compat antiga
                  titulo: str = '', job_id: str = '') -> str:
    """PDF executivo CO-BRANDED (logo + cor do escritório + nome cliente).

    branding dict (preferido): {project_name, architect_name, client_name,
      company, logo_local_path, brand_color, job_id}

    Estrutura:
      - Página 1: Capa cor da marca + hero + dados
      - Página 2: Gantt (foco principal F1)
      - Página 3: Curva S de avanço previsto
      - Página 4: Caminho crítico (tabela executiva)
      - Página 5: Marcos normativos + ressalvas
    """
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import (BaseDocTemplate, Frame, PageTemplate,
                                     Paragraph, Spacer, Image, PageBreak,
                                     Table, TableStyle, KeepTogether,
                                     FrameBreak, NextPageTemplate)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_RIGHT, TA_CENTER
    from reportlab.lib import colors
    from reportlab.lib.units import cm, mm
    from reportlab.pdfgen import canvas

    # Normaliza branding (compat com chamada antiga só com titulo+job_id)
    b = branding or {}
    project_name = b.get('project_name') or titulo or 'Projeto sem nome'
    architect_name = b.get('architect_name', '')
    client_name = b.get('client_name', '')
    company = b.get('company', '')
    logo_path = b.get('logo_local_path')
    brand_color = b.get('brand_color') or COLOR_INDIGO
    ref_job = b.get('job_id') or job_id

    tmp_dir = tempfile.mkdtemp()
    gantt_png = os.path.join(tmp_dir, 'gantt.png')
    curva_png = os.path.join(tmp_dir, 'curva.png')
    gerar_gantt_png(cronograma, gantt_png)
    gerar_curva_s_png(cronograma, curva_png)

    resumo = cronograma.get('resumo', {})
    data_emissao = _dt.now().strftime('%d/%m/%Y')

    # Tamanho A4 paisagem
    PAGE_W, PAGE_H = landscape(A4)
    MARGIN = 2 * cm

    # Tipografia — defino estilos antes
    styles = getSampleStyleSheet()
    s_title = ParagraphStyle('T', parent=styles['Title'], fontName='Helvetica-Bold',
                              fontSize=32, leading=38,
                              textColor=colors.HexColor(COLOR_WHITE),
                              alignment=TA_LEFT, spaceAfter=0)
    # Subtítulo + capa-meta: cores ajustadas pra capa BRANCA (não mais
    # fullbleed colorida). Antes era lavanda (#C7D2FE/#E0E7FF) — só
    # funcionava com fundo indigo. Agora slate suave.
    s_subtitle = ParagraphStyle('ST', parent=styles['Heading2'],
                                 fontName='Helvetica', fontSize=15, leading=20,
                                 textColor=colors.HexColor(COLOR_GRAY_TX),
                                 alignment=TA_LEFT, spaceAfter=0)
    s_capa_meta = ParagraphStyle('CM', parent=styles['BodyText'],
                                  fontName='Helvetica', fontSize=11, leading=16,
                                  textColor=colors.HexColor(COLOR_GRAY_MID),
                                  alignment=TA_LEFT)
    s_h1 = ParagraphStyle('H1', parent=styles['Heading1'],
                           fontName='Helvetica-Bold', fontSize=24, leading=30,
                           textColor=colors.HexColor(COLOR_DARK),
                           alignment=TA_LEFT, spaceAfter=4, spaceBefore=0)
    s_eyebrow = ParagraphStyle('EB', parent=styles['BodyText'],
                                fontName='Helvetica-Bold', fontSize=9, leading=12,
                                textColor=colors.HexColor(COLOR_INDIGO),
                                alignment=TA_LEFT, spaceAfter=4)
    s_body = ParagraphStyle('B', parent=styles['BodyText'],
                             fontName='Helvetica', fontSize=10, leading=15,
                             textColor=colors.HexColor(COLOR_DARK_2),
                             alignment=TA_LEFT)
    s_caption = ParagraphStyle('CAP', parent=styles['BodyText'],
                                fontName='Helvetica-Oblique', fontSize=8,
                                leading=11,
                                textColor=colors.HexColor(COLOR_GRAY_LIGHT),
                                alignment=TA_LEFT)
    s_meta_val = ParagraphStyle('MV', parent=styles['BodyText'],
                                 fontName='Helvetica-Bold', fontSize=18,
                                 leading=22,
                                 textColor=colors.HexColor(COLOR_DARK),
                                 alignment=TA_LEFT)
    s_meta_lbl = ParagraphStyle('ML', parent=styles['BodyText'],
                                 fontName='Helvetica', fontSize=9, leading=11,
                                 textColor=colors.HexColor(COLOR_GRAY_TX),
                                 alignment=TA_LEFT)

    # ─── Header/Footer helpers (canvas) — usam branding ──
    def _draw_logo(canv, x, y, max_w=40*mm, max_h=8*mm):
        """Desenha logo do escritório se houver, senão texto 'AI.arq'."""
        if logo_path and os.path.exists(logo_path):
            try:
                from reportlab.lib.utils import ImageReader
                img = ImageReader(logo_path)
                iw, ih = img.getSize()
                # Mantém aspect ratio dentro do max_w/max_h
                ratio = min(max_w/iw, max_h/ih)
                w, h = iw*ratio, ih*ratio
                canv.drawImage(logo_path, x, y - h/2 + 2*mm,
                                width=w, height=h, mask='auto',
                                preserveAspectRatio=True)
                return w
            except Exception:
                pass
        # Fallback texto
        canv.setFont('Helvetica-Bold', 11)
        canv.setFillColor(colors.HexColor(brand_color))
        canv.drawString(x, y, company or 'AI.arq')
        return canv.stringWidth(company or 'AI.arq', 'Helvetica-Bold', 11)

    def header_footer(canv, doc):
        canv.saveState()
        # Footer linha sutil
        canv.setStrokeColor(colors.HexColor(COLOR_BORDER))
        canv.setLineWidth(0.4)
        canv.line(MARGIN, MARGIN - 6*mm, PAGE_W - MARGIN, MARGIN - 6*mm)
        # Logo escritório à esquerda topo
        logo_w = _draw_logo(canv, MARGIN, PAGE_H - MARGIN + 3*mm,
                            max_w=35*mm, max_h=7*mm)
        # Tagline depois do logo
        canv.setFont('Helvetica', 9)
        canv.setFillColor(colors.HexColor(COLOR_GRAY_TX))
        canv.drawString(MARGIN + logo_w + 6*mm,
                        PAGE_H - MARGIN + 3*mm,
                        '· Cronograma da obra')
        # Carimbo direita topo
        canv.setFont('Helvetica', 8)
        canv.setFillColor(colors.HexColor(COLOR_GRAY_LIGHT))
        carimbo = f'emitido {data_emissao}'
        if architect_name:
            carimbo = f'{architect_name}  ·  {carimbo}'
        canv.drawRightString(PAGE_W - MARGIN, PAGE_H - MARGIN + 3*mm, carimbo)
        # Footer: título projeto + nº página
        canv.setFont('Helvetica', 9)
        canv.setFillColor(colors.HexColor(COLOR_GRAY_LIGHT))
        canv.drawString(MARGIN, MARGIN - 12*mm, project_name[:70])
        page_num = canv.getPageNumber()
        canv.drawRightString(PAGE_W - MARGIN, MARGIN - 12*mm, f'{page_num}')
        # "Powered by AI.arq" discreto no centro
        canv.setFont('Helvetica-Oblique', 7)
        canv.setFillColor(colors.HexColor(COLOR_GRAY_LIGHT))
        canv.drawCentredString(PAGE_W/2, MARGIN - 12*mm,
                                'powered by AI.arq · ai.arq.br')
        canv.restoreState()

    def capa_canvas(canv, doc):
        """Capa minimal — fundo branco, tipografia generosa, hairline da
        cor da marca como acento. Profissional, sem fullbleed colorido.

        Estrutura:
        - Topo: hairline horizontal cor da marca (6pt)
        - Logo no canto sup esq + carimbo no canto sup dir
        - Centro-baixo: "Cronograma da obra" (eyebrow) + nome do projeto
          (display tipográfico grande) + cliente
        - Footer: linha + powered by + página
        """
        canv.saveState()
        # Fundo branco (não preencher — default já é branco)
        # Hairline da cor da marca no topo (só accent)
        canv.setFillColor(colors.HexColor(brand_color))
        canv.rect(0, PAGE_H - 6, PAGE_W, 6, fill=1, stroke=0)

        # Logo escritório (canto sup esq) — em cor preta/escura
        if logo_path and os.path.exists(logo_path):
            try:
                from reportlab.lib.utils import ImageReader
                img = ImageReader(logo_path)
                iw, ih = img.getSize()
                max_w, max_h = 50*mm, 12*mm
                ratio = min(max_w/iw, max_h/ih)
                w, h = iw*ratio, ih*ratio
                canv.drawImage(logo_path, MARGIN,
                                PAGE_H - MARGIN - h + 4*mm,
                                width=w, height=h, mask='auto',
                                preserveAspectRatio=True)
            except Exception:
                canv.setFont('Helvetica-Bold', 13)
                canv.setFillColor(colors.HexColor(COLOR_DARK))
                canv.drawString(MARGIN, PAGE_H - MARGIN, company or 'AI.arq')
        else:
            canv.setFont('Helvetica-Bold', 13)
            canv.setFillColor(colors.HexColor(COLOR_DARK))
            canv.drawString(MARGIN, PAGE_H - MARGIN, company or 'AI.arq')

        # Carimbo direita (data + arquiteto) — slate suave
        canv.setFont('Helvetica', 9)
        canv.setFillColor(colors.HexColor(COLOR_GRAY_TX))
        carimbo = f'Emitido em {data_emissao}'
        if architect_name:
            carimbo = f'{architect_name}  ·  {carimbo}'
        canv.drawRightString(PAGE_W - MARGIN, PAGE_H - MARGIN, carimbo)

        # Conteúdo principal — eyebrow + título do projeto
        # Eyebrow pequeno (letterspaced)
        canv.setFont('Helvetica-Bold', 9)
        canv.setFillColor(colors.HexColor(brand_color))
        canv.drawString(MARGIN, PAGE_H/2 + 30*mm, 'CRONOGRAMA FÍSICO DA OBRA')

        # Hairline divider sutil abaixo do eyebrow
        canv.setStrokeColor(colors.HexColor(COLOR_BORDER))
        canv.setLineWidth(0.5)
        canv.line(MARGIN, PAGE_H/2 + 25*mm, MARGIN + 60*mm, PAGE_H/2 + 25*mm)

        # Nome do projeto — display grande, cor escura
        canv.setFont('Helvetica-Bold', 36)
        canv.setFillColor(colors.HexColor(COLOR_DARK))
        # Trunca se for muito longo
        proj_display = project_name if len(project_name) <= 50 else project_name[:48] + '…'
        canv.drawString(MARGIN, PAGE_H/2 + 5*mm, proj_display)

        # Cliente (se houver) — label + valor abaixo do título
        if client_name:
            canv.setFont('Helvetica', 9)
            canv.setFillColor(colors.HexColor(COLOR_GRAY_LIGHT))
            canv.drawString(MARGIN, PAGE_H/2 - 12*mm, 'CLIENTE')
            canv.setFont('Helvetica', 16)
            canv.setFillColor(colors.HexColor(COLOR_DARK_2))
            canv.drawString(MARGIN, PAGE_H/2 - 20*mm, client_name)

        # Mini cards de meta (4 números-chave do resumo).
        # Fix 2026-06-09 (P0): antes eram desenhados no `story` em branco
        # (#FFFFFF) sobre fundo branco = invisíveis, com rótulos lavanda
        # (#A5B4FC/#C7D2FE) resíduo do design fullbleed antigo. Agora vão
        # aqui no canvas, posicionados abaixo do bloco cliente, em slate
        # escuro legível (valor) + slate-600 (rótulo).
        meta_y = (PAGE_H/2 - 34*mm) if client_name else (PAGE_H/2 - 16*mm)
        metas = [
            (f"{_format_br(resumo.get('data_inicio'))}", 'INÍCIO PREVISTO'),
            (f"{_format_br(resumo.get('data_fim'))}", 'TÉRMINO PREVISTO'),
            (f"{resumo.get('duracao_dias_reais', '—')} dias", 'DURAÇÃO TOTAL'),
            (f"{resumo.get('n_fases', 0)} fases", 'DISCIPLINAS'),
        ]
        col_w = 62 * mm
        for i, (val, lbl) in enumerate(metas):
            cx = MARGIN + i * col_w
            canv.setFont('Helvetica-Bold', 20)
            canv.setFillColor(colors.HexColor(COLOR_DARK))
            canv.drawString(cx, meta_y, str(val))
            canv.setFont('Helvetica', 9)
            canv.setFillColor(colors.HexColor(COLOR_GRAY_MID))
            canv.drawString(cx, meta_y - 7*mm, lbl)

        # Footer — linha + powered by + página
        canv.setStrokeColor(colors.HexColor(COLOR_BORDER))
        canv.setLineWidth(0.4)
        canv.line(MARGIN, MARGIN - 6*mm, PAGE_W - MARGIN, MARGIN - 6*mm)
        canv.setFont('Helvetica', 8)
        canv.setFillColor(colors.HexColor(COLOR_GRAY_LIGHT))
        canv.drawString(MARGIN, MARGIN - 12*mm, 'powered by AI.arq · ai.arq.br')
        canv.drawRightString(PAGE_W - MARGIN, MARGIN - 12*mm, 'Página 1 de 5')
        canv.restoreState()

    # ─── Doc com 2 page templates ──
    doc = BaseDocTemplate(
        output_path, pagesize=landscape(A4),
        leftMargin=MARGIN, rightMargin=MARGIN,
        topMargin=MARGIN + 8*mm, bottomMargin=MARGIN + 6*mm,
        title='Cronograma da obra — AI.arq', author='AI.arq')

    frame_capa = Frame(MARGIN, MARGIN + 10*mm,
                       PAGE_W - 2*MARGIN, PAGE_H - 2*MARGIN - 20*mm,
                       id='capa', showBoundary=0,
                       leftPadding=0, rightPadding=0,
                       topPadding=0, bottomPadding=0)
    frame_normal = Frame(MARGIN, MARGIN,
                         PAGE_W - 2*MARGIN, PAGE_H - 2*MARGIN - 6*mm,
                         id='normal', showBoundary=0,
                         leftPadding=0, rightPadding=0,
                         topPadding=0, bottomPadding=0)

    tpl_capa = PageTemplate(id='capa', frames=[frame_capa],
                             onPage=capa_canvas)
    tpl_normal = PageTemplate(id='normal', frames=[frame_normal],
                               onPage=header_footer)
    doc.addPageTemplates([tpl_capa, tpl_normal])

    story = []

    # ═══ PÁGINA 1 — CAPA ═══
    # Fix 2026-06-09 (P0): TODO o conteúdo visível da capa (eyebrow, título do
    # projeto, cliente e os 4 meta cards) é desenhado por `capa_canvas` em
    # slate escuro, com posicionamento absoluto. Antes o `story` redesenhava
    # título/cliente em branco por cima (texto duplicado/borrado) e os 4
    # números em branco-sobre-branco (invisíveis). Aqui só fica um Spacer pra
    # o frame da capa não ficar vazio (Platypus exige conteúdo no frame).
    story.append(Spacer(1, PAGE_H - 2*MARGIN - 20*mm))

    # ═══ PÁGINA 2 — GANTT ═══
    # Muda pro template normal (insere ANTES do PageBreak)
    story.append(NextPageTemplate('normal'))
    story.append(PageBreak())

    story.append(Paragraph('GANTT', s_eyebrow))
    story.append(Paragraph('Cronograma físico das disciplinas', s_h1))
    story.append(Spacer(1, 0.3*cm))
    # Imagem Gantt
    gantt_w = PAGE_W - 2*MARGIN
    story.append(Image(gantt_png, width=gantt_w, height=PAGE_H - 7*cm,
                       kind='proportional'))
    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph(
        '<b>Como ler:</b> cada barra representa uma disciplina; a posição '
        'horizontal mostra início e fim previstos; a largura indica duração '
        'em dias. Cores diferenciam categorias construtivas. Atraso em '
        'fases do caminho crítico atrasa a obra inteira (ver página 4).',
        s_caption))

    # ═══ PÁGINA 3 — CURVA S ═══
    story.append(PageBreak())
    story.append(Paragraph('CURVA S · AVANÇO PREVISTO', s_eyebrow))
    story.append(Paragraph('Como a obra progride no tempo', s_h1))
    story.append(Spacer(1, 0.3*cm))
    story.append(Image(curva_png, width=gantt_w, height=PAGE_H - 7.5*cm,
                       kind='proportional'))
    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph(
        '<b>Modelo sigmoidal (logístico, k=10).</b> Reflete distribuição real '
        'de obra — arranque suave (canteiro+fundação), pico de execução no '
        'meio (estrutura+vedação+instalações), finalização gradual '
        '(acabamentos). Marcos 25/50/75/100% sinalizados pra acompanhamento '
        'mensal (Lei 14.133 Art 117).',
        s_caption))

    # ═══ PÁGINA 4 — CAMINHO CRÍTICO ═══
    cc = resumo.get('caminho_critico', [])
    if cc:
        story.append(PageBreak())
        story.append(Paragraph('CAMINHO CRÍTICO', s_eyebrow))
        story.append(Paragraph('Top 5 fases mais longas', s_h1))
        story.append(Spacer(1, 0.1*cm))
        story.append(Paragraph(
            'Atraso em qualquer dessas fases atrasa a obra inteira. '
            'Priorize fornecedor, equipe e gestão de risco aqui.',
            s_body))
        story.append(Spacer(1, 0.5*cm))

        rows = [['#', 'Disciplina', 'Duração']]
        for i, item in enumerate(cc, 1):
            rows.append([
                str(i),
                item.get('label', ''),
                f"{item.get('dur_dias', '?')} dias"
            ])
        t = Table(rows, colWidths=[1.5*cm, 18*cm, 4*cm])
        t.setStyle(TableStyle([
            # Header
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(COLOR_INDIGO)),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('ALIGN', (0, 0), (0, 0), 'CENTER'),
            ('ALIGN', (-1, 0), (-1, 0), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('TOPPADDING', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            # Body
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 12),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.HexColor(COLOR_DARK)),
            ('ALIGN', (0, 1), (0, -1), 'CENTER'),
            ('ALIGN', (-1, 1), (-1, -1), 'CENTER'),
            ('FONTNAME', (-1, 1), (-1, -1), 'Helvetica-Bold'),
            ('TEXTCOLOR', (-1, 1), (-1, -1), colors.HexColor(COLOR_INDIGO)),
            ('TOPPADDING', (0, 1), (-1, -1), 12),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 12),
            ('LEFTPADDING', (1, 0), (1, -1), 14),
            ('LINEBELOW', (0, 0), (-1, -2), 0.5,
             colors.HexColor(COLOR_BORDER)),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1),
             [colors.HexColor('#F8FAFC'), colors.white]),
        ]))
        story.append(t)

    # ═══ PÁGINA 5 — MARCOS NORMATIVOS + RESSALVAS ═══
    story.append(PageBreak())
    story.append(Paragraph('REFERÊNCIAS', s_eyebrow))
    story.append(Paragraph('Marcos normativos e ressalvas', s_h1))
    story.append(Spacer(1, 0.5*cm))

    marcos = cronograma.get('marcos_legais', [])
    if marcos:
        story.append(Paragraph(
            '<b>Marcos considerados na elaboração deste cronograma:</b>',
            s_body))
        story.append(Spacer(1, 0.2*cm))
        # Tabela com 2 colunas de marcos
        n = len(marcos)
        mid = (n + 1) // 2
        col1 = marcos[:mid]
        col2 = marcos[mid:]
        # Empareia até maior tamanho
        rows = []
        for i in range(max(len(col1), len(col2))):
            a = col1[i] if i < len(col1) else ''
            b = col2[i] if i < len(col2) else ''
            rows.append([
                Paragraph(f'▸ {a}', s_body) if a else '',
                Paragraph(f'▸ {b}', s_body) if b else '',
            ])
        marc_table = Table(rows, colWidths=[12*cm, 12*cm])
        marc_table.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 0),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
        ]))
        story.append(marc_table)
        story.append(Spacer(1, 0.8*cm))

    # Box ressalva amarelo
    story.append(Paragraph('<b>Ressalvas</b>', s_body))
    story.append(Spacer(1, 0.2*cm))
    ressalva_text = cronograma.get('ressalva', '')
    ressalva_box = Table(
        [[Paragraph(f'<b>⚠ Cronograma de referência.</b> {ressalva_text}',
                    s_body)]],
        colWidths=[PAGE_W - 2*MARGIN])
    ressalva_box.setStyle(TableStyle([
        # Box sóbrio: fundo cinza muito claro + borda slate (era amarelo
        # gritante FEF3C7 + amber F59E0B — colorido demais pro contexto)
        ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor(COLOR_BG_SOFT)),
        ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor(COLOR_BORDER)),
        ('LINEBEFORE', (0, 0), (0, -1), 3, colors.HexColor(COLOR_GRAY_TX)),
        ('LEFTPADDING', (0, 0), (-1, -1), 16),
        ('RIGHTPADDING', (0, 0), (-1, -1), 16),
        ('TOPPADDING', (0, 0), (-1, -1), 14),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 14),
    ]))
    story.append(ressalva_box)
    story.append(Spacer(1, 1*cm))

    # Linha fim
    story.append(Paragraph(
        'Gerado por AI.arq · ai.arq.br · Quantitativo com IA pra '
        'arquitetos brasileiros', s_caption))

    doc.build(story)

    try:
        os.remove(gantt_png); os.remove(curva_png); os.rmdir(tmp_dir)
    except Exception:
        pass
    return output_path


# ═════════════════════════════════════════════════════════════════
#  PPTX — Layout executivo 16:9
# ═════════════════════════════════════════════════════════════════

def exportar_pptx(cronograma: Dict, output_path: str,
                   branding: Optional[Dict] = None,
                   # compat antiga
                   titulo: str = '', job_id: str = '') -> str:
    """PPT 16:9 CO-BRANDED com 5 slides executivos."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    b = branding or {}
    project_name = b.get('project_name') or titulo or 'Projeto sem nome'
    architect_name = b.get('architect_name', '')
    client_name = b.get('client_name', '')
    company = b.get('company', '')
    logo_path = b.get('logo_local_path')
    brand_color_hex = b.get('brand_color') or COLOR_INDIGO
    ref_job = b.get('job_id') or job_id

    tmp_dir = tempfile.mkdtemp()
    gantt_png = os.path.join(tmp_dir, 'gantt.png')
    curva_png = os.path.join(tmp_dir, 'curva.png')
    gerar_gantt_png(cronograma, gantt_png)
    gerar_curva_s_png(cronograma, curva_png)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Paleta como RGBColor — sóbria, slate puro. rgb_brand vem da config
    # do escritório (cor de identidade do user) e entra só como acento
    # (hairline da capa, eyebrow), nunca como background fullbleed.
    _r, _g, _b = _hex_to_rgb_int(brand_color_hex)
    rgb_brand = RGBColor(_r, _g, _b)
    # Tons slate (substituem indigo/cyan/amber gritantes do design antigo)
    rgb_indigo = RGBColor(0x47, 0x55, 0x69)        # slate-600 (era 4F46E5)
    rgb_indigo_dark = RGBColor(0x33, 0x41, 0x55)   # slate-700 (era 3730A3)
    rgb_cyan = RGBColor(0x64, 0x74, 0x8B)          # slate-500 (era 06B6D4)
    rgb_dark = RGBColor(0x0F, 0x17, 0x2A)
    rgb_dark2 = RGBColor(0x1E, 0x29, 0x3B)
    rgb_gray = RGBColor(0x47, 0x55, 0x69)
    rgb_gray_light = RGBColor(0x94, 0xA3, 0xB8)
    rgb_gray_mid = RGBColor(0x64, 0x74, 0x8B)
    rgb_white = RGBColor(0xFF, 0xFF, 0xFF)
    rgb_indigo_text = RGBColor(0x64, 0x74, 0x8B)   # slate-500 (era C7D2FE)
    rgb_border = RGBColor(0xE2, 0xE8, 0xF0)
    rgb_bg_soft = RGBColor(0xF8, 0xFA, 0xFC)       # slate-50 — cards
    rgb_amber_bg = RGBColor(0xF5, 0xF5, 0xF4)      # stone-100 (era FEF3C7 amarelo gritante)

    resumo = cronograma.get('resumo', {})
    data_emissao = _dt.now().strftime('%d/%m/%Y')

    def add_text(slide, text, left, top, width, height, size=14, bold=False,
                  color=rgb_dark, align='left', font='Calibri',
                  anchor='top'):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = {
            'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE,
            'bottom': MSO_ANCHOR.BOTTOM
        }[anchor]
        p = tf.paragraphs[0]
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                       'right': PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
        return tb

    def add_rect(slide, left, top, width, height, fill_rgb, line=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if not line:
            shape.line.fill.background()
        return shape

    def add_header(slide):
        """Header com logo escritório no canto + carimbo direita."""
        # Logo do escritório (se houver)
        if logo_path and os.path.exists(logo_path):
            try:
                slide.shapes.add_picture(logo_path, Inches(0.5),
                                          Inches(0.22), height=Inches(0.45))
                txt_x = Inches(2.3)
            except Exception:
                add_text(slide, company or 'AI.arq', Inches(0.5),
                          Inches(0.25), Inches(2), Inches(0.4),
                          size=14, bold=True, color=rgb_brand)
                txt_x = Inches(2)
        else:
            add_text(slide, company or 'AI.arq', Inches(0.5),
                      Inches(0.25), Inches(2.5), Inches(0.4),
                      size=14, bold=True, color=rgb_brand)
            txt_x = Inches(2.5)
        add_text(slide, '· Cronograma da obra',
                  txt_x, Inches(0.27),
                  Inches(5), Inches(0.4), size=11, color=rgb_gray)
        # Carimbo: arquiteto + data
        carimbo = f'emitido {data_emissao}'
        if architect_name:
            carimbo = f'{architect_name} · {carimbo}'
        add_text(slide, carimbo, Inches(0.5), Inches(0.27),
                  Inches(12.3), Inches(0.4), size=9, color=rgb_gray_light,
                  align='right')

    def add_footer(slide, n_slide):
        # Linha sutil
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.5), Inches(7.0),
                                       Inches(12.3), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb_border
        line.line.fill.background()
        add_text(slide, project_name[:70], Inches(0.5), Inches(7.1),
                  Inches(8), Inches(0.3), size=9, color=rgb_gray_light)
        add_text(slide, 'powered by AI.arq · ai.arq.br',
                  Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
                  size=8, color=rgb_gray_light, align='center')
        add_text(slide, f'{n_slide} / 5', Inches(0.5), Inches(7.1),
                  Inches(12.3), Inches(0.3), size=9, color=rgb_gray_light,
                  align='right')

    # ═══ SLIDE 1 — CAPA MINIMAL ═══
    # Refatorada 2026-05-25: antes era fullbleed colorido + faixa cyan
    # (visual carnavalesco). Agora é capa branca minimal com hairline da
    # cor da marca + tipografia generosa + cliente discreto.
    s1 = prs.slides.add_slide(blank)

    # Hairline horizontal no topo (cor da marca, fina) — single accent
    add_rect(s1, 0, 0, prs.slide_width, Emu(60000), rgb_brand)

    # Logo escritório (canto sup esq) em escala discreta
    if logo_path and os.path.exists(logo_path):
        try:
            s1.shapes.add_picture(logo_path, Inches(0.6), Inches(0.6),
                                   height=Inches(0.55))
        except Exception:
            add_text(s1, company or 'AI.arq', Inches(0.6), Inches(0.55),
                      Inches(4), Inches(0.4), size=13, bold=True,
                      color=rgb_dark)
    else:
        add_text(s1, company or 'AI.arq', Inches(0.6), Inches(0.55),
                  Inches(4), Inches(0.4), size=13, bold=True,
                  color=rgb_dark)

    # Carimbo direita — slate suave
    carimbo = f'Emitido em {data_emissao}'
    if architect_name:
        carimbo = f'{architect_name} · {carimbo}'
    add_text(s1, carimbo, Inches(0.6), Inches(0.6),
              Inches(12.1), Inches(0.4), size=10, color=rgb_gray,
              align='right')

    # Eyebrow letterspaced — usa cor da marca como única identidade visual
    add_text(s1, 'CRONOGRAMA FÍSICO DA OBRA', Inches(0.6), Inches(3.0),
              Inches(12), Inches(0.4), size=10, bold=True, color=rgb_brand)

    # Hairline divider abaixo do eyebrow
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.6), Inches(3.45),
                                Inches(2), Emu(8000))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb_border
    line.line.fill.background()

    # Nome do projeto — display grande slate escuro
    proj_display = project_name if len(project_name) <= 50 else project_name[:48] + '…'
    add_text(s1, proj_display, Inches(0.6), Inches(3.7),
              Inches(12.1), Inches(1.2), size=44, bold=True,
              color=rgb_dark)

    # Cliente (se houver) — label + valor
    if client_name:
        add_text(s1, 'CLIENTE', Inches(0.6), Inches(5.4),
                  Inches(8), Inches(0.3), size=9, bold=True,
                  color=rgb_gray_light)
        add_text(s1, client_name, Inches(0.6), Inches(5.7),
                  Inches(12.1), Inches(0.5), size=18, bold=True,
                  color=rgb_dark2)

    # Metadados em 4 colunas
    # Fix 2026-06-09 (P0): os 4 valores eram desenhados em rgb_white (branco)
    # sobre o slide branco = invisíveis, resíduo do design fullbleed antigo.
    # Agora vão em slate escuro (valor) + slate-500 (rótulo) e descem pra
    # Inches(6.0) pra não sobrepor o bloco cliente (que vai até ~5.7).
    meta_top = 6.0 if client_name else 5.5
    metas = [
        (_format_br(resumo.get('data_inicio')), 'INÍCIO PREVISTO'),
        (_format_br(resumo.get('data_fim')), 'TÉRMINO PREVISTO'),
        (f"{resumo.get('duracao_dias_reais', '—')} dias", 'DURAÇÃO TOTAL'),
        (f"{resumo.get('n_fases', 0)} fases", 'DISCIPLINAS'),
    ]
    for i, (val, lbl) in enumerate(metas):
        x = Inches(0.6 + i * 3.05)
        add_text(s1, val, x, Inches(meta_top), Inches(3), Inches(0.5),
                  size=22, bold=True, color=rgb_dark)
        add_text(s1, lbl, x, Inches(meta_top + 0.5), Inches(3), Inches(0.35),
                  size=9, bold=True, color=rgb_gray_mid)

    # Footer capa — slate suave (era lavanda #A5B4FC invisível/baixo contraste)
    add_text(s1, 'powered by AI.arq · ai.arq.br',
              Inches(0.6), Inches(7.1), Inches(10), Inches(0.3),
              size=9, color=rgb_gray_light)
    add_text(s1, '1 / 5', Inches(0.6), Inches(7.1),
              Inches(12.1), Inches(0.3), size=9,
              color=rgb_gray_light, align='right')

    # ═══ SLIDE 2 — GANTT ═══
    s2 = prs.slides.add_slide(blank)
    add_header(s2)
    add_text(s2, 'GANTT', Inches(0.5), Inches(0.85),
              Inches(4), Inches(0.4), size=10, bold=True, color=rgb_brand)
    add_text(s2, 'Cronograma físico das disciplinas',
              Inches(0.5), Inches(1.15), Inches(12), Inches(0.6),
              size=24, bold=True, color=rgb_dark)
    s2.shapes.add_picture(gantt_png, Inches(0.5), Inches(2.0),
                           width=Inches(12.3))
    add_text(s2,
              'Cada barra = uma disciplina. Posição horizontal = início/fim. '
              'Largura = duração. Cores diferenciam categorias.',
              Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
              size=9, color=rgb_gray_light)
    add_footer(s2, 2)

    # ═══ SLIDE 3 — CURVA S ═══
    s3 = prs.slides.add_slide(blank)
    add_header(s3)
    add_text(s3, 'CURVA S · AVANÇO PREVISTO', Inches(0.5), Inches(0.85),
              Inches(8), Inches(0.4), size=10, bold=True, color=rgb_brand)
    add_text(s3, 'Como a obra progride no tempo',
              Inches(0.5), Inches(1.15), Inches(12), Inches(0.6),
              size=24, bold=True, color=rgb_dark)
    s3.shapes.add_picture(curva_png, Inches(0.5), Inches(2.0),
                           width=Inches(12.3))
    add_text(s3,
              'Modelo sigmoidal (k=10). Reflete distribuição real de obra: '
              'arranque suave, pico no meio, finalização gradual.',
              Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
              size=9, color=rgb_gray_light)
    add_footer(s3, 3)

    # ═══ SLIDE 4 — CAMINHO CRÍTICO ═══
    s4 = prs.slides.add_slide(blank)
    add_header(s4)
    add_text(s4, 'CAMINHO CRÍTICO', Inches(0.5), Inches(0.85),
              Inches(6), Inches(0.4), size=10, bold=True, color=rgb_brand)
    add_text(s4, 'Top 5 fases mais longas',
              Inches(0.5), Inches(1.15), Inches(12), Inches(0.6),
              size=24, bold=True, color=rgb_dark)
    add_text(s4,
              'Atraso em qualquer dessas fases atrasa a obra inteira. '
              'Priorize fornecedor, equipe e gestão de risco aqui.',
              Inches(0.5), Inches(1.85), Inches(12), Inches(0.4),
              size=12, color=rgb_gray)

    cc = resumo.get('caminho_critico', [])
    y = Inches(2.6)
    for i, item in enumerate(cc, 1):
        # Card row
        card = add_rect(s4, Inches(0.5), y, Inches(12.3), Inches(0.75),
                         RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0
                         else rgb_white)
        # Borda
        card.line.color.rgb = rgb_border
        card.line.width = Pt(0.5)
        # Número grande
        add_text(s4, str(i), Inches(0.7), y, Inches(0.8), Inches(0.75),
                  size=24, bold=True, color=rgb_brand, anchor='middle')
        # Label
        add_text(s4, item.get('label', ''),
                  Inches(1.6), y, Inches(8.5), Inches(0.75),
                  size=16, bold=True, color=rgb_dark, anchor='middle')
        # Duração à direita
        add_text(s4, f"{item.get('dur_dias', '?')} dias",
                  Inches(0.5), y, Inches(12.3), Inches(0.75),
                  size=14, bold=True, color=rgb_brand,
                  align='right', anchor='middle')
        y += Inches(0.85)

    add_footer(s4, 4)

    # ═══ SLIDE 5 — MARCOS + RESSALVAS ═══
    s5 = prs.slides.add_slide(blank)
    add_header(s5)
    add_text(s5, 'REFERÊNCIAS', Inches(0.5), Inches(0.85),
              Inches(6), Inches(0.4), size=10, bold=True, color=rgb_brand)
    add_text(s5, 'Marcos normativos e ressalvas',
              Inches(0.5), Inches(1.15), Inches(12), Inches(0.6),
              size=24, bold=True, color=rgb_dark)

    # Marcos em 2 colunas
    marcos = cronograma.get('marcos_legais', [])
    add_text(s5, 'Marcos considerados:',
              Inches(0.5), Inches(2.05), Inches(12), Inches(0.4),
              size=13, bold=True, color=rgb_dark)
    n = len(marcos)
    mid = (n + 1) // 2
    y_col = Inches(2.55)
    for i, m in enumerate(marcos):
        col = 0 if i < mid else 1
        idx_in_col = i if col == 0 else i - mid
        x = Inches(0.7 + col * 6.4)
        y = y_col + Inches(idx_in_col * 0.35)
        add_text(s5, f'▸ {m}', x, y, Inches(6), Inches(0.35),
                  size=10, color=rgb_dark2)

    # Box ressalva amarela
    y_ressalva = Inches(5.0)
    ressalva_box = add_rect(s5, Inches(0.5), y_ressalva,
                             Inches(12.3), Inches(1.6), rgb_amber_bg)
    ressalva_box.line.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
    ressalva_box.line.width = Pt(1)
    add_text(s5, '⚠ Cronograma de referência',
              Inches(0.75), y_ressalva + Inches(0.15),
              Inches(11.8), Inches(0.4),
              size=13, bold=True, color=RGBColor(0x92, 0x40, 0x0E))
    add_text(s5, cronograma.get('ressalva', ''),
              Inches(0.75), y_ressalva + Inches(0.55),
              Inches(11.8), Inches(1.0),
              size=10, color=RGBColor(0x78, 0x35, 0x0F))

    add_footer(s5, 5)

    prs.save(output_path)
    try:
        os.remove(gantt_png); os.remove(curva_png); os.rmdir(tmp_dir)
    except Exception:
        pass
    return output_path
